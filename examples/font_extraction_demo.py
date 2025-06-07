#!/usr/bin/env python3
"""
Demo script showing font and bullet hierarchy extraction from PowerPoint templates.

This script demonstrates the T-37 Font & Bullet Hierarchy Extraction functionality
implemented in the TemplateParser.
"""

import tempfile
from pathlib import Path
from pptx import Presentation

# Add the open_lilli package to the path
import sys
sys.path.insert(0, str(Path(__file__).parent.parent))

from open_lilli.template_parser import TemplateParser
from open_lilli.models import TemplateStyle, FontInfo, BulletInfo, PlaceholderStyleInfo


def create_sample_template() -> str:
    """Create a sample PowerPoint template for demonstration."""
    prs = Presentation()
    
    # Add a title slide
    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    if title_slide.shapes.title:
        title_slide.shapes.title.text = "Font Extraction Demo"
    
    # Add content slides if available
    if len(prs.slide_layouts) > 1:
        content_layout = prs.slide_layouts[1]
        content_slide = prs.slides.add_slide(content_layout)
        
        if content_slide.shapes.title:
            content_slide.shapes.title.text = "Content Slide"
        
        # Add some body text with bullets
        for shape in content_slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame and not shape.shapes.title:
                shape.text = "• First bullet point\n  ○ Sub bullet point\n• Second bullet point"
                break
    
    # Save to temporary file
    temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
    temp_path = temp_file.name
    temp_file.close()
    prs.save(temp_path)
    
    return temp_path


def demonstrate_font_extraction():
    """Demonstrate the font and bullet hierarchy extraction."""
    print("🎨 Font & Bullet Hierarchy Extraction Demo")
    print("=" * 50)
    
    # Create a sample template
    print("\n📝 Creating sample template...")
    template_path = create_sample_template()
    
    try:
        # Initialize the template parser
        print("🔍 Analyzing template...")
        parser = TemplateParser(template_path)
        
        # Get the template style
        style = parser.get_template_style()
        print(f"✅ Template style extracted successfully!")
        
        # Display template information
        template_info = parser.get_template_info()
        style_info = template_info['template_style']
        
        print(f"\n📊 Template Overview:")
        print(f"   • Layouts: {template_info['total_layouts']}")
        print(f"   • Placeholder styles: {style_info['placeholder_styles_count']}")
        print(f"   • Has master font: {style_info['has_master_font']}")
        print(f"   • Styled placeholder types: {style_info['placeholder_types_with_styles']}")
        
        # Display theme fonts
        print(f"\n🎨 Theme Fonts:")
        for role, font_name in style.theme_fonts.items():
            print(f"   • {role.capitalize()}: {font_name}")
        
        # Display master font
        if style.master_font:
            font = style.master_font
            print(f"\n📝 Master Font:")
            print(f"   • Name: {font.name}")
            print(f"   • Size: {font.size or 'default'}")
            print(f"   • Weight: {font.weight or 'default'}")
            print(f"   • Color: {font.color or 'default'}")
        
        # Display placeholder-specific fonts
        print(f"\n🔤 Placeholder Fonts:")
        placeholder_type_names = {1: "TITLE", 2: "BODY", 3: "SUBTITLE", 7: "OBJECT"}
        
        for ph_type in [1, 2, 3, 7]:  # Common placeholder types
            font = parser.get_font_for_placeholder_type(ph_type)
            type_name = placeholder_type_names.get(ph_type, f"TYPE_{ph_type}")
            
            if font:
                print(f"   • {type_name}: {font.name}")
                if font.size:
                    print(f"     Size: {font.size}pt")
                if font.weight and font.weight != "normal":
                    print(f"     Weight: {font.weight}")
            else:
                print(f"   • {type_name}: No specific font (uses master font)")
        
        # Display bullet styles
        print(f"\n🔘 Bullet Styles:")
        for level in range(3):  # Show first 3 levels
            bullet = parser.get_bullet_style_for_level(2, level)  # BODY placeholder
            if bullet:
                print(f"   • Level {level}: '{bullet.character}'")
                if bullet.font:
                    print(f"     Font: {bullet.font.name} ({bullet.font.size or 'default'}pt)")
            else:
                print(f"   • Level {level}: No bullet style defined")
        
        # Test style accessor methods
        print(f"\n🔧 Style Accessor Methods:")
        
        # Test get_font_for_placeholder_type
        title_font = style.get_font_for_placeholder_type(1)
        if title_font:
            print(f"   • Title font: {title_font.name}")
        
        # Test get_bullet_style_for_level
        bullet = style.get_bullet_style_for_level(2, 0)
        if bullet:
            print(f"   • Level 0 bullet: '{bullet.character}'")
        
        # Test get_placeholder_style
        placeholder_style = style.get_placeholder_style(2)
        if placeholder_style:
            print(f"   • BODY placeholder has {len(placeholder_style.bullet_styles)} bullet styles")
        
        print(f"\n✅ Font extraction demo completed successfully!")
        print(f"\n💡 Key Features Demonstrated:")
        print(f"   • Template style extraction from .pptx files")
        print(f"   • Font hierarchy detection (master, placeholder-specific)")
        print(f"   • Bullet character and styling extraction")
        print(f"   • Theme font identification")
        print(f"   • Comprehensive style information access")
        
    except Exception as e:
        print(f"❌ Error during demonstration: {e}")
        raise
    
    finally:
        # Clean up
        Path(template_path).unlink()
        print(f"\n🧹 Temporary template cleaned up")


if __name__ == "__main__":
    demonstrate_font_extraction()