import unittest
from pathlib import Path
import logging
from unittest.mock import MagicMock, patch, PropertyMock

from pptx import Presentation # Added for template creation

from open_lilli.models import (
    SlidePlan, Outline, ContentFitConfig, TemplateStyle,
    FontInfo, PlaceholderStyleInfo, GenerationConfig, ContentFitResult,
    ContentDensityAnalysis, LayoutRecommendation
)
from open_lilli.slide_planner import SlidePlanner
from open_lilli.template_parser import TemplateParser
from open_lilli.content_fit_analyzer import ContentFitAnalyzer # Added

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class TestSlidePlannerIntegrationT76(unittest.TestCase):
    def setUp(self):
        self.test_dir = Path(__file__).parent
        self.test_data_dir = self.test_dir / "test_data"
        self.test_data_dir.mkdir(exist_ok=True)
        self.minimal_template_path = self.test_data_dir / "minimal_template_t76.pptx"

        if not self.minimal_template_path.exists():
            prs = Presentation()
            # Ensure layouts 0 (Title), 1 (Content), 3 (Two Content) exist
            # Default python-pptx presentation has these:
            # Layout 0: Title Slide
            # Layout 1: Title and Content
            # Layout 2: Section Header
            # Layout 3: Two Content
            # Layout 4: Comparison
            # Layout 5: Title Only
            # Layout 6: Blank
            if len(prs.slide_layouts) < 7:
                 logger.warning(f"Default template has only {len(prs.slide_layouts)} layouts. Some tests might be affected if expected layouts (0,1,3) are not standard.")
            # Add a slide for each layout type we might need to ensure they are "used" if template is truly minimal
            for i in range(min(7, len(prs.slide_layouts))):
                 try:
                    slide_layout = prs.slide_layouts[i]
                    prs.slides.add_slide(slide_layout)
                 except IndexError:
                    logger.warning(f"Could not access slide layout {i} in minimal template.")
            prs.save(self.minimal_template_path)
            logger.info(f"Created minimal template at {self.minimal_template_path}")

        self.template_parser = TemplateParser(str(self.minimal_template_path))

        # Override layout_map for predictable testing
        # Ensure these layout IDs exist in the minimal_template.pptx (python-pptx default layouts)
        self.template_parser.layout_map = {
            'title': 0,
            'content': 1,
            'section': 2,
            'two_column': 3,
            'comparison': 4,
            'title_only': 5,
            'blank': 6
        }
        self.template_parser.reverse_layout_map = {v: k for k, v in self.template_parser.layout_map.items()}


        self.content_fit_config = ContentFitConfig(
            characters_per_line=50,
            lines_per_placeholder=5,
            font_tune_threshold=1.01,
            rewrite_threshold=1.1,
            split_threshold=1.2,
            proportional_shrink_factor=0.9,
            max_proportional_shrink_cap_factor=0.85,
            min_font_size=10
        )

        self.mock_template_style = MagicMock(spec=TemplateStyle)
        mock_body_font = MagicMock(spec=FontInfo, name='Arial', size=18.0, weight='normal', color='#000000')

        mock_placeholder_style_body = MagicMock(spec=PlaceholderStyleInfo)
        mock_placeholder_style_body.default_font = mock_body_font

        self.mock_template_style.get_placeholder_style = MagicMock(return_value=mock_placeholder_style_body)
        self.mock_template_style.get_font_for_placeholder_type = MagicMock(return_value=mock_body_font)
        self.mock_template_style.language_specific_fonts = {}
        self.template_parser.template_style = self.mock_template_style # Ensure parser uses this mock style for analyzer

        self.openai_client_mock = MagicMock()

        self.layout_recommender_mock = MagicMock()
        def mock_recommend_layout(slide, available_layouts):
            layout_id_to_return = slide.layout_id if slide.layout_id is not None else self.template_parser.get_layout_index('content')
            rec_type = self.template_parser.get_layout_type_by_id(layout_id_to_return)
            if rec_type is None:
                rec_type = 'content'
                layout_id_to_return = self.template_parser.get_layout_index('content')
            return LayoutRecommendation(
                slide_type=rec_type,
                layout_id=layout_id_to_return,
                confidence=1.0,
                reasoning="mocked"
            )
        self.layout_recommender_mock.recommend_layout = MagicMock(side_effect=mock_recommend_layout)


        self.slide_planner = SlidePlanner(
            template_parser=self.template_parser,
            openai_client=self.openai_client_mock,
            content_fit_config=self.content_fit_config
        )
        # Ensure the planner's content_fit_analyzer also uses the mocked template_style if it's not passed down
        # This is important if optimize_slide_content in SlidePlanner doesn't pass its template_style to analyzer.analyze_slide_density
        # However, it seems SlidePlanner._optimize_content_fit does: template_style = getattr(self.template_parser, 'template_style', None)
        # And then passes it to self.content_fit_analyzer.optimize_slide_content(slide, template_style)
        # And ContentFitAnalyzer.optimize_slide_content passes it to analyze_slide_density & recommend_font_adjustment
        # So self.template_parser.template_style = self.mock_template_style is the key.

        self.slide_planner.layout_recommender = self.layout_recommender_mock
        self.slide_planner.enable_ml_layouts = True

        self.generation_config = GenerationConfig(max_bullets_per_slide=10) # Allow many bullets to test overflow logic

    def test_layout_upshift_resolves_overflow(self):
        logger.info("Starting test_layout_upshift_resolves_overflow")
        original_bullets = ["This is a very long line of text. " * 3] * 3 # Content that will require action
        initial_slide_plan = SlidePlan(
            index=0, slide_type="content", title="Test Upshift",
            bullets=original_bullets, layout_id=self.template_parser.get_layout_index("content")
        )
        outline = Outline(language="en", title="Upshift Test", slides=[initial_slide_plan])

        # --- Mocking for ContentFitAnalyzer ---
        # 1. Initial optimize_slide_content call
        summarized_bullets = ["Summarized: This is a long line. " * 2] * 3
        summarized_slide_plan = initial_slide_plan.model_copy()
        summarized_slide_plan.bullets = summarized_bullets
        summarized_slide_plan.summarized_by_llm = True

        # This is the density of the summarized content on the ORIGINAL 'content' layout
        # We want this to still require action to trigger upshift.
        density_after_summary_on_content_layout = ContentDensityAnalysis(
            total_characters=sum(len(b) for b in summarized_bullets) + len(summarized_slide_plan.title),
            estimated_lines=10, # Added dummy value
            placeholder_capacity=200, # Mocked: Summarized content (e.g. 300 chars) / 200 = 1.5 (needs action)
            density_ratio=1.5, # This will trigger split_slide if no upshift
            requires_action=True,
            recommended_action="split_slide" # Action if upshift wasn't available or failed
        )
        initial_fit_result = ContentFitResult(
            slide_index=0,
            density_analysis=density_after_summary_on_content_layout,
            final_action="rewrite_content", # Original action was rewrite
            modified_slide_plan=summarized_slide_plan # This is the key output
        )

        # 2. Subsequent analyze_slide_density call for the UPSHIFTED 'two_column' layout
        # This should show that the summarized content FITS in the new layout.
        density_on_upshifted_layout = ContentDensityAnalysis(
            total_characters=sum(len(b) for b in summarized_bullets) + len(summarized_slide_plan.title),
            estimated_lines=5, # Added dummy value
            placeholder_capacity=1000, # Mocked: Summarized content (e.g. 300 chars) / 1000 = 0.3 (fits)
            density_ratio=0.3,
            requires_action=False,
            recommended_action="no_action"
        )

        mock_optimize_content = MagicMock(return_value=initial_fit_result)
        mock_analyze_density_for_upshift = MagicMock(return_value=density_on_upshifted_layout)

        with patch.object(self.slide_planner.content_fit_analyzer, 'optimize_slide_content', mock_optimize_content), \
             patch.object(self.slide_planner.content_fit_analyzer, 'analyze_slide_density', mock_analyze_density_for_upshift):

            planned_slides = self.slide_planner.plan_slides(outline, self.generation_config)

        self.assertEqual(len(planned_slides), 1)
        final_slide = planned_slides[0]

        self.assertEqual(final_slide.layout_id, self.template_parser.get_layout_index("two_column"))
        self.assertTrue(final_slide.summarized_by_llm, "Slide should be marked as summarized")
        self.assertEqual(final_slide.bullets, summarized_bullets)

        # Check call to optimize_content_fit
        mock_optimize_content.assert_called_once()
        args_opt, _ = mock_optimize_content.call_args # template_style is positional
        called_slide_plan_opt = args_opt[0]
        called_template_style_opt = args_opt[1]
        self.assertEqual(called_slide_plan_opt.index, initial_slide_plan.index) # Check a few key fields
        self.assertEqual(called_slide_plan_opt.bullets, initial_slide_plan.bullets) # Before internal modifications by _plan_individual_slide
        self.assertEqual(called_template_style_opt, self.mock_template_style)

        # analyze_slide_density is called for the upshift check
        mock_analyze_density_for_upshift.assert_called_once()
        args_analyze, _ = mock_analyze_density_for_upshift.call_args # template_style is positional
        called_slide_plan_analyze = args_analyze[0]
        called_template_style_analyze = args_analyze[1]
        self.assertEqual(called_slide_plan_analyze.layout_id, self.template_parser.get_layout_index("two_column"))
        self.assertEqual(called_template_style_analyze, self.mock_template_style)

        self.assertNotIn("font_adjustment", final_slide.__dict__, "No font adjustment should be on final upshifted slide if it fits directly")
        logger.info("Finished test_layout_upshift_resolves_overflow")

    def test_slide_split_if_upshift_fails_or_unavailable(self):
        logger.info("Starting test_slide_split_if_upshift_fails_or_unavailable")
        original_bullets = ["This is an extremely long line. " * 10] * 5 # Very long
        initial_slide_plan = SlidePlan(
            index=0, slide_type="content", title="Test Split After Upshift Fail",
            bullets=original_bullets, layout_id=self.template_parser.get_layout_index("content")
        )
        outline = Outline(language="en", title="Split Test", slides=[initial_slide_plan])

        # --- Mocking for ContentFitAnalyzer ---
        # 1. Initial optimize_slide_content call
        summarized_bullets = ["Summarized: Extremely long line. " * 8] * 5 # Still very long
        summarized_slide_plan = initial_slide_plan.model_copy()
        summarized_slide_plan.bullets = summarized_bullets
        summarized_slide_plan.summarized_by_llm = True

        # Density of summarized content on ORIGINAL 'content' layout - requires action
        density_after_summary_on_content_layout = ContentDensityAnalysis(
            total_characters=sum(len(b) for b in summarized_bullets) + len(summarized_slide_plan.title),
            estimated_lines=10, # Added dummy value
            placeholder_capacity=300, # Summarized (e.g. 1000 chars) / 300 = 3.33 (needs split)
            density_ratio=3.33,
            requires_action=True,
            recommended_action="split_slide"
        )
        initial_fit_result = ContentFitResult(
            slide_index=0,
            density_analysis=density_after_summary_on_content_layout,
            final_action="rewrite_content", # Original action was rewrite
            modified_slide_plan=summarized_slide_plan
        )

        # 2. Subsequent analyze_slide_density call for UPSHIFTED 'two_column' layout
        # This should show that summarized content STILL DOES NOT FIT in the new layout.
        density_on_upshifted_layout_still_overflows = ContentDensityAnalysis(
            total_characters=sum(len(b) for b in summarized_bullets) + len(summarized_slide_plan.title),
            estimated_lines=5, # Added dummy value
            placeholder_capacity=500, # Summarized (e.g. 1000 chars) / 500 = 2.0 (still needs split)
            density_ratio=2.0,
            requires_action=True,
            recommended_action="split_slide" # Important: recommends split
        )

        mock_optimize_content = MagicMock(return_value=initial_fit_result)
        # Mock for analyze_slide_density called on the upshifted layout
        mock_analyze_density_for_upshift_fail = MagicMock(return_value=density_on_upshifted_layout_still_overflows)

        # Mock split_slide_content to verify it's called with the summarized content
        # Make it return a list of slides to satisfy the calling code
        split_slide_part1 = summarized_slide_plan.model_copy()
        split_slide_part1.title += " (Part 1)"
        split_slide_part1.bullets = summarized_bullets[:len(summarized_bullets)//2]
        split_slide_part2 = summarized_slide_plan.model_copy()
        split_slide_part2.title += " (Part 2)"
        split_slide_part2.bullets = summarized_bullets[len(summarized_bullets)//2:]
        mock_split_content_return = [split_slide_part1, split_slide_part2]
        mock_split_content = MagicMock(return_value=mock_split_content_return)


        with patch.object(self.slide_planner.content_fit_analyzer, 'optimize_slide_content', mock_optimize_content), \
             patch.object(self.slide_planner.content_fit_analyzer, 'analyze_slide_density', mock_analyze_density_for_upshift_fail), \
             patch.object(self.slide_planner.content_fit_analyzer, 'split_slide_content', mock_split_content):

            planned_slides = self.slide_planner.plan_slides(outline, self.generation_config)

        self.assertTrue(len(planned_slides) > 1, f"Slide should have been split. Got {len(planned_slides)} slides.")
        mock_split_content.assert_called_once()
        # Check that split_slide_content was called with the summarized slide plan
        args_split, _ = mock_split_content.call_args
        slide_passed_to_split = args_split[0]
        self.assertTrue(slide_passed_to_split.summarized_by_llm)
        self.assertEqual(slide_passed_to_split.bullets, summarized_bullets) # Should be called with the full summarized content

        self.assertTrue(planned_slides[0].summarized_by_llm, "First split slide should reflect summarization")
        self.assertTrue("Summarized: Extremely long line." in planned_slides[0].bullets[0]) # Check content from mock_split_content_return
        self.assertTrue("(Part 1)" in planned_slides[0].title, "Title of first split slide should indicate part 1")
        logger.info("Finished test_slide_split_if_upshift_fails_or_unavailable")

if __name__ == '__main__':
    unittest.main()
