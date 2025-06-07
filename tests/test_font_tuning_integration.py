import unittest
from pathlib import Path
import tempfile
from pptx import Presentation
from pptx.util import Pt
import logging

# Assuming open_lilli models and classes are structured for such imports
from open_lilli.models import SlidePlan, Outline, ContentFitConfig, TemplateStyle, FontInfo, PlaceholderStyleInfo, FontAdjustment, ContentDensityAnalysis
from open_lilli.content_fit_analyzer import ContentFitAnalyzer
from open_lilli.slide_assembler import SlideAssembler
from open_lilli.template_parser import TemplateParser
from unittest.mock import MagicMock, patch

# It's good practice to set up logging for tests to see messages from the library
logging.basicConfig(level=logging.INFO) # Use INFO for less verbose test output, DEBUG for dev
logger = logging.getLogger(__name__)

class TestFontTuningIntegration(unittest.TestCase):
    def setUp(self):
        self.test_dir = Path(__file__).parent
        self.test_data_dir = self.test_dir / "test_data"
        self.test_data_dir.mkdir(exist_ok=True)
        self.minimal_template_path = self.test_data_dir / "minimal_template.pptx"

        if not self.minimal_template_path.exists():
            prs = Presentation()
            # Layout 1 is 'Title and Content' which has a title and a body placeholder
            # Title placeholder is usually prs.slides[x].shapes.title or placeholders[0]
            # Body placeholder is usually prs.slides[x].placeholders[1]
            # Add a slide with layout 1 to ensure the layout exists
            slide_layout = prs.slide_layouts[1]
            prs.slides.add_slide(slide_layout)
            prs.save(self.minimal_template_path)
            logger.info(f"Created minimal template at {self.minimal_template_path}")

        # Use a real TemplateParser for SlideAssembler to load layouts
        self.template_parser_for_assembler = TemplateParser(str(self.minimal_template_path))

        # For ContentFitAnalyzer, we want to control the template style precisely
        self.mock_template_style = MagicMock(spec=TemplateStyle)
        mock_body_font = MagicMock(spec=FontInfo)
        mock_body_font.size = 18.0  # Original body font size for template cap calculation
        mock_placeholder_style_body = MagicMock(spec=PlaceholderStyleInfo)
        mock_placeholder_style_body.default_font = mock_body_font

        # Make get_placeholder_style return the mock for body (type 2)
        # and also for any other type if needed, or a generic one.
        def get_style_side_effect(ph_type):
            if ph_type == 2: # BODY placeholder type
                return mock_placeholder_style_body
            # Add other types if your code might query them
            # For title (type 1), let's return a default font too
            if ph_type == 1: # TITLE placeholder type
                mock_title_font = MagicMock(spec=FontInfo)
                mock_title_font.size = 24.0
                mock_placeholder_style_title = MagicMock(spec=PlaceholderStyleInfo)
                mock_placeholder_style_title.default_font = mock_title_font
                return mock_placeholder_style_title
            return None
        self.mock_template_style.get_placeholder_style = MagicMock(side_effect=get_style_side_effect)
        self.mock_template_style.get_font_for_placeholder_type = MagicMock(return_value=mock_body_font) # Default fallback
        self.mock_template_style.language_specific_fonts = {}


        self.content_fit_config = ContentFitConfig(
            proportional_shrink_factor=0.8,
            max_proportional_shrink_cap_factor=0.85,
            min_font_size=10.0,
            font_tune_threshold=1.01, # Low threshold to ensure tuning is attempted
            rewrite_threshold=1.3, # ensure this is higher than font_tune if only testing font
            split_threshold=1.5,
            characters_per_line=50, # Added for density calculation
            lines_per_placeholder=8   # Added for density calculation
        )
        self.analyzer = ContentFitAnalyzer(config=self.content_fit_config)
        # SlideAssembler uses the real parser to get layout objects
        self.assembler = SlideAssembler(self.template_parser_for_assembler)
        # Disable strict style validation for this integration test to focus on font tuning logic
        self.assembler.validation_config.enabled = False # Option 1: Disable entirely
        # self.assembler.validation_config.mode = "lenient" # Option 2: Make lenient
        self.temp_files = []

    def tearDown(self):
        for tmpfile in self.temp_files:
            Path(tmpfile).unlink(missing_ok=True)
            logger.info(f"Deleted temporary file {tmpfile}")

    def test_font_tuning_e2e(self):
        logger.info("Starting test_font_tuning_e2e...")
        # --- Non-RTL Slide ---
        non_rtl_bullets = ["This is a long line of text that will surely cause an overflow situation and require font tuning."] * 8
        non_rtl_slide_plan = SlidePlan(index=0, slide_type="content", title="Non-RTL Overflow", bullets=non_rtl_bullets, layout_id=1)

        # Expected font size calculation:
        # Current: 18.0. Proportional target: 18.0 * 0.8 = 14.4.
        # Template cap: 18.0 * 0.85 = 15.3. Min_font_size: 10.0.
        # Recommended = max(14.4, 15.3, 10.0) = 15.3. Rounded = 15.0.
        expected_non_rtl_font_size = 15.0

        # Manually set density analysis to trigger font adjustment
        # This bypasses the actual text length calculation for more direct testing of recommend_font_adjustment
        density_non_rtl = MagicMock(spec=ContentDensityAnalysis)
        density_non_rtl.density_ratio = 1.15 # To ensure it's within font_tune range
        density_non_rtl.overflow_severity = "mild" # Consistent with ratio
        density_non_rtl.recommended_action = "adjust_font" # Directly set for testing this path


        font_adj_non_rtl = self.analyzer.recommend_font_adjustment(density_non_rtl, self.mock_template_style, current_font_size=18.0)
        self.assertIsNotNone(font_adj_non_rtl, "Font adjustment should be recommended for Non-RTL slide")
        self.assertAlmostEqual(font_adj_non_rtl.recommended_size, expected_non_rtl_font_size, delta=0.5)
        non_rtl_slide_plan.__dict__['font_adjustment'] = font_adj_non_rtl

        # --- RTL Slide ---
        rtl_bullets = ["هذا سطرمحتوى طويل جدًا سيؤدي بالتأكيد إلى تجاوزالحدود ويتطلب ضبط الخط."] * 8
        rtl_slide_plan = SlidePlan(index=1, slide_type="content", title="RTL Overflow", bullets=rtl_bullets, layout_id=1)

        density_rtl = MagicMock(spec=ContentDensityAnalysis)
        density_rtl.density_ratio = 1.15 # Similar overflow
        density_rtl.overflow_severity = "mild"
        density_rtl.recommended_action = "adjust_font"


        font_adj_rtl = self.analyzer.recommend_font_adjustment(density_rtl, self.mock_template_style, current_font_size=18.0)
        self.assertIsNotNone(font_adj_rtl, "Font adjustment should still be recommended by analyzer for RTL slide")
        # Analyzer recommends, Assembler skips applying
        rtl_slide_plan.__dict__['font_adjustment'] = font_adj_rtl

        # --- Assemble ---
        outline = Outline(language="en", title="Font Tuning Test", slides=[non_rtl_slide_plan, rtl_slide_plan])

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            self.temp_files.append(tmp.name)
            output_path = Path(tmp.name) # Ensure output_path is a Path object

        original_add_slide = self.assembler._add_slide

        # Using a list to track calls if needed, or just for the side_effect
        mock_calls = []
        def mocked_add_slide_for_lang_test(prs, slide_plan_arg, slide_visuals, language_arg_original_call):
            effective_language = "ar" if slide_plan_arg.index == 1 else "en"
            mock_calls.append({'index': slide_plan_arg.index, 'lang': effective_language, 'adj': slide_plan_arg.__dict__.get('font_adjustment')})
            logger.info(f"Mocked _add_slide: Slide {slide_plan_arg.index}, Language {effective_language}, Original Lang: {language_arg_original_call}, Adjustment: {slide_plan_arg.__dict__.get('font_adjustment')}")
            # Call the original method with the potentially overridden language
            return original_add_slide(prs, slide_plan_arg, slide_visuals, effective_language)

        with patch.object(self.assembler, '_add_slide', side_effect=mocked_add_slide_for_lang_test) as mock_add_slide_method:
            logger.info(f"Calling assemble to path: {output_path}")
            # Assemble now takes Path object for output_path
            self.assembler.assemble(outline, [non_rtl_slide_plan, rtl_slide_plan], output_path=output_path)

        # --- Verify ---
        logger.info(f"Verifying output PPTX: {output_path}")
        prs_output = Presentation(str(output_path)) # Presentation constructor takes str or file-like
        self.assertEqual(len(prs_output.slides), 2, "Number of slides in output should be 2")

        # Non-RTL slide (Slide 0), body placeholder is placeholders[1] for layout_id 1 (0-indexed title and content)
        self.assertTrue(len(prs_output.slides[0].placeholders) >= 2, f"Non-RTL slide missing body placeholder. Found: {len(prs_output.slides[0].placeholders)}")
        body_ph_non_rtl = prs_output.slides[0].placeholders[1]
        self.assertTrue(body_ph_non_rtl.has_text_frame, "Non-RTL body placeholder has no text frame")
        # Check if there's text, implying content was added
        self.assertTrue(len(body_ph_non_rtl.text_frame.text.strip()) > 0, "Non-RTL body placeholder has no text")
        self.assertTrue(len(body_ph_non_rtl.text_frame.paragraphs) > 0, "Non-RTL body placeholder has no paragraphs")
        self.assertTrue(len(body_ph_non_rtl.text_frame.paragraphs[0].runs) > 0, "Non-RTL body placeholder has no runs")
        self.assertAlmostEqual(body_ph_non_rtl.text_frame.paragraphs[0].runs[0].font.size.pt, expected_non_rtl_font_size, delta=0.5)

        # RTL slide (Slide 1)
        self.assertTrue(len(prs_output.slides[1].placeholders) >= 2, f"RTL slide missing body placeholder. Found: {len(prs_output.slides[1].placeholders)}")
        body_ph_rtl = prs_output.slides[1].placeholders[1]
        self.assertTrue(body_ph_rtl.has_text_frame, "RTL body placeholder has no text frame")
        self.assertTrue(len(body_ph_rtl.text_frame.text.strip()) > 0, "RTL body placeholder has no text")
        self.assertTrue(len(body_ph_rtl.text_frame.paragraphs) > 0, "RTL body placeholder has no paragraphs")
        self.assertTrue(len(body_ph_rtl.text_frame.paragraphs[0].runs) > 0, "RTL body placeholder has no runs")
        self.assertAlmostEqual(body_ph_rtl.text_frame.paragraphs[0].runs[0].font.size.pt, 18.0, delta=0.5) # Should be original 18pt
        logger.info("test_font_tuning_e2e completed.")

if __name__ == '__main__':
    unittest.main()
