"""Tests for slide assembler's content splitting functionalities."""

import os
import shutil
from pathlib import Path
import pytest # Using pytest style for new tests

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import PP_PLACEHOLDER as PP_PLACEHOLDER_ENUM

from open_lilli.models import (
    Outline, SlidePlan, BulletItem,
    StyleValidationConfig, TextOverflowConfig, FontInfo
)
from open_lilli.slide_assembler import SlideAssembler
from open_lilli.template_parser import TemplateParser

# Utility function to create a minimal template for testing (module level)
def create_minimal_template_for_splitting_tests(filepath: str) -> str:
    prs = Presentation()
    # Layout 1: Title and Content (common layout, index 1 in slide_layouts)
    # Default layouts are: 0:Title, 1:Title+Content, 5:Title Only, 6:Blank
    title_content_layout = prs.slide_layouts[1]

    dir_path = Path(filepath).parent
    if not dir_path.exists():
        dir_path.mkdir(parents=True, exist_ok=True)

    prs.save(filepath)
    return filepath

class TestSlideAssemblerSplitting:
    TEST_DATA_DIR = Path(__file__).parent / "test_data_splitting"
    TEMPLATE_FILE = TEST_DATA_DIR / "minimal_template_for_splitting.pptx"
    OUTPUT_DIR = Path(__file__).parent / "test_output_splitting"

    _created_files_for_cleanup = []

    @classmethod
    def setup_class(cls):
        if not cls.TEST_DATA_DIR.exists():
            cls.TEST_DATA_DIR.mkdir(parents=True, exist_ok=True)
        if not cls.OUTPUT_DIR.exists():
            cls.OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

        create_minimal_template_for_splitting_tests(str(cls.TEMPLATE_FILE))

    def setup_method(self, method):
        self.template_parser = TemplateParser(template_path=str(self.TEMPLATE_FILE))
        # Ensure template_style and its language_specific_fonts are initialized
        if not hasattr(self.template_parser, 'template_style') or self.template_parser.template_style is None:
            # A basic mock for template_style if not fully loaded by TemplateParser minimal init
            mock_style = Mock()
            mock_style.language_specific_fonts = {}
            mock_style.master_font = FontInfo(name="Calibri", size=12.0)
            self.template_parser.template_style = mock_style
        elif not hasattr(self.template_parser.template_style, 'language_specific_fonts'):
            self.template_parser.template_style.language_specific_fonts = {}


        self.outline = Outline(title="Test Splitting Presentation", slides=[], language="en")
        self._created_files_for_cleanup = []

    def teardown_method(self, method):
        for filepath in self._created_files_for_cleanup:
            if Path(filepath).exists():
                try:
                    os.remove(filepath)
                except OSError as e:
                    print(f"Error removing file {filepath} in teardown: {e}")

    @classmethod
    def teardown_class(cls):
        if cls.TEMPLATE_FILE.exists():
            try:
                os.remove(str(cls.TEMPLATE_FILE))
            except OSError as e:
                print(f"Error removing template file {cls.TEMPLATE_FILE}: {e}")

        if cls.OUTPUT_DIR.exists():
            try:
                shutil.rmtree(str(cls.OUTPUT_DIR))
            except OSError as e:
                 print(f"Error removing output directory {cls.OUTPUT_DIR}: {e}")

        if cls.TEST_DATA_DIR.exists() and not any(cls.TEST_DATA_DIR.iterdir()):
             try:
                cls.TEST_DATA_DIR.rmdir()
             except OSError as e:
                print(f"Error removing test_data directory {cls.TEST_DATA_DIR}: {e}")


    def _get_output_path(self, filename: str) -> str:
        path = str(self.OUTPUT_DIR / filename)
        self._created_files_for_cleanup.append(path)
        return path

    def _count_bullets_in_body_placeholder(self, slide, layout_type_for_idx_reference: int = 1) -> int:
        count = 0
        target_placeholder = None

        # Try common body placeholder types/indices for layout 1 ("Title and Content")
        # Placeholder index 1 is typical for body/content in this layout.
        if layout_type_for_idx_reference == 1:
            try:
                ph = slide.placeholders[1]
                if ph.placeholder_format.type in [PP_PLACEHOLDER_ENUM.BODY, PP_PLACEHOLDER_ENUM.CONTENT, PP_PLACEHOLDER_ENUM.OBJECT]:
                    target_placeholder = ph
            except IndexError: # Not enough placeholders
                pass

        if not target_placeholder: # Fallback scan if specific index didn't work or different layout
            for ph in slide.placeholders:
                if ph.placeholder_format.type in [PP_PLACEHOLDER_ENUM.BODY, PP_PLACEHOLDER_ENUM.CONTENT, PP_PLACEHOLDER_ENUM.OBJECT]:
                    if ph.placeholder_format.type != PP_PLACEHOLDER_ENUM.TITLE and getattr(ph, 'name', '').lower() != 'title':
                        target_placeholder = ph
                        break

        if target_placeholder and target_placeholder.has_text_frame:
            for para in target_placeholder.text_frame.paragraphs:
                if para.text.strip():
                    count += 1
        return count

    def test_bullet_overflow_splitting_hierarchical(self):
        output_path = self._get_output_path("test_splitting_hierarchical.pptx")

        text_overflow_config = TextOverflowConfig(
            enable_bullet_splitting=True, max_lines_per_placeholder=7, split_slide_title_suffix="(Cont.)"
        )
        style_config = StyleValidationConfig(text_overflow_config=text_overflow_config)
        # Initialize assembler with the real parser for these tests
        assembler = SlideAssembler(self.template_parser, validation_config=style_config)

        bullet_items = [BulletItem(text=f"This is bullet {i+1}, level 0", level=0) for i in range(10)]
        slide_plan1 = SlidePlan(index=0, title="Hierarchical Split Test", slide_type="content",
                                bullet_hierarchy=bullet_items, layout_id=1) # Layout 1: Title and Content
        slides_input = [slide_plan1]
        self.outline.slides = slides_input

        assembler.assemble(self.outline, slides_input, output_path=output_path)

        prs = Presentation(output_path)
        assert len(prs.slides) == 2, "Presentation should have 2 slides after hierarchical splitting."

        assert prs.slides[0].shapes.title is not None and prs.slides[0].shapes.title.has_text_frame
        assert prs.slides[0].shapes.title.text == "Hierarchical Split Test"

        assert prs.slides[1].shapes.title is not None and prs.slides[1].shapes.title.has_text_frame
        assert prs.slides[1].shapes.title.text == f"Hierarchical Split Test {text_overflow_config.split_slide_title_suffix}"

        bullets_on_slide1 = self._count_bullets_in_body_placeholder(prs.slides[0], layout_type_for_idx_reference=1)
        bullets_on_slide2 = self._count_bullets_in_body_placeholder(prs.slides[1], layout_type_for_idx_reference=1)

        assert bullets_on_slide1 == 7
        assert bullets_on_slide2 == 3

    def test_bullet_overflow_splitting_legacy(self):
        output_path = self._get_output_path("test_splitting_legacy.pptx")

        text_overflow_config = TextOverflowConfig(enable_bullet_splitting=True, max_lines_per_placeholder=6, split_slide_title_suffix="(Continued)")
        style_config = StyleValidationConfig(text_overflow_config=text_overflow_config)
        assembler = SlideAssembler(self.template_parser, validation_config=style_config)

        legacy_bullets = [f"Legacy bullet point {i+1}" for i in range(9)]
        slide_plan1 = SlidePlan(index=0, title="Legacy Split Test", slide_type="content",
                                bullets=legacy_bullets, layout_id=1)
        slides_input = [slide_plan1]
        self.outline.slides = slides_input

        assembler.assemble(self.outline, slides_input, output_path=output_path)

        prs = Presentation(output_path)
        assert len(prs.slides) == 2

        assert prs.slides[0].shapes.title is not None and prs.slides[0].shapes.title.has_text_frame
        assert prs.slides[0].shapes.title.text == "Legacy Split Test"
        assert prs.slides[1].shapes.title is not None and prs.slides[1].shapes.title.has_text_frame
        assert prs.slides[1].shapes.title.text == f"Legacy Split Test {text_overflow_config.split_slide_title_suffix}"

        bullets_on_slide1 = self._count_bullets_in_body_placeholder(prs.slides[0], 1)
        bullets_on_slide2 = self._count_bullets_in_body_placeholder(prs.slides[1], 1)

        assert bullets_on_slide1 == 6
        assert bullets_on_slide2 == 3

    def test_bullet_overflow_splitting_disabled(self):
        output_path = self._get_output_path("test_splitting_disabled.pptx")

        text_overflow_config = TextOverflowConfig(enable_bullet_splitting=False, max_lines_per_placeholder=5)
        style_config = StyleValidationConfig(text_overflow_config=text_overflow_config)
        assembler = SlideAssembler(self.template_parser, validation_config=style_config)

        bullet_items = [BulletItem(text=f"Bullet {i+1}", level=0) for i in range(10)]
        slide_plan1 = SlidePlan(index=0, title="No Split Test", slide_type="content",
                                bullet_hierarchy=bullet_items, layout_id=1)
        slides_input = [slide_plan1]
        self.outline.slides = slides_input

        assembler.assemble(self.outline, slides_input, output_path=output_path)

        prs = Presentation(output_path)
        assert len(prs.slides) == 1
        assert prs.slides[0].shapes.title is not None and prs.slides[0].shapes.title.has_text_frame
        assert prs.slides[0].shapes.title.text == "No Split Test"
        bullets_on_slide1 = self._count_bullets_in_body_placeholder(prs.slides[0], 1)
        assert bullets_on_slide1 == 10

# To run these tests with pytest, ensure pytest is installed and run `pytest` in the terminal
# from the root directory of the project.
# Example:
# poetry run pytest tests/test_slide_splitting.py
