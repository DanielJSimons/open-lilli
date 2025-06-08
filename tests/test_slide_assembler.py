"""Tests for slide assembler."""

import tempfile
from pathlib import Path
from unittest.mock import Mock, patch
from typing import List

import pytest
from pptx import Presentation

from open_lilli.models import Outline, SlidePlan
from open_lilli.slide_assembler import SlideAssembler
from open_lilli.template_parser import TemplateParser


class TestSlideAssembler:
    """Tests for SlideAssembler class."""

    def setup_method(self):
        """Set up test fixtures."""
        # Create a basic test template
        self.temp_dir = Path(tempfile.mkdtemp())
        self.template_path = self.temp_dir / "test_template.pptx"
        
        # Create minimal template
        prs = Presentation()
        prs.save(str(self.template_path))
        
        # Mock template parser
        self.mock_template_parser = Mock(spec=TemplateParser)
        self.mock_template_parser.template_path = self.template_path
        self.mock_template_parser.get_layout_index.return_value = 0
        
        # Mock presentation with layouts
        self.mock_template_parser.prs = Mock()
        self.mock_template_parser.prs.slide_layouts = [Mock() for _ in range(3)]
        
        self.assembler = SlideAssembler(self.mock_template_parser)

    def teardown_method(self):
        """Clean up test fixtures."""
        import shutil
        shutil.rmtree(self.temp_dir, ignore_errors=True)

    def create_test_outline(self) -> Outline:
        """Create a test outline."""
        return Outline(
            title="Test Presentation",
            subtitle="A test presentation",
            slides=[]
        )

    def create_test_slides(self) -> List[SlidePlan]:
        """Create test slides."""
        return [
            SlidePlan(
                index=0,
                slide_type="title",
                title="Presentation Title",
                bullets=[],
                layout_id=0,
                speaker_notes="Welcome to the presentation"
            ),
            SlidePlan(
                index=1,
                slide_type="content",
                title="Content Slide",
                bullets=["Point 1", "Point 2", "Point 3"],
                layout_id=1,
                speaker_notes="Discuss the main points"
            )
        ]

    def test_init(self):
        """Test SlideAssembler initialization."""
        assert self.assembler.template_parser == self.mock_template_parser
        assert self.assembler.max_title_length == 60
        assert self.assembler.max_bullet_length == 120

    def test_assemble_basic(self):
        """Test basic presentation assembly."""
        outline = self.create_test_outline()
        slides = self.create_test_slides()
        output_path = self.temp_dir / "output.pptx"
        
        # Mock the presentation methods we'll use
        with patch('pptx.Presentation') as mock_prs_class:
            mock_prs = Mock()
            mock_prs.slides = Mock()
            mock_prs.slides.__len__ = Mock(return_value=0)  # No existing slides
            mock_prs.slide_layouts = [Mock(), Mock()]
            mock_prs_class.return_value = mock_prs
            
            # Mock slide creation
            mock_slide = Mock()
            mock_prs.slides.add_slide.return_value = mock_slide
            mock_slide.shapes.title = Mock()
            mock_slide.placeholders = []
            
            result_path = self.assembler.assemble(outline, slides, {}, output_path)
            
            assert result_path == output_path
            assert mock_prs.save.called

    def test_add_title(self):
        """Test adding title to slide."""
        mock_slide = Mock()
        mock_title_shape = Mock()
        mock_slide.shapes.title = mock_title_shape
        
        self.assembler._add_title(mock_slide, "Test Title")
        
        assert mock_title_shape.text == "Test Title"

    def test_add_title_too_long(self):
        """Test title truncation when too long."""
        mock_slide = Mock()
        mock_title_shape = Mock()
        mock_slide.shapes.title = mock_title_shape
        
        long_title = "A" * 100  # Longer than max_title_length
        self.assembler._add_title(mock_slide, long_title)
        
        # Should be truncated
        assert len(mock_title_shape.text) <= self.assembler.max_title_length

    def test_add_title_no_placeholder(self):
        """Test handling when no title placeholder exists."""
        mock_slide = Mock()
        mock_slide.shapes.title = None
        
        # Should not raise exception
        self.assembler._add_title(mock_slide, "Test Title")

    def test_add_bullet_content(self):
        """Test adding bullet points to slide."""
        mock_slide = Mock()
        mock_placeholder = Mock()
        mock_placeholder.placeholder_format.type = 2  # BODY type
        mock_slide.placeholders = [mock_placeholder]
        
        mock_text_frame = Mock()
        mock_placeholder.text_frame = mock_text_frame
        mock_paragraph = Mock()
        mock_text_frame.paragraphs = [mock_paragraph]
        mock_text_frame.add_paragraph.return_value = mock_paragraph
        
        bullets = ["Point 1", "Point 2", "Point 3"]
        self.assembler._add_bullet_content(mock_slide, bullets)
        
        # Should have called clear and add_paragraph
        mock_text_frame.clear.assert_called_once()
        assert mock_text_frame.add_paragraph.call_count == 2  # First uses existing paragraph

    def test_add_bullet_content_no_placeholder(self):
        """Test bullet content when no content placeholder exists."""
        mock_slide = Mock()
        mock_slide.placeholders = []
        
        bullets = ["Point 1", "Point 2"]
        
        # Should not raise exception
        self.assembler._add_bullet_content(mock_slide, bullets)

    def test_add_chart_image_with_placeholder(self):
        """Test adding chart image with picture placeholder."""
        mock_slide = Mock()
        mock_placeholder = Mock()
        mock_placeholder.placeholder_format.type = 18  # PICTURE type
        mock_slide.placeholders = [mock_placeholder]
        
        # Create a test image file
        test_image = self.temp_dir / "test_chart.png"
        test_image.write_bytes(b"fake image data")
        
        self.assembler._add_chart_image(mock_slide, str(test_image))
        
        mock_placeholder.insert_picture.assert_called_once_with(str(test_image))

    def test_add_chart_image_no_placeholder(self):
        """Test adding chart image without picture placeholder."""
        mock_slide = Mock()
        mock_slide.placeholders = []
        mock_slide.shapes = Mock()
        mock_slide.part.presentation.slide_width = 10000000  # EMUs
        mock_slide.part.presentation.slide_height = 7500000
        
        # Create a test image file
        test_image = self.temp_dir / "test_chart.png"
        test_image.write_bytes(b"fake image data")
        
        self.assembler._add_chart_image(mock_slide, str(test_image))
        
        mock_slide.shapes.add_picture.assert_called_once()

    def test_add_chart_image_missing_file(self):
        """Test handling of missing chart file."""
        mock_slide = Mock()
        mock_slide.placeholders = []
        
        # Should not raise exception
        self.assembler._add_chart_image(mock_slide, "nonexistent.png")

    def test_add_speaker_notes(self):
        """Test adding speaker notes to slide."""
        mock_slide = Mock()
        mock_notes_slide = Mock()
        mock_text_frame = Mock()
        
        mock_slide.notes_slide = mock_notes_slide
        mock_notes_slide.notes_text_frame = mock_text_frame
        
        notes = "These are speaker notes"
        self.assembler._add_speaker_notes(mock_slide, notes)
        
        assert mock_text_frame.text == notes

    def test_add_fallback_slide(self):
        """Test adding fallback slide when normal creation fails."""
        mock_prs = Mock()
        mock_layout = Mock()
        mock_prs.slide_layouts = [mock_layout]
        mock_slide = Mock()
        mock_prs.slides.add_slide.return_value = mock_slide
        mock_slide.shapes.title = Mock()
        
        slide_plan = SlidePlan(
            index=1,
            slide_type="content",
            title="Test Slide",
            bullets=[]
        )
        
        self.assembler._add_fallback_slide(mock_prs, slide_plan)
        
        assert mock_slide.shapes.title.text == "Test Slide"

    def test_apply_metadata(self):
        """Test applying presentation metadata."""
        mock_prs = Mock()
        mock_core_properties = Mock()
        mock_prs.core_properties = mock_core_properties
        
        outline = Outline(
            title="Test Presentation",
            subtitle="Test Subtitle",
            slides=[]
        )
        
        self.assembler._apply_metadata(mock_prs, outline)
        
        assert mock_core_properties.title == "Test Presentation"
        assert mock_core_properties.subject == "Test Subtitle"
        assert mock_core_properties.author == "Open Lilli AI"

    def test_analyze_slide_placeholders(self):
        """Test slide placeholder analysis."""
        mock_slide = Mock()
        mock_placeholder = Mock()
        mock_placeholder.placeholder_format.type = 2
        mock_placeholder.name = "Content Placeholder"
        mock_placeholder.shape_type = 14
        mock_slide.placeholders = [mock_placeholder]
        
        analysis = self.assembler.analyze_slide_placeholders(mock_slide)
        
        assert analysis["total_placeholders"] == 1
        assert len(analysis["placeholders"]) == 1
        assert analysis["placeholders"][0]["type"] == 2
        assert analysis["placeholders"][0]["name"] == "Content Placeholder"

    def test_get_assembly_statistics(self):
        """Test assembly statistics generation."""
        slides = self.create_test_slides()
        visuals = {
            1: {"chart": "chart.png", "image": "image.jpg"}
        }
        
        stats = self.assembler.get_assembly_statistics(slides, visuals)
        
        assert stats["total_slides"] == 2
        assert stats["slides_with_visuals"] == 1
        assert stats["total_bullets"] == 3
        assert stats["slides_with_notes"] == 2
        assert stats["slide_types"]["title"] == 1
        assert stats["slide_types"]["content"] == 1
        assert stats["visual_types"]["charts"] == 1
        assert stats["visual_types"]["images"] == 1

    def test_validate_slides_before_assembly(self):
        """Test slide validation before assembly."""
        # Create slides with various issues
        slides = [
            SlidePlan(
                index=0,
                slide_type="title",
                title="",  # Missing title
                bullets=[],
                layout_id=0
            ),
            SlidePlan(
                index=1,
                slide_type="content", 
                title="Valid Title",
                bullets=["Point"] * 15,  # Too many bullets
                layout_id=99  # Invalid layout ID
            ),
            SlidePlan(
                index=2,
                slide_type="content",
                title="A" * 100,  # Title too long
                bullets=[],
                layout_id=1
            )
        ]
        
        issues = self.assembler.validate_slides_before_assembly(slides)
        
        assert len(issues) >= 4  # Should find multiple issues
        assert any("no title" in issue.lower() for issue in issues)
        assert any("too many bullets" in issue.lower() for issue in issues)
        assert any("invalid layout id" in issue.lower() for issue in issues)
        assert any("title is too long" in issue.lower() for issue in issues)

    def test_validate_slides_empty_list(self):
        """Test validation with empty slide list."""
        issues = self.assembler.validate_slides_before_assembly([])
        
        assert len(issues) == 1
        assert "No slides provided" in issues[0]

    def test_validate_slides_valid(self):
        """Test validation with valid slides."""
        slides = self.create_test_slides()
        
        issues = self.assembler.validate_slides_before_assembly(slides)
        
        assert len(issues) == 0

    def test_create_slide_from_layout(self):
        """Test creating slide from specific layout."""
        mock_prs = Mock()
        mock_layout = Mock()
        mock_prs.slide_layouts = [mock_layout]
        mock_slide = Mock()
        mock_prs.slides.add_slide.return_value = mock_slide
        mock_slide.shapes.title = Mock()
        
        self.assembler.create_slide_from_layout(
            mock_prs, "content", "Test Title", ["Point 1", "Point 2"]
        )
        
        mock_prs.slides.add_slide.assert_called_once_with(mock_layout)
        assert mock_slide.shapes.title.text == "Test Title"

    def test_assemble_with_visuals(self):
        """Test assembly with visual elements."""
        outline = self.create_test_outline()
        slides = self.create_test_slides()
        
        # Create test visual files
        chart_file = self.temp_dir / "chart.png"
        chart_file.write_bytes(b"fake chart")
        image_file = self.temp_dir / "image.jpg"
        image_file.write_bytes(b"fake image")
        
        visuals = {
            1: {
                "chart": str(chart_file),
                "image": str(image_file)
            }
        }
        
        output_path = self.temp_dir / "output_with_visuals.pptx"
        
        with patch('pptx.Presentation') as mock_prs_class:
            mock_prs = Mock()
            mock_prs.slides = Mock()
            mock_prs.slides.__len__ = Mock(return_value=0)
            mock_prs.slide_layouts = [Mock(), Mock()]
            mock_prs_class.return_value = mock_prs
            
            mock_slide = Mock()
            mock_prs.slides.add_slide.return_value = mock_slide
            mock_slide.shapes.title = Mock()
            mock_slide.placeholders = []
            mock_slide.shapes = Mock()
            mock_slide.part.presentation.slide_width = 10000000
            mock_slide.part.presentation.slide_height = 7500000
            
            result_path = self.assembler.assemble(outline, slides, visuals, output_path)
            
            assert result_path == output_path

# Tests for T-55 (RTL & Font Fallback)
class TestSlideAssemblerInternationalization:

    def setup_method(self):
        """Set up test fixtures for internationalization tests."""
        self.mock_template_parser = Mock(spec=TemplateParser)
        self.mock_template_style = Mock()

        # Configure language_specific_fonts for tests
        self.mock_template_style.language_specific_fonts = {
            "ar": "Test Arabic Font",
            "he": "Test Hebrew Font",
            # "en" is not here, will use default
        }
        self.mock_template_parser.template_style = self.mock_template_style

        # Mock methods from template_parser that _get_expected_font might call
        self.mock_template_parser.get_bullet_style_for_level.return_value = None
        self.mock_template_parser.get_font_for_placeholder_type.return_value = Mock(
            name="Calibri", size=18, weight="normal", color="#000000"
        )
        self.mock_template_style.master_font = Mock(
            name="Calibri", size=12, weight="normal", color="#000000"
        )


        self.assembler = SlideAssembler(template_parser=self.mock_template_parser)

        # Mock for PP_ALIGN
        self.PP_ALIGN_RIGHT = Mock() # Simulate PP_ALIGN.RIGHT
        self.PP_ALIGN_LEFT = Mock()  # Simulate PP_ALIGN.LEFT

        # Patch PP_ALIGN where it's imported in slide_assembler
        patcher = patch('open_lilli.slide_assembler.PP_ALIGN')
        self.mock_pp_align = patcher.start()
        self.mock_pp_align.RIGHT = self.PP_ALIGN_RIGHT
        self.mock_pp_align.LEFT = self.PP_ALIGN_LEFT # And any other used values

    def teardown_method(self):
        patch.stopall()

    def create_mock_slide_with_text_frame(self):
        """Creates a mock slide with title and body placeholders having text_frames and paragraphs."""
        mock_slide = Mock()

        # Title shape
        mock_title_shape = Mock()
        mock_title_shape.text_frame = Mock()
        mock_title_paragraph = Mock()
        mock_title_paragraph.runs = [Mock()] # Ensure run exists
        mock_title_shape.text_frame.paragraphs = [mock_title_paragraph]
        mock_slide.shapes.title = mock_title_shape

        # Body placeholder for bullets
        mock_body_placeholder = Mock()
        mock_body_placeholder.placeholder_format.type = 2 # BODY
        mock_body_placeholder.text_frame = Mock()

        # Mock paragraphs for bullet content
        # _add_bullet_content clears and then uses paragraphs[0] or adds new ones
        mock_bullet_para1 = Mock()
        mock_bullet_para1.runs = [Mock()]
        mock_bullet_para1.text = "Bullet 1" # for the strip() check in RTL alignment

        mock_body_placeholder.text_frame.paragraphs = [mock_bullet_para1]
        def add_paragraph_mock():
            new_para = Mock()
            new_para.runs = [Mock()]
            return new_para
        mock_body_placeholder.text_frame.add_paragraph.side_effect = add_paragraph_mock

        mock_slide.placeholders = [mock_body_placeholder]

        return mock_slide, mock_title_paragraph, mock_bullet_para1


    @pytest.mark.parametrize("language, should_align_rtl", [
        ("ar", True), ("he", True), ("fa", True), ("en", False), ("fr", False)
    ])
    def test_rtl_text_alignment(self, language, should_align_rtl):
        """Test _add_title and _add_bullet_content for RTL text alignment."""
        mock_slide, mock_title_para, mock_bullet_para = self.create_mock_slide_with_text_frame()

        # Test _add_title
        self.assembler._add_title(mock_slide, "Test Title", language)
        if should_align_rtl:
            assert mock_title_para.alignment == self.PP_ALIGN_RIGHT
        else:
            # Ensure it wasn't set to RIGHT if not RTL (could be LEFT, or None if not touched)
            if mock_title_para.alignment == self.PP_ALIGN_RIGHT:
                 # This case means it was set to RIGHT when it shouldn't have been
                assert not should_align_rtl # Fail the test

        # Reset alignment for bullet test
        mock_title_para.alignment = None
        mock_bullet_para.alignment = None

        # Test _add_bullet_content
        self.assembler._add_bullet_content(mock_slide, ["Bullet 1"], language)
        if should_align_rtl:
            # This check needs to be on the actual paragraph object that received text
            # For a single bullet, it's the first paragraph in the (mocked) text_frame
            assert mock_slide.placeholders[0].text_frame.paragraphs[0].alignment == self.PP_ALIGN_RIGHT
        else:
            if mock_slide.placeholders[0].text_frame.paragraphs[0].alignment == self.PP_ALIGN_RIGHT:
                assert not should_align_rtl


    @pytest.mark.parametrize("language, expected_font", [
        ("ar", "Test Arabic Font"),
        ("he", "Test Hebrew Font"),
        ("en", None) # No specific font for English in mock
    ])
    def test_proactive_font_selection(self, language, expected_font):
        """Test proactive font selection in _add_title and _add_bullet_content."""
        mock_slide, mock_title_para, mock_bullet_para = self.create_mock_slide_with_text_frame()

        title_run_font = mock_title_para.runs[0].font
        bullet_run_font = mock_bullet_para.runs[0].font

        # Test _add_title
        self.assembler._add_title(mock_slide, "Test Title", language)
        if expected_font:
            assert title_run_font.name == expected_font
        else:
            # Ensure it wasn't set if no specific font, or check it was set to a default if that's the logic
            # Current logic: only sets if specific_font_name is found.
             assert title_run_font.name != "Test Arabic Font" # Example, ensure it's not accidentally set
             assert title_run_font.name != "Test Hebrew Font"


        # Reset font name for bullet test
        title_run_font.name = None
        bullet_run_font.name = None

        # Test _add_bullet_content
        self.assembler._add_bullet_content(mock_slide, ["Bullet 1"], language)
        if expected_font:
            assert bullet_run_font.name == expected_font
        else:
            assert bullet_run_font.name != "Test Arabic Font"
            assert bullet_run_font.name != "Test Hebrew Font"


    def test_style_validation_aware_of_font_override(self):
        """Test _validate_run_style's awareness of language-specific fonts."""
        mock_run = Mock()
        mock_run.font.name = "Test Arabic Font" # Actual font used

        # Expected font from template (before override)
        expected_font_from_template = Mock(spec=SlidePlan().model_fields['title'].default_factory().__class__) # Dummy FontInfo like object
        expected_font_from_template.name = "Calibri"
        expected_font_from_template.size = 18
        expected_font_from_template.weight = "normal"
        expected_font_from_template.color = "#000000"

        # Enable font name enforcement for test
        self.assembler.validation_config.enforce_font_name = True

        # Test case 1: Language is "ar", font is the correct override "Test Arabic Font"
        # _get_expected_font will return a FontInfo with name "Test Arabic Font"
        # So _validate_run_style should receive an expected_font.name that matches run.font.name

        # We need to mock _get_expected_font to simulate its new behavior correctly,
        # or trust its implementation and test _validate_run_style directly with crafted expected_font.
        # The implementation of _validate_run_style itself has the override check logic.

        violations = self.assembler._validate_run_style(
            mock_run,
            expected_font_from_template, # This is what _validate_paragraph_style would pass after calling _get_expected_font without lang initially
            slide_index=0, para_index=0, run_index=0,
            language="ar" # This language is key for the internal override check
        )
        font_name_violations = [v for v in violations if v['type'] == 'font_name']
        assert not font_name_violations, f"Validation failed for correct override: {font_name_violations}"

        # Test case 2: Language is "ar", but font is something unexpected
        mock_run.font.name = "Comic Sans MS"
        violations_unexpected = self.assembler._validate_run_style(
            mock_run, expected_font_from_template,
            slide_index=0, para_index=0, run_index=0, language="ar"
        )
        font_name_violations_unexpected = [v for v in violations_unexpected if v['type'] == 'font_name']
        assert font_name_violations_unexpected, "Validation did not catch incorrect font with language override active."
        assert font_name_violations_unexpected[0]['actual'] == "Comic Sans MS"
        # The expected from `expected_font_from_template` is "Calibri", but the logic inside
        # _validate_run_style should ideally state what it expected considering the override language.
        # However, the current implementation of _validate_run_style's override check just suppresses the error.
        # It doesn't change the `expected` field in the violation report if an override was possible but not matched.

        # Test case 3: Language is "en" (no override), font is "Test Arabic Font" (mismatch)
        mock_run.font.name = "Test Arabic Font"
        violations_en = self.assembler._validate_run_style(
            mock_run, expected_font_from_template,
            slide_index=0, para_index=0, run_index=0, language="en"
        )
        font_name_violations_en = [v for v in violations_en if v['type'] == 'font_name']
        assert font_name_violations_en, "Validation failed for non-overridden language."
        assert font_name_violations_en[0]['expected'] == "Calibri"
        assert font_name_violations_en[0]['actual'] == "Test Arabic Font"


class TestSubtitlePlaceholderGuard:
    """Test cases for T-93: Subtitle Placeholder Guard."""
    
    def setup_method(self):
        """Set up test fixtures."""
        self.mock_template_parser = Mock(spec=TemplateParser)
        self.assembler = SlideAssembler(self.mock_template_parser)
    
    def test_t93_hide_empty_subtitle_placeholder(self):
        """Test T-93: Empty subtitle placeholder is hidden to avoid style warnings."""
        # Create mock slide with subtitle placeholder
        mock_slide = Mock()
        mock_subtitle_placeholder = Mock()
        mock_subtitle_placeholder.placeholder_format.type = 3  # SUBTITLE
        mock_subtitle_placeholder.text_frame = Mock()
        mock_subtitle_placeholder.text_frame.text = ""  # Empty
        
        # Mock other placeholders
        mock_title_placeholder = Mock()
        mock_title_placeholder.placeholder_format.type = 1  # TITLE
        mock_title_placeholder.text_frame = Mock()
        mock_title_placeholder.text_frame.text = "Title Text"
        
        mock_slide.placeholders = [mock_title_placeholder, mock_subtitle_placeholder]
        
        # Create slide plan without subtitle
        slide_plan = SlidePlan(
            index=0,
            slide_type="title",
            title="Test Title",
            bullets=[]  # No bullets, so no subtitle
        )
        
        # Mock the hide method to track calls
        with patch.object(self.assembler, '_hide_empty_placeholder') as mock_hide:
            self.assembler._add_title_slide_content(mock_slide, slide_plan, "en")
            
            # Should have called hide for empty subtitle
            mock_hide.assert_called_once_with(mock_subtitle_placeholder, "subtitle")
    
    def test_t93_no_hide_when_subtitle_has_content(self):
        """Test that subtitle placeholder is not hidden when it has content."""
        # Create mock slide with subtitle placeholder
        mock_slide = Mock()
        mock_subtitle_placeholder = Mock()
        mock_subtitle_placeholder.placeholder_format.type = 3  # SUBTITLE
        
        mock_slide.placeholders = [mock_subtitle_placeholder]
        
        # Create slide plan with subtitle content
        slide_plan = SlidePlan(
            index=0,
            slide_type="title",
            title="Test Title",
            bullets=["Subtitle content", "More content"]
        )
        
        # Mock the hide method to track calls
        with patch.object(self.assembler, '_hide_empty_placeholder') as mock_hide:
            self.assembler._add_title_slide_content(mock_slide, slide_plan, "en")
            
            # Should NOT have called hide since subtitle has content
            mock_hide.assert_not_called()
            
            # Should have set subtitle text
            expected_subtitle = "Subtitle content • More content"
            assert mock_subtitle_placeholder.text == expected_subtitle
    
    def test_t93_remove_empty_placeholders_from_slide(self):
        """Test removal of multiple empty placeholders from a slide."""
        # Create mock slide with various placeholders
        mock_slide = Mock()
        
        # Title placeholder (should never be hidden)
        mock_title = Mock()
        mock_title.placeholder_format.type = 1  # TITLE
        mock_title.text_frame = Mock()
        mock_title.text_frame.text = "Title"
        
        # Empty subtitle placeholder (should be hidden)
        mock_subtitle = Mock()
        mock_subtitle.placeholder_format.type = 3  # SUBTITLE
        mock_subtitle.text_frame = Mock()
        mock_subtitle.text_frame.text = ""
        
        # Empty footer placeholder (should be hidden)
        mock_footer = Mock()
        mock_footer.placeholder_format.type = 7  # FOOTER
        mock_footer.text_frame = Mock()
        mock_footer.text_frame.text = "   "  # Whitespace only
        
        # Non-empty body placeholder (should not be hidden)
        mock_body = Mock()
        mock_body.placeholder_format.type = 2  # BODY
        mock_body.text_frame = Mock()
        mock_body.text_frame.text = "Body content"
        
        mock_slide.placeholders = [mock_title, mock_subtitle, mock_footer, mock_body]
        
        # Mock the hide method to track calls
        with patch.object(self.assembler, '_hide_empty_placeholder') as mock_hide:
            hidden_count = self.assembler._remove_empty_placeholders_from_slide(mock_slide)
            
            # Should have hidden 2 placeholders (subtitle and footer)
            assert hidden_count == 2
            assert mock_hide.call_count == 2
            
            # Verify specific placeholders were hidden
            call_args = [call[0] for call in mock_hide.call_args_list]
            hidden_placeholders = [args[0] for args in call_args]
            assert mock_subtitle in hidden_placeholders
            assert mock_footer in hidden_placeholders
            assert mock_title not in hidden_placeholders  # Title should never be hidden
            assert mock_body not in hidden_placeholders   # Body has content
    
    def test_t93_hide_placeholder_methods(self):
        """Test various methods of hiding placeholders."""
        mock_placeholder = Mock()
        
        # Test Method 1: visible attribute
        mock_placeholder.visible = True
        self.assembler._hide_empty_placeholder(mock_placeholder, "test")
        assert mock_placeholder.visible is False
        
        # Test Method 3: positioning (if visible attribute doesn't exist)
        mock_placeholder_pos = Mock()
        del mock_placeholder_pos.visible  # Remove visible attribute
        mock_placeholder_pos.left = 100
        mock_placeholder_pos.top = 100
        
        self.assembler._hide_empty_placeholder(mock_placeholder_pos, "test")
        # Should have moved off-slide
        assert mock_placeholder_pos.left.value < 0  # Negative position
        assert mock_placeholder_pos.top.value < 0   # Negative position


class TestTwoColumnBulletDistribution:
    """Test cases for T-94: Two-Column Bullet Distribution."""
    
    def setup_method(self):
        """Set up test fixtures."""
        self.mock_template_parser = Mock(spec=TemplateParser)
        self.mock_template_parser.template_style = Mock()
        self.mock_template_parser.template_style.language_specific_fonts = {}
        self.assembler = SlideAssembler(self.mock_template_parser)
    
    def test_t94_two_column_bullet_distribution(self):
        """Test T-94: Bullets are distributed evenly across two BODY placeholders."""
        # Create mock slide with two BODY placeholders
        mock_slide = Mock()
        
        # First BODY placeholder
        mock_body1 = Mock()
        mock_body1.placeholder_format.type = 2  # BODY
        mock_body1.text_frame = Mock()
        mock_body1.text_frame.paragraphs = [Mock()]  # Initial paragraph
        mock_body1.text_frame.add_paragraph = Mock(return_value=Mock())
        
        # Second BODY placeholder
        mock_body2 = Mock()
        mock_body2.placeholder_format.type = 2  # BODY
        mock_body2.text_frame = Mock()
        mock_body2.text_frame.paragraphs = [Mock()]  # Initial paragraph
        mock_body2.text_frame.add_paragraph = Mock(return_value=Mock())
        
        mock_slide.placeholders = [mock_body1, mock_body2]
        
        # Test with 6 bullets - should be split 3-3
        bullets = [
            "Bullet 1", "Bullet 2", "Bullet 3", 
            "Bullet 4", "Bullet 5", "Bullet 6"
        ]
        
        # Mock the individual placeholder method to track calls
        with patch.object(self.assembler, '_add_bullets_to_placeholder') as mock_add_bullets:
            self.assembler._add_bullet_content(mock_slide, bullets, "en")
            
            # Should have called _add_bullets_to_placeholder twice
            assert mock_add_bullets.call_count == 2
            
            # Verify the first call (first column)
            first_call_args = mock_add_bullets.call_args_list[0]
            assert first_call_args[0][0] == mock_body1  # First placeholder
            assert first_call_args[0][1] == ["Bullet 1", "Bullet 2", "Bullet 3"]  # First 3 bullets
            
            # Verify the second call (second column)
            second_call_args = mock_add_bullets.call_args_list[1]
            assert second_call_args[0][0] == mock_body2  # Second placeholder
            assert second_call_args[0][1] == ["Bullet 4", "Bullet 5", "Bullet 6"]  # Last 3 bullets
    
    def test_t94_uneven_bullet_distribution(self):
        """Test T-94: Uneven bullet count distributed correctly across columns."""
        # Create mock slide with two BODY placeholders
        mock_slide = Mock()
        mock_body1 = Mock()
        mock_body1.placeholder_format.type = 2
        mock_body2 = Mock()
        mock_body2.placeholder_format.type = 2
        mock_slide.placeholders = [mock_body1, mock_body2]
        
        # Test with 5 bullets - should be split 3-2 (extra bullet goes to first column)
        bullets = ["Bullet 1", "Bullet 2", "Bullet 3", "Bullet 4", "Bullet 5"]
        
        with patch.object(self.assembler, '_add_bullets_to_placeholder') as mock_add_bullets:
            self.assembler._add_bullet_content(mock_slide, bullets, "en")
            
            # Verify distribution: 3 bullets in first column, 2 in second
            first_call_bullets = mock_add_bullets.call_args_list[0][0][1]
            second_call_bullets = mock_add_bullets.call_args_list[1][0][1]
            
            assert len(first_call_bullets) == 3  # First column gets extra bullet
            assert len(second_call_bullets) == 2
            assert first_call_bullets == ["Bullet 1", "Bullet 2", "Bullet 3"]
            assert second_call_bullets == ["Bullet 4", "Bullet 5"]
    
    def test_t94_single_body_placeholder_fallback(self):
        """Test that single BODY placeholder uses standard distribution."""
        # Create mock slide with only one BODY placeholder
        mock_slide = Mock()
        mock_body = Mock()
        mock_body.placeholder_format.type = 2  # BODY
        mock_slide.placeholders = [mock_body]
        
        bullets = ["Bullet 1", "Bullet 2", "Bullet 3"]
        
        with patch.object(self.assembler, '_add_bullets_to_placeholder') as mock_add_bullets:
            self.assembler._add_bullet_content(mock_slide, bullets, "en")
            
            # Should call _add_bullets_to_placeholder once with all bullets
            assert mock_add_bullets.call_count == 1
            call_args = mock_add_bullets.call_args_list[0]
            assert call_args[0][0] == mock_body
            assert call_args[0][1] == bullets  # All bullets in single column
    
    def test_t94_single_bullet_no_distribution(self):
        """Test that single bullet doesn't trigger two-column distribution."""
        # Create mock slide with two BODY placeholders
        mock_slide = Mock()
        mock_body1 = Mock()
        mock_body1.placeholder_format.type = 2
        mock_body2 = Mock()
        mock_body2.placeholder_format.type = 2
        mock_slide.placeholders = [mock_body1, mock_body2]
        
        # Single bullet should go to first placeholder only
        bullets = ["Single Bullet"]
        
        with patch.object(self.assembler, '_add_bullets_to_placeholder') as mock_add_bullets:
            self.assembler._add_bullet_content(mock_slide, bullets, "en")
            
            # Should use standard single-column distribution
            assert mock_add_bullets.call_count == 1
            call_args = mock_add_bullets.call_args_list[0]
            assert call_args[0][0] == mock_body1  # First placeholder
            assert call_args[0][1] == bullets
    
    def test_t94_distribute_bullets_across_columns_method(self):
        """Test the _distribute_bullets_across_columns method directly."""
        mock_body1 = Mock()
        mock_body2 = Mock()
        mock_body3 = Mock()  # Third placeholder to test 2-column limit
        
        body_placeholders = [mock_body1, mock_body2, mock_body3]
        bullets = ["A", "B", "C", "D", "E", "F", "G"]
        
        with patch.object(self.assembler, '_add_bullets_to_placeholder') as mock_add_bullets:
            self.assembler._distribute_bullets_across_columns(body_placeholders, bullets, "en")
            
            # Should limit to 2 columns even with 3 placeholders available
            assert mock_add_bullets.call_count == 2
            
            # Verify distribution: 4 bullets in first column, 3 in second
            first_call_bullets = mock_add_bullets.call_args_list[0][0][1]
            second_call_bullets = mock_add_bullets.call_args_list[1][0][1]
            
            assert len(first_call_bullets) == 4  # 7//2 + 1 (extra)
            assert len(second_call_bullets) == 3  # 7//2
            assert first_call_bullets == ["A", "B", "C", "D"]
            assert second_call_bullets == ["E", "F", "G"]