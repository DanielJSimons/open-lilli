import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE # For potential future use in helper

from open_lilli.slide_assembler import SlideAssembler
from open_lilli.template_parser import TemplateParser
from open_lilli.models import Outline, SlidePlan, StyleValidationConfig

# Path to the template presentation - ensure this is correct relative to repo root
TEMPLATE_PATH = "tests/fixtures/strategy-deck.pptx"

@pytest.fixture(scope="module")
def template_parser():
    """Fixture to provide a TemplateParser instance. Scoped to module for efficiency."""
    try:
        return TemplateParser(template_path=TEMPLATE_PATH)
    except Exception as e:
        # Attempt to locate the template if the initial path fails, common in different CWD contexts
        try:
            import os
            # Assuming the script runs from repo root or tests/ directory
            alt_path = os.path.join(os.path.dirname(__file__), "fixtures", "strategy-deck.pptx")
            if os.path.exists(alt_path):
                return TemplateParser(template_path=alt_path)
            # If still not found, try one level up from tests/ (repo root)
            alt_path_repo_root = os.path.join(os.path.dirname(os.path.dirname(__file__)), "tests", "fixtures", "strategy-deck.pptx")
            if os.path.exists(alt_path_repo_root) and "tests" in alt_path_repo_root : # simple check to avoid going too far up
                 return TemplateParser(template_path=alt_path_repo_root)

            pytest.fail(f"Failed to initialize TemplateParser with {TEMPLATE_PATH} or alternates: {e}")
        except Exception as e2:
            pytest.fail(f"Failed to initialize TemplateParser with {TEMPLATE_PATH} and alternates. Main error: {e}, Alt error: {e2}")


@pytest.fixture
def slide_assembler(template_parser):
    """Fixture to provide a SlideAssembler instance."""
    return SlideAssembler(template_parser=template_parser, validation_config=StyleValidationConfig())

@pytest.fixture
def new_presentation_with_layout(template_parser):
    """
    Helper fixture to get a new presentation and a function to add slides.
    This avoids issues with reusing the Presentation object from template_parser directly for slide addition tests.
    """
    def _get_presentation_and_add_slide(layout_idx):
        prs = Presentation(template_parser.template_path) # Load fresh from template path
        if layout_idx >= len(prs.slide_layouts):
            pytest.skip(f"Layout index {layout_idx} out of range for template {template_parser.template_path}")
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        return prs, slide
    return _get_presentation_and_add_slide


def is_placeholder_effectively_hidden(placeholder, prs):
    """
    Checks if a placeholder is hidden using various techniques employed by _hide_empty_placeholder.
    Args:
        placeholder: The placeholder shape to check.
        prs: The Presentation object, used for slide dimensions in off-slide checks.
    """
    if not placeholder:
        return False

    # Check 1: XML 'hidden' attribute (via element.set('hidden', '1'))
    # python-pptx Shape object doesn't have a direct 'visible' property reflecting all XML states.
    # The common way to hide a shape from appearing is to set its p:cNvPr 'hidden' attribute for selection pane hiding,
    # or by moving it off-slide / making it tiny for visual hiding.
    # SlideAssembler's _hide_empty_placeholder uses placeholder.element.set('hidden', '1')
    if hasattr(placeholder, 'element') and placeholder.element.get('hidden') == '1':
        return True # This is one of the attempts in _hide_empty_placeholder

    # Check 2: Moved off-slide
    if hasattr(placeholder, 'left') and hasattr(placeholder, 'top') and prs and \
       hasattr(prs, 'slide_width') and hasattr(prs, 'slide_height'):
        # Specific check for Inches(-10) if that's a magic number used
        if placeholder.left == Inches(-10) and placeholder.top == Inches(-10):
            return True
        # General off-slide check
        if (placeholder.left + placeholder.width <= Inches(0.01) or
            placeholder.top + placeholder.height <= Inches(0.01) or
            placeholder.left >= prs.slide_width or
            placeholder.top >= prs.slide_height):
            return True

    # Check 3: Resized to be tiny
    if hasattr(placeholder, 'width') and hasattr(placeholder, 'height'):
        if placeholder.width <= Inches(0.01) and placeholder.height <= Inches(0.01):
            return True

    return False

def test_non_placeholder_shape_removal(slide_assembler, new_presentation_with_layout):
    prs, slide = new_presentation_with_layout(5) # Layout 5 is "Blank" in strategy-deck

    left, top, width, height = Inches(1), Inches(1), Inches(2), Inches(1)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.name = "TestTextboxToRemove"

    assert any(shape.name == "TestTextboxToRemove" for shape in slide.shapes), "Textbox not added initially."

    slide_assembler._postprocess_slide(slide, shape_whitelist=None)

    found_shape = any(shape.name == "TestTextboxToRemove" for shape in slide.shapes)
    assert not found_shape, "Non-placeholder shape was not removed by _postprocess_slide."

def test_whitelisted_non_placeholder_shape_preservation(slide_assembler, new_presentation_with_layout):
    prs, slide = new_presentation_with_layout(5) # Blank layout

    textbox1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    textbox1.name = "WhitelistedShape"
    textbox2 = slide.shapes.add_textbox(Inches(4), Inches(1), Inches(2), Inches(1))
    textbox2.name = "ShapeToRemove"

    slide_assembler._postprocess_slide(slide, shape_whitelist=["WhitelistedShape"])

    final_shape_names = {shape.name for shape in slide.shapes if shape.name}
    assert "WhitelistedShape" in final_shape_names, "Whitelisted shape was removed."
    assert "ShapeToRemove" not in final_shape_names, "Non-whitelisted shape was not removed."


def test_empty_text_placeholder_hidden(slide_assembler, new_presentation_with_layout):
    # Layout 0: "Title Slide" (Title, Subtitle)
    prs, test_slide = new_presentation_with_layout(0)

    if test_slide.shapes.title: # Populate title
        test_slide.shapes.title.text = "Test Title"

    subtitle_placeholder = None
    for ph in test_slide.placeholders:
        if ph.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE: # SUBTITLE (3)
            subtitle_placeholder = ph
            break

    assert subtitle_placeholder is not None, "Could not find Subtitle placeholder on Layout 0."
    # Ensure it's empty
    if hasattr(subtitle_placeholder, "text_frame") and subtitle_placeholder.text_frame:
        subtitle_placeholder.text_frame.text = ""
        assert not subtitle_placeholder.text_frame.text.strip(), "Subtitle placeholder should be empty."

    slide_assembler._postprocess_slide(test_slide, shape_whitelist=None)

    # Re-fetch placeholder by its original element or idx to check its state
    refetched_subtitle_placeholder = None
    for shp in test_slide.placeholders:
        if hasattr(shp, 'placeholder_format') and shp.placeholder_format.idx == subtitle_placeholder.placeholder_format.idx:
            refetched_subtitle_placeholder = shp
            break

    assert refetched_subtitle_placeholder is not None, "Could not re-fetch subtitle placeholder."
    assert is_placeholder_effectively_hidden(refetched_subtitle_placeholder, prs), \
        f"Empty Subtitle placeholder (idx={subtitle_placeholder.placeholder_format.idx}) was not hidden."


def test_empty_picture_placeholder_hidden(slide_assembler, new_presentation_with_layout):
    # Layout 7: "Picture with Caption" (Title, Picture, Text)
    # Or Layout 1: "Title and Content" if 7 is not available
    prs, test_slide = new_presentation_with_layout(7 if len(Presentation(TEMPLATE_PATH).slide_layouts) > 7 else 1)

    if test_slide.shapes.title:
        test_slide.shapes.title.text = "Empty Picture Test"

    picture_placeholder = None
    for ph in test_slide.placeholders:
        if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE: # PICTURE (18)
            picture_placeholder = ph
            break
        elif ph.placeholder_format.type == PP_PLACEHOLDER.OBJECT: # OBJECT (10) can also be for pictures
             # Heuristic: if layout 1 (Title and Content), this is likely the one.
            if test_slide.slide_layout == prs.slide_layouts[1] and ph.placeholder_format.idx == 1:
                 picture_placeholder = ph
                 break

    assert picture_placeholder is not None, "Could not find a suitable Picture or Object placeholder."
    # Ensure it's empty (no image)
    try:
        if picture_placeholder.image:
             pytest.fail("Picture placeholder should be empty for this test (already has an image).")
    except ValueError:
        pass # Expected if no image

    slide_assembler._postprocess_slide(test_slide, shape_whitelist=None)

    refetched_pic_placeholder = None
    for shp in test_slide.placeholders:
        if hasattr(shp, 'placeholder_format') and shp.placeholder_format.idx == picture_placeholder.placeholder_format.idx:
            refetched_pic_placeholder = shp
            break

    assert refetched_pic_placeholder is not None, "Could not re-fetch picture placeholder."
    assert is_placeholder_effectively_hidden(refetched_pic_placeholder, prs), \
        f"Empty Picture/Object placeholder (idx={picture_placeholder.placeholder_format.idx}) was not hidden."


def test_content_placeholders_remain_unaffected(slide_assembler, new_presentation_with_layout):
    # Layout 1: "Title and Content"
    prs, test_slide = new_presentation_with_layout(1)

    title_text = "Full Slide Test"
    body_text = "This is some body text.\nWith multiple lines."

    title_ph = test_slide.shapes.title
    if title_ph:
        title_ph.text = title_text

    body_ph = None
    for ph in test_slide.placeholders: # Body placeholder is usually idx 1 on this layout
        if ph.placeholder_format.type == PP_PLACEHOLDER.BODY and ph.placeholder_format.idx == 1:
            body_ph = ph
            break
    assert body_ph is not None, "Could not find Body placeholder (idx 1) on Layout 1."
    body_ph.text_frame.text = body_text

    slide_assembler._postprocess_slide(test_slide, shape_whitelist=None)

    assert title_ph.text == title_text, "Title text was altered."
    assert not is_placeholder_effectively_hidden(title_ph, prs), "Title placeholder was hidden."
    assert body_ph.text_frame.text == body_text, "Body text was altered."
    assert not is_placeholder_effectively_hidden(body_ph, prs), "Body placeholder was hidden."


def test_filled_picture_placeholder_remains(slide_assembler, new_presentation_with_layout, tmp_path):
    pil_image = pytest.importorskip("PIL.Image")

    # Layout 7: "Picture with Caption" or Layout 1 as fallback
    prs, test_slide = new_presentation_with_layout(7 if len(Presentation(TEMPLATE_PATH).slide_layouts) > 7 else 1)

    if test_slide.shapes.title:
        test_slide.shapes.title.text = "Filled Picture Test"

    picture_placeholder = None
    for ph in test_slide.placeholders:
        if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            picture_placeholder = ph
            break
        elif ph.placeholder_format.type == PP_PLACEHOLDER.OBJECT:
             if test_slide.slide_layout == prs.slide_layouts[1] and ph.placeholder_format.idx == 1: # Content placeholder on layout 1
                 picture_placeholder = ph
                 break
    assert picture_placeholder is not None, "Could not find Picture/Object placeholder."

    dummy_image_path = tmp_path / "dummy.png"
    img = pil_image.new('RGB', (60, 30), color = 'red')
    img.save(dummy_image_path)

    picture_placeholder.insert_picture(str(dummy_image_path))
    assert picture_placeholder.image is not None # Verify image was inserted

    slide_assembler._postprocess_slide(test_slide, shape_whitelist=None)

    # Re-fetch to be sure
    refetched_pic_placeholder = None
    for shp in test_slide.placeholders:
        if hasattr(shp, 'placeholder_format') and shp.placeholder_format.idx == picture_placeholder.placeholder_format.idx:
            refetched_pic_placeholder = shp
            break

    assert refetched_pic_placeholder is not None
    assert refetched_pic_placeholder.image is not None, "Image was removed from filled picture placeholder."
    assert not is_placeholder_effectively_hidden(refetched_pic_placeholder, prs), "Filled picture placeholder was hidden."
