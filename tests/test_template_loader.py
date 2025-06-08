"""Unit tests for TemplateLoader hardened template loading."""

import pytest
import tempfile
from pathlib import Path
from unittest.mock import Mock, patch, MagicMock

from pptx import Presentation
from pptx.slide import SlideLayout

from open_lilli.template_loader import TemplateLoader


class TestTemplateLoader:
    """Test cases for TemplateLoader class."""
    
    def test_template_loader_initialization_file_not_found(self):
        """Test that FileNotFoundError is raised for non-existent template."""
        with pytest.raises(FileNotFoundError, match="Template file not found"):
            TemplateLoader("non_existent_template.pptx")
    
    def test_template_loader_initialization_invalid_extension(self):
        """Test that ValueError is raised for non-pptx files."""
        with tempfile.NamedTemporaryFile(suffix=".txt") as temp_file:
            with pytest.raises(ValueError, match="Template must be a .pptx file"):
                TemplateLoader(temp_file.name)
    
    @patch('open_lilli.template_loader.Presentation')
    def test_slides_stripped_on_load(self, mock_presentation_class):
        """Test acceptance criteria: calling loader leaves prs.slides empty."""
        # Create mock presentation with slides
        mock_prs = Mock()
        mock_prs.slides = Mock()
        mock_prs.slides._sld_lst = [Mock(), Mock(), Mock()]  # 3 mock slides
        mock_prs.slide_layouts = [Mock(), Mock()]  # 2 mock layouts
        mock_prs.part = Mock()
        
        # Configure slide layouts for classification
        for i, layout in enumerate(mock_prs.slide_layouts):
            layout.placeholders = []
            if i == 0:
                # Title layout: TITLE + SUBTITLE
                title_ph = Mock()
                title_ph.placeholder_format.type = 1  # TITLE
                subtitle_ph = Mock()
                subtitle_ph.placeholder_format.type = 3  # SUBTITLE
                layout.placeholders = [title_ph, subtitle_ph]
            else:
                # Content layout: TITLE + BODY
                title_ph = Mock()
                title_ph.placeholder_format.type = 1  # TITLE
                body_ph = Mock()
                body_ph.placeholder_format.type = 2  # BODY
                layout.placeholders = [title_ph, body_ph]
        
        mock_presentation_class.return_value = mock_prs
        
        # Create temp file to satisfy path validation
        with tempfile.NamedTemporaryFile(suffix=".pptx") as temp_file:
            loader = TemplateLoader(temp_file.name)
            
            # Verify slides were stripped
            assert len(loader.prs.slides._sld_lst) == 0, "Slides should be empty after loading"
            
            # Verify layouts are preserved
            assert len(loader.prs.slide_layouts) == 2, "Layouts should be preserved"
    
    @patch('open_lilli.template_loader.Presentation')
    def test_get_layout_index_content_layout(self, mock_presentation_class):
        """Test acceptance criteria: get_layout_index("content") returns correct index."""
        mock_prs = Mock()
        mock_prs.slides = Mock()
        mock_prs.slides._sld_lst = []
        mock_prs.part = Mock()
        
        # Create mock layouts
        title_layout = Mock()
        title_layout.placeholders = []
        title_ph = Mock()
        title_ph.placeholder_format.type = 1  # TITLE
        subtitle_ph = Mock()
        subtitle_ph.placeholder_format.type = 3  # SUBTITLE
        title_layout.placeholders = [title_ph, subtitle_ph]
        
        content_layout = Mock()
        content_layout.placeholders = []
        title_ph2 = Mock()
        title_ph2.placeholder_format.type = 1  # TITLE
        body_ph = Mock()
        body_ph.placeholder_format.type = 2  # BODY
        content_layout.placeholders = [title_ph2, body_ph]
        
        mock_prs.slide_layouts = [title_layout, content_layout]
        mock_presentation_class.return_value = mock_prs
        
        with tempfile.NamedTemporaryFile(suffix=".pptx") as temp_file:
            loader = TemplateLoader(temp_file.name)
            
            # Test that content layout is correctly identified
            content_index = loader.get_layout_index("content")
            assert content_index == 1, f"Content layout should be at index 1, got {content_index}"
            
            # Test that title layout is correctly identified
            title_index = loader.get_layout_index("title")
            assert title_index == 0, f"Title layout should be at index 0, got {title_index}"
    
    @patch('open_lilli.template_loader.Presentation')
    def test_placeholder_matching_validation(self, mock_presentation_class):
        """Test that layouts match expected placeholder patterns."""
        mock_prs = Mock()
        mock_prs.slides = Mock()
        mock_prs.slides._sld_lst = []
        mock_prs.part = Mock()
        
        # Create content layout with TITLE + BODY placeholders
        content_layout = Mock()
        content_layout.placeholders = []
        
        title_ph = Mock()
        title_ph.placeholder_format.type = 1  # TITLE
        title_ph.left = 100
        title_ph.top = 100
        title_ph.width = 800
        title_ph.height = 100
        
        body_ph = Mock()
        body_ph.placeholder_format.type = 2  # BODY
        body_ph.left = 100
        body_ph.top = 250
        body_ph.width = 800
        body_ph.height = 500
        
        content_layout.placeholders = [title_ph, body_ph]
        mock_prs.slide_layouts = [content_layout]
        mock_presentation_class.return_value = mock_prs
        
        with tempfile.NamedTemporaryFile(suffix=".pptx") as temp_file:
            loader = TemplateLoader(temp_file.name)
            
            # Test placeholder validation
            assert loader.validate_placeholder_match("content", {"TITLE", "BODY"}), \
                "Content layout should match title + body placeholders"
            
            assert not loader.validate_placeholder_match("content", {"TITLE", "PICTURE"}), \
                "Content layout should not match title + picture placeholders"
    
    @patch('open_lilli.template_loader.Presentation')
    def test_comprehensive_layout_classification(self, mock_presentation_class):
        """Test classification of various layout types."""
        mock_prs = Mock()
        mock_prs.slides = Mock()
        mock_prs.slides._sld_lst = []
        mock_prs.part = Mock()
        
        # Helper function to create placeholder mock
        def create_placeholder(ph_type):
            ph = Mock()
            ph.placeholder_format.type = ph_type
            ph.left = 100
            ph.top = 100
            ph.width = 400
            ph.height = 300
            return ph
        
        # Create various layout types
        layouts = []
        
        # 1. Title layout (TITLE + SUBTITLE)
        title_layout = Mock()
        title_layout.placeholders = [
            create_placeholder(1),   # TITLE
            create_placeholder(3)    # SUBTITLE
        ]
        layouts.append(title_layout)
        
        # 2. Content layout (TITLE + BODY)
        content_layout = Mock()
        content_layout.placeholders = [
            create_placeholder(1),   # TITLE
            create_placeholder(2)    # BODY
        ]
        layouts.append(content_layout)
        
        # 3. Two column layout (TITLE + BODY + BODY)
        two_column_layout = Mock()
        two_column_layout.placeholders = [
            create_placeholder(1),   # TITLE
            create_placeholder(2),   # BODY
            create_placeholder(2)    # BODY
        ]
        layouts.append(two_column_layout)
        
        # 4. Image content layout (TITLE + PICTURE + BODY)
        image_content_layout = Mock()
        image_content_layout.placeholders = [
            create_placeholder(1),   # TITLE
            create_placeholder(18),  # PICTURE
            create_placeholder(2)    # BODY
        ]
        layouts.append(image_content_layout)
        
        # 5. Chart layout (TITLE + CHART)
        chart_layout = Mock()
        chart_layout.placeholders = [
            create_placeholder(1),   # TITLE
            create_placeholder(14)   # CHART
        ]
        layouts.append(chart_layout)
        
        # 6. Blank layout (no placeholders)
        blank_layout = Mock()
        blank_layout.placeholders = []
        layouts.append(blank_layout)
        
        mock_prs.slide_layouts = layouts
        mock_presentation_class.return_value = mock_prs
        
        with tempfile.NamedTemporaryFile(suffix=".pptx") as temp_file:
            loader = TemplateLoader(temp_file.name)
            
            # Test all expected layout types are classified correctly
            expected_layouts = {
                "title": 0,
                "content": 1,
                "two_column": 2,
                "image_content": 3,
                "chart": 4,
                "blank": 5
            }
            
            for layout_name, expected_index in expected_layouts.items():
                actual_index = loader.get_layout_index(layout_name)
                assert actual_index == expected_index, \
                    f"Layout '{layout_name}' should be at index {expected_index}, got {actual_index}"
    
    @patch('open_lilli.template_loader.Presentation')
    def test_essential_layout_fallbacks(self, mock_presentation_class):
        """Test that essential layouts have proper fallback mechanisms."""
        mock_prs = Mock()
        mock_prs.slides = Mock()
        mock_prs.slides._sld_lst = []
        mock_prs.part = Mock()
        
        # Create only a basic content layout
        content_layout = Mock()
        title_ph = Mock()
        title_ph.placeholder_format.type = 1  # TITLE
        body_ph = Mock()
        body_ph.placeholder_format.type = 2  # BODY
        content_layout.placeholders = [title_ph, body_ph]
        
        mock_prs.slide_layouts = [content_layout]
        mock_presentation_class.return_value = mock_prs
        
        with tempfile.NamedTemporaryFile(suffix=".pptx") as temp_file:
            loader = TemplateLoader(temp_file.name)
            
            # Essential layouts should all fall back to the content layout
            assert loader.get_layout_index("content") == 0
            assert loader.get_layout_index("title") == 0, "Title should fallback to content"
            assert loader.get_layout_index("section") == 0, "Section should fallback to content" 
            assert loader.get_layout_index("blank") == 0, "Blank should fallback to content"
    
    @patch('open_lilli.template_loader.Presentation')
    def test_get_layout_info_details(self, mock_presentation_class):
        """Test detailed layout information retrieval."""
        mock_prs = Mock()
        mock_prs.slides = Mock()
        mock_prs.slides._sld_lst = []
        mock_prs.part = Mock()
        
        # Create content layout with detailed placeholder info
        content_layout = Mock()
        
        title_ph = Mock()
        title_ph.placeholder_format.type = 1  # TITLE
        title_ph.left = 1000
        title_ph.top = 2000
        title_ph.width = 8000
        title_ph.height = 1500
        
        body_ph = Mock()
        body_ph.placeholder_format.type = 2  # BODY
        body_ph.left = 1000
        body_ph.top = 4000
        body_ph.width = 8000
        body_ph.height = 5000
        
        content_layout.placeholders = [title_ph, body_ph]
        mock_prs.slide_layouts = [content_layout]
        mock_presentation_class.return_value = mock_prs
        
        with tempfile.NamedTemporaryFile(suffix=".pptx") as temp_file:
            loader = TemplateLoader(temp_file.name)
            
            # Get detailed layout info
            info = loader.get_layout_info("content")
            
            assert info is not None, "Layout info should be available"
            assert info["semantic_name"] == "content"
            assert info["layout_index"] == 0
            assert info["total_placeholders"] == 2
            assert len(info["placeholders"]) == 2
            
            # Check placeholder details
            title_placeholder = info["placeholders"][0]
            assert title_placeholder["type"] == 1
            assert title_placeholder["type_name"] == "TITLE"
            assert title_placeholder["position"]["left"] == 1000
            
            body_placeholder = info["placeholders"][1]
            assert body_placeholder["type"] == 2
            assert body_placeholder["type_name"] == "BODY"
            assert body_placeholder["position"]["top"] == 4000
    
    @patch('open_lilli.template_loader.Presentation')
    def test_available_layouts_list(self, mock_presentation_class):
        """Test getting list of available semantic layout names."""
        mock_prs = Mock()
        mock_prs.slides = Mock()
        mock_prs.slides._sld_lst = []
        mock_prs.part = Mock()
        
        # Create two different layouts
        title_layout = Mock()
        title_ph = Mock()
        title_ph.placeholder_format.type = 1  # TITLE
        subtitle_ph = Mock()
        subtitle_ph.placeholder_format.type = 3  # SUBTITLE
        title_layout.placeholders = [title_ph, subtitle_ph]
        
        content_layout = Mock()
        title_ph2 = Mock()
        title_ph2.placeholder_format.type = 1  # TITLE
        body_ph = Mock()
        body_ph.placeholder_format.type = 2  # BODY
        content_layout.placeholders = [title_ph2, body_ph]
        
        mock_prs.slide_layouts = [title_layout, content_layout]
        mock_presentation_class.return_value = mock_prs
        
        with tempfile.NamedTemporaryFile(suffix=".pptx") as temp_file:
            loader = TemplateLoader(temp_file.name)
            
            available = loader.get_available_layouts()
            
            assert "title" in available, "Title layout should be available"
            assert "content" in available, "Content layout should be available"
            assert "section" in available, "Section should be available (fallback to title)"
            assert "blank" in available, "Blank should be available (fallback to content)"
    
    @patch('open_lilli.template_loader.Presentation')
    def test_get_layout_map_copy(self, mock_presentation_class):
        """Test that get_layout_map returns a copy of the internal mapping."""
        mock_prs = Mock()
        mock_prs.slides = Mock()
        mock_prs.slides._sld_lst = []
        mock_prs.part = Mock()
        
        content_layout = Mock()
        title_ph = Mock()
        title_ph.placeholder_format.type = 1  # TITLE
        body_ph = Mock()
        body_ph.placeholder_format.type = 2  # BODY
        content_layout.placeholders = [title_ph, body_ph]
        
        mock_prs.slide_layouts = [content_layout]
        mock_presentation_class.return_value = mock_prs
        
        with tempfile.NamedTemporaryFile(suffix=".pptx") as temp_file:
            loader = TemplateLoader(temp_file.name)
            
            layout_map = loader.get_layout_map()
            
            # Modify the returned copy
            layout_map["test"] = 999
            
            # Verify internal map is unchanged
            assert "test" not in loader._layout_map, "Internal layout map should not be modified"
            assert loader.get_layout_index("test") is None, "Test layout should not exist in internal map"
    
    @patch('open_lilli.template_loader.Presentation')
    def test_nonexistent_layout_returns_none(self, mock_presentation_class):
        """Test that requesting non-existent layout returns None."""
        mock_prs = Mock()
        mock_prs.slides = Mock()
        mock_prs.slides._sld_lst = []
        mock_prs.part = Mock()
        mock_prs.slide_layouts = []
        
        mock_presentation_class.return_value = mock_prs
        
        with tempfile.NamedTemporaryFile(suffix=".pptx") as temp_file:
            loader = TemplateLoader(temp_file.name)
            
            assert loader.get_layout_index("nonexistent") is None
            assert loader.get_layout_info("nonexistent") is None
            assert not loader.validate_placeholder_match("nonexistent", {"TITLE"})