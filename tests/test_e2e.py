"""End-to-end integration tests."""

import tempfile
from pathlib import Path
from unittest.mock import Mock, patch

import pytest
from pptx import Presentation

from open_lilli.content_processor import ContentProcessor
from open_lilli.models import GenerationConfig, Outline, SlidePlan
from open_lilli.slide_assembler import SlideAssembler
from open_lilli.template_parser import TemplateParser


@pytest.mark.e2e
class TestEndToEnd:
    """End-to-end integration tests for the complete pipeline."""

    def setup_method(self):
        """Set up test fixtures."""
        self.temp_dir = Path(tempfile.mkdtemp())
        
        # Create test content
        self.test_content = """
        # Business Review Presentation
        
        ## Executive Summary
        Our company has shown strong performance this quarter with significant growth in key metrics.
        Revenue increased by 25% compared to last quarter, driven by expansion into new markets.
        
        ## Financial Performance
        - Revenue: $2.5M (up 25%)
        - Profit margin: 15% (up from 12%)
        - Customer acquisition cost decreased by 10%
        
        ## Market Analysis
        The market conditions remain favorable for continued growth.
        Competition has increased but we maintain our competitive advantage.
        
        ## Future Outlook
        We expect continued growth in the next quarter.
        Key focus areas include product development and market expansion.
        """
        
        self.content_file = self.temp_dir / "test_content.txt"
        self.content_file.write_text(self.test_content)
        
        # Create a basic template
        self.template_file = self.temp_dir / "test_template.pptx"
        prs = Presentation()
        prs.save(str(self.template_file))

    def teardown_method(self):
        """Clean up test fixtures."""
        import shutil
        shutil.rmtree(self.temp_dir, ignore_errors=True)

    def test_content_processing_pipeline(self):
        """Test the content processing part of the pipeline."""
        processor = ContentProcessor()
        
        # Test file extraction
        content = processor.extract_text(self.content_file)
        assert len(content) > 0
        assert "Business Review Presentation" in content
        
        # Test section extraction
        sections = processor.extract_sections(content)
        assert len(sections) > 0
        assert "Executive Summary" in str(sections)
        
        # Test statistics
        word_count = processor.get_word_count(content)
        assert word_count > 50
        
        reading_time = processor.get_reading_time(content)
        assert reading_time >= 1

    def test_template_analysis_pipeline(self):
        """Test template analysis pipeline."""
        parser = TemplateParser(str(self.template_file))
        
        # Test basic parsing
        assert parser.template_path.exists()
        assert len(parser.layout_map) > 0
        
        # Test layout information
        layouts = parser.list_available_layouts()
        assert len(layouts) > 0
        
        # Test getting layouts
        for layout_type in layouts:
            layout = parser.get_layout(layout_type)
            assert layout is not None
        
        # Test template info
        info = parser.get_template_info()
        assert info["total_layouts"] > 0
        assert "slide_dimensions" in info

    def test_slide_assembly_pipeline(self):
        """Test slide assembly with real components."""
        # Create test slides
        slides = [
            SlidePlan(
                index=0,
                slide_type="title",
                title="Test Presentation",
                bullets=[],
                layout_id=0,
                speaker_notes="Welcome"
            ),
            SlidePlan(
                index=1,
                slide_type="content",
                title="Main Points",
                bullets=["Point 1", "Point 2", "Point 3"],
                layout_id=1,
                speaker_notes="Discuss main points"
            )
        ]
        
        # Create outline
        outline = Outline(title="Test Presentation", slides=slides)
        
        # Test assembly
        parser = TemplateParser(str(self.template_file))
        assembler = SlideAssembler(parser)
        
        # Validate slides
        issues = assembler.validate_slides_before_assembly(slides)
        assert isinstance(issues, list)
        
        # Assemble presentation
        output_file = self.temp_dir / "output.pptx"
        result_path = assembler.assemble(outline, slides, {}, output_file)
        
        assert result_path.exists()
        assert result_path.suffix == ".pptx"
        
        # Verify the output is a valid PowerPoint file
        try:
            test_prs = Presentation(str(result_path))
            assert len(test_prs.slides) == 2
        except Exception as e:
            pytest.fail(f"Generated PowerPoint file is invalid: {e}")

    def test_configuration_integration(self):
        """Test that configuration objects work with all components."""
        config = GenerationConfig(
            max_slides=10,
            max_bullets_per_slide=4,
            tone="professional",
            complexity_level="intermediate",
            include_images=True,
            include_charts=True
        )
        
        # Test config validation
        assert config.max_slides == 10
        assert config.tone == "professional"
        assert config.include_images is True
        
        # Test config with different values
        casual_config = GenerationConfig(
            tone="casual",
            complexity_level="basic",
            include_images=False
        )
        
        assert casual_config.tone == "casual"
        assert casual_config.include_images is False

    @patch("open_lilli.outline_generator.OpenAI")
    def test_outline_generation_integration(self, mock_openai):
        """Test outline generation with mocked OpenAI."""
        from open_lilli.outline_generator import OutlineGenerator
        
        # Mock OpenAI response
        mock_client = Mock()
        mock_response = Mock()
        mock_response.choices = [Mock()]
        mock_response.choices[0].message.content = '''
        {
            "language": "en",
            "title": "Business Review",
            "slides": [
                {
                    "index": 0,
                    "slide_type": "title",
                    "title": "Business Review",
                    "bullets": [],
                    "image_query": null,
                    "chart_data": null
                },
                {
                    "index": 1,
                    "slide_type": "content",
                    "title": "Executive Summary",
                    "bullets": ["Strong performance", "25% growth"],
                    "image_query": "business growth",
                    "chart_data": null
                }
            ]
        }
        '''
        mock_client.chat.completions.create.return_value = mock_response
        mock_openai.return_value = mock_client
        
        # Test outline generation
        processor = ContentProcessor()
        content = processor.extract_text(self.content_file)
        
        generator = OutlineGenerator(mock_client)
        outline = generator.generate_outline(content)
        
        assert isinstance(outline, Outline)
        assert outline.title == "Business Review"
        assert len(outline.slides) == 2
        assert outline.slides[0].slide_type == "title"

    def test_error_handling_integration(self):
        """Test error handling across components."""
        # Test with invalid template path
        with pytest.raises(FileNotFoundError):
            TemplateParser("nonexistent.pptx")
        
        # Test with invalid content path
        processor = ContentProcessor()
        with pytest.raises(FileNotFoundError):
            processor.extract_text("nonexistent.txt")
        
        # Test assembly with invalid slides
        parser = TemplateParser(str(self.template_file))
        assembler = SlideAssembler(parser)
        
        invalid_slides = [
            SlidePlan(
                index=0,
                slide_type="content",
                title="",  # Invalid: empty title
                bullets=[],
                layout_id=99  # Invalid: layout doesn't exist
            )
        ]
        
        issues = assembler.validate_slides_before_assembly(invalid_slides)
        assert len(issues) > 0
        assert any("no title" in issue.lower() for issue in issues)

    def test_file_format_support(self):
        """Test support for different file formats."""
        processor = ContentProcessor()
        
        # Test .txt file
        txt_file = self.temp_dir / "test.txt"
        txt_file.write_text("Test content")
        content = processor.extract_text(txt_file)
        assert content == "Test content"
        
        # Test .md file
        md_file = self.temp_dir / "test.md"
        md_file.write_text("# Header\nContent")
        content = processor.extract_text(md_file)
        assert "Header" in content
        assert "Content" in content
        
        # Test unsupported file
        pdf_file = self.temp_dir / "test.pdf"
        pdf_file.write_bytes(b"fake pdf")
        
        with pytest.raises(ValueError, match="Unsupported file type"):
            processor.extract_text(pdf_file)

    def test_memory_and_performance(self):
        """Test memory usage and performance with larger content."""
        # Create larger content
        large_content = "Large content section. " * 1000
        large_file = self.temp_dir / "large_content.txt"
        large_file.write_text(large_content)
        
        processor = ContentProcessor()
        content = processor.extract_text(large_file)
        
        # Should handle large content without issues
        assert len(content) > 10000
        word_count = processor.get_word_count(content)
        assert word_count > 2000

    def test_internationalization_support(self):
        """Test support for different languages."""
        # Test content with non-ASCII characters
        multilingual_content = """
        # Présentation d'Affaires
        
        ## Résumé Exécutif
        Notre entreprise a montré de fortes performances ce trimestre.
        
        ## Análisis de Mercado
        Las condiciones del mercado siguen siendo favorables.
        
        ## 市场分析
        市场状况继续有利于持续增长。
        """
        
        ml_file = self.temp_dir / "multilingual.txt"
        ml_file.write_text(multilingual_content, encoding='utf-8')
        
        processor = ContentProcessor()
        content = processor.extract_text(ml_file)
        
        assert "Présentation" in content
        assert "Análisis" in content
        assert "市场分析" in content

    def test_data_flow_integrity(self):
        """Test that data flows correctly through the pipeline without corruption."""
        # Create test data with specific markers
        test_data = {
            "title": "Data Integrity Test",
            "content": "MARKER_START: This is test content with unique identifiers. MARKER_END",
            "bullets": ["BULLET_1: First point", "BULLET_2: Second point"]
        }
        
        content_text = f"# {test_data['title']}\n{test_data['content']}\n- {test_data['bullets'][0]}\n- {test_data['bullets'][1]}"
        test_file = self.temp_dir / "integrity_test.txt"
        test_file.write_text(content_text)
        
        # Process through content processor
        processor = ContentProcessor()
        processed_content = processor.extract_text(test_file)
        
        # Verify markers are preserved
        assert "MARKER_START" in processed_content
        assert "MARKER_END" in processed_content
        assert "BULLET_1" in processed_content
        assert "BULLET_2" in processed_content
        
        # Test section extraction preserves structure
        sections = processor.extract_sections(processed_content)
        assert test_data['title'] in str(sections)

    def test_output_validation(self):
        """Test that outputs meet quality standards."""
        # Create a complete test slide
        slide = SlidePlan(
            index=0,
            slide_type="content",
            title="Quality Test Slide",
            bullets=["Quality point 1", "Quality point 2"],
            speaker_notes="These are quality speaker notes",
            layout_id=1
        )
        
        # Test slide validation
        parser = TemplateParser(str(self.template_file))
        assembler = SlideAssembler(parser)
        
        issues = assembler.validate_slides_before_assembly([slide])
        assert len(issues) == 0  # Should be valid
        
        # Test assembly produces valid output
        outline = Outline(title="Quality Test", slides=[slide])
        output_file = self.temp_dir / "quality_test.pptx"
        
        result_path = assembler.assemble(outline, [slide], {}, output_file)
        assert result_path.exists()
        
        # Verify file can be opened
        prs = Presentation(str(result_path))
        assert len(prs.slides) == 1
        
        # Check slide has expected content
        slide_obj = prs.slides[0]
        assert hasattr(slide_obj.shapes, 'title')

    def test_resource_cleanup(self):
        """Test that resources are properly cleaned up."""
        # This test ensures we don't leave temporary files or open handles
        parser = TemplateParser(str(self.template_file))
        
        # Create multiple parser instances to test resource handling
        for i in range(5):
            temp_parser = TemplateParser(str(self.template_file))
            info = temp_parser.get_template_info()
            assert info is not None
        
        # Original parser should still work
        final_info = parser.get_template_info()
        assert final_info is not None