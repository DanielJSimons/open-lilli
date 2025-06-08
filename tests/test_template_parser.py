"""Tests for template parser."""

import tempfile
from pathlib import Path
from unittest.mock import Mock, patch, MagicMock

import pytest
from pptx import Presentation

from open_lilli.template_parser import TemplateParser
from open_lilli.models import TemplateStyle, FontInfo, BulletInfo, PlaceholderStyleInfo, DesignPattern


class TestTemplateParser:
    """Tests for TemplateParser class."""

    def create_test_template(self) -> str:
        """Create a minimal test template."""
        prs = Presentation()
        
        # Create a temporary file
        temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        temp_path = temp_file.name
        temp_file.close()
        
        # Save the presentation
        prs.save(temp_path)
        return temp_path

    def test_init_valid_template(self):
        """Test initialization with valid template."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            assert parser.template_path == Path(template_path)
            assert parser.prs is not None
            assert isinstance(parser.layout_map, dict)
            assert isinstance(parser.palette, dict)
            
        finally:
            Path(template_path).unlink()


    # --- Tests for analyze_design_pattern ---

    @pytest.fixture
    def design_mock_parser(self) -> TemplateParser:
        """
        Creates a TemplateParser instance with mocks tailored for design pattern analysis.
        This fixture avoids using parametrize with the more general mock_parser if test setup is complex.
        """
        template_path = self.create_test_template()
        parser = TemplateParser(template_path) # Real parser for basic structure

        # Mock template_style and its methods
        parser.template_style = MagicMock(spec=TemplateStyle)

        # Default mocks for fonts - these can be overridden in each test
        default_title_font = FontInfo(name="Arial", size=32.0, weight="bold", color="#000000")
        default_body_font = FontInfo(name="Arial", size=18.0, weight="normal", color="#000000")

        def mock_get_font(ph_type):
            if ph_type == 1 or ph_type == 13: # TITLE or CENTERED_TITLE
                return getattr(parser.template_style, '_mock_title_font', default_title_font)
            if ph_type == 2 or ph_type == 7: # BODY or OBJECT
                return getattr(parser.template_style, '_mock_body_font', default_body_font)
            return None
        parser.template_style.get_font_for_placeholder_type = MagicMock(side_effect=mock_get_font)

        # Default mocks for palette and layout_map - override in tests as needed
        parser.palette = {"dk1": "#111111", "lt1": "#EEEEEE", "acc1": "#FF0000"}
        parser.layout_map = {"title":0, "content":1, "section":2, "image_content":3, "blank":4}
        parser.reverse_layout_map = {v: k for k, v in parser.layout_map.items()}

        yield parser

        Path(template_path).unlink(missing_ok=True)


    def test_design_pattern_minimalist(self, design_mock_parser: TemplateParser):
        # Configure for minimalist
        design_mock_parser.template_style._mock_title_font = FontInfo(name="Segoe UI", size=24.0)
        design_mock_parser.template_style._mock_body_font = FontInfo(name="Segoe UI", size=18.0) # Low ratio: 24/18 = 1.33
        design_mock_parser.palette = {"dk1": "#202020", "lt1": "#F0F0F0"} # Only base colors, no accents
        design_mock_parser.layout_map = {"title": 0, "content": 1, "section": 2, "blank": 3} # Simple layouts

        pattern = design_mock_parser.analyze_design_pattern()

        assert pattern is not None
        assert pattern.name == "minimalist"
        assert pattern.primary_intent == "readability"
        assert pattern.font_scale_ratio == pytest.approx(1.33)
        assert pattern.color_complexity_score == pytest.approx(0.2) # Low
        assert pattern.layout_density_preference == "low"

    def test_design_pattern_vibrant(self, design_mock_parser: TemplateParser):
        # Configure for vibrant
        design_mock_parser.template_style._mock_title_font = FontInfo(name="Impact", size=48.0)
        design_mock_parser.template_style._mock_body_font = FontInfo(name="Arial", size=20.0) # High ratio: 48/20 = 2.4
        design_mock_parser.palette = {
            "dk1": "#000000", "lt1": "#FFFFFF",
            "acc1": "#FF00FF", "acc2": "#00FFFF", "acc3": "#FFFF00", "acc4": "#FF6600", "acc5": "#00FF66" # 5 distinct accents
        }
        design_mock_parser.layout_map = {"title":0, "content":1, "image_content":2, "two_column":3, "section":4} # Mixed, some dense

        pattern = design_mock_parser.analyze_design_pattern()

        assert pattern is not None
        assert pattern.name == "vibrant" # Or "bold & colorful" depending on exact thresholds
        assert pattern.primary_intent == "visual_impact"
        assert pattern.font_scale_ratio == pytest.approx(2.4)
        assert pattern.color_complexity_score == pytest.approx(0.9) # High
        assert pattern.layout_density_preference == "high" # 2 dense out of 5 is 40%

    def test_design_pattern_data_heavy(self, design_mock_parser: TemplateParser):
        # Configure for data-heavy
        design_mock_parser.template_style._mock_title_font = FontInfo(name="Calibri", size=28.0)
        design_mock_parser.template_style._mock_body_font = FontInfo(name="Calibri", size=16.0) # Moderate ratio: 28/16 = 1.75
        design_mock_parser.palette = {"dk1": "#333333", "lt1": "#FFFFFF", "acc1": "#0044CC", "acc2": "#0066DD"} # Few, professional accents
        design_mock_parser.layout_map = {
            "title": 0, "content": 1, "two_column": 2, "section": 3,
            "image_content":4, "layout_5_dense": 5 # Assuming layout_5_dense is also complex
        }
        # Manually classify "layout_5_dense" as dense for the test by adding it to dense_layout_keywords temporarily
        # This is a bit of a hack for testing; ideally the classification in TemplateParser would handle this.
        original_dense_keywords = design_mock_parser.dense_layout_keywords if hasattr(design_mock_parser, 'dense_layout_keywords') else ["two_column", "image_content"]
        design_mock_parser.dense_layout_keywords = original_dense_keywords + ["layout_5_dense"]


        pattern = design_mock_parser.analyze_design_pattern()

        # Restore original keywords if they existed
        if hasattr(design_mock_parser, 'dense_layout_keywords'):
             design_mock_parser.dense_layout_keywords = original_dense_keywords


        assert pattern is not None
        assert pattern.name == "data-heavy"
        assert pattern.primary_intent == "information_density"
        assert pattern.font_scale_ratio == pytest.approx(1.75)
        assert pattern.color_complexity_score == pytest.approx(0.6) # Medium
        assert pattern.layout_density_preference == "high" # 3 dense / 6 total layouts = 50%

    def test_design_pattern_standard_default(self, design_mock_parser: TemplateParser):
        # Configure for standard/default (values that don't strongly fit other patterns)
        design_mock_parser.template_style._mock_title_font = FontInfo(name="Arial", size=36.0)
        design_mock_parser.template_style._mock_body_font = FontInfo(name="Arial", size=20.0) # Ratio: 1.8
        design_mock_parser.palette = {"dk1": "#000000", "lt1": "#FFFFFF", "acc1": "#007ACC", "acc2": "#CCCCCC"} # 1-2 accents
        design_mock_parser.layout_map = {"title":0, "content":1, "section":2, "image":3, "blank":4} # Mostly simple

        pattern = design_mock_parser.analyze_design_pattern()

        assert pattern is not None
        assert pattern.name == "standard"
        assert pattern.primary_intent == "balanced" # Or "clarity" if color complexity is low enough
        assert pattern.font_scale_ratio == pytest.approx(1.8)
        assert pattern.color_complexity_score == pytest.approx(0.6) # Medium (acc1 is distinct, acc2 might be too close to lt1 or dk1 to count depending on implementation)
        assert pattern.layout_density_preference == "low" # No explicitly "dense" layouts listed

    def test_design_pattern_missing_font_info(self, design_mock_parser: TemplateParser):
        # Configure for missing font info
        design_mock_parser.template_style._mock_title_font = None # Title font missing
        design_mock_parser.template_style._mock_body_font = FontInfo(name="Times New Roman", size=16.0)
        design_mock_parser.palette = {"dk1": "#111", "lt1": "#EEE", "acc1": "#ABC"}
        design_mock_parser.layout_map = {"title": 0, "content": 1}

        pattern = design_mock_parser.analyze_design_pattern()

        assert pattern is not None
        assert pattern.name == "standard" # Should default gracefully
        assert pattern.font_scale_ratio == pytest.approx(1.8) # Default ratio when one font is missing
        assert pattern.color_complexity_score == pytest.approx(0.6) # Medium (1 distinct accent)
        assert pattern.layout_density_preference == "low"

    def test_design_pattern_both_fonts_missing(self, design_mock_parser: TemplateParser):
        design_mock_parser.template_style._mock_title_font = None
        design_mock_parser.template_style._mock_body_font = None

        pattern = design_mock_parser.analyze_design_pattern()
        assert pattern is not None
        assert pattern.name == "standard"
        assert pattern.font_scale_ratio == pytest.approx(1.8) # Default
        assert pattern.color_complexity_score == pytest.approx(0.6) # Based on default palette in fixture
        assert pattern.layout_density_preference == "low" # Based on default layout_map in fixture if not overridden

    def test_design_pattern_edge_case_no_accent_colors(self, design_mock_parser: TemplateParser):
        design_mock_parser.palette = {"dk1": "#000000", "lt1": "#FFFFFF"} # No accents
        pattern = design_mock_parser.analyze_design_pattern()
        assert pattern.color_complexity_score == pytest.approx(0.2) # Low

    def test_design_pattern_all_accent_colors_same_as_base(self, design_mock_parser: TemplateParser):
        design_mock_parser.palette = {
            "dk1": "#000000", "lt1": "#FFFFFF",
            "acc1": "#000000", "acc2": "#FFFFFF", "acc3": "#000001", "acc4": "#FFFFFE" # Last two are very similar
        }
        pattern = design_mock_parser.analyze_design_pattern()
        # acc3 and acc4 might be counted if the similarity check is loose, or not if strict.
        # Based on current logic (1 char diff allowed), they might not be distinct enough from base.
        # If #000001 is not distinct from #000000, and #FFFFFE not from #FFFFFF, then score is 0.2
        assert pattern.color_complexity_score == pytest.approx(0.2)

    def test_design_pattern_no_layouts(self, design_mock_parser: TemplateParser):
        design_mock_parser.layout_map = {}
        design_mock_parser.reverse_layout_map = {}
        pattern = design_mock_parser.analyze_design_pattern()
        assert pattern.layout_density_preference == "medium" # Default


    # --- Tests for TemplateCompatibilityReport ---

    @pytest.fixture
    def mock_parser(self, request) -> TemplateParser:
        """Creates a mock TemplateParser for compatibility tests."""
        # Create a real TemplateParser instance with a minimal template
        # to ensure all basic structures are initialized.
        template_path = self.create_test_template()
        parser = TemplateParser(template_path)

        # Default mocks - can be overridden by individual tests using request.param
        default_layout_map = {"title": 0, "content": 1, "section": 2, "image": 3, "blank": 4}
        default_palette = {
            "dk1": "#000000", "lt1": "#FFFFFF", "acc1": "#FF0000",
            "acc2": "#00FF00", "acc3": "#0000FF"
        }

        # Apply parameters if provided by the test
        layout_map_override = getattr(request, "param", {}).get("layout_map")
        palette_override = getattr(request, "param", {}).get("palette")

        parser.layout_map = layout_map_override if layout_map_override is not None else default_layout_map
        parser.palette = palette_override if palette_override is not None else default_palette

        # Ensure reverse map is also updated
        parser.reverse_layout_map = {v: k for k, v in parser.layout_map.items()}

        yield parser # Provide the parser to the test

        # Teardown: remove the temp template file
        Path(template_path).unlink(missing_ok=True)


    def test_compatibility_good_template(self, mock_parser: TemplateParser):
        """Test with a well-formed template."""
        report = mock_parser.check_template_compatibility()
        assert report.passed_all_checks is True
        assert not report.issues
        assert not report.missing_placeholders
        assert not report.color_scheme_warnings
        assert not report.contrast_issues
        assert "Template appears to meet basic compatibility checks." in report.suggestions


    @pytest.mark.parametrize("mock_parser", [{"layout_map": {"title": 0, "section": 1}}], indirect=True)
    def test_compatibility_missing_placeholders(self, mock_parser: TemplateParser):
        """Test with missing essential placeholders."""
        report = mock_parser.check_template_compatibility()
        assert report.passed_all_checks is False
        assert "content" in report.missing_placeholders
        assert "image" in report.missing_placeholders
        assert "Missing essential placeholder types: content, image." in report.issues[0]
        assert "Consider adding a 'content' layout for better versatility." in report.suggestions
        assert "Consider adding a 'image' layout for better versatility." in report.suggestions


    @pytest.mark.parametrize("mock_parser", [{"palette": {
        "dk1": "#808080", # Grey
        "lt1": "#A0A0A0", # Light Grey (low contrast with dk1)
        "acc1": "#B0B0B0", # Another Grey (low contrast with lt1)
    }}], indirect=True)
    def test_compatibility_low_contrast(self, mock_parser: TemplateParser):
        """Test with low contrast colors."""
        report = mock_parser.check_template_compatibility()
        assert report.passed_all_checks is False
        assert len(report.contrast_issues) > 0
        assert "Low contrast between 'dk1' (#808080) and 'lt1' (#A0A0A0)" in report.contrast_issues[0]
        # acc1 vs lt1 might also be low, depending on exact calculation and other palette colors
        assert "One or more color pairs have insufficient contrast" in report.issues[0]
        assert "Review theme colors to ensure text is clearly readable" in report.suggestions[0]

    @pytest.mark.parametrize("mock_parser", [{"palette": {
        "dk1": "#FFFFFF", # White
        "lt1": "#FFFFFF", # White (identical to dk1)
        "acc1": "#0000FF",
    }}], indirect=True)
    def test_compatibility_identical_dk1_lt1(self, mock_parser: TemplateParser):
        """Test with dk1 and lt1 being the same color."""
        report = mock_parser.check_template_compatibility()
        assert report.passed_all_checks is False
        assert "'dk1' and 'lt1' colors are identical, which will cause contrast issues." in report.color_scheme_warnings
        assert "Issues found with the template's color scheme definitions." in report.issues
        # Contrast check for dk1 vs lt1 should also fail spectacularly
        assert any("Low contrast between 'dk1' (#FFFFFF) and 'lt1' (#FFFFFF)" in issue for issue in report.contrast_issues)


    @pytest.mark.parametrize("mock_parser", [{"palette": {
        # dk1 is missing
        "lt1": "#FFFFFF",
        "acc1": "#0000FF",
    }}], indirect=True)
    def test_compatibility_missing_dk1(self, mock_parser: TemplateParser):
        """Test with missing dk1 color."""
        report = mock_parser.check_template_compatibility()
        assert report.passed_all_checks is False
        assert "Theme color 'dk1' (primary dark) is not defined." in report.color_scheme_warnings
        assert "Issues found with the template's color scheme definitions." in report.issues
        # Contrast checks involving dk1 should warn about missing color
        assert "Color 'dk1' not found in palette for contrast check." in report.color_scheme_warnings


    @pytest.mark.parametrize("mock_parser", [{"palette": {
        "dk1": "#000000",
        # lt1 is missing
        "acc1": "#0000FF",
    }}], indirect=True)
    def test_compatibility_missing_lt1(self, mock_parser: TemplateParser):
        """Test with missing lt1 color."""
        report = mock_parser.check_template_compatibility()
        assert report.passed_all_checks is False
        assert "Theme color 'lt1' (primary light) is not defined." in report.color_scheme_warnings
        assert "Issues found with the template's color scheme definitions." in report.issues
        # Contrast checks involving lt1 should warn about missing color
        assert "Color 'lt1' not found in palette for contrast check." in report.color_scheme_warnings


    def test_hex_to_rgb_conversion(self, mock_parser: TemplateParser):
        """Test _hex_to_rgb utility."""
        assert mock_parser._hex_to_rgb("#FF0000") == (255, 0, 0)
        assert mock_parser._hex_to_rgb("00FF00") == (0, 255, 0)
        assert mock_parser._hex_to_rgb("#00F") == (0, 0, 255) # Shorthand
        with pytest.raises(ValueError, match="Invalid hex color format: #12345"):
            mock_parser._hex_to_rgb("#12345")
        with pytest.raises(ValueError, match="Invalid character in hex color: GG0000"):
            mock_parser._hex_to_rgb("GG0000")

    def test_relative_luminance_calculation(self, mock_parser: TemplateParser):
        """Test _relative_luminance utility."""
        # Test pure white
        assert mock_parser._relative_luminance((255, 255, 255)) == pytest.approx(1.0)
        # Test pure black
        assert mock_parser._relative_luminance((0, 0, 0)) == pytest.approx(0.0)
        # Test a grey color
        assert mock_parser._relative_luminance((128, 128, 128)) == pytest.approx(0.21586)

    def test_contrast_ratio_calculation(self, mock_parser: TemplateParser):
        """Test _calculate_contrast_ratio utility."""
        # Black text on white background (max contrast)
        assert mock_parser._calculate_contrast_ratio("#000000", "#FFFFFF") == pytest.approx(21.0)
        # White text on black background (max contrast)
        assert mock_parser._calculate_contrast_ratio("#FFFFFF", "#000000") == pytest.approx(21.0)
        # Grey on white (example from WCAG)
        assert mock_parser._calculate_contrast_ratio("#767676", "#FFFFFF") == pytest.approx(4.5, abs=0.1)
        # Red on green (poor contrast)
        assert mock_parser._calculate_contrast_ratio("#FF0000", "#00FF00") == pytest.approx(2.91, abs=0.01)
        # Identical colors
        assert mock_parser._calculate_contrast_ratio("#ABCDEF", "#ABCDEF") == pytest.approx(1.0)
        # Invalid color should return 1.0 and log warning
        with patch.object(mock_parser.logger, 'warning') as mock_log:
            assert mock_parser._calculate_contrast_ratio("invalid", "#FFFFFF") == 1.0
            mock_log.assert_called_once()
            assert "Invalid color format for contrast calculation" in mock_log.call_args[0][0]

    def test_init_file_not_found(self):
        """Test initialization with non-existent file."""
        with pytest.raises(FileNotFoundError):
            TemplateParser("nonexistent.pptx")

    def test_init_invalid_extension(self):
        """Test initialization with invalid file extension."""
        with pytest.raises(ValueError, match="Template must be a .pptx file"):
            TemplateParser("template.txt")

    def test_init_invalid_file(self):
        """Test initialization with invalid PowerPoint file."""
        # Create a text file with .pptx extension
        with tempfile.NamedTemporaryFile(suffix='.pptx', mode='w', delete=False) as f:
            f.write("This is not a PowerPoint file")
            temp_path = f.name
        
        try:
            with pytest.raises(ValueError, match="Failed to load template"):
                TemplateParser(temp_path)
        finally:
            Path(temp_path).unlink()

    def test_classify_layout(self):
        """Test layout classification logic."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            # Test with actual layouts from the template
            for i, layout in enumerate(parser.prs.slide_layouts):
                layout_type = parser._classify_layout(layout, i)
                assert isinstance(layout_type, str)
                assert len(layout_type) > 0
                
        finally:
            Path(template_path).unlink()

    def test_ensure_basic_layouts(self):
        """Test that basic layouts are ensured."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            # These should always be present after _ensure_basic_layouts
            essential_layouts = ["title", "content", "section", "blank"]
            
            for layout_type in essential_layouts:
                assert layout_type in parser.layout_map
                assert isinstance(parser.layout_map[layout_type], int)
                
        finally:
            Path(template_path).unlink()

    def test_get_layout_valid(self):
        """Test getting a valid layout."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            # Should be able to get any layout in the map
            for layout_type in parser.layout_map:
                layout = parser.get_layout(layout_type)
                assert layout is not None
                
        finally:
            Path(template_path).unlink()

    def test_get_layout_invalid(self):
        """Test getting an invalid layout."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            with pytest.raises(ValueError, match="Layout type 'nonexistent' not found"):
                parser.get_layout("nonexistent")
                
        finally:
            Path(template_path).unlink()

    def test_get_layout_index(self):
        """Test getting layout index."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            # Valid layout type
            for layout_type in parser.layout_map:
                index = parser.get_layout_index(layout_type)
                assert isinstance(index, int)
                assert index >= 0
            
            # Invalid layout type should return default
            index = parser.get_layout_index("nonexistent")
            assert isinstance(index, int)
            assert index >= 0
            
        finally:
            Path(template_path).unlink()

    def test_list_available_layouts(self):
        """Test listing available layouts."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            layouts = parser.list_available_layouts()
            assert isinstance(layouts, list)
            assert len(layouts) > 0
            
            # Should match the layout_map keys
            assert set(layouts) == set(parser.layout_map.keys())
            
        finally:
            Path(template_path).unlink()

    def test_analyze_layout_placeholders(self):
        """Test placeholder analysis."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            # Test with first available layout
            layout_type = list(parser.layout_map.keys())[0]
            analysis = parser.analyze_layout_placeholders(layout_type)
            
            assert isinstance(analysis, dict)
            assert "total_placeholders" in analysis
            assert "placeholder_details" in analysis
            assert "has_title" in analysis
            assert "has_content" in analysis
            assert "has_image" in analysis
            assert "has_chart" in analysis
            
            assert isinstance(analysis["total_placeholders"], int)
            assert isinstance(analysis["placeholder_details"], list)
            assert isinstance(analysis["has_title"], bool)
            
        finally:
            Path(template_path).unlink()

    def test_get_theme_color(self):
        """Test theme color retrieval."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            # Should return a color for any name in palette
            for color_name in parser.palette:
                color = parser.get_theme_color(color_name)
                assert isinstance(color, str)
                assert color.startswith("#")
                assert len(color) == 7  # Hex color format
            
            # Should return default for unknown color
            color = parser.get_theme_color("unknown_color")
            assert color == "#000000"
            
        finally:
            Path(template_path).unlink()

    def test_get_slide_size(self):
        """Test slide size retrieval."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            width, height = parser.get_slide_size()
            assert isinstance(width, int)
            assert isinstance(height, int)
            assert width > 0
            assert height > 0
            
        finally:
            Path(template_path).unlink()

    def test_get_template_info(self):
        """Test comprehensive template info."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            info = parser.get_template_info()
            
            assert isinstance(info, dict)
            assert "template_path" in info
            assert "total_layouts" in info
            assert "available_layout_types" in info
            assert "layout_mapping" in info
            assert "slide_dimensions" in info
            assert "theme_colors" in info
            assert "template_style" in info
            
            # Verify slide dimensions structure
            dimensions = info["slide_dimensions"]
            assert "width" in dimensions
            assert "height" in dimensions
            assert "width_inches" in dimensions
            assert "height_inches" in dimensions
            
            # Verify template style structure
            template_style = info["template_style"]
            assert "placeholder_styles_count" in template_style
            assert "theme_fonts" in template_style
            assert "has_master_font" in template_style
            assert "placeholder_types_with_styles" in template_style
            assert isinstance(template_style["placeholder_styles_count"], int)
            assert isinstance(template_style["has_master_font"], bool)
            assert isinstance(template_style["placeholder_types_with_styles"], list)
            
        finally:
            Path(template_path).unlink()

    def test_placeholder_type_name(self):
        """Test placeholder type name conversion."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            # Test known types
            assert parser._placeholder_type_name(1) == "TITLE"
            assert parser._placeholder_type_name(2) == "BODY"
            assert parser._placeholder_type_name(18) == "PICTURE"
            
            # Test unknown type
            assert parser._placeholder_type_name(999) == "TYPE_999"
            
        finally:
            Path(template_path).unlink()

    def test_get_layout_info(self):
        """Test layout info generation."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            # Test with first layout
            layout = parser.prs.slide_layouts[0]
            info = parser._get_layout_info(layout)
            
            assert isinstance(info, str)
            assert info.startswith("[")
            assert info.endswith("]")
            
        finally:
            Path(template_path).unlink()

    def test_get_theme_colors_method(self):
        """Test the get_theme_colors method exists and returns dict."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            # Method should exist and return a dictionary
            colors = parser.get_theme_colors()
            assert isinstance(colors, dict)
            
            # Should contain standard theme color keys when extraction works
            # (Will be empty for basic template without proper theme)
            for key in colors:
                assert isinstance(key, str)
                if colors[key]:  # If color is not None/empty
                    assert isinstance(colors[key], str)
                    
        finally:
            Path(template_path).unlink()

    def test_palette_contains_standard_colors(self):
        """Test that palette contains standard theme colors."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            # Should have standard theme color names
            expected_colors = ['dk1', 'lt1', 'acc1', 'acc2', 'acc3', 'acc4', 'acc5', 'acc6']
            
            for color_name in expected_colors:
                assert color_name in parser.palette
                color_value = parser.palette[color_name]
                assert isinstance(color_value, str)
                assert color_value.startswith("#")
                assert len(color_value) == 7  # Hex format #RRGGBB
                
        finally:
            Path(template_path).unlink()

    def test_get_theme_color_fallback(self):
        """Test theme color retrieval with fallback behavior."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            # Should return a valid color for standard names
            for color_name in ['dk1', 'lt1', 'acc1', 'acc2', 'acc3', 'acc4', 'acc5', 'acc6']:
                color = parser.get_theme_color(color_name)
                assert isinstance(color, str)
                assert color.startswith("#")
                assert len(color) == 7
            
            # Should return black for unknown color (existing test)
            color = parser.get_theme_color("unknown_color")
            assert color == "#000000"
            
        finally:
            Path(template_path).unlink()

    def test_template_style_extraction(self):
        """Test that template style is extracted during initialization."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            # Check that template_style attribute exists
            assert hasattr(parser, 'template_style')
            assert isinstance(parser.template_style, TemplateStyle)
            
            # Check that get_template_style method works
            style = parser.get_template_style()
            assert isinstance(style, TemplateStyle)
            assert style == parser.template_style
            
        finally:
            Path(template_path).unlink()

    def test_font_extraction_from_placeholder(self):
        """Test font extraction from placeholders."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            # Test getting font for different placeholder types
            title_font = parser.get_font_for_placeholder_type(1)  # TITLE
            if title_font:
                assert isinstance(title_font, FontInfo)
                assert isinstance(title_font.name, str)
                assert title_font.name != ""
                
            body_font = parser.get_font_for_placeholder_type(2)  # BODY
            if body_font:
                assert isinstance(body_font, FontInfo)
                assert isinstance(body_font.name, str)
                
        finally:
            Path(template_path).unlink()

    def test_bullet_style_extraction(self):
        """Test bullet style extraction."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            # Test getting bullet style for body placeholders
            bullet_style = parser.get_bullet_style_for_level(2, 0)  # BODY, level 0
            if bullet_style:
                assert isinstance(bullet_style, BulletInfo)
                assert isinstance(bullet_style.character, str)
                assert bullet_style.indent_level == 0
                
                if bullet_style.font:
                    assert isinstance(bullet_style.font, FontInfo)
                
        finally:
            Path(template_path).unlink()

    def test_template_style_placeholder_styles(self):
        """Test that placeholder styles are extracted correctly."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            style = parser.get_template_style()
            
            # Check placeholder_styles structure
            assert isinstance(style.placeholder_styles, dict)
            
            # Check that placeholder style info is properly structured
            for ph_type, style_info in style.placeholder_styles.items():
                assert isinstance(ph_type, int)
                assert isinstance(style_info, PlaceholderStyleInfo)
                assert style_info.placeholder_type == ph_type
                assert isinstance(style_info.type_name, str)
                assert isinstance(style_info.bullet_styles, list)
                
                # Check bullet styles if present
                for bullet in style_info.bullet_styles:
                    assert isinstance(bullet, BulletInfo)
                    assert isinstance(bullet.character, str)
                    assert isinstance(bullet.indent_level, int)
                    assert bullet.indent_level >= 0
                
        finally:
            Path(template_path).unlink()

    def test_theme_fonts_extraction(self):
        """Test theme fonts extraction."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            style = parser.get_template_style()
            
            # Check theme_fonts structure
            assert isinstance(style.theme_fonts, dict)
            
            # Should have at least default entries
            if style.theme_fonts:
                for font_role, font_name in style.theme_fonts.items():
                    assert isinstance(font_role, str)
                    assert isinstance(font_name, str)
                    assert font_name != ""
                
        finally:
            Path(template_path).unlink()

    def test_master_font_extraction(self):
        """Test master font extraction."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            style = parser.get_template_style()
            
            # Master font should be present (at least default)
            if style.master_font:
                assert isinstance(style.master_font, FontInfo)
                assert isinstance(style.master_font.name, str)
                assert style.master_font.name != ""
                
                # Size should be positive if present
                if style.master_font.size:
                    assert style.master_font.size > 0
                
                # Weight should be valid
                if style.master_font.weight:
                    assert style.master_font.weight in ["normal", "bold"]
                    
                # Color should be valid hex if present
                if style.master_font.color:
                    assert style.master_font.color.startswith("#")
                    assert len(style.master_font.color) == 7
                
        finally:
            Path(template_path).unlink()

    def test_font_fallback_behavior(self):
        """Test font fallback when extraction fails."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            # Test with non-existent placeholder type
            font = parser.get_font_for_placeholder_type(999)
            # Should return master font or None gracefully
            if font:
                assert isinstance(font, FontInfo)
                
        finally:
            Path(template_path).unlink()

    def test_bullet_level_fallback(self):
        """Test bullet style fallback for different levels."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            
            # Test multiple indentation levels
            for level in range(5):
                bullet = parser.get_bullet_style_for_level(2, level)  # BODY type
                if bullet:
                    assert isinstance(bullet, BulletInfo)
                    assert bullet.indent_level == level
                    assert bullet.character in ["•", "○", "▪", "‒", "►", ""]
                    
        finally:
            Path(template_path).unlink()

    def test_style_info_methods(self):
        """Test TemplateStyle helper methods."""
        template_path = self.create_test_template()
        
        try:
            parser = TemplateParser(template_path)
            style = parser.get_template_style()
            
            # Test get_font_for_placeholder_type method
            for ph_type in [1, 2, 3]:  # Common placeholder types
                font = style.get_font_for_placeholder_type(ph_type)
                if font:
                    assert isinstance(font, FontInfo)
            
            # Test get_bullet_style_for_level method
            bullet = style.get_bullet_style_for_level(2, 0)
            if bullet:
                assert isinstance(bullet, BulletInfo)
                
            # Test get_placeholder_style method
            placeholder_style = style.get_placeholder_style(2)
            if placeholder_style:
                assert isinstance(placeholder_style, PlaceholderStyleInfo)
                
        finally:
            Path(template_path).unlink()

    def create_test_template_with_fonts(self) -> str:
        """Create a test template with specific font configuration."""
        prs = Presentation()
        
        # Add a slide with title and content
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title text to trigger font extraction
        if slide.shapes.title:
            slide.shapes.title.text = "Test Title"
            
        # Add content slide
        if len(prs.slide_layouts) > 1:
            content_layout = prs.slide_layouts[1]  # Content slide
            content_slide = prs.slides.add_slide(content_layout)
            
            if content_slide.shapes.title:
                content_slide.shapes.title.text = "Content Title"
                
            # Try to add body text
            for shape in content_slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    shape.text = "• Bullet point 1\n• Bullet point 2"
                    break
        
        # Create temporary file
        temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        temp_path = temp_file.name
        temp_file.close()
        
        # Save the presentation
        prs.save(temp_path)
        return temp_path

    def test_font_extraction_with_content(self):
        """Test font extraction with actual content."""
        template_path = self.create_test_template_with_fonts()
        
        try:
            parser = TemplateParser(template_path)
            
            # Should be able to extract fonts from populated template
            assert hasattr(parser, 'template_style')
            style = parser.get_template_style()
            
            # Check that we have some placeholder styles
            assert isinstance(style.placeholder_styles, dict)
            
            # Master font should exist
            if style.master_font:
                assert isinstance(style.master_font, FontInfo)
                assert style.master_font.name != ""
                
        finally:
            Path(template_path).unlink()

    def test_title_font_detection_specific(self):
        """Test that title font detection works correctly on sample template."""
        template_path = self.create_test_template_with_fonts()
        
        try:
            parser = TemplateParser(template_path)
            
            # Test title font detection (placeholder type 1)
            title_font = parser.get_font_for_placeholder_type(1)
            
            # Should have some font information for title
            if title_font:
                assert isinstance(title_font, FontInfo)
                assert title_font.name is not None
                assert title_font.name != ""
                # Title fonts are typically larger
                if title_font.size:
                    assert title_font.size >= 12  # Reasonable minimum
            else:
                # If no specific title font, should fall back to master font
                master_font = parser.get_template_style().master_font
                assert master_font is not None
                assert isinstance(master_font, FontInfo)
                
        finally:
            Path(template_path).unlink()