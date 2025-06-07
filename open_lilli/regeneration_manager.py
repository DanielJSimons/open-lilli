"""Regeneration manager for selective slide regeneration."""

import logging
import re
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.text.text import _Paragraph

from .models import SlidePlan, Outline, GenerationConfig
from .template_parser import TemplateParser
from .content_generator import ContentGenerator
from .slide_assembler import SlideAssembler
from .exceptions import ValidationConfigError

logger = logging.getLogger(__name__)


class RegenerationManager:
    """Manages selective slide regeneration for existing presentations."""

    def __init__(
        self,
        template_parser: TemplateParser,
        content_generator: ContentGenerator,
        slide_assembler: SlideAssembler
    ):
        """
        Initialize the regeneration manager.
        
        Args:
            template_parser: TemplateParser instance for template access
            content_generator: ContentGenerator for regenerating slide content
            slide_assembler: SlideAssembler for patching presentations
        """
        self.template_parser = template_parser
        self.content_generator = content_generator
        self.slide_assembler = slide_assembler
        
        logger.info("RegenerationManager initialized")

    def extract_slides_from_presentation(self, pptx_path: Path) -> Tuple[Outline, List[SlidePlan]]:
        """
        Extract slide content from an existing PowerPoint presentation.
        
        Args:
            pptx_path: Path to the existing PowerPoint file
            
        Returns:
            Tuple of (Outline, List[SlidePlan]) extracted from the presentation
            
        Raises:
            FileNotFoundError: If the presentation file doesn't exist
            ValueError: If the presentation structure is invalid
        """
        if not pptx_path.exists():
            raise FileNotFoundError(f"Presentation file not found: {pptx_path}")
        
        logger.info(f"Extracting slides from presentation: {pptx_path}")
        
        try:
            prs = Presentation(str(pptx_path))
            
            # Extract presentation metadata
            title = self._extract_presentation_title(prs)
            slides = []
            
            # Extract each slide
            for slide_idx, slide in enumerate(prs.slides):
                slide_plan = self._extract_slide_content(slide, slide_idx)
                slides.append(slide_plan)
                logger.debug(f"Extracted slide {slide_idx}: {slide_plan.title}")
            
            # Create outline
            outline = Outline(
                title=title,
                slides=slides,
                language="en",  # Default, could be detected
                style_guidance="Extracted from existing presentation"
            )
            
            logger.info(f"Successfully extracted {len(slides)} slides from presentation")
            return outline, slides
            
        except Exception as e:
            logger.error(f"Failed to extract slides from presentation: {e}")
            raise ValueError(f"Invalid presentation structure: {e}") from e

    def select_slides_for_regeneration(
        self, 
        slides: List[SlidePlan], 
        target_indices: List[int]
    ) -> List[SlidePlan]:
        """
        Select specific slides for regeneration.
        
        Args:
            slides: List of all slide plans
            target_indices: List of slide indices to regenerate (0-based)
            
        Returns:
            List of slide plans to regenerate
            
        Raises:
            ValueError: If any target index is invalid
        """
        max_index = len(slides) - 1
        invalid_indices = [idx for idx in target_indices if idx < 0 or idx > max_index]
        
        if invalid_indices:
            raise ValueError(
                f"Invalid slide indices: {invalid_indices}. "
                f"Valid range is 0-{max_index} for {len(slides)} slides."
            )
        
        selected_slides = [slides[idx] for idx in target_indices]
        logger.info(f"Selected {len(selected_slides)} slides for regeneration: {target_indices}")
        
        return selected_slides

    def coordinate_selective_regeneration(
        self,
        input_pptx: Path,
        target_indices: List[int],
        config: GenerationConfig,
        feedback: Optional[str] = None,
        language: str = "en"
    ) -> Tuple[Outline, List[SlidePlan]]:
        """
        Coordinate the selective regeneration process.
        
        Args:
            input_pptx: Path to existing presentation
            target_indices: List of slide indices to regenerate
            config: Generation configuration
            feedback: Optional feedback for targeted improvements
            language: Language code for content generation
            
        Returns:
            Tuple of (updated_outline, updated_slides)
        """
        logger.info(f"Starting selective regeneration for slides {target_indices}")
        
        # Extract existing slides
        outline, all_slides = self.extract_slides_from_presentation(input_pptx)
        
        # Select slides for regeneration
        target_slides = self.select_slides_for_regeneration(all_slides, target_indices)
        
        # Regenerate content for selected slides
        regenerated_slides = self.content_generator.regenerate_specific_slides(
            target_slides, config, outline.style_guidance, language, feedback
        )
        
        # Update the slides list with regenerated content
        updated_slides = all_slides.copy()
        for regenerated_slide in regenerated_slides:
            updated_slides[regenerated_slide.index] = regenerated_slide
        
        # Update outline with new slides
        updated_outline = outline.model_copy()
        updated_outline.slides = updated_slides
        
        logger.info(f"Selective regeneration completed for {len(regenerated_slides)} slides")
        return updated_outline, updated_slides

    def _extract_presentation_title(self, prs: Presentation) -> str:
        """
        Extract the presentation title from the first slide or properties.
        
        Args:
            prs: PowerPoint presentation object
            
        Returns:
            Extracted presentation title
        """
        # Try to get title from first slide
        if len(prs.slides) > 0:
            first_slide = prs.slides[0]
            title = self._extract_slide_title(first_slide)
            if title and title != "Untitled Slide":
                return title
        
        # Fallback to generic title
        return "Extracted Presentation"

    def _extract_slide_content(self, slide, slide_idx: int) -> SlidePlan:
        """
        Extract content from a single slide.
        
        Args:
            slide: PowerPoint slide object
            slide_idx: Index of the slide
            
        Returns:
            SlidePlan object with extracted content
        """
        title = self._extract_slide_title(slide)
        bullets = self._extract_slide_bullets(slide)
        slide_type = self._determine_slide_type(slide, title, bullets)
        
        return SlidePlan(
            index=slide_idx,
            slide_type=slide_type,
            title=title or f"Slide {slide_idx + 1}",
            bullets=bullets,
            layout_id=self._determine_layout_id(slide_type)
        )

    def _extract_slide_title(self, slide) -> Optional[str]:
        """
        Extract title from a slide.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            Extracted title or None
        """
        try:
            # Check if slide has a title shape
            if hasattr(slide, 'shapes') and slide.shapes.title:
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    return title_text
            
            # Search for title-like text in all shapes
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    text = shape.text.strip()
                    # Consider it a title if it's short and not bullet-like
                    if text and len(text) < 100 and not text.startswith('•'):
                        return text
                        
        except Exception as e:
            logger.debug(f"Error extracting slide title: {e}")
        
        return None

    def _extract_slide_bullets(self, slide) -> List[str]:
        """
        Extract bullet points from a slide.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            List of bullet point strings
        """
        bullets = []
        
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text_frame = shape.text_frame
                    
                    for paragraph in text_frame.paragraphs:
                        if paragraph.text.strip():
                            # Clean up bullet text
                            bullet_text = paragraph.text.strip()
                            # Remove common bullet characters
                            bullet_text = re.sub(r'^[•·▪▫◦‣⁃]\s*', '', bullet_text)
                            bullet_text = re.sub(r'^[-*]\s*', '', bullet_text)
                            
                            if bullet_text and len(bullet_text) > 3:  # Avoid very short fragments
                                bullets.append(bullet_text)
                                
        except Exception as e:
            logger.debug(f"Error extracting slide bullets: {e}")
        
        return bullets

    def _determine_slide_type(self, slide, title: Optional[str], bullets: List[str]) -> str:
        """
        Determine the slide type based on content.
        
        Args:
            slide: PowerPoint slide object
            title: Extracted title
            bullets: Extracted bullets
            
        Returns:
            Slide type string
        """
        # Check for images
        has_images = any(
            shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
            for shape in slide.shapes
        )
        
        # Check for charts
        has_charts = any(
            shape.shape_type == 3  # MSO_SHAPE_TYPE.CHART
            for shape in slide.shapes
        )
        
        # Determine type based on content
        if has_charts:
            return "chart"
        elif has_images:
            return "image"
        elif bullets:
            return "content"
        elif title and not bullets:
            return "title"
        else:
            return "content"

    def _determine_layout_id(self, slide_type: str) -> int:
        """
        Map slide type to layout ID.
        
        Args:
            slide_type: Type of slide
            
        Returns:
            Layout ID for the slide type
        """
        layout_mapping = {
            "title": 0,
            "content": 1,
            "chart": 1,
            "image": 5,  # Assuming image layout exists
            "two_column": 3
        }
        
        return layout_mapping.get(slide_type, 1)  # Default to content layout