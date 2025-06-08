"""Template loader for hardened template loading with content stripping."""

import logging
from pathlib import Path
from typing import Dict, List, Optional, Set

from pptx import Presentation
from pptx.slide import SlideLayout
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


class TemplateLoader:
    """
    Hardened template loader that strips content slides and caches semantic layout mappings.
    
    This class ensures that when load_template(path) is called, all content slides are removed,
    preserving only masters and layouts. It builds and caches semantic layout name mappings.
    """
    
    def __init__(self, template_path: str):
        """
        Initialize the template loader and perform hardened loading.
        
        Args:
            template_path: Path to the .pptx template file
            
        Raises:
            FileNotFoundError: If template file doesn't exist
            ValueError: If template file is invalid
        """
        self.template_path = Path(template_path)
        
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template file not found: {template_path}")
        
        if not self.template_path.suffix.lower() == '.pptx':
            raise ValueError(f"Template must be a .pptx file, got: {self.template_path.suffix}")
        
        logger.info(f"Loading template with hardened mode: {self.template_path}")
        
        try:
            # Load the presentation
            self.prs = Presentation(str(self.template_path))
            logger.info(f"Template loaded with {len(self.prs.slide_layouts)} layouts and {len(self.prs.slides)} content slides")
            
            # Strip all content slides - this is the key hardening step
            self._strip_content_slides()
            
            # Build semantic layout mappings
            self._layout_map = self._build_semantic_layout_map()
            
            logger.info(f"Template hardening complete. Stripped to {len(self.prs.slides)} slides, mapped {len(self._layout_map)} semantic layouts")
            
        except Exception as e:
            raise ValueError(f"Failed to load template: {e}")
    
    def _strip_content_slides(self) -> None:
        """
        Remove all content slides from the presentation, preserving only masters/layouts.
        
        This ensures the loaded template contains no actual slide content, only the structure.
        """
        # Remove slides in reverse order to avoid index shifting issues
        slides_to_remove = list(range(len(self.prs.slides)))
        for slide_index in reversed(slides_to_remove):
            slide_part = self.prs.slides._sld_lst[slide_index]
            self.prs.part.drop_rel(slide_part.rId)
            del self.prs.slides._sld_lst[slide_index]
        
        logger.info(f"Stripped {len(slides_to_remove)} content slides from template")
    
    def _build_semantic_layout_map(self) -> Dict[str, int]:
        """
        Build mapping from semantic layout names to layout indices by analyzing placeholders.
        
        Returns:
            Dictionary mapping semantic names (e.g. "title", "content") to layout indices
        """
        layout_map = {}
        
        for i, layout in enumerate(self.prs.slide_layouts):
            semantic_name = self._classify_layout_by_placeholders(layout, i)
            
            # Store the first occurrence of each semantic type
            if semantic_name not in layout_map:
                layout_map[semantic_name] = i
                logger.debug(f"Layout {i}: '{semantic_name}' - {self._get_placeholder_summary(layout)}")
        
        # Ensure we have essential layouts with fallbacks
        self._ensure_essential_layouts(layout_map)
        
        return layout_map
    
    def _classify_layout_by_placeholders(self, layout: SlideLayout, index: int) -> str:
        """
        Classify a layout based on its placeholder structure.
        
        Args:
            layout: SlideLayout to classify
            index: Layout index for fallback naming
            
        Returns:
            Semantic name for the layout
        """
        placeholders = layout.placeholders
        placeholder_types = [ph.placeholder_format.type for ph in placeholders]
        
        # Count different types of placeholders
        title_count = sum(1 for t in placeholder_types if t in (1, 13))  # TITLE, CENTERED_TITLE
        content_count = sum(1 for t in placeholder_types if t in (2, 7))  # BODY, OBJECT
        subtitle_count = sum(1 for t in placeholder_types if t == 3)  # SUBTITLE
        picture_count = sum(1 for t in placeholder_types if t == 18)  # PICTURE
        chart_count = sum(1 for t in placeholder_types if t == 14)  # CHART
        table_count = sum(1 for t in placeholder_types if t == 15)  # TABLE
        
        total_placeholders = len(placeholders)
        
        logger.debug(f"Layout {index}: {total_placeholders} placeholders - "
                    f"title:{title_count}, content:{content_count}, "
                    f"subtitle:{subtitle_count}, picture:{picture_count}, "
                    f"chart:{chart_count}, table:{table_count}")
        
        # Classification logic based on placeholder combinations
        if title_count >= 1 and subtitle_count >= 1 and total_placeholders <= 3:
            return "title"
        elif title_count >= 1 and content_count >= 2:
            return "two_column"  
        elif title_count >= 1 and picture_count >= 1 and content_count >= 1:
            return "image_content"
        elif title_count >= 1 and chart_count >= 1:
            return "chart"
        elif title_count >= 1 and table_count >= 1:
            return "table"
        elif title_count >= 1 and picture_count >= 1:
            return "image"
        elif title_count >= 1 and content_count == 1:
            return "content"
        elif title_count >= 1 and content_count == 0 and total_placeholders <= 2:
            return "section"
        elif total_placeholders == 0:
            return "blank"
        else:
            # Fallback to generic naming
            return f"layout_{index}"
    
    def _get_placeholder_summary(self, layout: SlideLayout) -> str:
        """Get a summary of placeholders in the layout for debugging."""
        placeholders = layout.placeholders
        type_counts = {}
        
        for ph in placeholders:
            ph_type = ph.placeholder_format.type
            type_name = self._placeholder_type_name(ph_type)
            type_counts[type_name] = type_counts.get(type_name, 0) + 1
        
        summary_parts = [f"{name}({count})" for name, count in type_counts.items()]
        return f"[{', '.join(summary_parts)}]"
    
    def _placeholder_type_name(self, ph_type: int) -> str:
        """Convert placeholder type number to readable name."""
        type_names = {
            1: "TITLE",
            2: "BODY", 
            3: "SUBTITLE",
            7: "OBJECT",
            13: "CENTERED_TITLE",
            14: "CHART",
            15: "TABLE",
            16: "CLIP_ART",
            17: "MEDIA_CLIP",
            18: "PICTURE"
        }
        return type_names.get(ph_type, f"TYPE_{ph_type}")
    
    def _ensure_essential_layouts(self, layout_map: Dict[str, int]) -> None:
        """
        Ensure we have mappings for essential layout types with fallbacks.
        
        Args:
            layout_map: Dictionary to update with essential layouts
        """
        essential_layouts = {
            "title": ["title", "section", "content"],
            "content": ["content", "two_column", "blank"],
            "section": ["section", "title", "content"],
            "blank": ["blank", "content"]
        }
        
        for essential_type, fallback_chain in essential_layouts.items():
            if essential_type not in layout_map:
                # Try each fallback in order
                for fallback in fallback_chain:
                    if fallback in layout_map:
                        layout_map[essential_type] = layout_map[fallback]
                        logger.info(f"Mapped essential layout '{essential_type}' to '{fallback}' (index {layout_map[fallback]})")
                        break
                else:
                    # Use the first available layout as last resort
                    if layout_map:
                        first_index = list(layout_map.values())[0]
                        layout_map[essential_type] = first_index
                        logger.warning(f"No suitable fallback for '{essential_type}', using layout {first_index}")
    
    def get_layout_index(self, semantic_name: str) -> Optional[int]:
        """
        Get layout index by semantic name.
        
        Args:
            semantic_name: Semantic layout name (e.g. "title", "content", "two_column")
            
        Returns:
            Layout index or None if not found
        """
        return self._layout_map.get(semantic_name)
    
    def get_available_layouts(self) -> List[str]:
        """
        Get list of available semantic layout names.
        
        Returns:
            List of semantic layout names
        """
        return list(self._layout_map.keys())
    
    def get_layout_map(self) -> Dict[str, int]:
        """
        Get the complete layout mapping.
        
        Returns:
            Dictionary mapping semantic names to layout indices
        """
        return self._layout_map.copy()
    
    def validate_placeholder_match(self, semantic_name: str, expected_placeholders: Set[str]) -> bool:
        """
        Validate that a layout contains expected placeholder types.
        
        Args:
            semantic_name: Semantic layout name to validate
            expected_placeholders: Set of expected placeholder type names
            
        Returns:
            True if layout contains all expected placeholders
        """
        layout_index = self.get_layout_index(semantic_name)
        if layout_index is None:
            return False
        
        layout = self.prs.slide_layouts[layout_index]
        actual_placeholders = set()
        
        for ph in layout.placeholders:
            ph_type = ph.placeholder_format.type
            type_name = self._placeholder_type_name(ph_type)
            actual_placeholders.add(type_name)
        
        return expected_placeholders.issubset(actual_placeholders)
    
    def get_layout_info(self, semantic_name: str) -> Optional[Dict]:
        """
        Get detailed information about a specific layout.
        
        Args:
            semantic_name: Semantic layout name
            
        Returns:
            Dictionary with layout information or None if not found
        """
        layout_index = self.get_layout_index(semantic_name)
        if layout_index is None:
            return None
        
        layout = self.prs.slide_layouts[layout_index]
        placeholders = layout.placeholders
        
        placeholder_info = []
        for i, ph in enumerate(placeholders):
            ph_type = ph.placeholder_format.type
            type_name = self._placeholder_type_name(ph_type)
            
            placeholder_info.append({
                "index": i,
                "type": ph_type,
                "type_name": type_name,
                "position": {
                    "left": ph.left,
                    "top": ph.top,
                    "width": ph.width,
                    "height": ph.height
                }
            })
        
        return {
            "semantic_name": semantic_name,
            "layout_index": layout_index,
            "total_placeholders": len(placeholders),
            "placeholders": placeholder_info,
            "summary": self._get_placeholder_summary(layout)
        }