"""Template parser for analyzing PowerPoint templates."""

import logging
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.slide import SlideLayout
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE

from .models import (
    TemplateStyle, FontInfo, BulletInfo, PlaceholderStyleInfo,
    TemplateCompatibilityReport, DesignPattern
)

logger = logging.getLogger(__name__)


class TemplateParser:
    """Analyzes PowerPoint templates and catalogues slide layouts."""

    def __init__(self, template_path: str):
        """
        Initialize the template parser.
        
        Args:
            template_path: Path to the .pptx template file
            
        Raises:
            FileNotFoundError: If template file doesn't exist
            ValueError: If template file is invalid
        """
        self.template_path = Path(template_path)
        
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template file not found: {template_path}")
        
        if not self.template_path.suffix.lower() == '.pptx':
            raise ValueError(f"Template must be a .pptx file, got: {self.template_path.suffix}")
        
        logger.info(f"Loading template: {self.template_path}")
        
        try:
            self.prs = Presentation(str(self.template_path))
            logger.info(f"Template loaded successfully with {len(self.prs.slide_layouts)} layouts")
        except Exception as e:
            raise ValueError(f"Failed to load template: {e}")
        
        # Analyze the template
        self.layout_map = self._index_layouts()
        self.reverse_layout_map = {v: k for k, v in self.layout_map.items()} # Added reverse map
        self.palette = self._extract_theme_colors()
        self.template_style = self._extract_template_style()

        if self.template_style:
            default_lang_fonts = {
                "ar": "Arial Unicode MS",
                "he": "David Libre",
                "fa": "Tahoma",
                "ja": "MS Gothic",
                "ko": "Malgun Gothic",
                "zh": "SimSun"
            }
            if not self.template_style.language_specific_fonts:
                self.template_style.language_specific_fonts = default_lang_fonts
            else:
                merged_fonts = default_lang_fonts.copy()
                merged_fonts.update(self.template_style.language_specific_fonts)
                self.template_style.language_specific_fonts = merged_fonts
        
        logger.info(f"Template analysis complete. Layout map: {self.layout_map}")
        logger.info(f"Template style extracted with {len(self.template_style.placeholder_styles)} placeholder styles")

    def _index_layouts(self) -> Dict[str, int]:
        """
        Analyze slide layouts and create a mapping from semantic names to indices.
        
        Returns:
            Dictionary mapping layout types to layout indices
        """
        layout_map = {}
        
        for i, layout in enumerate(self.prs.slide_layouts):
            layout_type = self._classify_layout(layout, i)
            
            # Store the first occurrence of each type
            if layout_type not in layout_map:
                layout_map[layout_type] = i
                logger.debug(f"Layout {i}: {layout_type} - {self._get_layout_info(layout)}")
        
        # Ensure we have at least basic layouts
        self._ensure_basic_layouts(layout_map)
        
        return layout_map

    def _classify_layout(self, layout: SlideLayout, index: int) -> str:
        """
        Classify a layout based on its placeholders and structure.
        
        Args:
            layout: SlideLayout to classify
            index: Layout index for fallback naming
            
        Returns:
            Semantic name for the layout
        """
        placeholders = layout.placeholders
        placeholder_types = [ph.placeholder_format.type for ph in placeholders]
        
        # Count different types of placeholders
        title_count = sum(1 for t in placeholder_types if t in (1, 13))  # TITLE, CENTERED_TITLE
        content_count = sum(1 for t in placeholder_types if t in (2, 7))  # BODY, OBJECT
        subtitle_count = sum(1 for t in placeholder_types if t == 3)  # SUBTITLE
        picture_count = sum(1 for t in placeholder_types if t == 18)  # PICTURE
        
        total_placeholders = len(placeholders)
        
        logger.debug(f"Layout {index}: {total_placeholders} placeholders - "
                    f"title:{title_count}, content:{content_count}, "
                    f"subtitle:{subtitle_count}, picture:{picture_count}")
        
        # Enhanced classification logic with more layout types
        if title_count >= 1 and subtitle_count >= 1 and total_placeholders <= 3:
            return "title"
        elif title_count >= 1 and content_count >= 3:
            return "three_column"
        elif title_count >= 1 and content_count >= 2:
            # Check if it's a comparison layout based on placeholder positioning
            if total_placeholders >= 4:  # Title + 2+ content + possibly other elements
                return "comparison"
            else:
                return "two_column"
        elif title_count >= 1 and picture_count >= 1 and content_count >= 1:
            return "image_content"
        elif title_count >= 1 and picture_count >= 1:
            return "image"
        elif title_count >= 1 and content_count == 1:
            # Distinguish between regular content and dense content layouts
            if total_placeholders >= 3:  # Title + content + additional elements
                return "content_dense"
            else:
                return "content"
        elif title_count >= 1 and content_count == 0 and total_placeholders <= 2:
            return "section"
        elif total_placeholders == 0:
            return "blank"
        else:
            # Try to infer from total placeholder count
            if total_placeholders >= 5:
                return "content_dense"
            elif total_placeholders == 4:
                return "two_column"
            else:
                return "content"  # Better default than generic naming

    def _ensure_basic_layouts(self, layout_map: Dict[str, int]) -> None:
        """Ensure we have mappings for essential layout types."""
        essentials = {
            "title": "title",
            "content": "content", 
            "section": "section",
            "blank": "blank"
        }
        
        for layout_type, fallback in essentials.items():
            if layout_type not in layout_map:
                # Try to find a suitable fallback
                if fallback in layout_map:
                    layout_map[layout_type] = layout_map[fallback]
                elif layout_map:
                    # Use the first available layout as last resort
                    layout_map[layout_type] = list(layout_map.values())[0]
                    logger.warning(f"No {layout_type} layout found, using layout {layout_map[layout_type]} as fallback")

    def _get_layout_info(self, layout: SlideLayout) -> str:
        """Get human-readable info about a layout."""
        placeholders = layout.placeholders
        ph_info = []
        
        for ph in placeholders:
            ph_type = ph.placeholder_format.type
            type_name = self._placeholder_type_name(ph_type)
            ph_info.append(f"{type_name}({ph_type})")
        
        return f"[{', '.join(ph_info)}]"

    def _placeholder_type_name(self, ph_type: int) -> str:
        """Convert placeholder type number to readable name."""
        type_names = {
            1: "TITLE",
            2: "BODY", 
            3: "SUBTITLE",
            7: "OBJECT",
            13: "CENTERED_TITLE",
            18: "PICTURE",
            14: "CHART",
            15: "TABLE",
            16: "CLIP_ART",
            17: "MEDIA_CLIP"
        }
        return type_names.get(ph_type, f"TYPE_{ph_type}")

    def _extract_theme_colors(self) -> Dict[str, str]:
        """
        Extract theme colors from the presentation.
        
        Returns:
            Dictionary of color names to hex values
        """
        try:
            # Extract theme colors from theme1.xml
            theme_colors = self.get_theme_colors()
            if theme_colors:
                logger.debug(f"Extracted theme colors: {list(theme_colors.keys())}")
                return theme_colors
        except Exception as e:
            logger.warning(f"Failed to extract theme colors from XML: {e}")
        
        # Fallback to default colors if extraction fails
        default_palette = {
            "dk1": "#000000",      # Dark 1 (usually black)
            "lt1": "#FFFFFF",      # Light 1 (usually white)
            "acc1": "#1F497D",     # Accent 1 (dark blue)
            "acc2": "#4F81BD",     # Accent 2 (medium blue)
            "acc3": "#9BBB59",     # Accent 3 (green)
            "acc4": "#F79646",     # Accent 4 (orange)
            "acc5": "#8064A2",     # Accent 5 (purple)
            "acc6": "#4BACC6",     # Accent 6 (light blue)
        }
        
        logger.debug(f"Using default color palette: {list(default_palette.keys())}")
        return default_palette

    def get_layout(self, slide_type: str) -> SlideLayout:
        """
        Get a slide layout by semantic type.
        
        Args:
            slide_type: Semantic layout type (title, content, image, etc.)
            
        Returns:
            SlideLayout object
            
        Raises:
            ValueError: If layout type is not found
        """
        if slide_type not in self.layout_map:
            available = list(self.layout_map.keys())
            raise ValueError(f"Layout type '{slide_type}' not found. Available: {available}")
        
        layout_index = self.layout_map[slide_type]
        return self.prs.slide_layouts[layout_index]

    def get_layout_index(self, slide_type: str) -> int:
        """
        Get layout index by semantic type.
        
        Args:
            slide_type: Semantic layout type
            
        Returns:
            Layout index
        """
        if slide_type not in self.layout_map:
            # Return a sensible default
            return self.layout_map.get("content", 0)
        
        return self.layout_map[slide_type]

    def list_available_layouts(self) -> List[str]:
        """
        Get list of available layout types.
        
        Returns:
            List of layout type names
        """
        return list(self.layout_map.keys())

    def analyze_layout_placeholders(self, slide_type: str) -> Dict[str, any]:
        """
        Analyze the placeholders in a specific layout.
        
        Args:
            slide_type: Layout type to analyze
            
        Returns:
            Dictionary with placeholder analysis
        """
        layout = self.get_layout(slide_type)
        placeholders = layout.placeholders
        
        analysis = {
            "total_placeholders": len(placeholders),
            "placeholder_details": [],
            "has_title": False,
            "has_content": False,
            "has_image": False,
            "has_chart": False
        }
        
        for i, ph in enumerate(placeholders):
            ph_type = ph.placeholder_format.type
            type_name = self._placeholder_type_name(ph_type)
            
            ph_detail = {
                "index": i,
                "type": ph_type,
                "type_name": type_name,
                "position": {
                    "left": ph.left,
                    "top": ph.top,
                    "width": ph.width,
                    "height": ph.height
                }
            }
            
            analysis["placeholder_details"].append(ph_detail)
            
            # Set flags based on placeholder types
            if ph_type in (1, 13):  # TITLE types
                analysis["has_title"] = True
            elif ph_type in (2, 7):  # BODY, OBJECT types
                analysis["has_content"] = True
            elif ph_type == 18:  # PICTURE
                analysis["has_image"] = True
            elif ph_type == 14:  # CHART
                analysis["has_chart"] = True
        
        return analysis
    
    def extract_complete_layout_visual_data(self, layout: SlideLayout, layout_index: int) -> Dict:
        """
        Extract comprehensive visual layout information including all shapes, positioning,
        and spatial relationships for LLM-based template selection.
        
        Args:
            layout: SlideLayout to analyze
            layout_index: Index of the layout
            
        Returns:
            Complete visual analysis including shapes, relationships, and design patterns
        """
        try:
            visual_data = {
                "layout_index": layout_index,
                "layout_name": getattr(layout, 'name', f"Layout_{layout_index}"),
                "shapes": self._extract_all_shapes(layout),
                "spatial_analysis": self._analyze_spatial_relationships(layout),
                "design_patterns": self._detect_design_patterns(layout),
                "content_zones": self._map_content_zones(layout),
                "visual_summary": self._generate_visual_summary(layout),
                "recommended_content_types": self._infer_content_types(layout)
            }
            
            logger.debug(f"Extracted visual data for layout {layout_index}: {len(visual_data['shapes']['all_shapes'])} total shapes")
            return visual_data
            
        except Exception as e:
            logger.error(f"Failed to extract visual data for layout {layout_index}: {e}")
            return {"layout_index": layout_index, "error": str(e)}
    
    def _extract_all_shapes(self, layout: SlideLayout) -> Dict:
        """Extract all shapes from layout with complete visual information."""
        shapes_data = {
            "all_shapes": [],
            "placeholders": [],
            "text_elements": [],
            "visual_elements": [],
            "background_elements": []
        }
        
        try:
            for shape in layout.shapes:
                shape_info = {
                    "id": getattr(shape, 'shape_id', None),
                    "type": str(shape.shape_type),
                    "name": getattr(shape, 'name', ''),
                    "position": {
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height,
                        "center_x": shape.left + shape.width // 2,
                        "center_y": shape.top + shape.height // 2
                    },
                    "visual_properties": self._extract_shape_visual_properties(shape),
                    "content_type": self._classify_shape_content(shape)
                }
                
                shapes_data["all_shapes"].append(shape_info)
                
                # Categorize shapes
                if hasattr(shape, 'placeholder_format'):
                    shape_info["placeholder_type"] = shape.placeholder_format.type
                    shapes_data["placeholders"].append(shape_info)
                elif hasattr(shape, 'text_frame') and shape.text_frame:
                    shapes_data["text_elements"].append(shape_info)
                elif shape.shape_type in [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM]:
                    shapes_data["visual_elements"].append(shape_info)
                else:
                    shapes_data["background_elements"].append(shape_info)
                    
        except Exception as e:
            logger.warning(f"Error extracting shapes: {e}")
            
        return shapes_data
    
    def _extract_shape_visual_properties(self, shape) -> Dict:
        """Extract visual styling properties from a shape."""
        properties = {"fill": None, "line": None, "effects": None}
        
        try:
            # Fill properties
            if hasattr(shape, 'fill'):
                fill = shape.fill
                if hasattr(fill, 'type') and fill.type:
                    properties["fill"] = {
                        "type": str(fill.type),
                        "color": self._extract_color_info(fill) if hasattr(fill, 'fore_color') else None
                    }
            
            # Line properties
            if hasattr(shape, 'line'):
                line = shape.line
                properties["line"] = {
                    "color": self._extract_color_info(line) if hasattr(line, 'color') else None,
                    "width": getattr(line, 'width', None)
                }
                
        except Exception as e:
            logger.debug(f"Could not extract visual properties: {e}")
            
        return properties
    
    def _extract_color_info(self, color_obj) -> Optional[str]:
        """Extract color information as hex string."""
        try:
            if hasattr(color_obj, 'rgb'):
                rgb = color_obj.rgb
                return f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
        except:
            pass
        return None
    
    def _classify_shape_content(self, shape) -> str:
        """Classify what type of content this shape is intended for."""
        if hasattr(shape, 'placeholder_format'):
            ph_type = shape.placeholder_format.type
            if ph_type in (1, 13):  # TITLE types
                return "title"
            elif ph_type in (2, 7):  # BODY, OBJECT
                return "content"
            elif ph_type == 18:  # PICTURE
                return "image"
            elif ph_type == 14:  # CHART
                return "chart"
            else:
                return f"placeholder_{ph_type}"
        elif hasattr(shape, 'text_frame') and shape.text_frame:
            return "text"
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "image"
        else:
            return "decoration"
    
    def _analyze_spatial_relationships(self, layout: SlideLayout) -> Dict:
        """Analyze how shapes relate spatially to understand layout structure."""
        shapes = list(layout.shapes)
        
        return {
            "alignment_patterns": self._detect_alignment_patterns(shapes),
            "layout_grid": self._detect_grid_structure(shapes),
            "content_zones": self._identify_content_zones(shapes),
            "visual_balance": self._analyze_visual_balance(shapes)
        }
    
    def _detect_alignment_patterns(self, shapes) -> Dict:
        """Detect common alignment patterns in the layout."""
        left_edges = [shape.left for shape in shapes]
        top_edges = [shape.top for shape in shapes]
        
        # Group shapes by similar positions (within tolerance)
        tolerance = 50000  # EMUs (about 1.8 inches)
        
        aligned_left = self._group_by_proximity(left_edges, tolerance)
        aligned_top = self._group_by_proximity(top_edges, tolerance)
        
        return {
            "vertical_alignments": len(aligned_left),
            "horizontal_alignments": len(aligned_top),
            "has_center_alignment": self._check_center_alignment(shapes),
            "has_grid_layout": len(aligned_left) > 1 and len(aligned_top) > 1
        }
    
    def _group_by_proximity(self, values: List[int], tolerance: int) -> List[List[int]]:
        """Group values that are within tolerance of each other."""
        groups = []
        sorted_values = sorted(values)
        
        current_group = [sorted_values[0]] if sorted_values else []
        
        for value in sorted_values[1:]:
            if abs(value - current_group[-1]) <= tolerance:
                current_group.append(value)
            else:
                if len(current_group) > 1:  # Only count as group if multiple items
                    groups.append(current_group)
                current_group = [value]
        
        if len(current_group) > 1:
            groups.append(current_group)
            
        return groups
    
    def _check_center_alignment(self, shapes) -> bool:
        """Check if shapes are center-aligned."""
        try:
            slide_width = self.prs.slide_width
            centers = [shape.left + shape.width // 2 for shape in shapes]
            slide_center = slide_width // 2
            
            # Check if any shapes are close to slide center
            tolerance = slide_width // 10  # 10% tolerance
            return any(abs(center - slide_center) < tolerance for center in centers)
        except:
            return False
    
    def _detect_grid_structure(self, shapes) -> Dict:
        """Detect if layout follows a grid structure."""
        try:
            # Analyze placeholder positions to infer grid
            placeholders = [s for s in shapes if hasattr(s, 'placeholder_format')]
            
            if len(placeholders) < 2:
                return {"type": "single", "columns": 1, "rows": 1}
            
            # Group by horizontal position (columns)
            left_positions = [p.left for p in placeholders]
            unique_lefts = sorted(set(left_positions))
            
            # Group by vertical position (rows)  
            top_positions = [p.top for p in placeholders]
            unique_tops = sorted(set(top_positions))
            
            cols = len(unique_lefts)
            rows = len(unique_tops)
            
            grid_type = "single"
            if cols > 1 and rows == 1:
                grid_type = "horizontal"
            elif cols == 1 and rows > 1:
                grid_type = "vertical"
            elif cols > 1 and rows > 1:
                grid_type = "grid"
                
            return {
                "type": grid_type,
                "columns": cols,
                "rows": rows,
                "total_cells": len(placeholders)
            }
            
        except Exception as e:
            logger.debug(f"Grid detection failed: {e}")
            return {"type": "unknown", "columns": 1, "rows": 1}
    
    def _detect_design_patterns(self, layout: SlideLayout) -> Dict:
        """Detect common design patterns and layout intentions."""
        shapes = list(layout.shapes)
        placeholders = [s for s in shapes if hasattr(s, 'placeholder_format')]
        
        patterns = {
            "layout_style": "unknown",
            "content_orientation": "mixed",
            "visual_complexity": "medium",
            "primary_purpose": "general"
        }
        
        try:
            # Determine layout style
            if len(placeholders) <= 2:
                patterns["layout_style"] = "minimal"
            elif len(placeholders) >= 5:
                patterns["layout_style"] = "complex"
            else:
                patterns["layout_style"] = "standard"
            
            # Analyze content orientation
            title_shapes = [s for s in placeholders if hasattr(s, 'placeholder_format') and s.placeholder_format.type in (1, 13)]
            content_shapes = [s for s in placeholders if hasattr(s, 'placeholder_format') and s.placeholder_format.type in (2, 7)]
            
            if len(content_shapes) > 1:
                # Check if content is arranged horizontally or vertically
                content_tops = [s.top for s in content_shapes]
                content_lefts = [s.left for s in content_shapes]
                
                if max(content_tops) - min(content_tops) > max(content_lefts) - min(content_lefts):
                    patterns["content_orientation"] = "vertical"
                else:
                    patterns["content_orientation"] = "horizontal"
            
            # Determine primary purpose
            image_placeholders = [s for s in placeholders if hasattr(s, 'placeholder_format') and s.placeholder_format.type == 18]
            
            if len(image_placeholders) > 0 and len(content_shapes) > 0:
                patterns["primary_purpose"] = "image_content"
            elif len(image_placeholders) > 0:
                patterns["primary_purpose"] = "image_focused"
            elif len(content_shapes) > 2:
                patterns["primary_purpose"] = "content_heavy"
            elif len(content_shapes) == 0:
                patterns["primary_purpose"] = "title_only"
            else:
                patterns["primary_purpose"] = "balanced"
                
        except Exception as e:
            logger.debug(f"Pattern detection failed: {e}")
            
        return patterns
    
    def _map_content_zones(self, layout: SlideLayout) -> Dict:
        """Map out distinct content zones in the layout."""
        try:
            slide_width = self.prs.slide_width
            slide_height = self.prs.slide_height
            
            zones = {
                "header": [],
                "main_content": [],
                "sidebar": [],
                "footer": []
            }
            
            for shape in layout.shapes:
                if hasattr(shape, 'placeholder_format'):
                    # Classify zone based on position
                    rel_top = shape.top / slide_height
                    rel_left = shape.left / slide_width
                    rel_width = shape.width / slide_width
                    
                    shape_data = {
                        "type": shape.placeholder_format.type,
                        "position": (rel_left, rel_top),
                        "size": (rel_width, shape.height / slide_height)
                    }
                    
                    if rel_top < 0.2:  # Top 20%
                        zones["header"].append(shape_data)
                    elif rel_top > 0.8:  # Bottom 20%
                        zones["footer"].append(shape_data)
                    elif rel_left > 0.7:  # Right 30%
                        zones["sidebar"].append(shape_data)
                    else:
                        zones["main_content"].append(shape_data)
                        
            return zones
            
        except Exception as e:
            logger.debug(f"Zone mapping failed: {e}")
            return {"header": [], "main_content": [], "sidebar": [], "footer": []}
    
    def _generate_visual_summary(self, layout: SlideLayout) -> str:
        """Generate a human-readable description of the layout for LLM analysis."""
        try:
            shapes = list(layout.shapes)
            placeholders = [s for s in shapes if hasattr(s, 'placeholder_format')]
            
            # Count different types
            title_count = len([s for s in placeholders if s.placeholder_format.type in (1, 13)])
            content_count = len([s for s in placeholders if s.placeholder_format.type in (2, 7)])
            image_count = len([s for s in placeholders if s.placeholder_format.type == 18])
            
            # Analyze layout structure
            grid_info = self._detect_grid_structure(shapes)
            patterns = self._detect_design_patterns(layout)
            
            # Generate description
            summary_parts = []
            
            if title_count > 0:
                summary_parts.append(f"{title_count} title area{'s' if title_count > 1 else ''}")
            
            if content_count > 0:
                if content_count == 1:
                    summary_parts.append("single content area")
                elif content_count == 2:
                    summary_parts.append("two-column content layout")
                else:
                    summary_parts.append(f"{content_count}-column content layout")
            
            if image_count > 0:
                summary_parts.append(f"{image_count} image placeholder{'s' if image_count > 1 else ''}")
            
            layout_desc = f"Layout with {', '.join(summary_parts)}" if summary_parts else "Minimal layout"
            
            # Add structural information
            if grid_info["type"] != "single":
                layout_desc += f", arranged in {grid_info['type']} structure ({grid_info['columns']}x{grid_info['rows']})"
            
            # Add design pattern
            layout_desc += f", {patterns['layout_style']} style, optimized for {patterns['primary_purpose']} content"
            
            return layout_desc
            
        except Exception as e:
            logger.debug(f"Summary generation failed: {e}")
            return "Layout analysis unavailable"
    
    def _infer_content_types(self, layout: SlideLayout) -> List[str]:
        """Infer what types of content this layout is best suited for."""
        try:
            shapes = list(layout.shapes)
            placeholders = [s for s in shapes if hasattr(s, 'placeholder_format')]
            
            content_types = []
            
            # Analyze placeholder composition
            title_count = len([s for s in placeholders if s.placeholder_format.type in (1, 13)])
            content_count = len([s for s in placeholders if s.placeholder_format.type in (2, 7)])
            image_count = len([s for s in placeholders if s.placeholder_format.type == 18])
            
            if content_count >= 3:
                content_types.extend(["comparison", "detailed_analysis", "multi_point_discussion"])
            elif content_count == 2:
                content_types.extend(["comparison", "before_after", "two_column_data"])
            elif content_count == 1:
                if image_count > 0:
                    content_types.extend(["image_explanation", "process_diagram", "visual_content"])
                else:
                    content_types.extend(["bullet_points", "key_messages", "overview"])
            
            if image_count > 0 and content_count == 0:
                content_types.extend(["image_showcase", "visual_impact", "photo_gallery"])
            
            if title_count > 0 and content_count == 0 and image_count == 0:
                content_types.extend(["section_divider", "title_slide", "chapter_intro"])
                
            return content_types[:5]  # Limit to top 5 recommendations
            
        except Exception as e:
            logger.debug(f"Content type inference failed: {e}")
            return ["general_content"]
    
    def create_layout_descriptions_for_llm(self) -> Dict[str, str]:
        """
        Create detailed descriptions of all layouts for LLM-based template selection.
        
        Returns:
            Dictionary mapping layout names to detailed descriptions
        """
        layout_descriptions = {}
        
        try:
            for i, layout in enumerate(self.prs.slide_layouts):
                # Extract comprehensive visual data
                visual_data = self.extract_complete_layout_visual_data(layout, i)
                
                # Create semantic layout name
                semantic_name = self._classify_layout(layout, i)
                
                # Generate detailed description for LLM
                description = self._create_detailed_layout_description(visual_data, semantic_name)
                
                layout_descriptions[semantic_name] = description
                
                logger.debug(f"Created description for layout {i} ({semantic_name}): {len(description)} chars")
                
        except Exception as e:
            logger.error(f"Failed to create layout descriptions: {e}")
            
        return layout_descriptions
    
    def _create_detailed_layout_description(self, visual_data: Dict, semantic_name: str) -> str:
        """
        Create a detailed, structured description of a layout for LLM analysis.
        
        Args:
            visual_data: Complete visual analysis data
            semantic_name: Semantic name of the layout
            
        Returns:
            Detailed description string
        """
        try:
            if "error" in visual_data:
                return f"Layout {semantic_name}: Analysis failed - {visual_data['error']}"
            
            description_parts = [
                f"LAYOUT: {semantic_name}",
                f"VISUAL SUMMARY: {visual_data.get('visual_summary', 'No summary available')}",
            ]
            
            # Add shape inventory
            shapes = visual_data.get('shapes', {})
            if shapes:
                total_shapes = len(shapes.get('all_shapes', []))
                placeholders = len(shapes.get('placeholders', []))
                visual_elements = len(shapes.get('visual_elements', []))
                
                description_parts.append(
                    f"ELEMENTS: {total_shapes} total shapes ({placeholders} placeholders, {visual_elements} visual elements)"
                )
            
            # Add spatial analysis
            spatial = visual_data.get('spatial_analysis', {})
            if spatial:
                grid = spatial.get('layout_grid', {})
                alignment = spatial.get('alignment_patterns', {})
                
                grid_desc = f"{grid.get('type', 'unknown')} layout"
                if grid.get('columns', 1) > 1 or grid.get('rows', 1) > 1:
                    grid_desc += f" ({grid.get('columns', 1)}x{grid.get('rows', 1)} grid)"
                
                description_parts.append(f"STRUCTURE: {grid_desc}")
                
                if alignment.get('has_center_alignment'):
                    description_parts.append("ALIGNMENT: Center-aligned elements")
                elif alignment.get('vertical_alignments', 0) > 1:
                    description_parts.append("ALIGNMENT: Vertical column structure")
            
            # Add design patterns
            patterns = visual_data.get('design_patterns', {})
            if patterns:
                style = patterns.get('layout_style', 'unknown')
                purpose = patterns.get('primary_purpose', 'general')
                orientation = patterns.get('content_orientation', 'mixed')
                
                description_parts.append(f"STYLE: {style} complexity, {purpose} purpose, {orientation} orientation")
            
            # Add recommended content types
            content_types = visual_data.get('recommended_content_types', [])
            if content_types:
                description_parts.append(f"BEST FOR: {', '.join(content_types[:3])}")
            
            # Add content zones
            zones = visual_data.get('content_zones', {})
            if zones:
                zone_summary = []
                for zone_name, zone_content in zones.items():
                    if zone_content:
                        zone_summary.append(f"{zone_name}({len(zone_content)})")
                
                if zone_summary:
                    description_parts.append(f"ZONES: {', '.join(zone_summary)}")
            
            return " | ".join(description_parts)
            
        except Exception as e:
            logger.error(f"Failed to create detailed description: {e}")
            return f"Layout {semantic_name}: Description generation failed"

    def get_theme_color(self, color_name: str) -> str:
        """
        Get a theme color by name.
        
        Args:
            color_name: Name of the color (primary, secondary, etc.)
            
        Returns:
            Hex color code
        """
        return self.palette.get(color_name, "#000000")

    def get_slide_size(self) -> Tuple[int, int]:
        """
        Get the slide dimensions.
        
        Returns:
            Tuple of (width, height) in EMUs (English Metric Units)
        """
        return (self.prs.slide_width, self.prs.slide_height)

    def get_template_info(self) -> Dict[str, any]:
        """
        Get comprehensive information about the template.
        
        Returns:
            Dictionary with template analysis
        """
        width, height = self.get_slide_size()
        
        info = {
            "template_path": str(self.template_path),
            "total_layouts": len(self.prs.slide_layouts),
            "available_layout_types": self.list_available_layouts(),
            "layout_mapping": self.layout_map.copy(),
            "slide_dimensions": {
                "width": width,
                "height": height,
                "width_inches": width / 914400,  # Convert EMU to inches
                "height_inches": height / 914400
            },
            "theme_colors": self.palette.copy(),
            "template_style": {
                "placeholder_styles_count": len(self.template_style.placeholder_styles),
                "theme_fonts": self.template_style.theme_fonts.copy(),
                "has_master_font": self.template_style.master_font is not None,
                "placeholder_types_with_styles": list(self.template_style.placeholder_styles.keys())
            }
        }
        
        return info

    def get_theme_colors(self) -> Dict[str, str]:
        """
        Extract theme colors from ppt/theme/theme1.xml.
        
        Returns:
            Dictionary mapping color names (dk1, lt1, acc1-6) to hex RGB values
        """
        try:
            with zipfile.ZipFile(self.template_path, 'r') as pptx_zip:
                # Read the theme XML file
                theme_xml = pptx_zip.read('ppt/theme/theme1.xml')
                root = ET.fromstring(theme_xml)
                
                # Define namespaces used in the theme XML
                namespaces = {
                    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
                }
                
                # Extract colors from the color scheme
                colors = {}
                color_scheme = root.find('.//a:clrScheme', namespaces)
                
                if color_scheme is not None:
                    # Map the standard color names to their XML elements
                    color_mappings = {
                        'dk1': 'a:dk1',
                        'lt1': 'a:lt1', 
                        'acc1': 'a:accent1',
                        'acc2': 'a:accent2',
                        'acc3': 'a:accent3',
                        'acc4': 'a:accent4',
                        'acc5': 'a:accent5',
                        'acc6': 'a:accent6'
                    }
                    
                    for color_name, xml_path in color_mappings.items():
                        color_elem = color_scheme.find(xml_path, namespaces)
                        if color_elem is not None:
                            hex_color = self._extract_color_value(color_elem, namespaces)
                            if hex_color:
                                colors[color_name] = hex_color
                
                return colors
                
        except (zipfile.BadZipFile, FileNotFoundError, ET.ParseError, KeyError) as e:
            logger.debug(f"Could not extract theme colors: {e}")
            return {}
    
    def _extract_color_value(self, color_elem: ET.Element, namespaces: Dict[str, str]) -> Optional[str]:
        """
        Extract hex color value from a color element.
        
        Args:
            color_elem: XML element containing color information
            namespaces: XML namespaces dictionary
            
        Returns:
            Hex color string or None if extraction fails
        """
        try:
            # Look for srgbClr (explicit RGB color)
            srgb_color = color_elem.find('.//a:srgbClr', namespaces)
            if srgb_color is not None:
                val = srgb_color.get('val')
                if val and len(val) == 6:
                    return f"#{val.upper()}"
            
            # Look for sysClr (system color) 
            sys_color = color_elem.find('.//a:sysClr', namespaces)
            if sys_color is not None:
                last_clr = sys_color.get('lastClr')
                if last_clr and len(last_clr) == 6:
                    return f"#{last_clr.upper()}"
            
            # Look for prstClr (preset color)
            preset_color = color_elem.find('.//a:prstClr', namespaces)
            if preset_color is not None:
                val = preset_color.get('val')
                # Convert common preset colors to hex
                preset_to_hex = {
                    'black': '#000000',
                    'white': '#FFFFFF',
                    'red': '#FF0000',
                    'green': '#008000',
                    'blue': '#0000FF',
                    'yellow': '#FFFF00',
                    'cyan': '#00FFFF',
                    'magenta': '#FF00FF',
                    'gray': '#808080',
                    'darkBlue': '#000080',
                    'darkGreen': '#008000',
                    'darkRed': '#800000'
                }
                return preset_to_hex.get(val, '#000000')
                
        except Exception as e:
            logger.debug(f"Error extracting color value: {e}")
            
        return None

    def _extract_template_style(self) -> TemplateStyle:
        """
        Extract comprehensive style information from the template.
        
        Returns:
            TemplateStyle object with font and bullet hierarchy information
        """
        try:
            # Extract theme fonts
            theme_fonts = self._extract_theme_fonts()
            
            # Extract master font information
            master_font = self._extract_master_font()
            
            # Extract placeholder-specific styles
            placeholder_styles = self._extract_placeholder_styles()
            
            template_style = TemplateStyle(
                master_font=master_font,
                placeholder_styles=placeholder_styles,
                theme_fonts=theme_fonts
            )
            
            logger.debug(f"Extracted template style with {len(placeholder_styles)} placeholder styles")
            return template_style
            
        except Exception as e:
            logger.warning(f"Failed to extract template style: {e}")
            # Return minimal style object
            return TemplateStyle()

    def _extract_theme_fonts(self) -> Dict[str, str]:
        """
        Extract theme font definitions from the presentation.
        
        Returns:
            Dictionary mapping font roles (major, minor) to font names
        """
        try:
            with zipfile.ZipFile(self.template_path, 'r') as pptx_zip:
                theme_xml = pptx_zip.read('ppt/theme/theme1.xml')
                root = ET.fromstring(theme_xml)
                
                namespaces = {
                    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
                }
                
                fonts = {}
                font_scheme = root.find('.//a:fontScheme', namespaces)
                
                if font_scheme is not None:
                    # Extract major font (typically used for headings)
                    major_font = font_scheme.find('.//a:majorFont/a:latin', namespaces)
                    if major_font is not None:
                        fonts['major'] = major_font.get('typeface', 'Calibri')
                    
                    # Extract minor font (typically used for body text)  
                    minor_font = font_scheme.find('.//a:minorFont/a:latin', namespaces)
                    if minor_font is not None:
                        fonts['minor'] = minor_font.get('typeface', 'Calibri')
                
                return fonts
                
        except Exception as e:
            logger.debug(f"Could not extract theme fonts: {e}")
            return {'major': 'Calibri', 'minor': 'Calibri'}

    def _extract_master_font(self) -> Optional[FontInfo]:
        """
        Extract default font information from the slide master.
        
        Returns:
            FontInfo object or None if extraction fails
        """
        try:
            with zipfile.ZipFile(self.template_path, 'r') as pptx_zip:
                # Try to read slide master XML
                master_xml = pptx_zip.read('ppt/slideMasters/slideMaster1.xml')
                root = ET.fromstring(master_xml)
                
                namespaces = {
                    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
                }
                
                # Look for default font in master text styles
                default_font = root.find('.//a:defRPr/a:latin', namespaces)
                if default_font is not None:
                    font_name = default_font.get('typeface', 'Calibri')
                    
                    # Look for font size in the same element
                    rpr = default_font.getparent()
                    font_size = None
                    if rpr is not None:
                        size_attr = rpr.get('sz')
                        if size_attr:
                            # Font size is in hundredths of a point
                            font_size = int(size_attr) // 100
                    
                    return FontInfo(
                        name=font_name,
                        size=font_size,
                        weight="normal",
                        color="#000000"
                    )
                
        except Exception as e:
            logger.debug(f"Could not extract master font: {e}")
            
        # Return default font if extraction fails
        return FontInfo(
            name="Calibri",
            size=12,
            weight="normal", 
            color="#000000"
        )

    def _extract_placeholder_styles(self) -> Dict[int, PlaceholderStyleInfo]:
        """
        Extract style information for each placeholder type from layouts.
        
        Returns:
            Dictionary mapping placeholder types to their style information
        """
        placeholder_styles = {}
        
        # Process each layout to find placeholder-specific styles
        for layout in self.prs.slide_layouts:
            try:
                layout_styles = self._extract_layout_placeholder_styles(layout)
                
                # Merge layout styles into main dictionary
                for ph_type, style_info in layout_styles.items():
                    if ph_type not in placeholder_styles:
                        placeholder_styles[ph_type] = style_info
                    else:
                        # If we already have this placeholder type, merge bullet styles
                        existing = placeholder_styles[ph_type]
                        for bullet_style in style_info.bullet_styles:
                            if not any(b.indent_level == bullet_style.indent_level for b in existing.bullet_styles):
                                existing.bullet_styles.append(bullet_style)
                        
            except Exception as e:
                logger.debug(f"Failed to extract styles from layout: {e}")
                continue
        
        return placeholder_styles

    def _extract_layout_placeholder_styles(self, layout: SlideLayout) -> Dict[int, PlaceholderStyleInfo]:
        """
        Extract placeholder styles from a specific layout.
        
        Args:
            layout: SlideLayout to analyze
            
        Returns:
            Dictionary mapping placeholder types to style information
        """
        styles = {}
        
        for placeholder in layout.placeholders:
            ph_type = placeholder.placeholder_format.type
            type_name = self._placeholder_type_name(ph_type)
            
            try:
                # Extract font information for this placeholder
                font_info = self._extract_placeholder_font(placeholder)
                
                # Extract font information for this placeholder
                font_info = self._extract_placeholder_font(placeholder)

                # Extract bullet styles for this placeholder
                bullet_styles = self._extract_placeholder_bullets(placeholder)
                
                # Extract fill color for this placeholder
                fill_color = self._extract_placeholder_fill_color(placeholder)

                styles[ph_type] = PlaceholderStyleInfo(
                    placeholder_type=ph_type,
                    type_name=type_name,
                    default_font=font_info,
                    bullet_styles=bullet_styles,
                    fill_color=fill_color
                )
                
            except Exception as e:
                logger.debug(f"Failed to extract style for placeholder type {ph_type}: {e}")
                continue
        
        return styles

    def _extract_placeholder_font(self, placeholder) -> Optional[FontInfo]:
        """
        Extract font information from a placeholder.
        
        Args:
            placeholder: Placeholder shape to analyze
            
        Returns:
            FontInfo object or None if extraction fails
        """
        try:
            # Try to get font information from placeholder's text frame
            if hasattr(placeholder, 'text_frame') and placeholder.text_frame:
                text_frame = placeholder.text_frame
                
                # Look at the first paragraph's font properties
                if len(text_frame.paragraphs) > 0:
                    paragraph = text_frame.paragraphs[0]
                    
                    if hasattr(paragraph, 'font'):
                        font = paragraph.font
                        
                        # Extract font properties
                        font_name = getattr(font, 'name', None) or 'Calibri'
                        font_size = getattr(font, 'size', None)
                        if font_size:
                            font_size = font_size.pt
                        
                        font_bold = getattr(font, 'bold', None)
                        font_weight = "bold" if font_bold else "normal"
                        
                        font_color = "#000000"  # Default color
                        if hasattr(font, 'color') and font.color:
                            try:
                                # Try to get RGB color
                                rgb = font.color.rgb
                                if rgb:
                                    font_color = f"#{rgb}"
                            except:
                                pass
                        
                        return FontInfo(
                            name=font_name,
                            size=int(font_size) if font_size else None,
                            weight=font_weight,
                            color=font_color
                        )
                        
        except Exception as e:
            logger.debug(f"Could not extract placeholder font: {e}")
        
        return None

    def _extract_placeholder_fill_color(self, placeholder) -> Optional[str]:
        """
        Extract fill color from a placeholder shape if it's a solid fill.

        Args:
            placeholder: Placeholder shape to analyze

        Returns:
            Hex color string (e.g., "#RRGGBB") or None if not a solid fill or error.
        """
        try:
            if hasattr(placeholder, 'fill') and placeholder.fill.type == MSO_FILL_TYPE.SOLID:
                if hasattr(placeholder.fill.fore_color, 'rgb'):
                    rgb = placeholder.fill.fore_color.rgb
                    return f"#{str(rgb).upper()}"
        except Exception as e:
            logger.debug(f"Could not extract placeholder fill color: {e}")
        return None

    def _extract_placeholder_bullets(self, placeholder) -> List[BulletInfo]:
        """
        Extract bullet style information from a placeholder.
        
        Args:
            placeholder: Placeholder shape to analyze
            
        Returns:
            List of BulletInfo objects for different indentation levels
        """
        bullet_styles = []
        
        try:
            if hasattr(placeholder, 'text_frame') and placeholder.text_frame:
                text_frame = placeholder.text_frame
                
                # Analyze each paragraph level (up to 5 levels typical in PowerPoint)
                for level in range(5):
                    try:
                        # Try to get bullet character for this level
                        bullet_char = self._get_bullet_character_for_level(text_frame, level)
                        
                        if bullet_char:
                            # Get font for this level
                            level_font = self._get_font_for_level(text_frame, level)
                            
                            bullet_info = BulletInfo(
                                character=bullet_char,
                                font=level_font,
                                indent_level=level
                            )
                            bullet_styles.append(bullet_info)
                            
                    except Exception as e:
                        logger.debug(f"Could not extract bullet for level {level}: {e}")
                        continue
                        
        except Exception as e:
            logger.debug(f"Could not extract bullets from placeholder: {e}")
        
        return bullet_styles

    def _get_bullet_character_for_level(self, text_frame, level: int) -> Optional[str]:
        """
        Get the bullet character for a specific indentation level.
        
        Args:
            text_frame: Text frame to analyze
            level: Indentation level (0-based)
            
        Returns:
            Bullet character string or None
        """
        try:
            # Default bullet characters for each level
            default_bullets = ["•", "○", "▪", "‒", "►"]
            
            # Return default bullet for the level
            if level < len(default_bullets):
                return default_bullets[level]
            else:
                return "•"  # Default to filled circle
                
        except Exception:
            return "•"  # Fallback

    def _get_font_for_level(self, text_frame, level: int) -> Optional[FontInfo]:
        """
        Get font information for a specific bullet level.
        
        Args:
            text_frame: Text frame to analyze  
            level: Indentation level (0-based)
            
        Returns:
            FontInfo object or None
        """
        try:
            # For now, return the same font for all levels
            # This could be enhanced to extract level-specific fonts
            if len(text_frame.paragraphs) > 0:
                paragraph = text_frame.paragraphs[0]
                if hasattr(paragraph, 'font'):
                    font = paragraph.font
                    
                    font_name = getattr(font, 'name', None) or 'Calibri'
                    font_size = getattr(font, 'size', None)
                    if font_size:
                        font_size = font_size.pt
                    
                    return FontInfo(
                        name=font_name,
                        size=int(font_size) if font_size else 14,
                        weight="normal",
                        color="#000000"
                    )
                    
        except Exception:
            pass
        
        # Return default font
        return FontInfo(
            name="Calibri",
            size=14,
            weight="normal", 
            color="#000000"
        )

    def get_template_style(self) -> TemplateStyle:
        """
        Get the extracted template style information.
        
        Returns:
            TemplateStyle object with font and bullet hierarchy
        """
        return self.template_style

    def get_font_for_placeholder_type(self, placeholder_type: int) -> Optional[FontInfo]:
        """
        Get font information for a specific placeholder type.
        
        Args:
            placeholder_type: PowerPoint placeholder type number
            
        Returns:
            FontInfo object or None if not found
        """
        return self.template_style.get_font_for_placeholder_type(placeholder_type)

    def get_bullet_style_for_level(self, placeholder_type: int, level: int = 0) -> Optional[BulletInfo]:
        """
        Get bullet style for a specific placeholder type and indentation level.
        
        Args:
            placeholder_type: PowerPoint placeholder type number
            level: Indentation level (0-based)
            
        Returns:
            BulletInfo object or None if not found
        """
        return self.template_style.get_bullet_style_for_level(placeholder_type, level)

    def get_layout_type_by_id(self, layout_id: int) -> Optional[str]:
        """
        Get the semantic layout type name by its ID/index.

        Args:
            layout_id: The ID/index of the layout.

        Returns:
            The semantic name of the layout type (e.g., "content", "title") or None if not found.
        """
        return self.reverse_layout_map.get(layout_id)

    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """Convert hex color string to RGB tuple."""
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 3: # Expand shorthand hex
            hex_color = "".join([c*2 for c in hex_color])
        if len(hex_color) != 6:
            raise ValueError(f"Invalid hex color format: {hex_color}")
        try:
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        except ValueError:
            raise ValueError(f"Invalid character in hex color: {hex_color}")

    def _relative_luminance(self, rgb_color: Tuple[int, int, int]) -> float:
        """Calculate relative luminance for an RGB color."""
        r, g, b = [x / 255.0 for x in rgb_color]
        r = r / 12.92 if r <= 0.03928 else ((r + 0.055) / 1.055) ** 2.4
        g = g / 12.92 if g <= 0.03928 else ((g + 0.055) / 1.055) ** 2.4
        b = b / 12.92 if b <= 0.03928 else ((b + 0.055) / 1.055) ** 2.4
        return 0.2126 * r + 0.7152 * g + 0.0722 * b

    def _calculate_contrast_ratio(self, color1_hex: str, color2_hex: str) -> float:
        """Calculate contrast ratio between two hex colors."""
        try:
            rgb1 = self._hex_to_rgb(color1_hex)
            rgb2 = self._hex_to_rgb(color2_hex)
        except ValueError as e:
            logger.warning(f"Invalid color format for contrast calculation: {e}")
            return 1.0 # Default to lowest contrast if color is invalid

        lum1 = self._relative_luminance(rgb1)
        lum2 = self._relative_luminance(rgb2)

        if lum1 > lum2:
            return (lum1 + 0.05) / (lum2 + 0.05)
        else:
            return (lum2 + 0.05) / (lum1 + 0.05)

    def check_template_compatibility(self) -> TemplateCompatibilityReport:
        """
        Checks the template for compatibility and best practices.

        Returns:
            A TemplateCompatibilityReport object.
        """
        report = TemplateCompatibilityReport()
        essential_placeholders = ["title", "content", "section", "image"]
        found_placeholders = self.layout_map.keys()

        # Placeholder Check
        for ph_type in essential_placeholders:
            if ph_type not in found_placeholders:
                report.missing_placeholders.append(ph_type)
                report.suggestions.append(f"Consider adding a '{ph_type}' layout for better versatility.")

        if report.missing_placeholders:
            report.issues.append(
                f"Missing essential placeholder types: {', '.join(report.missing_placeholders)}. "
                "This may limit the types of slides that can be automatically generated."
            )
            report.passed_all_checks = False

        # Color Scheme & Contrast Check
        dk1 = self.palette.get("dk1")
        lt1 = self.palette.get("lt1")

        if not dk1:
            report.color_scheme_warnings.append("Theme color 'dk1' (primary dark) is not defined.")
            report.passed_all_checks = False
        if not lt1:
            report.color_scheme_warnings.append("Theme color 'lt1' (primary light) is not defined.")
            report.passed_all_checks = False

        if dk1 and lt1 and dk1.upper() == lt1.upper():
            report.color_scheme_warnings.append("'dk1' and 'lt1' colors are identical, which will cause contrast issues.")
            report.passed_all_checks = False

        if report.color_scheme_warnings:
             report.issues.append("Issues found with the template's color scheme definitions.")


        # Contrast Checks for common pairs
        # Using a simplified threshold of 3:1 for basic check as per requirements
        # WCAG AA typically requires 4.5:1 for normal text.
        contrast_threshold = 3.0
        color_pairs_to_check = [
            ("dk1", "lt1"),
            ("acc1", "lt1"),
            ("dk1", "acc1"), # Check against accent if it's used for text on dark backgrounds
            ("acc2", "lt1"),
            ("acc3", "dk1"), # Example: accent on dark background
        ]

        for c1_name, c2_name in color_pairs_to_check:
            color1 = self.palette.get(c1_name)
            color2 = self.palette.get(c2_name)

            if color1 and color2:
                try:
                    ratio = self._calculate_contrast_ratio(color1, color2)
                    if ratio < contrast_threshold:
                        report.contrast_issues.append(
                            f"Low contrast between '{c1_name}' ({color1}) and '{c2_name}' ({color2}). "
                            f"Ratio: {ratio:.2f}:1. Minimum recommended: {contrast_threshold}:1."
                        )
                        report.passed_all_checks = False
                except ValueError as e: # Handles invalid hex codes passed to ratio calc
                     report.color_scheme_warnings.append(f"Could not calculate contrast for '{c1_name}' and '{c2_name}': {e}")
                     report.passed_all_checks = False # Treat as a failure if colors can't be parsed

            elif not color1:
                report.color_scheme_warnings.append(f"Color '{c1_name}' not found in palette for contrast check.")
            elif not color2:
                report.color_scheme_warnings.append(f"Color '{c2_name}' not found in palette for contrast check.")

        if report.contrast_issues:
            report.issues.append("One or more color pairs have insufficient contrast, potentially affecting readability.")
            report.suggestions.append("Review theme colors to ensure text is clearly readable against backgrounds. Aim for a contrast ratio of at least 3:1 (or 4.5:1 for WCAG AA).")

        # Overall Status
        if report.passed_all_checks:
            report.suggestions.append("Template appears to meet basic compatibility checks.")
        else:
            if not report.issues: # Ensure there's at least one issue if not passing
                 report.issues.append("Template has one or more compatibility issues requiring attention.")
            if not report.suggestions: # Ensure there's at least one suggestion
                report.suggestions.append("Review the reported issues and warnings to improve template compatibility.")

        return report

    def analyze_design_pattern(self) -> DesignPattern:
        """
        Analyzes the template to infer a design pattern.

        Returns:
            A DesignPattern object with inferred characteristics.
        """
        # Defaults
        font_scale_ratio = 1.8  # Typical default (e.g., 36pt title / 20pt body)
        color_complexity_score = 0.5
        layout_density_preference = "medium"
        name = "standard"
        primary_intent = "balanced"

        # --- Font Analysis ---
        title_font_info = self.template_style.get_font_for_placeholder_type(1) # TITLE
        if not title_font_info: # Try CENTERED_TITLE
            title_font_info = self.template_style.get_font_for_placeholder_type(13)

        body_font_info = self.template_style.get_font_for_placeholder_type(2) # BODY
        if not body_font_info: # Try OBJECT
            body_font_info = self.template_style.get_font_for_placeholder_type(7)

        if title_font_info and title_font_info.size and \
           body_font_info and body_font_info.size and body_font_info.size > 0:
            font_scale_ratio = round(title_font_info.size / body_font_info.size, 2)
        elif title_font_info and title_font_info.size:
            # Only title font is available, assume a default body size for ratio estimation
            # This is a heuristic. Common body font sizes are 10-12pt for notes, 18-24pt for content.
            # Let's assume a generic "smaller" body font.
            assumed_body_font_size = max(12.0, title_font_info.size / 2.0) # Ensure it's not excessively small
            font_scale_ratio = round(title_font_info.size / assumed_body_font_size, 2)
        elif body_font_info and body_font_info.size:
             # Only body font is available
            assumed_title_font_size = body_font_info.size * 1.8
            font_scale_ratio = round(assumed_title_font_size / body_font_info.size, 2)
        else:
            logger.debug("Could not determine both title and body font sizes for scale ratio.")
            # font_scale_ratio remains default 1.8

        # --- Color Analysis ---
        distinct_accent_colors = set()
        base_colors = {self.palette.get("dk1", "").upper(), self.palette.get("lt1", "").upper()}

        for i in range(1, 7): # acc1 to acc6
            accent_color = self.palette.get(f"acc{i}")
            if accent_color:
                accent_color_upper = accent_color.upper()
                # Check if it's different from base colors and other accents already counted
                # A simple threshold for "significantly different" might be needed if colors are very close
                # For now, exact match check.
                if accent_color_upper not in base_colors and accent_color_upper not in distinct_accent_colors:
                    # Rudimentary check for very similar colors (e.g. #FEFEFE vs #FFFFFF)
                    # This is simplistic; a proper color difference metric (delta E) would be better.
                    is_very_similar_to_base = False
                    for base_c in base_colors:
                        if base_c and len(base_c) == 7 and len(accent_color_upper) == 7:
                            # Count differing hex characters (simple difference)
                            diff = sum(1 for c1, c2 in zip(base_c[1:], accent_color_upper[1:]) if c1 != c2)
                            if diff <= 1: # Allow only 1 char difference (e.g. #F0F0F0 vs #F1F0F0)
                                is_very_similar_to_base = True
                                break
                    if not is_very_similar_to_base:
                         distinct_accent_colors.add(accent_color_upper)

        num_distinct_accents = len(distinct_accent_colors)
        if num_distinct_accents <= 1:
            color_complexity_score = 0.2  # Low
        elif num_distinct_accents <= 3:
            color_complexity_score = 0.6  # Medium
        else:
            color_complexity_score = 0.9  # High

        # --- Layout Density Analysis ---
        dense_layout_keywords = ["two_column", "image_content"] # Add more if other complex types are classified
        dense_layout_count = 0
        simple_layout_count = 0

        for layout_name in self.layout_map.keys():
            if any(keyword in layout_name for keyword in dense_layout_keywords):
                dense_layout_count += 1
            elif "content" in layout_name or "title" in layout_name or "section" in layout_name or "blank" in layout_name:
                # Only count if it's not also a dense layout (e.g. "image_content" is not simple)
                if not any(keyword in layout_name for keyword in dense_layout_keywords):
                    simple_layout_count +=1

        total_layouts = len(self.layout_map)
        if total_layouts > 0:
            dense_ratio = dense_layout_count / total_layouts
            if dense_ratio >= 0.4: # At least 40% of layouts are "dense"
                layout_density_preference = "high"
            elif dense_ratio <= 0.15 and simple_layout_count / total_layouts >= 0.5 : # Few dense, mostly simple
                layout_density_preference = "low"
            else:
                layout_density_preference = "medium"
        else:
            layout_density_preference = "medium" # Default if no layouts

        # --- Derive Name and Primary Intent ---
        # These rules are heuristics and can be expanded.
        if font_scale_ratio < 1.5 and color_complexity_score < 0.35 and layout_density_preference == "low":
            name = "minimalist"
            primary_intent = "readability"
        elif color_complexity_score >= 0.7 and font_scale_ratio >= 1.8:
            name = "vibrant"
            primary_intent = "visual_impact"
        elif layout_density_preference == "high" and font_scale_ratio < 2.0 : # Data-heavy might prefer slightly smaller titles relative to content
            name = "data-heavy"
            primary_intent = "information_density"
        elif font_scale_ratio >= 2.2 and color_complexity_score >= 0.5:
            name = "bold & colorful"
            primary_intent = "visual_impact"
        elif font_scale_ratio <= 1.3 and color_complexity_score <= 0.3 and layout_density_preference != "high":
            name = "subtle & clean"
            primary_intent = "readability"
        else: # Default "standard" if no strong indicators
            name = "standard"
            if font_scale_ratio > 2.0: # Standard but with large titles
                primary_intent = "strong_hierarchy"
            elif color_complexity_score < 0.4:
                primary_intent = "clarity"
            else:
                primary_intent = "balanced"


        return DesignPattern(
            name=name,
            font_scale_ratio=font_scale_ratio,
            color_complexity_score=color_complexity_score,
            layout_density_preference=layout_density_preference,
            primary_intent=primary_intent
        )