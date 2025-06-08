import unittest
from unittest.mock import MagicMock, patch
import logging
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from typing import Optional, List, Dict

from open_lilli.slide_assembler import SlideAssembler
from open_lilli.models import (
    StyleValidationConfig, SlidePlan, Outline, QualityGates, ReviewFeedback, DesignIssueType
)
from open_lilli.template_parser import TemplateParser
# Need to import VisualProofreader and DesignIssue for mocking and type hints in new tests
from open_lilli.visual_proofreader import VisualProofreader, DesignIssue, ProofreadingResult
from openai import OpenAI # Required for patching


class TestSlideAssemblerPlaceholderPopulation(unittest.TestCase):
    def setUp(self):
        self.mock_template_parser = MagicMock(spec=TemplateParser)
        self.mock_template_parser.template_style = MagicMock()
        self.mock_template_parser.template_style.language_specific_fonts = {}
        self.mock_prs = MagicMock(spec=Presentation)
        self.mock_prs.slide_layouts = [MagicMock()]
        self.mock_prs.slides = []
        self.patcher = patch('pptx.Presentation', return_value=self.mock_prs)
        self.mock_presentation_cls = self.patcher.start()

    def tearDown(self):
        self.patcher.stop()

    def _create_mock_slide(self, placeholders_config, slide_width=Inches(10), slide_height=Inches(7.5)):
        slide = MagicMock()
        mock_placeholders = []
        for i, (ph_type_enum, text_content, shape_type_enum) in enumerate(placeholders_config):
            placeholder = MagicMock()
            placeholder.placeholder_format = MagicMock()
            placeholder.placeholder_format.type = ph_type_enum
            placeholder.placeholder_format.idx = i
            placeholder.name = f"Placeholder_{ph_type_enum.name}_{i}"
            placeholder.shape_id = i + 100
            placeholder.width = Inches(5)
            placeholder.height = Inches(3)
            placeholder.left = Inches(1)
            placeholder.top = Inches(1)

            if shape_type_enum == MSO_SHAPE_TYPE.PICTURE:
                placeholder.shape_type = MSO_SHAPE_TYPE.PICTURE
                placeholder.text_frame = None
                if text_content:
                    placeholder.image = MagicMock()
                else:
                    placeholder.image = None
            else:
                placeholder.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
                placeholder.text_frame = MagicMock(spec=True)
                placeholder.text_frame.text = text_content if text_content else ""
                placeholder.text_frame.margin_top = Inches(0.1)
                placeholder.text_frame.margin_bottom = Inches(0.1)
                placeholder.text_frame.margin_left = Inches(0.1)
                placeholder.text_frame.margin_right = Inches(0.1)
                placeholder.text_frame.auto_size = MSO_AUTO_SIZE.NONE
                mock_paragraph = MagicMock()
                mock_run = MagicMock()
                mock_font = MagicMock()
                mock_font.size = Pt(18)
                mock_run.font = mock_font
                mock_paragraph.runs = [mock_run]
                placeholder.text_frame.paragraphs = [mock_paragraph]
            mock_placeholders.append(placeholder)

        slide.placeholders = mock_placeholders
        slide.shapes = mock_placeholders
        mock_presentation_part = MagicMock()
        mock_presentation_part.slide_width = slide_width
        mock_presentation_part.slide_height = slide_height
        slide.parent = mock_presentation_part
        return slide

    def test_valid_placeholder_population(self):
        validation_config = StyleValidationConfig(check_placeholder_population=True)
        assembler = SlideAssembler(self.mock_template_parser, validation_config)
        mock_slide = self._create_mock_slide([
            (PP_PLACEHOLDER.TITLE, "Valid Title", MSO_SHAPE_TYPE.TEXT_BOX),
            (PP_PLACEHOLDER.BODY, "Valid Body Content", MSO_SHAPE_TYPE.TEXT_BOX)
        ])
        self.mock_prs.slides.append(mock_slide)
        feedback_items = assembler._check_slide_placeholder_population(mock_slide, slide_index=0)
        self.assertEqual(len(feedback_items), 0, "Should be no violations for validly populated slide.")

    def test_valid_placeholder_population_with_picture_in_body(self):
        validation_config = StyleValidationConfig(check_placeholder_population=True)
        assembler = SlideAssembler(self.mock_template_parser, validation_config)
        mock_slide = self._create_mock_slide([
            (PP_PLACEHOLDER.TITLE, "Valid Title", MSO_SHAPE_TYPE.TEXT_BOX),
            (PP_PLACEHOLDER.BODY, "Picture is present", MSO_SHAPE_TYPE.PICTURE)
        ])
        self.mock_prs.slides.append(mock_slide)
        feedback_items = assembler._check_slide_placeholder_population(mock_slide, slide_index=0)
        self.assertEqual(len(feedback_items), 0, "Should be no violations when body has a picture.")

    def test_missing_title_placeholder(self):
        validation_config = StyleValidationConfig(check_placeholder_population=True)
        assembler = SlideAssembler(self.mock_template_parser, validation_config)
        mock_slide = self._create_mock_slide([
            (PP_PLACEHOLDER.BODY, "Body Content Only", MSO_SHAPE_TYPE.TEXT_BOX)
        ])
        self.mock_prs.slides.append(mock_slide)
        feedback_items = assembler._check_slide_placeholder_population(mock_slide, slide_index=0)
        self.assertEqual(len(feedback_items), 1, "Should be one violation for missing title.")
        self.assertEqual(feedback_items[0].category, 'placeholder_population')
        self.assertIn("'Title' placeholder is empty or missing", feedback_items[0].message)

    def test_empty_title_placeholder(self):
        validation_config = StyleValidationConfig(check_placeholder_population=True)
        assembler = SlideAssembler(self.mock_template_parser, validation_config)
        mock_slide = self._create_mock_slide([
            (PP_PLACEHOLDER.TITLE, "", MSO_SHAPE_TYPE.TEXT_BOX),
            (PP_PLACEHOLDER.BODY, "Body Content", MSO_SHAPE_TYPE.TEXT_BOX)
        ])
        self.mock_prs.slides.append(mock_slide)
        feedback_items = assembler._check_slide_placeholder_population(mock_slide, slide_index=0)
        self.assertEqual(len(feedback_items), 1, "Should be one violation for empty title.")
        self.assertEqual(feedback_items[0].category, 'placeholder_population')
        self.assertIn("'Title' placeholder is empty or missing", feedback_items[0].message)

    def test_missing_body_placeholder(self):
        validation_config = StyleValidationConfig(check_placeholder_population=True)
        assembler = SlideAssembler(self.mock_template_parser, validation_config)
        mock_slide = self._create_mock_slide([
            (PP_PLACEHOLDER.TITLE, "Title Only Content", MSO_SHAPE_TYPE.TEXT_BOX)
        ])
        self.mock_prs.slides.append(mock_slide)
        feedback_items = assembler._check_slide_placeholder_population(mock_slide, slide_index=0)
        self.assertEqual(len(feedback_items), 0, "Should be zero violations if body placeholder is completely missing (current behavior).")

    def test_empty_body_placeholder(self):
        validation_config = StyleValidationConfig(check_placeholder_population=True)
        assembler = SlideAssembler(self.mock_template_parser, validation_config)
        mock_slide = self._create_mock_slide([
            (PP_PLACEHOLDER.TITLE, "Title Content", MSO_SHAPE_TYPE.TEXT_BOX),
            (PP_PLACEHOLDER.BODY, "", MSO_SHAPE_TYPE.TEXT_BOX)
        ])
        self.mock_prs.slides.append(mock_slide)
        feedback_items = assembler._check_slide_placeholder_population(mock_slide, slide_index=0)
        self.assertEqual(len(feedback_items), 1, "Should be one violation for empty body.")
        self.assertEqual(feedback_items[0].category, 'placeholder_population')
        self.assertIn("'Body Content' placeholder is present but empty", feedback_items[0].message)

    def test_empty_body_placeholder_picture(self):
        validation_config = StyleValidationConfig(check_placeholder_population=True)
        assembler = SlideAssembler(self.mock_template_parser, validation_config)
        mock_slide = self._create_mock_slide([
            (PP_PLACEHOLDER.TITLE, "Title Content", MSO_SHAPE_TYPE.TEXT_BOX),
            (PP_PLACEHOLDER.BODY, None, MSO_SHAPE_TYPE.PICTURE)
        ])
        self.mock_prs.slides.append(mock_slide)
        feedback_items = assembler._check_slide_placeholder_population(mock_slide, slide_index=0)
        self.assertEqual(len(feedback_items), 1, "Should be one violation for an empty picture body placeholder.")
        self.assertEqual(feedback_items[0].category, 'placeholder_population')
        self.assertIn("'Body Content' placeholder is present but empty", feedback_items[0].message)

    def test_placeholder_check_disabled(self):
        validation_config = StyleValidationConfig(check_placeholder_population=False)
        assembler = SlideAssembler(self.mock_template_parser, validation_config)
        mock_slide = self._create_mock_slide([
            (PP_PLACEHOLDER.BODY, "", MSO_SHAPE_TYPE.TEXT_BOX)
        ])
        self.mock_prs.slides.append(mock_slide)
        feedback_items = assembler._check_slide_placeholder_population(mock_slide, slide_index=0)
        self.assertEqual(len(feedback_items), 0, "Should be no violations when check is disabled.")

    def test_both_title_and_body_empty(self):
        validation_config = StyleValidationConfig(check_placeholder_population=True)
        assembler = SlideAssembler(self.mock_template_parser, validation_config)
        mock_slide = self._create_mock_slide([
            (PP_PLACEHOLDER.TITLE, "", MSO_SHAPE_TYPE.TEXT_BOX),
            (PP_PLACEHOLDER.BODY, "", MSO_SHAPE_TYPE.TEXT_BOX)
        ])
        self.mock_prs.slides.append(mock_slide)
        feedback_items = assembler._check_slide_placeholder_population(mock_slide, slide_index=0)
        self.assertEqual(len(feedback_items), 2, "Should be two violations, one for empty title and one for empty body.")
        title_feedback_found = any(
            fb.category == 'placeholder_population' and "'Title' placeholder is empty or missing" in fb.message
            for fb in feedback_items
        )
        self.assertTrue(title_feedback_found, "Missing title feedback not found.")
        body_feedback_found = any(
            fb.category == 'placeholder_population' and "'Body Content' placeholder is present but empty" in fb.message
            for fb in feedback_items
        )
        self.assertTrue(body_feedback_found, "Empty body feedback not found.")


class TestSlideAssemblerTextOverflow(unittest.TestCase):
    def setUp(self):
        self.mock_template_parser = MagicMock(spec=TemplateParser)
        self.mock_template_parser.template_style = MagicMock()
        self.mock_template_parser.template_style.language_specific_fonts = {}
        self.base_validation_config = StyleValidationConfig(
            check_text_overflow=True,
            autofix_text_overflow=False
        )

    def _create_mock_shape(self, text_content, shape_height_emu, initial_font_size_pt=18, is_title=False, text_frame_margins_emu=0):
        shape = MagicMock(spec=True)
        shape.name = "TestShape"
        shape.shape_id = 123
        shape.width = Inches(5)
        shape.height = shape_height_emu
        shape.text_frame = MagicMock(spec=True)
        shape.text_frame.text = text_content
        shape.text_frame.margin_top = text_frame_margins_emu
        shape.text_frame.margin_bottom = text_frame_margins_emu
        shape.text_frame.margin_left = text_frame_margins_emu
        shape.text_frame.margin_right = text_frame_margins_emu
        shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        shape.text_frame.clear = MagicMock()
        mock_paragraph = MagicMock(spec=True)
        mock_run = MagicMock(spec=True)
        mock_font = MagicMock(spec=True)
        current_font_size_pt_obj = MagicMock(spec=Pt)
        current_font_size_pt_obj.pt = initial_font_size_pt
        mock_font.size = current_font_size_pt_obj
        mock_font.name = "Test Font"
        mock_font.bold = False
        mock_font.italic = False
        mock_font.color = MagicMock(spec=True)
        mock_font.color.rgb = RGBColor(0,0,0)
        mock_run.font = mock_font
        mock_run.text = text_content
        mock_paragraph.runs = [mock_run]
        mock_paragraph.text = text_content
        new_mock_run_after_clear = MagicMock(spec=True)
        new_mock_run_after_clear.font = mock_font
        mock_paragraph.add_run = MagicMock(return_value=new_mock_run_after_clear)
        shape.text_frame.paragraphs = [mock_paragraph]
        shape.placeholder_format = MagicMock(spec=True)
        shape.placeholder_format.idx = 0
        if is_title:
            shape.placeholder_format.type = PP_PLACEHOLDER.TITLE
        else:
            shape.placeholder_format.type = PP_PLACEHOLDER.BODY
        return shape

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_no_overflow(self, mock_logger):
        assembler = SlideAssembler(self.mock_template_parser, self.base_validation_config)
        shape = self._create_mock_shape("Short text", shape_height_emu=Inches(2), initial_font_size_pt=18)
        with patch.object(assembler, '_estimate_text_lines', return_value=2) as mock_estimate:
            feedback_items = assembler._check_text_frame_overflow(shape, slide_index=0, language="en")
            self.assertEqual(len(feedback_items), 0, "Should be no violations if text fits.")
            mock_estimate.assert_called_once()

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_overflow_detected_autofix_disabled(self, mock_logger):
        config = StyleValidationConfig(check_text_overflow=True, autofix_text_overflow=False)
        assembler = SlideAssembler(self.mock_template_parser, config)
        shape = self._create_mock_shape("This is a very long line of text that will surely overflow the small shape.",
                                        shape_height_emu=Inches(0.5), initial_font_size_pt=18)
        with patch.object(assembler, '_estimate_text_lines', return_value=10) as mock_estimate:
            feedback_items = assembler._check_text_frame_overflow(shape, slide_index=0, language="en")
            self.assertEqual(len(feedback_items), 1, "Should be one violation for overflow.")
            self.assertEqual(feedback_items[0].category, 'text_overflow')
            mock_estimate.assert_called_once()
            self.assertFalse(any("auto-fixed" in call.args[0].lower() for call in mock_logger.info.call_args_list))

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_overflow_check_disabled(self, mock_logger):
        config = StyleValidationConfig(check_text_overflow=False)
        assembler = SlideAssembler(self.mock_template_parser, config)
        shape = self._create_mock_shape("Overflowing text", shape_height_emu=Inches(0.1))
        feedback_items = assembler._check_text_frame_overflow(shape, slide_index=0, language="en")
        self.assertEqual(len(feedback_items), 0, "Should be no violations if check is disabled.")

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_overflow_autofix_font_reduction_succeeds(self, mock_logger):
        config = StyleValidationConfig(check_text_overflow=True, autofix_text_overflow=True)
        assembler = SlideAssembler(self.mock_template_parser, config)
        shape = self._create_mock_shape("Long text that needs shrinking", shape_height_emu=Inches(1), initial_font_size_pt=18)
        with patch.object(assembler, '_estimate_text_lines', side_effect=[10, 10, 2]) as mock_estimate_lines_method:
            feedback_items = assembler._check_text_frame_overflow(shape, slide_index=0, language="en")
            self.assertEqual(len(feedback_items), 0, "Should be no violations if autofix (font shrink) succeeds.")
            self.assertEqual(mock_estimate_lines_method.call_count, 3)
            self.assertEqual(shape.text_frame.paragraphs[0].runs[0].font.size.pt, 17)
            mock_logger.info.assert_any_call(f"Overflow fixed for shape '{shape.name}' on slide 0 by font size reduction (attempt 2).")

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_overflow_autofix_font_reduction_fails_min_size_reached(self, mock_logger):
        config = StyleValidationConfig(check_text_overflow=True, autofix_text_overflow=True)
        assembler = SlideAssembler(self.mock_template_parser, config)
        initial_font_pt = 11
        shape = self._create_mock_shape("Very very long text", shape_height_emu=Inches(0.2), initial_font_size_pt=initial_font_pt)
        mock_estimate_lines = MagicMock(return_value=20)
        with patch.object(assembler, '_estimate_text_lines', mock_estimate_lines):
            feedback_items = assembler._check_text_frame_overflow(shape, slide_index=0, language="en")
            self.assertEqual(len(feedback_items), 1, "Should be one violation if autofix (font shrink) fails.")
            self.assertEqual(feedback_items[0].category, 'text_overflow')
            self.assertEqual(shape.text_frame.paragraphs[0].runs[0].font.size.pt, 10)
            self.assertTrue(mock_estimate_lines.call_count >= 1)
            mock_logger.warning.assert_any_call(f"Failed to autofix text overflow for shape '{shape.name}' on slide 0.")

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_overflow_autofix_title_truncation_succeeds(self, mock_logger):
        config = StyleValidationConfig(check_text_overflow=True, autofix_text_overflow=True)
        assembler = SlideAssembler(self.mock_template_parser, config)
        assembler.max_title_length = 20
        long_title = "This is a very long title that will certainly be truncated"
        shape = self._create_mock_shape(long_title, shape_height_emu=Inches(0.2), initial_font_size_pt=12, is_title=True)
        shape.text_frame.paragraphs[0].runs[0].font.size.pt = 10
        mock_estimate_lines = MagicMock(return_value=20)
        with patch.object(assembler, '_estimate_text_lines', mock_estimate_lines) :
            feedback_items = assembler._check_text_frame_overflow(shape, slide_index=0, language="en")
            self.assertEqual(len(feedback_items), 0, "Should be no violations if title truncation succeeds.")
            max_len = assembler.max_title_length
            original_text_for_calc = long_title
            temp_safe_truncate_pos = (original_text_for_calc[:max_len - 3]).rfind(' ')
            if temp_safe_truncate_pos == -1 or temp_safe_truncate_pos < max_len // 2:
                expected_truncated_text = original_text_for_calc[:max_len - 3] + "..."
            else:
                expected_truncated_text = original_text_for_calc[:temp_safe_truncate_pos] + "..."
            self.assertEqual(shape.text_frame.paragraphs[0].add_run.return_value.text, expected_truncated_text)
            mock_logger.info.assert_any_call(f"Overflow in title shape '{shape.name}' fixed by truncation.")

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_overflow_autofix_non_title_no_truncation_font_fail(self, mock_logger):
        config = StyleValidationConfig(check_text_overflow=True, autofix_text_overflow=True)
        assembler = SlideAssembler(self.mock_template_parser, config)
        long_text = "This is a very long non-title text that overflows"
        shape = self._create_mock_shape(long_text, shape_height_emu=Inches(0.2), initial_font_size_pt=10)
        shape.text_frame.paragraphs[0].runs[0].font.size.pt = 10
        mock_estimate_lines = MagicMock(return_value=20)
        with patch.object(assembler, '_estimate_text_lines', mock_estimate_lines):
            feedback_items = assembler._check_text_frame_overflow(shape, slide_index=0, language="en")
            self.assertEqual(len(feedback_items), 1, "Should be one violation as non-title won't truncate and font fix fails.")
            self.assertEqual(feedback_items[0].category, 'text_overflow')
            self.assertEqual(shape.text_frame.paragraphs[0].runs[0].text, long_text)
            mock_logger.warning.assert_any_call(f"Failed to autofix text overflow for shape '{shape.name}' on slide 0.")


class TestSlideAssemblerShapeAlignment(unittest.TestCase):
    def setUp(self):
        self.mock_template_parser = MagicMock(spec=TemplateParser)
        self.slide_width_emu = Inches(10)
        self.slide_height_emu = Inches(7.5)
        self.default_margins = {"top": 0.5, "bottom": 0.5, "left": 0.5, "right": 0.5}
        self.assembler = None

    def _create_test_slide_and_shape(self, shape_left_in, shape_top_in, shape_width_in, shape_height_in,
                                     margins_in=None):
        if margins_in is None:
            margins_in = self.default_margins
        mock_slide = MagicMock(spec=True)
        mock_slide_parent = MagicMock(spec=True)
        mock_slide_parent.slide_width = self.slide_width_emu
        mock_slide_parent.slide_height = self.slide_height_emu
        mock_slide.parent = mock_slide_parent
        shape = MagicMock(spec=True)
        shape.left = Inches(shape_left_in)
        shape.top = Inches(shape_top_in)
        shape.width = Inches(shape_width_in)
        shape.height = Inches(shape_height_in)
        shape.name = "TestShape"
        shape.shape_id = 777
        shape.placeholder_format = MagicMock(spec=True)
        shape.placeholder_format.type = PP_PLACEHOLDER.OBJECT
        shape.placeholder_format.idx = 0
        mock_slide.shapes = [shape]
        return mock_slide, shape

    def _get_validation_config(self, check_alignment=True, autofix_alignment=False, margins_in=None):
        if margins_in is None:
            margins_in = self.default_margins
        actual_quality_gates = QualityGates()
        actual_quality_gates.slide_margin_top_inches = margins_in["top"]
        actual_quality_gates.slide_margin_bottom_inches = margins_in["bottom"]
        actual_quality_gates.slide_margin_left_inches = margins_in["left"]
        actual_quality_gates.slide_margin_right_inches = margins_in["right"]
        actual_quality_gates.enable_alignment_check = True
        config = StyleValidationConfig(
            check_alignment=check_alignment,
            autofix_alignment=autofix_alignment,
            quality_gates_config=actual_quality_gates.model_dump()
        )
        return config

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_shape_correctly_aligned(self, mock_logger):
        config = self._get_validation_config(check_alignment=True, autofix_alignment=False)
        assembler = SlideAssembler(self.mock_template_parser, config)
        slide, shape = self._create_test_slide_and_shape(1, 1, 2, 2)
        feedback_items = assembler._check_shape_alignment_against_margins(shape, slide, 0, "en")
        self.assertEqual(len(feedback_items), 0, "No violations for correctly aligned shape.")

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_alignment_check_disabled(self, mock_logger):
        config = self._get_validation_config(check_alignment=False)
        assembler = SlideAssembler(self.mock_template_parser, config)
        slide, shape = self._create_test_slide_and_shape(0, 0, 2, 2)
        feedback_items = assembler._check_shape_alignment_against_margins(shape, slide, 0, "en")
        self.assertEqual(len(feedback_items), 0, "No violations when alignment check is disabled.")

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_left_margin_violation_detected_no_autofix(self, mock_logger):
        config = self._get_validation_config(check_alignment=True, autofix_alignment=False)
        assembler = SlideAssembler(self.mock_template_parser, config)
        slide, shape = self._create_test_slide_and_shape(0.1, 1, 2, 2)
        original_left = shape.left
        feedback_items = assembler._check_shape_alignment_against_margins(shape, slide, 0, "en")
        self.assertEqual(len(feedback_items), 1)
        self.assertEqual(feedback_items[0].category, 'alignment')
        self.assertIn("extends beyond the left slide margin", feedback_items[0].message)
        self.assertEqual(shape.left, original_left, "Shape should not be moved if autofix is off.")

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_top_margin_violation_detected_no_autofix(self, mock_logger):
        config = self._get_validation_config(check_alignment=True, autofix_alignment=False)
        assembler = SlideAssembler(self.mock_template_parser, config)
        slide, shape = self._create_test_slide_and_shape(1, 0.1, 2, 2)
        original_top = shape.top
        feedback_items = assembler._check_shape_alignment_against_margins(shape, slide, 0, "en")
        self.assertEqual(len(feedback_items), 1)
        self.assertEqual(feedback_items[0].category, 'alignment')
        self.assertIn("extends beyond the top slide margin", feedback_items[0].message)
        self.assertEqual(shape.top, original_top, "Shape should not be moved if autofix is off.")

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_right_margin_violation_detected_no_autofix(self, mock_logger):
        config = self._get_validation_config(check_alignment=True, autofix_alignment=False)
        assembler = SlideAssembler(self.mock_template_parser, config)
        slide, shape = self._create_test_slide_and_shape(8, 1, 2, 2)
        original_width = shape.width
        original_left = shape.left
        feedback_items = assembler._check_shape_alignment_against_margins(shape, slide, 0, "en")
        self.assertEqual(len(feedback_items), 1)
        self.assertEqual(feedback_items[0].category, 'alignment')
        self.assertIn("extends beyond the right slide margin", feedback_items[0].message)
        self.assertEqual(shape.width, original_width, "Shape should not be resized if autofix is off.")
        self.assertEqual(shape.left, original_left, "Shape should not be moved if autofix is off.")

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_bottom_margin_violation_detected_no_autofix(self, mock_logger):
        config = self._get_validation_config(check_alignment=True, autofix_alignment=False)
        assembler = SlideAssembler(self.mock_template_parser, config)
        slide, shape = self._create_test_slide_and_shape(1, 6, 2, 1.5)
        original_height = shape.height
        original_top = shape.top
        feedback_items = assembler._check_shape_alignment_against_margins(shape, slide, 0, "en")
        self.assertEqual(len(feedback_items), 1)
        self.assertEqual(feedback_items[0].category, 'alignment')
        self.assertIn("extends beyond the bottom slide margin", feedback_items[0].message)
        self.assertEqual(shape.height, original_height, "Shape should not be resized if autofix is off.")
        self.assertEqual(shape.top, original_top, "Shape should not be moved if autofix is off.")

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_left_margin_violation_autofix_move(self, mock_logger):
        config = self._get_validation_config(check_alignment=True, autofix_alignment=True)
        assembler = SlideAssembler(self.mock_template_parser, config)
        slide, shape = self._create_test_slide_and_shape(0.1, 1, 2, 2)
        feedback_items = assembler._check_shape_alignment_against_margins(shape, slide, 0, "en")
        self.assertEqual(len(feedback_items), 0, "Violation should be fixed.")
        self.assertEqual(shape.left, Inches(self.default_margins["left"]))
        mock_logger.info.assert_any_call(f"Autofix: Moved shape '{assembler._get_placeholder_type_name(shape.placeholder_format.type.value)} (Idx: {shape.placeholder_format.idx}, Name: {shape.name})' on slide 0 from L:{Inches(0.1):.0f} to L:{Inches(0.5):.0f} (due to left margin).")

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_right_margin_violation_autofix_move(self, mock_logger):
        config = self._get_validation_config(check_alignment=True, autofix_alignment=True)
        assembler = SlideAssembler(self.mock_template_parser, config)
        slide, shape = self._create_test_slide_and_shape(8, 1, 2, 2)
        feedback_items = assembler._check_shape_alignment_against_margins(shape, slide, 0, "en")
        self.assertEqual(len(feedback_items), 0, "Violation should be fixed by moving left.")
        self.assertEqual(shape.left, Inches(10 - self.default_margins["right"] - 2))
        mock_logger.info.assert_any_call(f"Autofix: Moved shape '{assembler._get_placeholder_type_name(shape.placeholder_format.type.value)} (Idx: {shape.placeholder_format.idx}, Name: {shape.name})' on slide 0 left to fit right margin. New L:{shape.left:.0f}, R_Edge:{(shape.left + shape.width):.0f}.")

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_right_margin_violation_autofix_resize(self, mock_logger):
        config = self._get_validation_config(check_alignment=True, autofix_alignment=True)
        assembler = SlideAssembler(self.mock_template_parser, config)
        slide, shape = self._create_test_slide_and_shape(0.5, 1, 9.5, 2)
        feedback_items = assembler._check_shape_alignment_against_margins(shape, slide, 0, "en")
        self.assertEqual(len(feedback_items), 0, "Violation should be fixed by resizing width.")
        self.assertEqual(shape.width, Inches(10 - self.default_margins["left"] - self.default_margins["right"]))
        self.assertEqual(shape.left, Inches(self.default_margins["left"]))
        mock_logger.info.assert_any_call(f"Autofix: Resized width of shape '{assembler._get_placeholder_type_name(shape.placeholder_format.type.value)} (Idx: {shape.placeholder_format.idx}, Name: {shape.name})' on slide 0 to fit right margin. New W:{shape.width:.0f}, R_Edge:{(shape.left + shape.width):.0f}.")

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_bottom_margin_violation_autofix_resize(self, mock_logger):
        config = self._get_validation_config(check_alignment=True, autofix_alignment=True)
        assembler = SlideAssembler(self.mock_template_parser, config)
        slide, shape = self._create_test_slide_and_shape(1, 0.5, 2, 7.0)
        feedback_items = assembler._check_shape_alignment_against_margins(shape, slide, 0, "en")
        self.assertEqual(len(feedback_items), 0, "Violation should be fixed by resizing height.")
        self.assertEqual(shape.height, Inches(7.5 - self.default_margins["top"] - self.default_margins["bottom"]))
        self.assertEqual(shape.top, Inches(self.default_margins["top"]))
        mock_logger.info.assert_any_call(f"Autofix: Resized height of shape '{assembler._get_placeholder_type_name(shape.placeholder_format.type.value)} (Idx: {shape.placeholder_format.idx}, Name: {shape.name})' on slide 0 to fit bottom margin. New H:{shape.height:.0f}, B_Edge:{(shape.top + shape.height):.0f}.")

    @patch('open_lilli.slide_assembler.logger', autospec=True)
    def test_right_margin_autofix_resize_fails_too_small(self, mock_logger):
        custom_margins = self.default_margins.copy()
        custom_margins["right"] = 10.0 - 0.59
        config = self._get_validation_config(check_alignment=True, autofix_alignment=True, margins_in=custom_margins)
        assembler = SlideAssembler(self.mock_template_parser, config)
        slide, shape = self._create_test_slide_and_shape(0.5, 1, Inches(0.11), 2)
        slide.parent.slide_width = Inches(10)
        feedback_items = assembler._check_shape_alignment_against_margins(shape, slide, 0, "en")
        self.assertEqual(len(feedback_items), 1, "Violation should persist if resize makes shape too small.")
        self.assertEqual(feedback_items[0].category, 'alignment')
        self.assertIn("extends beyond the right slide margin", feedback_items[0].message)
        self.assertNotIn("after autofix", feedback_items[0].message)
        mock_logger.warning.assert_any_call(f"Autofix: Failed to fit shape '{assembler._get_placeholder_type_name(shape.placeholder_format.type.value)} (Idx: {shape.placeholder_format.idx}, Name: {shape.name})' on slide 0 within right margin without making width too small.")


class TestSlideAssemblerAccessibility(unittest.TestCase):
    def setUp(self):
        self.mock_template_parser = MagicMock(spec=TemplateParser)
        self.base_validation_config = StyleValidationConfig(check_alt_text_accessibility=True)
        self.assembler = SlideAssembler(self.mock_template_parser, self.base_validation_config)

    def _create_mock_image_shape(self, name_alt_text: Optional[str], shape_id=888):
        shape = MagicMock(spec=True)
        shape.name = name_alt_text
        shape.shape_id = shape_id
        shape.shape_type = MSO_SHAPE_TYPE.PICTURE
        shape.placeholder_format = MagicMock(spec=True)
        shape.placeholder_format.type = PP_PLACEHOLDER.PICTURE
        shape.placeholder_format.idx = 1
        return shape

    def test_image_with_good_alt_text(self):
        shape = self._create_mock_image_shape("A majestic mountain range at sunset, with vibrant orange and purple hues.")
        feedback_items = self.assembler._check_image_accessibility(shape, slide_index=0)
        self.assertEqual(len(feedback_items), 0, "Should be no violations for good alt text.")

    def test_image_missing_alt_text_empty_name(self):
        shape = self._create_mock_image_shape("")
        feedback_items = self.assembler._check_image_accessibility(shape, slide_index=0)
        self.assertEqual(len(feedback_items), 1)
        self.assertEqual(feedback_items[0].category, 'accessibility')
        self.assertIn("missing alt text (name property is empty)", feedback_items[0].message)

    def test_image_missing_alt_text_none_name(self):
        shape = self._create_mock_image_shape(None)
        feedback_items = self.assembler._check_image_accessibility(shape, slide_index=0)
        self.assertEqual(len(feedback_items), 1)
        self.assertEqual(feedback_items[0].category, 'accessibility')
        self.assertIn("missing alt text (name property is empty)", feedback_items[0].message)

    def test_image_generic_alt_text_exact_match(self):
        generic_names_to_test = ["picture", "image", "chart", "graph", "diagram", "img", "pic", "graphic"]
        for name in generic_names_to_test:
            with self.subTest(generic_name=name):
                shape = self._create_mock_image_shape(name)
                feedback_items = self.assembler._check_image_accessibility(shape, slide_index=0)
                self.assertEqual(len(feedback_items), 1, f"Violation expected for generic name '{name}'")
                self.assertEqual(feedback_items[0].category, 'accessibility')
                self.assertIn("alt text that is too generic", feedback_items[0].message)

    def test_image_generic_alt_text_default_pattern(self):
        patterns_to_test = ["Picture 1", "Image 23", "Chart 55", "graphic 100"]
        for name in patterns_to_test:
            with self.subTest(pattern_name=name):
                shape = self._create_mock_image_shape(name)
                feedback_items = self.assembler._check_image_accessibility(shape, slide_index=0)
                self.assertEqual(len(feedback_items), 1, f"Violation expected for pattern '{name}'")
                self.assertEqual(feedback_items[0].category, 'accessibility')
                self.assertIn("follows a default pattern", feedback_items[0].message)

    def test_image_alt_text_case_insensitivity(self):
        shape = self._create_mock_image_shape("PiCtUrE 12")
        feedback_items = self.assembler._check_image_accessibility(shape, slide_index=0)
        self.assertEqual(len(feedback_items), 1)
        self.assertEqual(feedback_items[0].category, 'accessibility')
        self.assertIn("follows a default pattern", feedback_items[0].message)

    def test_accessibility_check_disabled(self):
        config = StyleValidationConfig(check_alt_text_accessibility=False)
        assembler = SlideAssembler(self.mock_template_parser, config)
        shape = self._create_mock_image_shape(None)
        feedback_items = assembler._check_image_accessibility(shape, slide_index=0)
        self.assertEqual(len(feedback_items), 0, "No violations if accessibility check is off.")

    def test_non_picture_shape_not_checked_by_this_function(self):
        shape = self._create_mock_image_shape("Some name", shape_id=999)
        shape.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        feedback_items = self.assembler._check_image_accessibility(shape, slide_index=0)
        self.assertEqual(len(feedback_items), 0)


class TestSlideAssemblerVisualProofreaderIntegration(unittest.TestCase):
    def setUp(self):
        self.mock_template_parser = MagicMock(spec=TemplateParser)
        self.mock_template_parser.template_style = MagicMock()
        self.mock_template_parser.template_style.language_specific_fonts = {}

        self.mock_ppt_presentation = MagicMock(spec=Presentation)
        self.mock_ppt_presentation.slides = []

    @patch('open_lilli.slide_assembler.OpenAI')
    @patch('open_lilli.slide_assembler.VisualProofreader')
    def test_vp_initialized_when_enabled(self, MockVisualProofreader, MockOpenAI):
        config = StyleValidationConfig(enable_visual_proofreader=True, visual_proofreader_model="test-model", visual_proofreader_temperature=0.5)
        assembler = SlideAssembler(self.mock_template_parser, config)

        MockOpenAI.assert_called_once()
        MockVisualProofreader.assert_called_once_with(
            client=MockOpenAI.return_value,
            model="test-model",
            temperature=0.5
        )
        self.assertIsNotNone(assembler.visual_proofreader)

    @patch('open_lilli.slide_assembler.OpenAI')
    @patch('open_lilli.slide_assembler.VisualProofreader')
    def test_vp_not_initialized_when_disabled(self, MockVisualProofreader, MockOpenAI):
        config = StyleValidationConfig(enable_visual_proofreader=False)
        assembler = SlideAssembler(self.mock_template_parser, config)

        MockOpenAI.assert_not_called()
        MockVisualProofreader.assert_not_called()
        self.assertIsNone(assembler.visual_proofreader)

    @patch('open_lilli.slide_assembler.OpenAI') # Keep patch order consistent
    @patch('open_lilli.slide_assembler.VisualProofreader')
    @patch.object(SlideAssembler, '_validate_slide_style', return_value=[]) # Mock rule-based checks to return no issues
    def test_vp_invocation_and_feedback_merging(self, mock_validate_slide_style, MockVisualProofreader, MockOpenAI):
        # Setup VisualProofreader enabled
        config = StyleValidationConfig(
            enable_visual_proofreader=True,
            visual_proofreader_focus_areas=[DesignIssueType.CAPITALIZATION],
            visual_proofreader_enable_corrections=True,
            mode="lenient" # to prevent StyleError from being raised
        )

        # Mock OpenAI client and VisualProofreader instance
        mock_openai_instance = MockOpenAI.return_value
        mock_vp_instance = MockVisualProofreader.return_value

        assembler = SlideAssembler(self.mock_template_parser, config)
        # Ensure the assembler uses the instance we can assert on
        assembler.visual_proofreader = mock_vp_instance

        # Mock presentation with one slide
        mock_slide_obj = MagicMock(name="MockPPTXSlide")
        mock_shapes_collection = MagicMock(name="ShapesCollection")

        # Setup title access: slide.shapes.title
        mock_title_placeholder = MagicMock(name="TitlePlaceholder")
        mock_title_placeholder.has_text_frame = True
        mock_title_placeholder.text = "Test Slide Title From VP Test"
        mock_shapes_collection.title = mock_title_placeholder

        # Mock shapes for bullet extraction (simplified)
        mock_body_shape = MagicMock(name="BodyShape")
        mock_body_shape.has_text_frame = True
        mock_body_shape.is_placeholder = True
        mock_body_shape.placeholder_format = MagicMock()
        mock_body_shape.placeholder_format.type = PP_PLACEHOLDER.BODY.value # Ensure .value for enum comparison if needed
        mock_body_para = MagicMock()
        mock_body_para.text = "Test bullet point from VP test"
        mock_body_shape.text_frame = MagicMock()
        mock_body_shape.text_frame.paragraphs = [mock_body_para]

        # Make the shapes collection iterable
        mock_shapes_collection.__iter__ = MagicMock(return_value=iter([mock_body_shape]))
        mock_slide_obj.shapes = mock_shapes_collection
        mock_slide_obj.placeholders = [mock_body_shape] # Assuming placeholders might be accessed separately

        self.mock_ppt_presentation.slides = [mock_slide_obj]

        # VisualProofreader is expected to be called with SlidePlan objects.
        # The actual creation of these is simplified in the main code.
        # We are testing the interaction, not the SlidePlan creation fidelity here.

        # Mock VisualProofreader outputs
        mock_design_issue = DesignIssue(
            slide_index=0, issue_type=DesignIssueType.CAPITALIZATION, severity="medium",
            element="title", original_text="test", corrected_text="Test",
            description="Bad capitalization", confidence=0.9
        )
        mock_vp_result = ProofreadingResult(
            total_slides=1, issues_found=[mock_design_issue],
            processing_time_seconds=1.0, model_used="test-model"
        )
        mock_vp_instance.proofread_slides.return_value = mock_vp_result

        expected_vp_feedback = ReviewFeedback(
            slide_index=0, severity="medium", category="design", # from convert_to_review_feedback
            message="Bad capitalization", suggestion="Fix capitalization issue in title: change 'test' to 'Test'"
        )
        mock_vp_instance.convert_to_review_feedback.return_value = [expected_vp_feedback]

        # Rule-based check (mocked to return one issue for merging check)
        rule_based_feedback_item = ReviewFeedback(slide_index=0, severity="high", category="placeholder_population", message="Missing title", suggestion="Add title.")
        mock_validate_slide_style.return_value = [rule_based_feedback_item]

        # Call the main validation method
        final_feedback = assembler.validate_presentation_style(self.mock_ppt_presentation, language="en")

        mock_vp_instance.proofread_slides.assert_called_once()
        args, kwargs = mock_vp_instance.proofread_slides.call_args

        self.assertIn('slides', kwargs, "slides argument not found in kwargs")
        self.assertIsInstance(kwargs['slides'], list, "'slides' kwarg is not a list")
        if kwargs['slides']: # Check if the list is not empty
            self.assertIsInstance(kwargs['slides'][0], SlidePlan, "First item in 'slides' kwarg is not a SlidePlan")

        self.assertEqual(kwargs.get('focus_areas'), [DesignIssueType.CAPITALIZATION])
        self.assertEqual(kwargs.get('enable_corrections'), True)

        mock_vp_instance.convert_to_review_feedback.assert_called_once_with([mock_design_issue])

        self.assertIn(expected_vp_feedback, final_feedback, "VisualProofreader feedback missing from final result.")
        self.assertIn(rule_based_feedback_item, final_feedback, "Rule-based feedback missing from final result.")
        self.assertEqual(len(final_feedback), 2)


    @patch('open_lilli.slide_assembler.OpenAI', side_effect=Exception("OpenAI Init Error"))
    @patch('open_lilli.slide_assembler.logger')
    def test_vp_initialization_fails_gracefully(self, mock_logger, MockOpenAI):
        config = StyleValidationConfig(enable_visual_proofreader=True)
        assembler = SlideAssembler(self.mock_template_parser, config)

        self.assertIsNone(assembler.visual_proofreader)
        mock_logger.error.assert_any_call("Failed to initialize VisualProofreader: OpenAI Init Error")

    @patch('open_lilli.slide_assembler.OpenAI')
    @patch('open_lilli.slide_assembler.VisualProofreader')
    @patch.object(SlideAssembler, '_validate_slide_style', return_value=[])
    @patch('open_lilli.slide_assembler.logger')
    def test_vp_proofread_slides_fails_gracefully(self, mock_logger, mock_validate_slide_style, MockVisualProofreader, MockOpenAI):
        config = StyleValidationConfig(enable_visual_proofreader=True, mode="lenient")
        mock_vp_instance = MockVisualProofreader.return_value
        mock_vp_instance.proofread_slides.side_effect = Exception("VP Proofread Error")

        assembler = SlideAssembler(self.mock_template_parser, config)
        assembler.visual_proofreader = mock_vp_instance # Ensure it's using our mock

        # Setup a simple mock slide for the presentation
        mock_slide_obj_for_error_test = MagicMock(name="MockPPTXSlideForError")
        mock_shapes_collection_err = MagicMock(name="ShapesCollectionError")

        mock_title_placeholder_err = MagicMock(name="TitlePlaceholderError")
        mock_title_placeholder_err.has_text_frame = True
        mock_title_placeholder_err.text = "Error Test Slide"
        mock_shapes_collection_err.title = mock_title_placeholder_err

        mock_shapes_collection_err.__iter__ = MagicMock(return_value=iter([])) # No other shapes for iteration
        mock_slide_obj_for_error_test.shapes = mock_shapes_collection_err
        mock_slide_obj_for_error_test.placeholders = [] # Assuming no placeholders iterated for this test path
        self.mock_ppt_presentation.slides = [mock_slide_obj_for_error_test]


        feedback = assembler.validate_presentation_style(self.mock_ppt_presentation, language="en")

        mock_logger.error.assert_any_call("VisualProofreader failed: VP Proofread Error")
        # Check if an error feedback item was added
        self.assertTrue(any(fb.category == "visual_proofreader_error" for fb in feedback))


if __name__ == '__main__':
    unittest.main()
