"""Reviewer for AI-powered presentation critique and refinement."""

import json
import logging
import re
import time
import asyncio
import math # Added for APCA calculations if needed, though abs() and ** are built-in
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Union

import openai
from openai import OpenAI, AsyncOpenAI
from pydantic import ValidationError
from pptx.util import Inches
from pptx.enum.text import MSO_AUTO_SIZE
from pptx import Presentation as PptxPresentationType # For type hinting

from .models import QualityGateResult, QualityGates, ReviewFeedback, SlidePlan, TemplateStyle, PlaceholderStyleInfo, FontInfo

logger = logging.getLogger(__name__)


def calculate_readability_score(text: str) -> float:
    """
    Calculate a simplified readability score based on Flesch Reading Ease.
    
    This is a simplified version that estimates reading difficulty.
    Higher scores indicate more complex text.
    
    Args:
        text: Text to analyze
        
    Returns:
        Grade level (0-20+, where 9 = 9th grade level)
    """
    if not text or not text.strip():
        return 0.0
    
    # Clean text and split into sentences
    clean_text = re.sub(r'[^\w\s\.]', ' ', text)
    sentences = [s.strip() for s in re.split(r'[.!?]+', clean_text) if s.strip()]
    
    if not sentences:
        return 0.0
    
    # Count words and syllables
    total_words = 0
    total_syllables = 0
    
    for sentence in sentences:
        words = sentence.split()
        total_words += len(words)
        
        for word in words:
            # Simple syllable counting
            syllables = count_syllables(word.lower())
            total_syllables += syllables
    
    if total_words == 0:
        return 0.0
    
    # Calculate average sentence length and syllables per word
    avg_sentence_length = total_words / len(sentences)
    avg_syllables_per_word = total_syllables / total_words
    
    # Simplified Flesch-Kincaid Grade Level formula
    # Adjusted for presentation context (typically shorter sentences)
    grade_level = (0.39 * avg_sentence_length) + (11.8 * avg_syllables_per_word) - 15.59
    
    # Ensure reasonable bounds
    return max(0.0, min(20.0, grade_level))


def count_syllables(word: str) -> int:
    """
    Count syllables in a word using simple heuristics.
    
    Args:
        word: Word to count syllables for
        
    Returns:
        Number of syllables (minimum 1)
    """
    if not word:
        return 0
    
    word = word.lower()
    
    # Remove common endings that don't add syllables
    word = re.sub(r'[.,:;!?]', '', word)
    
    # Count vowel groups
    vowels = 'aeiouy'
    syllable_count = 0
    prev_was_vowel = False
    
    for char in word:
        is_vowel = char in vowels
        if is_vowel and not prev_was_vowel:
            syllable_count += 1
        prev_was_vowel = is_vowel
    
    # Handle silent 'e'
    if word.endswith('e') and syllable_count > 1:
        syllable_count -= 1
    
    # Ensure at least 1 syllable
    return max(1, syllable_count)


APCA_CONSTANTS = {
    'exponents': {'mainTRC': 2.4, 'normBG': 0.56, 'normTXT': 0.57, 'revTXT': 0.62, 'revBG': 0.65},
    'colorSpace': {'sRco': 0.2126729, 'sGco': 0.7151522, 'sBco': 0.0721750},
    'clamps': {'blkThrs': 0.022, 'blkClmp': 1.414, 'loClip': 0.1, 'deltaYmin': 0.0005},
    'scalers': {'scaleBoW': 1.14, 'loBoWoffset': 0.027, 'scaleWoB': 1.14, 'loWoBoffset': 0.027}
}


def _calculate_apca_luminance(hex_color: str) -> float:
    """
    Calculate screen luminance (Y_s) for APCA from a hex color string.
    """
    # Remove # if present
    if hex_color.startswith("#"):
        hex_color = hex_color[1:]

    # Convert hex to RGB (0-255)
    r_int = int(hex_color[0:2], 16)
    g_int = int(hex_color[2:4], 16)
    b_int = int(hex_color[4:6], 16)

    # Convert R, G, B to linear values (0-1 range)
    r_lin = (r_int / 255.0) ** APCA_CONSTANTS['exponents']['mainTRC']
    g_lin = (g_int / 255.0) ** APCA_CONSTANTS['exponents']['mainTRC']
    b_lin = (b_int / 255.0) ** APCA_CONSTANTS['exponents']['mainTRC']

    # Calculate Y_s (screen luminance)
    y_s = (r_lin * APCA_CONSTANTS['colorSpace']['sRco'] +
           g_lin * APCA_CONSTANTS['colorSpace']['sGco'] +
           b_lin * APCA_CONSTANTS['colorSpace']['sBco'])
    return y_s


def _soft_clamp_black_level(y_c: float) -> float:
    """
    Apply a soft clamp to the black level of a luminance value.
    """
    b_thresh = APCA_CONSTANTS['clamps']['blkThrs']
    b_exp = APCA_CONSTANTS['clamps']['blkClmp']

    if y_c >= b_thresh:
        return y_c
    else:
        # Ensure the base of the exponent is non-negative
        base = b_thresh - y_c
        if base < 0: # Should not happen if y_c < b_thresh
            base = 0
        return y_c + (base)**b_exp


def calculate_apca_contrast_value(fg_hex: str, bg_hex: str) -> float:
    """
    Calculate the APCA contrast value (Lc) between two hex colors.
    """
    y_s_fg = _calculate_apca_luminance(fg_hex)
    y_s_bg = _calculate_apca_luminance(bg_hex)

    # Apply soft clamp
    y_txt = _soft_clamp_black_level(y_s_fg)
    y_bg = _soft_clamp_black_level(y_s_bg)

    # Determine polarity and calculate S_pol
    norm_txt_exp = APCA_CONSTANTS['exponents']['normTXT']
    norm_bg_exp = APCA_CONSTANTS['exponents']['normBG']
    rev_txt_exp = APCA_CONSTANTS['exponents']['revTXT']
    rev_bg_exp = APCA_CONSTANTS['exponents']['revBG']

    # Ensure y_txt and y_bg are non-negative before exponentiation if they can be negative
    # However, luminance values (Y_s) and clamped values (Y_txt, Y_bg) should be >= 0.
    # Based on APCA, Y values are typically >= 0.

    if y_bg > y_txt:  # Normal polarity (perceivably dark text on light background)
        s_pol = (y_bg**norm_bg_exp) - (y_txt**norm_txt_exp)
    else:  # Reverse polarity (perceivably light text on dark background)
        s_pol = (y_bg**rev_bg_exp) - (y_txt**rev_txt_exp)

    # Clamp noise and scale (C)
    p_in = APCA_CONSTANTS['clamps']['deltaYmin']
    # scaleBoW and scaleWoB are the same in the provided constants
    r_scale = APCA_CONSTANTS['scalers']['scaleBoW']

    if abs(y_bg - y_txt) < p_in:
        c_val = 0.0
    else:
        c_val = s_pol * r_scale

    # Clamp minimum contrast and offset (S_apc)
    p_out = APCA_CONSTANTS['clamps']['loClip']
    # loBoWoffset and loWoBoffset are the same
    w_offset = APCA_CONSTANTS['scalers']['loBoWoffset']

    if abs(c_val) < p_out:
        s_apc = 0.0
    elif c_val > 0: # Positive contrast
        s_apc = c_val - w_offset
    else: # Negative contrast
        s_apc = c_val + w_offset # Effectively c_val - (-w_offset) for negative C

    # Calculate final Lc (Lightness contrast)
    lc = s_apc * 100
    return lc


def calculate_contrast_ratio(hex_color1: str, hex_color2: str) -> float:
    """
    Calculate the APCA contrast value (Lc) between two hex colors.
    This function now serves as an entry point for APCA calculation,
    replacing the previous WCAG contrast ratio logic.

    Args:
        hex_color1: The foreground hex color string (e.g., "#RRGGBB").
        hex_color2: The background hex color string (e.g., "#RRGGBB").

    Returns:
        The APCA Lc value.
    """
    # The APCA calculation is directional (foreground vs background matters).
    # We assume hex_color1 is foreground and hex_color2 is background.
    # If the original usage of calculate_contrast_ratio was for any two colors
    # without specific fg/bg designation, this change makes it directional.
    # The quality gate code seems to use it as text_color vs background_color,
    # so this should align.
    return calculate_apca_contrast_value(hex_color1, hex_color2)


def _is_shape_within_margins(shape_bbox: Dict[str, int], slide_dims: Dict[str, int], margins_emu: Dict[str, int]) -> bool:
    """Check if a shape's bounding box is within the defined slide margins."""
    if shape_bbox['left'] < margins_emu['left']:
        return False
    if shape_bbox['top'] < margins_emu['top']:
        return False
    if (shape_bbox['left'] + shape_bbox['width']) > (slide_dims['width'] - margins_emu['right']):
        return False
    if (shape_bbox['top'] + shape_bbox['height']) > (slide_dims['height'] - margins_emu['bottom']):
        return False
    return True

def _check_overset_text(shape) -> bool:
    """
    Heuristic check for overset text in a shape.
    Returns True if overset text is suspected, False otherwise.
    """
    if not shape.has_text_frame or not shape.text_frame.text or not shape.text_frame.text.strip():
        return False  # No text or only whitespace

    text_frame = shape.text_frame

    # If text auto-fits to shape or shape resizes to text, it's generally not overset in a hidden way.
    if text_frame.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT or \
       text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
        return False

    # Heuristic for MSO_AUTO_SIZE.NONE (or inherited NONE)
    # This is a simplified check. More sophisticated checks might involve rendering text.
    # Using a simple character count heuristic for now.
    # This threshold might need tuning based on typical content and shape sizes.
    if text_frame.auto_size == MSO_AUTO_SIZE.NONE or text_frame.auto_size is None:
        # A large amount of text in a shape not set to auto-size is a strong indicator.
        if len(text_frame.text) > 250:  # Arbitrary heuristic value
            # Further simple check: if word_wrap is False and there's a long paragraph, it's likely overset.
            if not text_frame.word_wrap:
                for para in text_frame.paragraphs:
                    if len(para.text) > 80: # Long paragraph without wrapping
                        return True
            return True # General flag for long text in non-autosizing shape

    # Placeholder for more advanced heuristics (e.g., comparing text box size to rendered text size)
    # For now, we rely on the above.
    return False


class Reviewer:
    """AI-powered presentation reviewer for quality feedback and iterative refinement."""

    def __init__(self,
                 client: OpenAI | AsyncOpenAI,
                 model: str = "gpt-4",
                 temperature: float = 0.2,
                 template_style: Optional[TemplateStyle] = None,
                 presentation: Optional[PptxPresentationType] = None): # Added presentation
        """
        Initialize the reviewer.
        
        Args:
            client: OpenAI client instance
            model: Model name to use for review
            temperature: Temperature for generation (lower for more consistent reviews)
            template_style: Optional TemplateStyle object for style-related checks.
            presentation: Optional pptx.Presentation object for layout-based checks.
        """
        self.client = client
        self.model = model
        self.temperature = temperature
        self.max_retries = 3
        self.retry_delay = 1.0
        self.template_style = template_style
        self.presentation = presentation # Stored presentation object

    def review_presentation(
        self,
        slides: List[SlidePlan],
        presentation_context: Optional[str] = None,
        review_criteria: Optional[Dict[str, any]] = None,
        include_quality_gates: bool = False,
        quality_gates: Optional[QualityGates] = None
    ) -> Union[List[ReviewFeedback], Tuple[List[ReviewFeedback], QualityGateResult]]:
        """
        Review an entire presentation and provide feedback.
        
        Args:
            slides: List of slides to review
            presentation_context: Context about the presentation (audience, purpose, etc.)
            review_criteria: Specific criteria to focus on during review
            include_quality_gates: If True, also evaluate quality gates and return tuple
            quality_gates: Quality gate configuration (uses defaults if None)
            
        Returns:
            List of feedback items for improvement, or tuple of (feedback, quality_gates_result)
            if include_quality_gates is True
        """
        logger.info(f"Reviewing presentation with {len(slides)} slides")
        
        review_criteria = review_criteria or self._get_default_criteria()
        
        # Create presentation summary for review
        presentation_summary = self._create_presentation_summary(slides, presentation_context)
        
        # Generate review prompt
        prompt = self._build_review_prompt(presentation_summary, review_criteria)
        
        try:
            # Get AI feedback
            feedback_data = self._call_openai_with_retries(prompt)
            
            # Parse and validate feedback
            feedback_list = self._parse_feedback_response(feedback_data)
            
            logger.info(f"Generated {len(feedback_list)} feedback items")
            
            # Evaluate quality gates if requested
            if include_quality_gates:
                quality_result = self.evaluate_quality_gates(slides, feedback_list, quality_gates)
                return feedback_list, quality_result
            
            return feedback_list
            
        except Exception as e:
            logger.error(f"Presentation review failed: {e}")
            if include_quality_gates:
                # Return empty feedback and a failing quality gates result
                empty_feedback = []
                default_gates = quality_gates or QualityGates()
                failing_result = QualityGateResult(
                    status="needs_fix",
                    gate_results={
                        "bullet_count": False,
                        "readability": False,
                        "style_errors": False,
                        "overall_score": False
                    },
                    violations=["Review failed due to an error"],
                    recommendations=["Please retry the review"],
                    metrics={}
                )
                return empty_feedback, failing_result
            return []

    async def review_presentation_async(
        self,
        slides: List[SlidePlan],
        presentation_context: Optional[str] = None,
        review_criteria: Optional[Dict[str, any]] = None,
        include_quality_gates: bool = False,
        quality_gates: Optional[QualityGates] = None,
    ) -> Union[List[ReviewFeedback], Tuple[List[ReviewFeedback], QualityGateResult]]:
        """Asynchronous version of ``review_presentation``."""
        logger.info(f"Reviewing presentation with {len(slides)} slides (async)")

        review_criteria = review_criteria or self._get_default_criteria()
        presentation_summary = self._create_presentation_summary(slides, presentation_context)
        prompt = self._build_review_prompt(presentation_summary, review_criteria)

        try:
            feedback_data = await self._call_openai_with_retries_async(prompt)
            feedback_list = self._parse_feedback_response(feedback_data)
            logger.info(f"Generated {len(feedback_list)} feedback items")
            if include_quality_gates:
                quality_result = self.evaluate_quality_gates(slides, feedback_list, quality_gates)
                return feedback_list, quality_result
            return feedback_list
        except Exception as e:
            logger.error(f"Presentation review failed: {e}")
            if include_quality_gates:
                empty_feedback = []
                default_gates = quality_gates or QualityGates()
                failing_result = QualityGateResult(
                    status="needs_fix",
                    gate_results={
                        "bullet_count": False,
                        "readability": False,
                        "style_errors": False,
                        "overall_score": False,
                    },
                    violations=["Review failed due to an error"],
                    recommendations=["Please retry the review"],
                    metrics={},
                )
                return empty_feedback, failing_result
            return []

    def review_individual_slide(
        self,
        slide: SlidePlan,
        slide_context: Optional[str] = None
    ) -> List[ReviewFeedback]:
        """
        Review an individual slide in detail.
        
        Args:
            slide: Slide to review
            slide_context: Additional context about this slide's role
            
        Returns:
            List of feedback items for this slide
        """
        logger.debug(f"Reviewing individual slide {slide.index}: {slide.title}")
        
        prompt = self._build_slide_review_prompt(slide, slide_context)
        
        try:
            feedback_data = self._call_openai_with_retries(prompt)
            feedback_list = self._parse_feedback_response(feedback_data)
            
            # Ensure all feedback is for this slide
            for feedback in feedback_list:
                feedback.slide_index = slide.index
            
            return feedback_list
            
        except Exception as e:
            logger.error(f"Slide review failed for slide {slide.index}: {e}")
            return []

    def check_presentation_flow(self, slides: List[SlidePlan]) -> List[ReviewFeedback]:
        """
        Check the logical flow and narrative structure of the presentation.
        
        Args:
            slides: List of slides to check flow for
            
        Returns:
            List of flow-related feedback
        """
        logger.info("Checking presentation flow and narrative structure")
        
        # Create flow summary
        flow_summary = self._create_flow_summary(slides)
        
        prompt = f"""You are an expert presentation consultant. Analyze the flow and narrative structure of this presentation.

PRESENTATION FLOW:
{flow_summary}

FOCUS AREAS:
1. Logical progression of ideas
2. Smooth transitions between sections
3. Clear narrative arc (beginning, middle, end)
4. Appropriate pacing and content distribution
5. Missing or redundant content

Return a JSON array of feedback items with this structure:
[
    {{
        "slide_index": 2,
        "severity": "medium",
        "category": "flow",
        "message": "Abrupt transition from introduction to detailed analysis",
        "suggestion": "Add an overview slide to bridge the gap"
    }}
]

Provide 3-5 specific, actionable feedback items focused on improving the presentation flow."""

        try:
            feedback_data = self._call_openai_with_retries(prompt)
            return self._parse_feedback_response(feedback_data)
            
        except Exception as e:
            logger.error(f"Flow check failed: {e}")
            return []

    def evaluate_quality_gates(
        self,
        slides: List[SlidePlan],
        feedback: List[ReviewFeedback],
        gates: Optional[QualityGates] = None
    ) -> QualityGateResult:
        """
        Evaluate presentation against quality gates to determine pass/fail status.
        
        Args:
            slides: List of slides to evaluate
            feedback: List of review feedback items
            gates: Quality gate configuration (uses defaults if None)
            
        Returns:
            QualityGateResult with pass/fail status and detailed metrics
        """
        logger.info("Evaluating presentation against quality gates")
        
        if gates is None:
            gates = QualityGates()
        
        # Initialize results
        gate_results = {}
        violations = []
        recommendations = []
        metrics = {}

        # --- Pre-computation for layout related checks ---
        slide_width_emu = 0
        slide_height_emu = 0
        can_perform_layout_checks = False

        if self.template_style and hasattr(self.template_style, 'slide_width') and hasattr(self.template_style, 'slide_height') and \
           self.template_style.slide_width and self.template_style.slide_height:
            slide_width_emu = self.template_style.slide_width
            slide_height_emu = self.template_style.slide_height
            if self.presentation:
                can_perform_layout_checks = True
            else:
                logger.warning("Presentation object (self.presentation) not available to Reviewer. Skipping alignment and overset text checks.")
        else:
            logger.warning("TemplateStyle or its slide dimensions not available. Skipping alignment and overset text checks.")

        if not can_perform_layout_checks:
            gates.enable_alignment_check = False
            gates.enable_overset_text_check = False
            # Ensure gate_results for these checks are set to True (skipped = passed) if they were initially enabled
            if "alignment_check" not in gate_results: gate_results["alignment_check"] = True
            if "overset_text_check" not in gate_results: gate_results["overset_text_check"] = True


        margins_emu = {
            'left': Inches(gates.slide_margin_left_inches),
            'top': Inches(gates.slide_margin_top_inches),
            'right': Inches(gates.slide_margin_right_inches),
            'bottom': Inches(gates.slide_margin_bottom_inches)
        }
        
        # --- Initialize counts for new checks ---
        misaligned_shapes_count = 0
        overset_text_shapes_count = 0
        alignment_violations_details = [] # To gather names of misaligned shapes for summary recommendation
        overset_violations_details = []   # To gather names of overset shapes for summary recommendation

        # 1. Check bullet count per slide
        max_bullets_found = 0
        bullet_violations = []
        
        for slide in slides:
            bullet_count = len(slide.bullets)
            if bullet_count > max_bullets_found:
                max_bullets_found = bullet_count
            
            if bullet_count > gates.max_bullets_per_slide:
                violation = f"Slide {slide.index + 1} has {bullet_count} bullets (exceeds limit of {gates.max_bullets_per_slide})"
                violations.append(violation)
                bullet_violations.append(slide.index)
        
        gate_results["bullet_count"] = len(bullet_violations) == 0
        metrics["max_bullets_found"] = max_bullets_found
        
        if bullet_violations:
            recommendations.append(f"Reduce bullet points on slides {', '.join(str(i+1) for i in bullet_violations)} to {gates.max_bullets_per_slide} or fewer")
        
        # 2. Check readability grade
        readability_scores = []
        readability_violations = []
        
        for slide in slides:
            # Combine title and bullet text for readability analysis
            slide_text = slide.title
            if slide.bullets:
                slide_text += ". " + ". ".join(slide.bullets)
            
            if slide_text.strip():
                readability_grade = calculate_readability_score(slide_text)
                readability_scores.append(readability_grade)
                
                if readability_grade > gates.max_readability_grade:
                    violation = f"Slide {slide.index + 1} has readability grade {readability_grade:.1f} (exceeds limit of {gates.max_readability_grade})"
                    violations.append(violation)
                    readability_violations.append(slide.index)
        
        gate_results["readability"] = len(readability_violations) == 0
        
        if readability_scores:
            metrics["avg_readability_grade"] = sum(readability_scores) / len(readability_scores)
            metrics["max_readability_grade"] = max(readability_scores)
        else:
            metrics["avg_readability_grade"] = 0.0
            metrics["max_readability_grade"] = 0.0
        
        if readability_violations:
            recommendations.append(f"Simplify language on slides {', '.join(str(i+1) for i in readability_violations)} to improve readability")
        
        # 3. Check style errors
        style_error_count = sum(1 for f in feedback if f.category == "design" or f.category == "consistency")
        gate_results["style_errors"] = style_error_count <= gates.max_style_errors
        metrics["style_error_count"] = style_error_count
        
        if style_error_count > gates.max_style_errors:
            violations.append(f"Found {style_error_count} style errors (exceeds limit of {gates.max_style_errors})")
            recommendations.append("Address style and consistency issues identified in feedback")
        
        # 4. Check overall score
        review_summary = self.get_review_summary(feedback)
        overall_score = review_summary["overall_score"]
        gate_results["overall_score"] = overall_score >= gates.min_overall_score
        metrics["overall_score"] = overall_score
        
        if overall_score < gates.min_overall_score:
            violations.append(f"Overall score {overall_score} is below minimum threshold of {gates.min_overall_score}")
            recommendations.append("Address critical and high severity feedback to improve overall score")

        # 5. Contrast Ratio Check (Now APCA Lc Value Check)
        # CRITICAL ASSUMPTION: self.template_style (TemplateStyle) is available on this Reviewer instance.
        # This is not currently handled by __init__ and needs to be addressed separately.
        # Also, this check simplifies by only considering a 'body' placeholder and its fill vs text/master font.
        # A full check would need to iterate all text elements on a slide and their actual backgrounds.
        contrast_violations = []
        all_lc_values = [] # Changed from all_ratios
        # APCA Recommended Levels (example for reference, actual interpretation may vary):
        # 60 Lc: Minimum for body text.
        # 75 Lc: Preferred for body text.
        # 45 Lc: Minimum for large text (e.g., headlines).
        # For this quality gate, we use a configurable threshold.
        apca_lc_threshold_pass = gates.min_apca_lc_for_body_text # Assuming this will be added to QualityGates model

        if hasattr(self, 'template_style') and self.template_style is not None:
            placeholder_type_body = 2  # Assuming BODY placeholder type index for main content

            for slide in slides:
                placeholder_style: Optional[PlaceholderStyleInfo] = self.template_style.get_placeholder_style(placeholder_type_body)

                text_color_hex: Optional[str] = None
                background_color_hex: Optional[str] = None

                if placeholder_style:
                    # Determine Text Color
                    if placeholder_style.default_font and placeholder_style.default_font.color:
                        text_color_hex = placeholder_style.default_font.color

                    # Determine Background Color (from placeholder fill first)
                    if placeholder_style.fill_color:
                        background_color_hex = placeholder_style.fill_color

                # Fallback for Text Color (if not found in placeholder style)
                if text_color_hex is None:
                    if self.template_style.master_font and self.template_style.master_font.color:
                        text_color_hex = self.template_style.master_font.color
                    else:
                        text_color_hex = "#000000"  # Default to black

                # Fallback for Background Color (if not found in placeholder fill)
                if background_color_hex is None:
                    background_color_hex = self.template_style.theme_colors.get('lt1', '#FFFFFF') # Default to theme light or white

                if text_color_hex and background_color_hex:
                    try:
                        # calculate_contrast_ratio now returns APCA Lc value
                        lc_value = calculate_contrast_ratio(text_color_hex, background_color_hex)
                        all_lc_values.append(lc_value)

                        # APCA Lc values can be negative. Higher absolute values are generally better.
                        # We check if the absolute Lc value meets the threshold.
                        if abs(lc_value) < apca_lc_threshold_pass:
                            violation_message = (
                                f"Slide {slide.index + 1} (Body Placeholder): APCA Lc is {lc_value:.2f} "
                                f"(Text: {text_color_hex}, Background: {background_color_hex}). "
                                f"Recommended minimum absolute Lc is {apca_lc_threshold_pass} for body text readability."
                            )
                            violations.append(violation_message)
                            contrast_violations.append(slide.index)
                    except Exception as e:
                        logger.warning(f"Could not calculate APCA contrast for slide {slide.index + 1}: {e}")
                else:
                    logger.warning(f"Could not determine text/background color for APCA contrast check on slide {slide.index + 1}.")

            gate_results["contrast_check"] = len(contrast_violations) == 0
            if all_lc_values:
                metrics["min_apca_lc_found"] = min(all_lc_values) if all_lc_values else 0.0
                metrics["max_apca_lc_found"] = max(all_lc_values) if all_lc_values else 0.0
                metrics["avg_apca_lc_found"] = sum(all_lc_values) / len(all_lc_values) if all_lc_values else 0.0

                abs_lc_values = [abs(lc) for lc in all_lc_values]
                metrics["min_abs_apca_lc_found"] = min(abs_lc_values) if abs_lc_values else 0.0
                metrics["avg_abs_apca_lc_found"] = sum(abs_lc_values) / len(abs_lc_values) if abs_lc_values else 0.0
            else: # Default if no Lc values calculated
                metrics["min_apca_lc_found"] = 0.0
                metrics["max_apca_lc_found"] = 0.0
                metrics["avg_apca_lc_found"] = 0.0
                metrics["min_abs_apca_lc_found"] = 0.0
                metrics["avg_abs_apca_lc_found"] = 0.0

            if contrast_violations:
                recommendations.append(
                    f"Improve text contrast on slides: {', '.join(str(i + 1) for i in sorted(list(set(contrast_violations))))}. "
                    f"Aim for an absolute APCA Lc value of {apca_lc_threshold_pass} or higher for body text."
                )
        else:
            logger.warning("Skipping contrast ratio check: TemplateStyle not found on Reviewer instance.")
            gate_results["contrast_check"] = False # Mark as failed if template_style is missing
            violations.append("Contrast check could not be performed: Template style information is missing.")
            recommendations.append("Ensure Reviewer is initialized with template style information to perform contrast checks.")

        # 6. Alignment Check
        # Ensure checks run only if enabled in gates AND possible to perform
        if gates.enable_alignment_check and can_perform_layout_checks:
            slide_dims_emu = {'width': slide_width_emu, 'height': slide_height_emu}
            for slide_plan in slides:
                try:
                    # self.presentation is confirmed available by can_perform_layout_checks
                    pptx_slide = self.presentation.slides[slide_plan.index]
                    for shape in pptx_slide.shapes:
                        if not hasattr(shape, 'left') or not hasattr(shape, 'top') or \
                           not hasattr(shape, 'width') or not hasattr(shape, 'height') or \
                           shape.width is None or shape.height is None: # Ensure width/height are not None
                            continue # Skip shapes without standard bounding box attributes

                        shape_bbox = {'left': shape.left, 'top': shape.top, 'width': shape.width, 'height': shape.height}
                        if not _is_shape_within_margins(shape_bbox, slide_dims_emu, margins_emu):
                            misaligned_shapes_count += 1
                            shape_name = shape.name or f"Shape ID {shape.shape_id}"
                            detail = f"Slide {slide_plan.index + 1}: Shape '{shape_name}' is outside defined margins."
                            violations.append(detail)
                            if shape_name not in alignment_violations_details: # Add only unique names for summary
                                alignment_violations_details.append(shape_name)

                            # Specific recommendation for logos (example)
                            if "logo" in shape_name.lower():
                                recommendations.append(f"Adjust position of logo '{shape_name}' on slide {slide_plan.index + 1} to be within slide margins.")
                except IndexError:
                    logger.warning(f"Could not access slide at index {slide_plan.index} in presentation for alignment check.")
                except Exception as e:
                    logger.error(f"Error during alignment check for slide {slide_plan.index + 1}: {e}")

            gate_results["alignment_check"] = misaligned_shapes_count == 0
            if not gate_results["alignment_check"] and alignment_violations_details:
                recommendations.append(f"Review alignment of shapes: {', '.join(sorted(list(set(alignment_violations_details))))} to ensure they are within margins.")
        elif "alignment_check" not in gate_results: # Ensure it's set if check was enabled but skipped
            gate_results["alignment_check"] = True


        # 7. Overset Text Check
        # Ensure checks run only if enabled in gates AND possible to perform
        if gates.enable_overset_text_check and can_perform_layout_checks:
            for slide_plan in slides:
                try:
                    # self.presentation is confirmed available by can_perform_layout_checks
                    pptx_slide = self.presentation.slides[slide_plan.index]
                    for shape in pptx_slide.shapes:
                        if _check_overset_text(shape):
                            overset_text_shapes_count += 1
                            shape_name = shape.name or f"Shape ID {shape.shape_id}"
                            detail = f"Slide {slide_plan.index + 1}: Shape '{shape_name}' may have overset/hidden text."
                            violations.append(detail)
                            summary_detail = f"'{shape_name}' on slide {slide_plan.index + 1}"
                            if summary_detail not in overset_violations_details:
                                overset_violations_details.append(summary_detail)
                except IndexError:
                    logger.warning(f"Could not access slide at index {slide_plan.index} in presentation for overset text check.")
                except Exception as e:
                    logger.error(f"Error during overset text check for slide {slide_plan.index + 1}: {e}")

            gate_results["overset_text_check"] = overset_text_shapes_count == 0
            if not gate_results["overset_text_check"] and overset_violations_details:
                recommendations.append(f"Review shapes for potential overset text: {', '.join(sorted(list(set(overset_violations_details))))}.")
        elif "overset_text_check" not in gate_results: # Ensure it's set if check was enabled but skipped
             gate_results["overset_text_check"] = True


        metrics["misaligned_shapes_count"] = misaligned_shapes_count
        metrics["overset_text_shapes_count"] = overset_text_shapes_count

        # Determine overall status
        all_gates_passed = all(gate_results.values())
        status = "pass" if all_gates_passed else "needs_fix"
        
        logger.info(f"Quality gates evaluation completed: {status} ({sum(gate_results.values())}/{len(gate_results)} gates passed)")
        
        return QualityGateResult(
            status=status,
            gate_results=gate_results,
            violations=violations,
            recommendations=recommendations,
            metrics=metrics
        )

    def _create_presentation_summary(
        self,
        slides: List[SlidePlan],
        context: Optional[str]
    ) -> str:
        """Create a summary of the presentation for review."""
        
        summary_parts = []
        
        if context:
            summary_parts.append(f"CONTEXT: {context}\n")
        
        summary_parts.append(f"TOTAL SLIDES: {len(slides)}\n")
        
        summary_parts.append("SLIDE BREAKDOWN:")
        for slide in slides:
            slide_summary = f"\nSlide {slide.index + 1} ({slide.slide_type}): {slide.title}"
            
            if slide.bullets:
                slide_summary += f"\n  • {len(slide.bullets)} bullet points"
                if len(slide.bullets) <= 3:
                    for bullet in slide.bullets:
                        slide_summary += f"\n    - {bullet[:60]}{'...' if len(bullet) > 60 else ''}"
            
            if slide.image_query:
                slide_summary += f"\n  • Image: {slide.image_query}"
            
            if slide.chart_data:
                slide_summary += f"\n  • Chart: {slide.chart_data.get('type', 'unknown')} chart"
            
            summary_parts.append(slide_summary)
        
        return "\n".join(summary_parts)

    def _create_flow_summary(self, slides: List[SlidePlan]) -> str:
        """Create a flow-focused summary of the presentation."""
        
        flow_parts = []
        
        flow_parts.append("PRESENTATION STRUCTURE:")
        
        current_section = None
        for i, slide in enumerate(slides):
            # Detect section changes
            if slide.slide_type in ["title", "section"]:
                current_section = slide.title
                flow_parts.append(f"\n[SECTION] {slide.title}")
            else:
                if current_section:
                    flow_parts.append(f"  {i+1}. {slide.title}")
                else:
                    flow_parts.append(f"{i+1}. {slide.title}")
            
            # Add brief content summary
            if slide.bullets:
                main_points = [bullet[:40] + "..." if len(bullet) > 40 else bullet 
                              for bullet in slide.bullets[:2]]
                flow_parts.append(f"     → {' | '.join(main_points)}")
        
        # Add transition analysis
        flow_parts.append("\nTRANSITION ANALYSIS:")
        for i in range(len(slides) - 1):
            current = slides[i]
            next_slide = slides[i + 1]
            
            transition = f"Slide {i+1} → Slide {i+2}: "
            transition += f"'{current.title}' to '{next_slide.title}'"
            flow_parts.append(transition)
        
        return "\n".join(flow_parts)

    def _build_review_prompt(
        self,
        presentation_summary: str,
        criteria: Dict[str, any]
    ) -> str:
        """Build the main review prompt."""
        
        criteria_text = ""
        for category, details in criteria.items():
            criteria_text += f"\n- {category.upper()}: {details['description']}"
            if details.get('weight'):
                criteria_text += f" (Weight: {details['weight']})"
        
        prompt = f"""You are an expert presentation consultant with experience reviewing corporate presentations. Provide constructive feedback to improve this presentation.

{presentation_summary}

REVIEW CRITERIA:{criteria_text}

INSTRUCTIONS:
1. Focus on actionable, specific improvements
2. Consider the overall narrative and flow
3. Identify slides that may be too dense or sparse
4. Check for consistency in tone and style
5. Look for missing transitions or unclear connections
6. Assess visual balance and content distribution

Return a JSON array of feedback items. Each item should have:
- slide_index: Which slide (0-based index, or -1 for presentation-wide issues)
- severity: "low", "medium", "high", or "critical"
- category: "content", "flow", "design", "clarity", "consistency"
- message: Clear description of the issue
- suggestion: Specific recommendation for improvement

Provide 5-10 most important feedback items that would significantly improve the presentation."""

        return prompt

    def _build_slide_review_prompt(
        self,
        slide: SlidePlan,
        context: Optional[str]
    ) -> str:
        """Build prompt for individual slide review."""
        
        context_text = f"Context: {context}\n" if context else ""
        
        prompt = f"""You are an expert presentation consultant. Review this individual slide for improvement opportunities.

{context_text}
SLIDE DETAILS:
- Index: {slide.index}
- Type: {slide.slide_type}
- Title: {slide.title}
- Bullets: {slide.bullets}
- Speaker Notes: {slide.speaker_notes or 'None'}
- Has Image: {'Yes' if slide.image_query else 'No'}
- Has Chart: {'Yes' if slide.chart_data else 'No'}

REVIEW FOCUS:
1. Title clarity and impact
2. Content density and readability
3. Bullet point structure and parallelism
4. Visual elements appropriateness
5. Speaker notes quality
6. Overall slide effectiveness

Return a JSON array of specific feedback items for this slide. Include 2-4 actionable suggestions."""

        return prompt

    def _get_default_criteria(self) -> Dict[str, any]:
        """Get default review criteria."""
        return {
            "clarity": {
                "description": "Content is clear, concise, and easy to understand",
                "weight": "high"
            },
            "flow": {
                "description": "Logical progression and smooth transitions between ideas",
                "weight": "high"
            },
            "content": {
                "description": "Appropriate content density and meaningful information",
                "weight": "medium"
            },
            "consistency": {
                "description": "Consistent tone, style, and formatting throughout",
                "weight": "medium"
            },
            "engagement": {
                "description": "Content engages audience and maintains interest",
                "weight": "medium"
            }
        }

    def _supports_json_mode(self) -> bool:
        """Check if the current model supports JSON mode."""
        json_mode_models = [
            "gpt-4", "gpt-4-0613", "gpt-4-1106-preview", "gpt-4-0125-preview",
            "gpt-4-turbo-preview", "gpt-4-turbo", "gpt-4o", "gpt-4o-mini",
            "gpt-3.5-turbo", "gpt-3.5-turbo-0613", "gpt-3.5-turbo-1106", "gpt-3.5-turbo-0125"
        ]
        return any(model in self.model for model in json_mode_models)

    def _call_openai_with_retries(self, prompt: str) -> dict:
        """Call OpenAI API with retry logic."""
        
        for attempt in range(self.max_retries):
            try:
                logger.debug(f"Review API call attempt {attempt + 1}")
                
                # Build request parameters
                request_params = {
                    "model": self.model,
                    "messages": [
                        {
                            "role": "system",
                            "content": "You are an expert presentation consultant. Always respond with valid JSON only."
                        },
                        {
                            "role": "user",
                            "content": prompt
                        }
                    ],
                    "temperature": self.temperature,
                    "max_tokens": 2000
                }
                
                # Only add response_format for models that support it
                if self._supports_json_mode():
                    request_params["response_format"] = {"type": "json_object"}
                
                response = self.client.chat.completions.create(**request_params)
                
                content = response.choices[0].message.content
                if not content:
                    raise ValueError("Empty response from OpenAI")
                
                try:
                    return json.loads(content)
                except json.JSONDecodeError as e:
                    logger.error(f"Failed to parse JSON response: {e}")
                    if attempt == self.max_retries - 1:
                        raise ValueError(f"Invalid JSON response: {e}")
                    continue
                
            except openai.RateLimitError:
                if attempt == self.max_retries - 1:
                    raise
                wait_time = self.retry_delay * (2 ** attempt)
                logger.warning(f"Rate limited, waiting {wait_time}s")
                time.sleep(wait_time)
                
            except openai.APIError as e:
                if attempt == self.max_retries - 1:
                    raise
                wait_time = self.retry_delay * (2 ** attempt)
                logger.warning(f"API error: {e}, waiting {wait_time}s")
                time.sleep(wait_time)
        
        raise ValueError("Failed to get review after all retries")

    async def _call_openai_with_retries_async(self, prompt: str) -> dict:
        """Asynchronous version of ``_call_openai_with_retries``."""
        for attempt in range(self.max_retries):
            try:
                logger.debug(f"Review API call attempt {attempt + 1} (async)")
                
                # Build request parameters
                request_params = {
                    "model": self.model,
                    "messages": [
                        {
                            "role": "system",
                            "content": "You are an expert presentation consultant. Always respond with valid JSON only.",
                        },
                        {"role": "user", "content": prompt},
                    ],
                    "temperature": self.temperature,
                    "max_tokens": 2000
                }
                
                # Only add response_format for models that support it
                if self._supports_json_mode():
                    request_params["response_format"] = {"type": "json_object"}
                
                response = await self.client.chat.completions.create(**request_params)

                content = response.choices[0].message.content
                if not content:
                    raise ValueError("Empty response from OpenAI")

                try:
                    return json.loads(content)
                except json.JSONDecodeError as e:
                    logger.error(f"Failed to parse JSON response: {e}")
                    if attempt == self.max_retries - 1:
                        raise ValueError(f"Invalid JSON response: {e}")
                    continue

            except openai.RateLimitError:
                if attempt == self.max_retries - 1:
                    raise
                wait_time = self.retry_delay * (2 ** attempt)
                logger.warning(f"Rate limited, waiting {wait_time}s")
                await asyncio.sleep(wait_time)

            except openai.APIError as e:
                if attempt == self.max_retries - 1:
                    raise
                wait_time = self.retry_delay * (2 ** attempt)
                logger.warning(f"API error: {e}, waiting {wait_time}s")
                await asyncio.sleep(wait_time)

        raise ValueError("Failed to get review after all retries")

    def _parse_feedback_response(self, response_data: dict) -> List[ReviewFeedback]:
        """Parse and validate feedback response from AI."""
        
        feedback_list = []
        
        # Handle different response formats
        if isinstance(response_data, list):
            feedback_data = response_data
        elif isinstance(response_data, dict) and "feedback" in response_data:
            feedback_data = response_data["feedback"]
        elif isinstance(response_data, dict) and "items" in response_data:
            feedback_data = response_data["items"]
        else:
            logger.error(f"Unexpected response format: {response_data}")
            return []
        
        for item in feedback_data:
            try:
                # Ensure required fields
                feedback_item = ReviewFeedback(
                    slide_index=item.get("slide_index", -1),
                    severity=item.get("severity", "medium"),
                    category=item.get("category", "general"),
                    message=item.get("message", ""),
                    suggestion=item.get("suggestion")
                )
                feedback_list.append(feedback_item)
                
            except ValidationError as e:
                logger.warning(f"Invalid feedback item: {item}, error: {e}")
                continue
        
        return feedback_list

    def prioritize_feedback(
        self,
        feedback_list: List[ReviewFeedback]
    ) -> List[ReviewFeedback]:
        """
        Prioritize feedback items by severity and category.
        
        Args:
            feedback_list: List of feedback items to prioritize
            
        Returns:
            Sorted list with highest priority items first
        """
        severity_order = {"critical": 4, "high": 3, "medium": 2, "low": 1}
        category_order = {"flow": 4, "clarity": 3, "content": 2, "consistency": 1, "design": 1}
        
        def priority_score(feedback: ReviewFeedback) -> tuple:
            severity_score = severity_order.get(feedback.severity, 0)
            category_score = category_order.get(feedback.category, 0)
            return (severity_score, category_score)
        
        return sorted(feedback_list, key=priority_score, reverse=True)

    def filter_feedback(
        self,
        feedback_list: List[ReviewFeedback],
        min_severity: str = "low",
        categories: Optional[List[str]] = None,
        slide_indices: Optional[List[int]] = None
    ) -> List[ReviewFeedback]:
        """
        Filter feedback based on criteria.
        
        Args:
            feedback_list: List of feedback to filter
            min_severity: Minimum severity level to include
            categories: List of categories to include (None for all)
            slide_indices: List of slide indices to include (None for all)
            
        Returns:
            Filtered list of feedback
        """
        severity_levels = {"low": 1, "medium": 2, "high": 3, "critical": 4}
        min_level = severity_levels.get(min_severity, 1)
        
        filtered = []
        for feedback in feedback_list:
            # Check severity
            feedback_level = severity_levels.get(feedback.severity, 1)
            if feedback_level < min_level:
                continue
            
            # Check category
            if categories and feedback.category not in categories:
                continue
            
            # Check slide index
            if slide_indices and feedback.slide_index not in slide_indices:
                continue
            
            filtered.append(feedback)
        
        return filtered

    def generate_improvement_plan(
        self,
        feedback_list: List[ReviewFeedback]
    ) -> Dict[str, any]:
        """
        Generate an improvement plan based on feedback.
        
        Args:
            feedback_list: List of feedback items
            
        Returns:
            Structured improvement plan
        """
        plan = {
            "total_issues": len(feedback_list),
            "by_severity": {},
            "by_category": {},
            "by_slide": {},
            "quick_wins": [],
            "major_improvements": [],
            "action_items": []
        }
        
        # Categorize feedback
        for feedback in feedback_list:
            # By severity
            plan["by_severity"][feedback.severity] = plan["by_severity"].get(feedback.severity, 0) + 1
            
            # By category
            plan["by_category"][feedback.category] = plan["by_category"].get(feedback.category, 0) + 1
            
            # By slide
            slide_key = f"slide_{feedback.slide_index}"
            if slide_key not in plan["by_slide"]:
                plan["by_slide"][slide_key] = []
            plan["by_slide"][slide_key].append(feedback.message)
            
            # Categorize by effort level
            if feedback.severity in ["low", "medium"] and len(feedback.message) < 100:
                plan["quick_wins"].append({
                    "slide": feedback.slide_index,
                    "issue": feedback.message,
                    "fix": feedback.suggestion
                })
            elif feedback.severity in ["high", "critical"]:
                plan["major_improvements"].append({
                    "slide": feedback.slide_index,
                    "issue": feedback.message,
                    "fix": feedback.suggestion
                })
        
        # Generate action items
        prioritized = self.prioritize_feedback(feedback_list)
        for i, feedback in enumerate(prioritized[:10]):  # Top 10 items
            plan["action_items"].append({
                "priority": i + 1,
                "slide": feedback.slide_index,
                "category": feedback.category,
                "action": feedback.suggestion or "Review and improve based on feedback",
                "severity": feedback.severity
            })
        
        return plan

    def get_review_summary(
        self,
        feedback_list: List[ReviewFeedback]
    ) -> Dict[str, any]:
        """
        Get a summary of the review results.
        
        Args:
            feedback_list: List of feedback items
            
        Returns:
            Review summary statistics
        """
        if not feedback_list:
            return {
                "total_feedback": 0,
                "overall_score": 10,  # Perfect if no issues
                "summary": "No issues found - presentation looks good!"
            }
        
        severity_counts = {}
        category_counts = {}
        slides_with_issues = set()
        
        for feedback in feedback_list:
            severity_counts[feedback.severity] = severity_counts.get(feedback.severity, 0) + 1
            category_counts[feedback.category] = category_counts.get(feedback.category, 0) + 1
            if feedback.slide_index >= 0:
                slides_with_issues.add(feedback.slide_index)
        
        # Calculate overall score (0-10 scale)
        severity_weights = {"critical": -3, "high": -2, "medium": -1, "low": -0.5}
        total_deduction = sum(severity_weights.get(sev, 0) * count 
                             for sev, count in severity_counts.items())
        overall_score = max(0, min(10, 10 + total_deduction))
        
        # Generate summary text
        critical_count = severity_counts.get("critical", 0)
        high_count = severity_counts.get("high", 0)
        
        if critical_count > 0:
            summary = f"{critical_count} critical issues need immediate attention"
        elif high_count > 0:
            summary = f"{high_count} high-priority improvements recommended"
        elif len(feedback_list) > 5:
            summary = "Several minor improvements would enhance the presentation"
        else:
            summary = "Good presentation with some minor refinements possible"
        
        return {
            "total_feedback": len(feedback_list),
            "severity_breakdown": severity_counts,
            "category_breakdown": category_counts,
            "slides_with_issues": len(slides_with_issues),
            "overall_score": round(overall_score, 1),
            "summary": summary
        }