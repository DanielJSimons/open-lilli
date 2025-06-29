"""Slide assembler for building PowerPoint presentations."""

import logging
import re
from pathlib import Path
from typing import Dict, List, Optional, Union

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from openai import OpenAI # For VisualProofreader initialization

from .models import (
    Outline, SlidePlan, StyleValidationConfig, FontInfo, BulletInfo, FontAdjustment,
    NativeChartData, BulletItem, ReviewFeedback, DesignIssueType, TextOverflowConfig
)
from .template_parser import TemplateParser
from .exceptions import StyleError, ValidationConfigError
from .visual_proofreader import VisualProofreader # Assuming this path is correct

logger = logging.getLogger(__name__)

RTL_LANGUAGES = ["ar", "he", "fa"]

class SlideAssembler:
    """Assembles PowerPoint presentations from templates and generated content."""

    def __init__(
        self, 
        template_parser: TemplateParser,
        validation_config: Optional[StyleValidationConfig] = None
    ):
        """
        Initialize the slide assembler.
        
        Args:
            template_parser: TemplateParser instance for template access
            validation_config: Style validation configuration
        """
        self.template_parser = template_parser
        self.max_title_length = 60
        self.max_bullet_length = 120
        self.validation_config = validation_config or StyleValidationConfig()
        
        # Initialize text overflow configuration
        if self.validation_config.text_overflow_config is None:
            self.overflow_config = TextOverflowConfig()
            logger.info("TextOverflowConfig not provided, initialized with default settings.")
        else:
            self.overflow_config = self.validation_config.text_overflow_config
            logger.info("TextOverflowConfig loaded from validation_config.")

        logger.info(f"Text overflow - Enable bullet splitting: {self.overflow_config.enable_bullet_splitting}")
        logger.info(f"Text overflow - Max lines per placeholder: {self.overflow_config.max_lines_per_placeholder}")

        # Conditionally initialize the VisualProofreader based on configuration.
        # This allows for LLM-based checks if enabled.
        self.visual_proofreader = None
        if self.validation_config.enable_visual_proofreader:
            try:
                # Note: OpenAI client initialization. For production, consider a shared client
                # instance managed at a higher application level. API keys are expected
                # to be available in the environment (e.g., OPENAI_API_KEY).
                openai_client = OpenAI()
                self.visual_proofreader = VisualProofreader(
                    client=openai_client,
                    model=self.validation_config.visual_proofreader_model,
                    temperature=self.validation_config.visual_proofreader_temperature
                    # enable_corrections for VisualProofreader is passed during the proofread_slides call
                )
                logger.info(f"VisualProofreader initialized with model '{self.validation_config.visual_proofreader_model}'.")
            except Exception as e:
                # Gracefully handle VisualProofreader initialization errors.
                # Validation will proceed with only rule-based checks if VP fails to initialize.
                logger.error(f"Failed to initialize VisualProofreader: {e}. Visual proofreading will be disabled.")
                self.visual_proofreader = None

        logger.info("SlideAssembler initialized")
        logger.debug(f"Style validation: {self.validation_config.enabled}, mode: {self.validation_config.mode}")
        logger.debug(f"Visual Proofreader enabled: {self.validation_config.enable_visual_proofreader}, instance: {'present' if self.visual_proofreader else 'None'}")

    def assemble(
        self,
        outline: Outline,
        slides: List[SlidePlan],
        visuals: Optional[Dict[int, Dict[str, str]]] = None,
        output_path: Union[str, Path] = "output.pptx"
    ) -> Path:
        """
        Assemble a complete PowerPoint presentation.
        
        Args:
            outline: Presentation outline with metadata
            slides: List of slide plans with content
            visuals: Dictionary mapping slide indices to visual file paths
            output_path: Path to save the output presentation
            
        Returns:
            Path to the assembled presentation file
        """
        output_path = Path(output_path)
        visuals = visuals or {}
        language = outline.language # Get language from outline
        
        logger.info(f"Assembling presentation with {len(slides)} slides in language: {language}")
        logger.info(f"Output path: {output_path}")
        
        # Create new presentation from template
        prs = Presentation(str(self.template_parser.template_path))
        
        # Clear any existing slides from the template
        # Iterate backwards to safely delete
        for i in range(len(prs.slides) - 1, -1, -1):
            r_id = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[i]
        
        # Create a mutable list of slide plans to process
        slide_plans_to_process: List[SlidePlan] = list(slides) # Operate on a copy

        current_idx = 0
        while current_idx < len(slide_plans_to_process):
            current_slide_plan = slide_plans_to_process[current_idx]

            logger.info(f"Processing slide plan for original index {current_slide_plan.index}, current title: '{current_slide_plan.title}' (List length: {len(slide_plans_to_process)})")

            try:
                # _add_slide now returns a list of new SlidePlan objects for *additional* slides.
                # The current_slide_plan is processed and added to prs within _add_slide.
                additional_slide_plans = self._add_slide(
                    prs,
                    current_slide_plan,
                    visuals.get(current_slide_plan.index, {}), # Use original index for visuals
                    language
                )

                if additional_slide_plans:
                    logger.info(f"Original slide {current_slide_plan.index} ('{current_slide_plan.title}') generated {len(additional_slide_plans)} additional slide(s) due to content splitting.")
                    for i, new_plan in enumerate(additional_slide_plans):
                        # Insert new plans into the processing list right after the current one.
                        # Their 'index' field is initially -1 or based on the parent.
                        # The actual slide index in the final presentation is determined by its order in prs.slides.
                        slide_plans_to_process.insert(current_idx + 1 + i, new_plan)
                        logger.debug(f"Inserted new slide plan '{new_plan.title}' to be processed next.")
                    logger.debug(f"Slide plans to process now has {len(slide_plans_to_process)} items.")

            except Exception as e:
                logger.error(f"Failed to process or add slide from plan (original index {current_slide_plan.index}, title '{current_slide_plan.title}'): {e}")
                # Optionally, add a fallback slide to prs if _add_slide failed catastrophically before creating a pptx slide object.
                # However, _add_slide is expected to handle its own errors and add a basic slide.
                # If _add_slide itself throws, it means the pptx slide object wasn't even created.
                # Consider if a global fallback is needed here, or if _add_slide's internal fallback is sufficient.
                # For now, we assume _add_slide handles its own errors including adding a fallback.

            current_idx += 1
        
        # Apply presentation metadata
        self._apply_metadata(prs, outline)
        
        # Validate presentation style before saving
        self.validate_presentation_style(prs)
        
        # Save presentation
        prs.save(str(output_path))
        
        logger.info(f"Successfully assembled presentation: {output_path}")
        return output_path

    def patch_existing_presentation(
        self,
        input_pptx: Path,
        updated_slides: List[SlidePlan],
        target_indices: List[int],
        language: str, # Add language parameter
        visuals: Optional[Dict[int, Dict[str, str]]] = None,
        output_path: Union[str, Path] = "updated.pptx"
    ) -> Path:
        """
        Patch specific slides in an existing presentation.
        
        Args:
            input_pptx: Path to existing presentation
            updated_slides: List of all slide plans (with some updated)
            target_indices: List of slide indices that were updated
            language: Language code for the presentation
            visuals: Dictionary mapping slide indices to visual file paths
            output_path: Path to save the patched presentation
            
        Returns:
            Path to the patched presentation file
        """
        output_path = Path(output_path)
        visuals = visuals or {}
        
        logger.info(f"Patching presentation with {len(target_indices)} updated slides in language: {language}")
        logger.info(f"Target indices: {target_indices}")
        logger.info(f"Output path: {output_path}")
        
        # Load existing presentation
        prs = Presentation(str(input_pptx))
        
        # Verify slide count matches
        if len(prs.slides) != len(updated_slides):
            logger.warning(
                f"Slide count mismatch: existing {len(prs.slides)}, "
                f"updated {len(updated_slides)}"
            )
        
        # Replace only the target slides
        for target_idx in target_indices:
            if target_idx >= len(prs.slides):
                logger.error(f"Target index {target_idx} exceeds slide count {len(prs.slides)}")
                continue
                
            if target_idx >= len(updated_slides):
                logger.error(f"Target index {target_idx} exceeds updated slides count {len(updated_slides)}")
                continue
            
            try:
                # Get the slide plan for this index
                slide_plan = updated_slides[target_idx]
                
                # Replace the slide at this index
                self._replace_slide_at_index(
                    prs, target_idx, slide_plan, visuals.get(target_idx, {}), language # Pass language
                )
                
                logger.debug(f"Replaced slide {target_idx}: {slide_plan.title}")
                
            except Exception as e:
                logger.error(f"Failed to replace slide {target_idx}: {e}")
        
        # Validate patched presentation style
        self.validate_presentation_style(prs)
        
        # Save patched presentation
        prs.save(str(output_path))
        
        logger.info(f"Successfully patched presentation: {output_path}")
        return output_path

    def _replace_slide_at_index(
        self,
        prs: Presentation,
        slide_idx: int,
        slide_plan: SlidePlan,
        slide_visuals: Dict[str, str],
        language: str, # Add language parameter
        shape_whitelist: Optional[List[str]] = None
    ) -> None:
        """
        Replace a slide at a specific index with new content.
        
        Args:
            prs: PowerPoint presentation object
            slide_idx: Index of slide to replace
            slide_plan: New slide plan with content
            slide_visuals: Visual assets for the slide
            language: Language code for the presentation
            shape_whitelist: Optional list of shape names to preserve during clearing
        """
        # Get the existing slide
        existing_slide = prs.slides[slide_idx]
        
        # Get the layout that matches our slide plan
        layout_index = slide_plan.layout_id or 0
        if layout_index >= len(prs.slide_layouts):
            logger.warning(f"Layout index {layout_index} out of range, using existing layout")
            layout = existing_slide.slide_layout
        else:
            layout = prs.slide_layouts[layout_index]
        
        # Clear existing content while preserving slide structure
        self._clear_slide_content(existing_slide, shape_whitelist=shape_whitelist)
        
        # Apply new layout if different
        if layout != existing_slide.slide_layout:
            # Note: Changing layout is complex in python-pptx
            # For now, we'll work with the existing layout
            logger.debug(f"Layout change requested but using existing layout for slide {slide_idx}")
        
        # Add new content
        self._add_title(existing_slide, slide_plan.title, language) # Pass language
        
        # Add content based on slide type
        if slide_plan.slide_type == "title":
            self._add_title_slide_content(existing_slide, slide_plan, language) # Pass language
        elif slide_plan.bullet_hierarchy is not None or slide_plan.bullets:
            additional_plans_from_bullets: List[SlidePlan] = []
            if slide_plan.bullet_hierarchy is not None:
                additional_plans_from_bullets = self._add_hierarchical_bullet_content(existing_slide, slide_plan, language)
            elif slide_plan.bullets: # Ensure bullets exist before calling
                additional_plans_from_bullets = self._add_bullet_content(existing_slide, slide_plan, language)

            if additional_plans_from_bullets:
                logger.warning(
                    f"Text overflow in slide {slide_idx} during patch operation would result in new slides. "
                    "This is not fully supported in patch mode. Extra content will be on the new slide plan "
                    "but not inserted into this patched presentation."
                )
        
        # Add visuals
        if "native_chart" in slide_visuals and slide_plan.chart_data:
            self._add_native_chart(existing_slide, slide_plan.chart_data)
        elif "chart" in slide_visuals:
            chart_alt_text = None
            if slide_plan.chart_data:
                if isinstance(slide_plan.chart_data, NativeChartData) and slide_plan.chart_data.title:
                    chart_alt_text = slide_plan.chart_data.title
                elif isinstance(slide_plan.chart_data, dict) and slide_plan.chart_data.get("title"):
                    chart_alt_text = slide_plan.chart_data.get("title")
            if not chart_alt_text:
                chart_alt_text = slide_plan.image_alt_text or "Chart" # Fallback
            self._add_chart_image(existing_slide, slide_visuals["chart"], alt_text=chart_alt_text)
        
        if "process_flow" in slide_visuals:
            # Assuming process flow is like an image, use image_alt_text or a generic one
            self._add_process_flow(existing_slide, slide_visuals["process_flow"], alt_text=slide_plan.image_alt_text or "Process flow diagram")
        
        if "image" in slide_visuals:
            self._add_image(existing_slide, slide_visuals["image"], alt_text=slide_plan.image_alt_text)
        
        # Update speaker notes
        if slide_plan.speaker_notes:
            self._add_speaker_notes(existing_slide, slide_plan.speaker_notes)

    def _clear_slide_content(self, slide, shape_whitelist: Optional[List[str]] = None) -> None:
        """
        Clear content from a slide while preserving structure.
        
        Args:
            slide: PowerPoint slide object to clear
            shape_whitelist: Optional list of shape names to preserve
        """
        try:
            # Clear text content from placeholders
            for placeholder in slide.placeholders:
                if hasattr(placeholder, 'text_frame') and placeholder.text_frame:
                    placeholder.text_frame.clear()
                elif hasattr(placeholder, 'text'):
                    placeholder.text = ""
            
            # Remove non-placeholder shapes (like added images/charts)
            shapes_to_remove = []
            for shape in slide.shapes:
                # Only remove shapes that are not placeholders
                if not hasattr(shape, 'placeholder_format'):
                    # If a whitelist is provided, don't remove shapes in the whitelist
                    if shape_whitelist and shape.name in shape_whitelist:
                        logger.debug(f"Skipping removal of whitelisted shape: {shape.name}")
                        continue
                    shapes_to_remove.append(shape)
            
            # Remove collected shapes
            for shape in shapes_to_remove:
                try:
                    slide.shapes._spTree.remove(shape._element)
                except Exception as e:
                    logger.debug(f"Could not remove shape: {e}")
            
            # Clear speaker notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if hasattr(notes_slide, 'notes_text_frame'):
                    notes_slide.notes_text_frame.clear()
            
            logger.debug("Cleared slide content")
            
        except Exception as e:
            logger.error(f"Failed to clear slide content: {e}")

    def _add_slide(
        self,
        prs: Presentation,
        slide_plan: SlidePlan,
        slide_visuals: Dict[str, str],
        language: str # Add language parameter
    ) -> List[SlidePlan]:
        """
        Adds a single slide to the presentation based on slide_plan.
        Returns a list of any additional SlidePlan objects created due to splitting.
        """
        all_newly_created_slide_plans: List[SlidePlan] = []
        
        # Get the appropriate layout
        layout_index = slide_plan.layout_id or 0
        if layout_index >= len(prs.slide_layouts):
            logger.warning(f"Layout index {layout_index} out of range, using 0")
            layout_index = 0
        
        layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(layout)
        
        # Log if content was summarized by LLM
        if hasattr(slide_plan, 'summarized_by_llm') and slide_plan.summarized_by_llm:
            logger.info(f"Slide {slide_plan.index} content was summarized by LLM.")

        # Add title
        self._add_title(slide, slide_plan.title, language) # Pass language
        
        # Add content based on slide type
        if slide_plan.slide_type == "title":
            # Assuming _add_title_slide_content also needs language for RTL subtitle
            self._add_title_slide_content(slide, slide_plan, language) # Pass language
        elif slide_plan.bullet_hierarchy is not None or slide_plan.bullets:
            new_plans: List[SlidePlan] = []
            # T-100: Use hierarchical bullet content if available, otherwise fall back to legacy
            if slide_plan.bullet_hierarchy is not None:
                new_plans = self._add_hierarchical_bullet_content(slide, slide_plan, language)
            elif slide_plan.bullets: # Ensure bullets exist before calling _add_bullet_content
                # Pass the slide_plan object to _add_bullet_content
                new_plans = self._add_bullet_content(slide, slide_plan, language)

            if new_plans:
                all_newly_created_slide_plans.extend(new_plans)
        
        # Add visuals
        if "native_chart" in slide_visuals and slide_plan.chart_data:
            self._add_native_chart(slide, slide_plan.chart_data)
        elif "chart" in slide_visuals:
            chart_alt_text = None
            if slide_plan.chart_data:
                if isinstance(slide_plan.chart_data, NativeChartData) and slide_plan.chart_data.title:
                    chart_alt_text = slide_plan.chart_data.title
                elif isinstance(slide_plan.chart_data, dict) and slide_plan.chart_data.get("title"):
                    chart_alt_text = slide_plan.chart_data.get("title")
            if not chart_alt_text: # Fallback if no specific chart title
                chart_alt_text = slide_plan.image_alt_text or "Chart"
            self._add_chart_image(slide, slide_visuals["chart"], alt_text=chart_alt_text)
        
        if "process_flow" in slide_visuals:
            # Assuming process flow is like an image, use image_alt_text or a generic one
            self._add_process_flow(slide, slide_visuals["process_flow"], alt_text=slide_plan.image_alt_text or "Process flow diagram")
        
        if "image" in slide_visuals:
            self._add_image(slide, slide_visuals["image"], alt_text=slide_plan.image_alt_text)
        
        # Check for RTL language and skip font adjustment if necessary
        if language.lower() in RTL_LANGUAGES:
            logger.info(f"Slide {slide_plan.index} is in RTL language ({language}), skipping font adjustment.")
            # Clear any existing font adjustment to prevent it from being applied
            if hasattr(slide_plan, '__dict__') and 'font_adjustment' in slide_plan.__dict__:
                slide_plan.__dict__['font_adjustment'] = None

        # Perform post-processing: remove unwanted shapes and hide empty placeholders
        # Passing None for shape_whitelist as this is a new slide.
        self._postprocess_slide(slide, shape_whitelist=None)

        # Apply font adjustments if needed (will be skipped if cleared above)
        self._apply_font_adjustments(slide, slide_plan)
        
        # Add speaker notes
        if slide_plan.speaker_notes:
            self._add_speaker_notes(slide, slide_plan.speaker_notes)
        
        return all_newly_created_slide_plans

    def _postprocess_slide(self, slide, shape_whitelist: Optional[List[str]] = None) -> None:
        """
        Perform post-processing on a slide after content has been added.
        This includes removing unwanted non-placeholder shapes and hiding empty placeholders.

        Args:
            slide: PowerPoint slide object to post-process.
            shape_whitelist: Optional list of shape names to preserve from removal.
        """
        logger.debug(f"Starting post-processing for slide (ID: {slide.slide_id if hasattr(slide, 'slide_id') else 'N/A'})...")

        # 1. Remove non-placeholder shapes (like added images/charts not in whitelist)
        # This logic is adapted from _clear_slide_content but only for non-placeholders
        shapes_to_remove = []
        for shape in slide.shapes:
            if not hasattr(shape, 'placeholder_format'): # It's a non-placeholder shape
                if shape_whitelist and shape.name in shape_whitelist:
                    logger.debug(f"Post-process: Skipping removal of whitelisted non-placeholder shape: {shape.name}")
                    continue
                shapes_to_remove.append(shape)
                logger.debug(f"Post-process: Queuing non-placeholder shape for removal: {shape.name or 'Unnamed shape'}")

        removed_count = 0
        for shape in shapes_to_remove:
            try:
                slide.shapes._spTree.remove(shape._element)
                removed_count +=1
            except Exception as e:
                logger.debug(f"Post-process: Could not remove shape {shape.name or 'Unnamed shape'}: {e}")

        if removed_count > 0:
            logger.debug(f"Post-process: Removed {removed_count} non-placeholder shapes.")

        # 2. Remove/hide any empty placeholders to avoid style warnings
        hidden_placeholder_count = self._remove_empty_placeholders_from_slide(slide)
        if hidden_placeholder_count > 0:
            logger.debug(f"Post-process: Hid {hidden_placeholder_count} empty placeholders on slide.")

        logger.debug(f"Finished post-processing for slide (ID: {slide.slide_id if hasattr(slide, 'slide_id') else 'N/A'}).")


    def _add_title(self, slide, title: str, language: str) -> None: # Add language
        """Add title to slide."""
        try:
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                title_shape = slide.shapes.title
                
                # Truncate title if too long
                if len(title) > self.max_title_length:
                    title = title[:self.max_title_length - 3] + "..."
                
                title_shape.text = title
                
                # Apply basic formatting
                if title_shape.text_frame.paragraphs:
                    paragraph = title_shape.text_frame.paragraphs[0]
                    if paragraph.runs:
                        run = paragraph.runs[0]
                        run.font.bold = True
                        
                        # Proactive font selection
                        if self.template_parser.template_style:
                            specific_font_name = self.template_parser.template_style.language_specific_fonts.get(language.lower())
                            if specific_font_name:
                                run.font.name = specific_font_name
                                logger.debug(f"Applied language-specific font '{specific_font_name}' for title in language '{language.lower()}'")

                # RTL alignment for title
                if language.lower() in RTL_LANGUAGES:
                    if title_shape.text_frame and title_shape.text_frame.paragraphs:
                        for para in title_shape.text_frame.paragraphs:
                            para.alignment = PP_ALIGN.RIGHT

                # Apply auto-fit text if enabled
                if self.overflow_config.enable_auto_fit_text:
                    if hasattr(title_shape, 'text_frame') and title_shape.text_frame:
                        title_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        logger.debug(f"Applied TEXT_TO_FIT_SHAPE to title placeholder due to enable_auto_fit_text=True.")

                logger.debug(f"Added title: {title}")
            else:
                logger.warning("No title placeholder found in slide")
                
        except Exception as e:
            logger.error(f"Failed to add title '{title}': {e}")

    def _add_title_slide_content(self, slide, slide_plan: SlidePlan, language: str) -> None: # Add language
        """Add content specific to title slides."""
        try:
            # Try to find subtitle placeholder
            subtitle_placeholder = None
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:  # SUBTITLE
                    subtitle_placeholder = placeholder
                    break
            
            if subtitle_placeholder:
                subtitle_text = ""
                if slide_plan.bullets: # Assuming bullets might be used as subtitle parts
                    subtitle_text = " • ".join(slide_plan.bullets[:2])
                elif hasattr(slide_plan, 'subtitle') and slide_plan.subtitle:
                    subtitle_text = slide_plan.subtitle
                
                if subtitle_text:
                    subtitle_placeholder.text = subtitle_text
                    # Apply RTL alignment for subtitle if needed
                    if language.lower() in RTL_LANGUAGES:
                        if subtitle_placeholder.text_frame and subtitle_placeholder.text_frame.paragraphs:
                            for para in subtitle_placeholder.text_frame.paragraphs:
                                para.alignment = PP_ALIGN.RIGHT
                    logger.debug(f"Added subtitle: {subtitle_text}")
                else:
                    # T-93: Hide empty subtitle placeholder to avoid style warning
                    self._hide_empty_placeholder(subtitle_placeholder, "subtitle")
                    logger.debug("Hidden empty subtitle placeholder to avoid style warning")
                    
        except Exception as e:
            logger.error(f"Failed to add title slide content: {e}")

    def _add_bullet_content(self, slide, slide_plan: SlidePlan, language: str) -> List[SlidePlan]: # Add language
        """
        Add bullet points to slide using intelligent content distribution.
        This method handles legacy flat list of bullets with improved placeholder utilization.
        """
        newly_created_slides: List[SlidePlan] = []
        
        # Use intelligent content distribution first
        distribution_result = self._intelligently_distribute_content(slide, slide_plan, language)
        
        # Check if content was successfully distributed
        if distribution_result["bullets_placed"]:
            logger.info(f"Content successfully distributed across available placeholders")
            return newly_created_slides
        
        # Fallback to original logic if intelligent distribution failed
        logger.warning("Intelligent distribution failed, falling back to legacy method")
        
        bullets: List[str] = slide_plan.bullets # Get bullets from slide_plan

        # Find BODY placeholders (legacy fallback)
        body_placeholders = []
        for placeholder in slide.placeholders:
            ph_type = placeholder.placeholder_format.type
            if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                body_placeholders.append(placeholder)

        if not body_placeholders:
            logger.warning("No BODY/OBJECT placeholder found, looking for any text placeholder")
            for placeholder in slide.placeholders:
                if hasattr(placeholder, 'text_frame'):
                    body_placeholders.append(placeholder)
                    break

        if not body_placeholders:
            logger.warning(f"Slide {slide_plan.index}: No suitable placeholder found for bullet content.")
            return []

        # Original logic path for no splitting or if splitting is disabled
        def _apply_original_bullet_logic():
            if len(body_placeholders) >= 2 and len(bullets) > 1:
                logger.debug(f"Slide {slide_plan.index}: Applying two-column distribution for {len(bullets)} bullets.")
                self._distribute_bullets_across_columns(body_placeholders, bullets, language)
            else:
                logger.debug(f"Slide {slide_plan.index}: Applying single-column distribution for {len(bullets)} bullets.")
                self._add_bullets_to_placeholder(body_placeholders[0], bullets, language)
            return []

        if not self.overflow_config.enable_bullet_splitting:
            logger.info(f"Slide {slide_plan.index}: Bullet splitting is disabled. Applying original bullet logic.")
            return _apply_original_bullet_logic()

        if not bullets:
            return []

        bullet_items = [BulletItem(text=b, level=0) for b in bullets]
        primary_placeholder = body_placeholders[0]
        
        placeholder_width_inches = (primary_placeholder.width / 914400.0) if primary_placeholder.width else 6.0
        base_font_size_pt = 18.0
        max_lines_allowed = self.overflow_config.max_lines_per_placeholder

        estimated_total_lines = self._estimate_lines_for_bullets(
            bullet_items, placeholder_width_inches, base_font_size_pt
        )
        logger.info(f"Slide {slide_plan.index} (_add_bullet_content): Estimated total lines for {len(bullet_items)} bullets: {estimated_total_lines}. Max allowed: {max_lines_allowed}.")

        # If content fits OR if there are multiple body placeholders (implying a two-column layout where native distribution is preferred unless single column overflows)
        # The condition `len(body_placeholders) >= 2` ensures that if it's a two-column layout,
        # we prefer the original distribution logic unless the content is so large it would overflow even one column.
        if estimated_total_lines <= max_lines_allowed or len(body_placeholders) >= 2:
            if len(body_placeholders) >= 2:
                 logger.debug(f"Slide {slide_plan.index}: Content fits single column estimate or is two-column layout. Applying original distribution logic.")
            else:
                 logger.debug(f"Slide {slide_plan.index}: Content fits single column estimate. Adding all bullets.")
            return _apply_original_bullet_logic()
        else:
            # Splitting is needed for a single column equivalent
            logger.info(f"Slide {slide_plan.index}: Splitting needed for {len(bullets)} legacy bullets (single column focus).")
            current_slide_bullets_text: List[str] = []
            lines_on_current_slide = 0
            split_idx = 0

            for i, bullet_item in enumerate(bullet_items):
                lines_for_this_bullet = self._estimate_lines_for_bullets(
                    [bullet_item], placeholder_width_inches, base_font_size_pt
                )
                if not current_slide_bullets_text or (lines_on_current_slide + lines_for_this_bullet <= max_lines_allowed):
                    current_slide_bullets_text.append(bullets[i]) # Use original string bullet
                    lines_on_current_slide += lines_for_this_bullet
                    split_idx = i + 1
                else:
                    break
            
            if current_slide_bullets_text:
                self._add_bullets_to_placeholder(primary_placeholder, current_slide_bullets_text, language)
                logger.info(f"Slide {slide_plan.index}: Added {len(current_slide_bullets_text)} bullets ({lines_on_current_slide} lines) to current slide's primary placeholder.")
            
            remaining_bullets_text = bullets[split_idx:]
            if remaining_bullets_text:
                logger.info(f"Slide {slide_plan.index}: {len(remaining_bullets_text)} legacy bullets remaining for new slide(s).")
                # Try to get current slide's layout ID for the new slide
                try:
                    current_layout_id = slide.slide_layout.slide_master.slide_layouts.index(slide.slide_layout)
                except ValueError: # If layout not found in master's list (should not happen)
                    logger.warning(f"Could not determine layout ID for slide {slide_plan.index}, defaulting to original slide_plan.layout_id or 0.")
                    current_layout_id = slide_plan.layout_id or 0


                new_slide_plan_data = {
                    "index": -1, # Placeholder
                    "title": f"{slide_plan.title} {self.overflow_config.split_slide_title_suffix}",
                    "layout_id": current_layout_id,
                    "slide_type": slide_plan.slide_type, # Keep original type or use "content"
                    "bullets": remaining_bullets_text,   # Legacy string list for the new slide
                    "bullet_hierarchy": None,            # Explicitly None
                    "image_query": slide_plan.image_query,
                    "chart_data": slide_plan.chart_data,
                    "speaker_notes": slide_plan.speaker_notes or "Content continued from previous slide.",
                    "summarized_by_llm": slide_plan.summarized_by_llm,
                    "font_adjustment": None,
                    "image_alt_text": slide_plan.image_alt_text,
                    "needs_splitting": False
                }
                try:
                    new_slide = SlidePlan(**new_slide_plan_data)
                    newly_created_slides.append(new_slide)
                    logger.info(f"Created new SlidePlan for continued legacy bullet content from slide {slide_plan.index} with title '{new_slide.title}'.")
                except Exception as e:
                    logger.error(f"Error creating new SlidePlan (legacy bullets) during split: {e}. Data: {new_slide_plan_data}")

            return newly_created_slides

    def _add_hierarchical_bullet_content(self, slide, slide_plan: SlidePlan, language: str) -> List[SlidePlan]:
        """
        Add hierarchical bullet content to slide (T-100), potentially splitting it
        across multiple slides if it overflows.

        Args:
            slide: PowerPoint slide object for the current slide.
            slide_plan: SlidePlan for the current slide.
            language: Language code for text formatting.
            
        Returns:
            List[SlidePlan]: A list of newly created SlidePlan objects if splitting occurred.
                             Returns an empty list if no splitting was necessary.
        """
        newly_created_slides: List[SlidePlan] = []

        # Try intelligent content distribution first (regardless of splitting setting)
        distribution_result = self._intelligently_distribute_content(slide, slide_plan, language)
        
        # If intelligent distribution succeeded, we're done
        if distribution_result["bullets_placed"]:
            logger.info(f"Slide {slide_plan.index}: Content successfully distributed using intelligent placement")
            return newly_created_slides
        
        # If intelligent distribution failed, fall back to original logic
        logger.warning(f"Slide {slide_plan.index}: Intelligent distribution failed, using legacy hierarchical bullet logic")
        
        if not self.overflow_config.enable_bullet_splitting:
            logger.info("Bullet splitting is disabled. Adding all bullets to current slide using fallback method.")
            # Find and add to placeholder without splitting (original logic path)
            bullet_items_no_split = slide_plan.get_effective_bullets()
            if not bullet_items_no_split:
                return []

            body_placeholders_no_split = []
            for placeholder_no_split in slide.placeholders:
                ph_type_no_split = placeholder_no_split.placeholder_format.type
                if ph_type_no_split in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                    body_placeholders_no_split.append(placeholder_no_split)

            if not body_placeholders_no_split: # Fallback if no BODY/OBJECT
                for placeholder_no_split in slide.placeholders:
                    if hasattr(placeholder_no_split, 'text_frame'):
                        body_placeholders_no_split.append(placeholder_no_split)
                        break
            
            if body_placeholders_no_split:
                self._add_hierarchical_bullets_to_placeholder(body_placeholders_no_split[0], bullet_items_no_split, language)
            else:
                logger.warning(f"Slide {slide_plan.index}: No suitable placeholder found for bullets when splitting is disabled.")
            return []

        bullet_items = slide_plan.get_effective_bullets()
        if not bullet_items:
            return []

        # Find the primary body placeholder
        placeholder_to_use = None
        for ph_candidate in slide.placeholders:
            ph_type_candidate = ph_candidate.placeholder_format.type
            if ph_type_candidate in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                placeholder_to_use = ph_candidate
                break
        if not placeholder_to_use: # Fallback
            for ph_candidate in slide.placeholders:
                if hasattr(ph_candidate, 'text_frame'):
                    placeholder_to_use = ph_candidate
                    break

        if not placeholder_to_use:
            logger.warning(f"Slide {slide_plan.index}: No suitable body placeholder found for hierarchical bullets. Cannot apply splitting logic.")
            # Add all bullets to the first available text frame as a last resort if no splitting can be determined
            if slide.shapes.placeholders: # Check if there are any placeholders
                 first_ph_with_tf = None
                 for ph_any in slide.shapes.placeholders:
                     if hasattr(ph_any, 'text_frame'):
                         first_ph_with_tf = ph_any
                         break
                 if first_ph_with_tf:
                    self._add_hierarchical_bullets_to_placeholder(first_ph_with_tf, bullet_items, language)
                 else: # Truly no text frame anywhere
                    logger.error(f"Slide {slide_plan.index}: No placeholder with a text_frame found. Bullets cannot be added.")
            return []

        # Placeholder and Font Information (rudimentary for now)
        # 1 inch = 914400 EMUs. python-pptx shape dimensions are in EMUs.
        placeholder_width_inches = (placeholder_to_use.width / 914400.0) if placeholder_to_use.width else 6.0 # Default if width is 0 or None
        # placeholder_height_inches = (placeholder_to_use.height / 914400.0) if placeholder_to_use.height else 4.0 # Default
        base_font_size_pt = 18.0 # Default base font size

        # Max lines based on configuration
        max_lines_allowed = self.overflow_config.max_lines_per_placeholder
        logger.debug(f"Slide {slide_plan.index}: Max lines allowed per placeholder (config): {max_lines_allowed}")

        estimated_total_lines = self._estimate_lines_for_bullets(
            bullet_items, placeholder_width_inches, base_font_size_pt
        )
        logger.info(f"Slide {slide_plan.index}: Estimated total lines for {len(bullet_items)} bullets: {estimated_total_lines}. Max allowed: {max_lines_allowed}.")

        if estimated_total_lines <= max_lines_allowed:
            self._add_hierarchical_bullets_to_placeholder(placeholder_to_use, bullet_items, language)
            logger.debug(f"Slide {slide_plan.index}: All {len(bullet_items)} bullets fit, no splitting needed.")
            return []
        else:
            logger.info(f"Slide {slide_plan.index}: Splitting needed for {len(bullet_items)} bullets.")
            bullets_for_current_slide: List[BulletItem] = []
            lines_on_current_slide = 0
            split_index = 0 # Marks the start of bullets for the *next* slide

            for i, bullet_item in enumerate(bullet_items):
                lines_for_this_bullet = self._estimate_lines_for_bullets(
                    [bullet_item], placeholder_width_inches, base_font_size_pt
                )

                # Ensure at least one bullet is on the slide, even if it exceeds max_lines.
                if not bullets_for_current_slide or (lines_on_current_slide + lines_for_this_bullet <= max_lines_allowed) :
                    bullets_for_current_slide.append(bullet_item)
                    lines_on_current_slide += lines_for_this_bullet
                    split_index = i + 1
                else:
                    # This bullet causes overflow for the current set.
                    # split_index is already correctly set to the start of this overflowing bullet.
                    break

            # Add the determined bullets to the current slide's placeholder
            if bullets_for_current_slide:
                self._add_hierarchical_bullets_to_placeholder(placeholder_to_use, bullets_for_current_slide, language)
                logger.info(f"Slide {slide_plan.index}: Added {len(bullets_for_current_slide)} bullets ({lines_on_current_slide} lines) to current slide.")
            else: # Should not happen if bullet_items is not empty due to "if not bullets_for_current_slide" logic
                logger.warning(f"Slide {slide_plan.index}: No bullets added to current slide during split. This indicates an issue.")


            remaining_bullets = bullet_items[split_index:]
            if remaining_bullets:
                logger.info(f"Slide {slide_plan.index}: {len(remaining_bullets)} bullets remaining for new slide(s).")
                new_slide_plan_data = {
                    "index": -1, # Placeholder index; to be managed by caller
                    "title": f"{slide_plan.title} {self.overflow_config.split_slide_title_suffix}",
                    "layout_id": slide_plan.layout_id,
                    "slide_type": slide_plan.slide_type, # Or a default like "content"
                    "bullet_hierarchy": remaining_bullets,
                    "bullets": [], # Clear legacy bullets field
                    "image_query": slide_plan.image_query, # Preserve other fields
                    "chart_data": slide_plan.chart_data,
                    "speaker_notes": slide_plan.speaker_notes or "Content continued from previous slide.",
                    "summarized_by_llm": slide_plan.summarized_by_llm,
                    "font_adjustment": None, # Typically None for a new slide from split
                    "image_alt_text": slide_plan.image_alt_text, # Preserve
                    "needs_splitting": False # Assume this new plan might or might not need further splitting
                }
                try:
                    new_slide = SlidePlan(**new_slide_plan_data)
                    newly_created_slides.append(new_slide)
                    logger.info(f"Created new SlidePlan for continued content from slide {slide_plan.index} with title '{new_slide.title}'.")
                except Exception as e:
                    logger.error(f"Error creating new SlidePlan during split: {e}. Data: {new_slide_plan_data}")
            
            return newly_created_slides
                
        # Fallback if logic somehow doesn't return earlier (should not happen)
        logger.debug(f"Slide {slide_plan.index}: Reached end of _add_hierarchical_bullet_content, returning {len(newly_created_slides)} new slides.")
        return newly_created_slides

    def _distribute_bullets_across_columns(
        self, 
        body_placeholders: List, 
        bullets: List[str], 
        language: str
    ) -> None:
        """
        Distribute bullets evenly across multiple BODY placeholders (T-94).
        
        Args:
            body_placeholders: List of BODY placeholder objects
            bullets: List of bullet text strings
            language: Language code for text formatting
        """
        num_columns = min(len(body_placeholders), 2)  # Limit to 2 columns for now
        bullets_per_column = len(bullets) // num_columns
        extra_bullets = len(bullets) % num_columns
        
        start_idx = 0
        for col_idx in range(num_columns):
            # Distribute extra bullets to first columns
            bullets_in_this_column = bullets_per_column + (1 if col_idx < extra_bullets else 0)
            end_idx = start_idx + bullets_in_this_column
            
            column_bullets = bullets[start_idx:end_idx]
            
            if column_bullets:  # Only add if there are bullets for this column
                self._add_bullets_to_placeholder(body_placeholders[col_idx], column_bullets, language)
                logger.debug(f"Added {len(column_bullets)} bullets to column {col_idx + 1}")
            
            start_idx = end_idx
        
        logger.info(f"Distributed {len(bullets)} bullets across {num_columns} columns")
    
    def _discover_all_content_placeholders(self, slide) -> Dict[str, List]:
        """
        Discover ALL available placeholders that can hold content.
        
        Args:
            slide: Slide object to analyze
            
        Returns:
            Dictionary categorizing all available content placeholders
        """
        placeholders = {
            "primary_content": [],    # BODY, OBJECT - main content areas
            "secondary_content": [],  # SUBTITLE, FOOTER, HEADER - supplementary content
            "title": [],              # TITLE placeholders
            "media": [],              # PICTURE, CHART placeholders
            "metadata": [],           # DATE, SLIDE_NUMBER placeholders
            "all_text_capable": []   # All placeholders that can hold text
        }
        
        try:
            for placeholder in slide.placeholders:
                ph_type = placeholder.placeholder_format.type
                
                # Categorize by function
                if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                    placeholders["primary_content"].append(placeholder)
                elif ph_type in (PP_PLACEHOLDER.SUBTITLE, 10, 9):  # SUBTITLE, FOOTER, HEADER
                    placeholders["secondary_content"].append(placeholder)
                elif ph_type in (PP_PLACEHOLDER.TITLE, 13):  # TITLE, CENTERED_TITLE
                    placeholders["title"].append(placeholder)
                elif ph_type in (PP_PLACEHOLDER.PICTURE, 14):  # PICTURE, CHART
                    placeholders["media"].append(placeholder)
                elif ph_type in (8, 11):  # DATE, SLIDE_NUMBER
                    placeholders["metadata"].append(placeholder)
                
                # Check if placeholder can hold text
                if hasattr(placeholder, 'text_frame') and placeholder.text_frame is not None:
                    placeholders["all_text_capable"].append(placeholder)
                    
        except Exception as e:
            logger.error(f"Error discovering placeholders: {e}")
            
        logger.debug(f"Discovered placeholders: {len(placeholders['primary_content'])} primary, "
                    f"{len(placeholders['secondary_content'])} secondary, "
                    f"{len(placeholders['all_text_capable'])} total text-capable")
        
        return placeholders
    
    def _intelligently_distribute_content(self, slide, slide_plan: SlidePlan, language: str) -> Dict[str, bool]:
        """
        Intelligently distribute slide content across ALL available placeholders.
        
        Args:
            slide: Slide object
            slide_plan: SlidePlan with content to distribute
            language: Language code
            
        Returns:
            Dictionary indicating what content was placed where
        """
        distribution_result = {
            "title_placed": False,
            "bullets_placed": False,
            "subtitle_used": False,
            "footer_used": False,
            "all_placeholders_utilized": False
        }
        
        try:
            # Discover all available placeholders
            placeholders = self._discover_all_content_placeholders(slide)
            
            # 1. Place title in title placeholder
            if slide_plan.title and placeholders["title"]:
                self._add_title(slide, slide_plan.title, language)
                distribution_result["title_placed"] = True
            
            # 2. Prepare content for distribution
            bullets = slide_plan.get_bullet_texts()
            total_content_items = len(bullets)
            
            if total_content_items == 0:
                logger.debug("No content to distribute")
                return distribution_result
            
            # 3. Create content distribution strategy
            content_strategy = self._create_content_distribution_strategy(
                placeholders, bullets, slide_plan
            )
            
            # 4. Execute distribution strategy
            self._execute_content_distribution(slide, content_strategy, language)
            
            # 5. Update distribution results
            distribution_result.update({
                "bullets_placed": len(content_strategy["primary_bullets"]) > 0,
                "subtitle_used": len(content_strategy.get("subtitle_content", [])) > 0,
                "footer_used": len(content_strategy.get("footer_content", [])) > 0,
                "all_placeholders_utilized": self._check_placeholder_utilization(placeholders)
            })
            
            logger.info(f"Content distribution completed: {sum(distribution_result.values())} of {len(distribution_result)} goals achieved")
            
        except Exception as e:
            logger.error(f"Content distribution failed: {e}")
            
        return distribution_result
    
    def _create_content_distribution_strategy(self, placeholders: Dict, bullets: List[str], slide_plan: SlidePlan) -> Dict:
        """
        Create an intelligent strategy for distributing content across available placeholders.
        
        Args:
            placeholders: Available placeholders by category
            bullets: List of bullet point texts
            slide_plan: Slide plan with additional context
            
        Returns:
            Distribution strategy dictionary
        """
        strategy = {
            "primary_bullets": [],
            "subtitle_content": [],
            "footer_content": [],
            "distributed_bullets": [],
            "placeholder_assignments": {}
        }
        
        try:
            primary_placeholders = placeholders["primary_content"]
            secondary_placeholders = placeholders["secondary_content"]
            
            # Analyze content for intelligent distribution
            if len(bullets) <= 3 and len(secondary_placeholders) > 0:
                # For short content, use subtitle for key takeaway
                strategy["primary_bullets"] = bullets
                    
            elif len(bullets) > 6 and len(primary_placeholders) >= 2:
                # For long content, distribute across multiple primary placeholders
                bullets_per_placeholder = len(bullets) // len(primary_placeholders)
                remainder = len(bullets) % len(primary_placeholders)
                
                distributed = []
                start_idx = 0
                
                for i, placeholder in enumerate(primary_placeholders):
                    bullets_for_this = bullets_per_placeholder + (1 if i < remainder else 0)
                    end_idx = start_idx + bullets_for_this
                    
                    placeholder_bullets = bullets[start_idx:end_idx]
                    distributed.append(placeholder_bullets)
                    strategy["placeholder_assignments"][i] = placeholder_bullets
                    
                    start_idx = end_idx
                
                strategy["distributed_bullets"] = distributed
                
            else:
                # Standard distribution
                strategy["primary_bullets"] = bullets
            
            # Always try to populate subtitle placeholders when they exist
            subtitle_placeholders = [p for p in secondary_placeholders 
                                   if p.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE]
            
            if subtitle_placeholders and len(strategy["subtitle_content"]) == 0:
                # Create appropriate subtitle content based on slide content
                if bullets:
                    if slide_plan.slide_type == "title":
                        # For title slides, use a summary or tagline
                        if len(bullets) == 1:
                            strategy["subtitle_content"] = [bullets[0]]
                        else:
                            # Create a brief summary from first bullet
                            strategy["subtitle_content"] = [bullets[0][:80] + "..." if len(bullets[0]) > 80 else bullets[0]]
                    else:
                        # For content slides, create a summary or key point
                        if len(bullets) <= 3:
                            strategy["subtitle_content"] = [f"Key takeaway: {bullets[0][:60]}{'...' if len(bullets[0]) > 60 else ''}"]
                        else:
                            # For longer content, create a brief overview
                            strategy["subtitle_content"] = [f"{len(bullets)} key points covering {slide_plan.title.lower()}"]
                elif slide_plan.slide_type == "title":
                    # For title slides without bullets, create a generic subtitle
                    strategy["subtitle_content"] = [""]  # Empty string to clear placeholder
                else:
                    # For non-title slides without bullets, create descriptive subtitle
                    strategy["subtitle_content"] = [f"Overview of {slide_plan.title.lower()}"]
            
            # Add footer content if available and appropriate
            if len(secondary_placeholders) > 1 and slide_plan.speaker_notes:
                # Use footer for condensed speaker notes
                strategy["footer_content"] = [slide_plan.speaker_notes[:100] + "..."]
            
            logger.debug(f"Created distribution strategy: {len(strategy['primary_bullets'])} primary bullets, "
                        f"{len(strategy['subtitle_content'])} subtitle items, "
                        f"{len(strategy['footer_content'])} footer items")
                        
        except Exception as e:
            logger.error(f"Strategy creation failed: {e}")
            # Fallback to simple strategy
            strategy["primary_bullets"] = bullets
            
        return strategy
    
    def _execute_content_distribution(self, slide, strategy: Dict, language: str) -> None:
        """
        Execute the content distribution strategy.
        
        Args:
            slide: Slide object
            strategy: Distribution strategy
            language: Language code
        """
        try:
            placeholders = self._discover_all_content_placeholders(slide)
            
            # Distribute primary bullets
            if strategy["distributed_bullets"]:
                # Multi-placeholder distribution
                for i, bullets in enumerate(strategy["distributed_bullets"]):
                    if i < len(placeholders["primary_content"]) and bullets:
                        self._add_bullets_to_placeholder(
                            placeholders["primary_content"][i], bullets, language
                        )
                        logger.debug(f"Added {len(bullets)} bullets to primary placeholder {i}")
                        
            elif strategy["primary_bullets"]:
                # Single placeholder distribution
                if placeholders["primary_content"]:
                    self._add_bullets_to_placeholder(
                        placeholders["primary_content"][0], strategy["primary_bullets"], language
                    )
                    logger.debug(f"Added {len(strategy['primary_bullets'])} bullets to primary placeholder")
            
            # Add subtitle content
            if strategy["subtitle_content"] and placeholders["secondary_content"]:
                for placeholder in placeholders["secondary_content"]:
                    if placeholder.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                        subtitle_text = strategy["subtitle_content"][0]
                        if subtitle_text:  # Only set text if not empty
                            placeholder.text = subtitle_text
                            logger.debug(f"Added subtitle content: {subtitle_text[:50]}...")
                        else:
                            # Hide empty subtitle placeholder
                            self._hide_empty_placeholder(placeholder, "subtitle")
                            logger.debug("Hidden empty subtitle placeholder")
                        break
            
            # Add footer content
            if strategy["footer_content"] and placeholders["secondary_content"]:
                for placeholder in placeholders["secondary_content"]:
                    if placeholder.placeholder_format.type == 10:  # FOOTER
                        placeholder.text = strategy["footer_content"][0]
                        logger.debug("Added content to footer placeholder")
                        break
            
            # Final cleanup: Ensure no subtitle placeholders are left with default text
            self._cleanup_empty_subtitle_placeholders(placeholders["secondary_content"])
                        
        except Exception as e:
            logger.error(f"Content distribution execution failed: {e}")
    
    def _check_placeholder_utilization(self, placeholders: Dict) -> bool:
        """
        Check if we've utilized most available text placeholders.
        
        Args:
            placeholders: Placeholder dictionary
            
        Returns:
            True if most placeholders are utilized
        """
        try:
            total_text_placeholders = len(placeholders["all_text_capable"])
            if total_text_placeholders == 0:
                return True
                
            utilized_count = 0
            for placeholder in placeholders["all_text_capable"]:
                if hasattr(placeholder, 'text') and placeholder.text.strip():
                    utilized_count += 1
                elif hasattr(placeholder, 'text_frame') and placeholder.text_frame:
                    if any(para.text.strip() for para in placeholder.text_frame.paragraphs):
                        utilized_count += 1
            
            utilization_rate = utilized_count / total_text_placeholders
            logger.debug(f"Placeholder utilization: {utilized_count}/{total_text_placeholders} ({utilization_rate:.1%})")
            
            return utilization_rate >= 0.7  # 70% utilization threshold
            
        except Exception as e:
            logger.error(f"Utilization check failed: {e}")
            return False

    def _add_bullets_to_placeholder(
        self, 
        placeholder, 
        bullets: List[str], 
        language: str
    ) -> None:
        """
        Add bullets to a specific placeholder (legacy format).
        
        Args:
            placeholder: Placeholder object to add bullets to
            bullets: List of bullet text strings
            language: Language code for text formatting
        """
        # Convert to BulletItem format for unified handling
        bullet_items = [BulletItem(text=bullet, level=0) for bullet in bullets]
        self._add_hierarchical_bullets_to_placeholder(placeholder, bullet_items, language)

    def _add_hierarchical_bullets_to_placeholder(
        self, 
        placeholder, 
        bullet_items: List[BulletItem], 
        language: str
    ) -> None:
        """
        Add hierarchical bullets to a specific placeholder (T-100).
        
        Args:
            placeholder: Placeholder object to add bullets to
            bullet_items: List of BulletItem objects with hierarchy information
            language: Language code for text formatting
        """
        if not hasattr(placeholder, 'text_frame'):
            logger.warning("Placeholder does not have text_frame attribute")
            return
        
        text_frame = placeholder.text_frame
        text_frame.clear()  # Clear existing content
        
        # Add bullets with hierarchy support
        for i, bullet_item in enumerate(bullet_items):
            bullet_text = bullet_item.text
            bullet_level = bullet_item.level
            
            # Truncate if too long
            if len(bullet_text) > self.max_bullet_length:
                bullet_text = bullet_text[:self.max_bullet_length - 3] + "..."
            
            if i == 0:
                # Use the first paragraph
                p = text_frame.paragraphs[0]
            else:
                # Add new paragraphs for subsequent bullets
                p = text_frame.add_paragraph()
            
            p.text = bullet_text
            p.level = min(bullet_level, 4)  # PowerPoint supports levels 0-4
            
            # Apply formatting (font size will be adjusted later if needed)
            if p.runs:
                run = p.runs[0]
                
                # Adjust font size based on hierarchy level (T-100)
                # If autofit is enabled, set a smaller, consistent initial size.
                # Otherwise, use the original level-based sizing.
                if self.overflow_config.enable_auto_fit_text:
                    run.font.size = Pt(14) # Smaller default for autofit
                    logger.debug(f"Set initial font size to 14pt for bullet level {bullet_level} (autofit enabled).")
                else:
                    base_size = 18
                    size_reduction = bullet_level * 2  # Reduce by 2pt per level
                    adjusted_size = max(12, base_size - size_reduction)  # Minimum 12pt
                    run.font.size = Pt(adjusted_size)
                    logger.debug(f"Applied bullet level {bullet_level} with font size {adjusted_size}pt (autofit disabled).")

                # Proactive font selection for bullets
                if self.template_parser.template_style:
                    specific_font_name = self.template_parser.template_style.language_specific_fonts.get(language.lower())
                    if specific_font_name:
                        run.font.name = specific_font_name
                        logger.debug(f"Applied language-specific font '{specific_font_name}' for bullets in language '{language.lower()}'")

                # Logger line for font size is now inside the conditional block above.

        # RTL alignment for bullets
        if language.lower() in RTL_LANGUAGES:
            for para in text_frame.paragraphs:
                if para.text and para.text.strip(): # Only align if paragraph has text
                    para.alignment = PP_ALIGN.RIGHT

        # Apply auto-fit text if enabled
        if self.overflow_config.enable_auto_fit_text:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            logger.debug(f"Applied TEXT_TO_FIT_SHAPE to bullet placeholder due to enable_auto_fit_text=True.")

    def find_placeholder(self, slide, placeholder_type: int) -> Optional[object]:
        """Find a placeholder of the specified type on the slide.
        
        Args:
            slide: PowerPoint slide object
            placeholder_type: PP_PLACEHOLDER constant (e.g., PP_PLACEHOLDER.PICTURE)
            
        Returns:
            Placeholder object if found, None otherwise
        """
        try:
            for placeholder in slide.placeholders:
                if (hasattr(placeholder, 'placeholder_format') and 
                    placeholder.placeholder_format.type == placeholder_type):
                    return placeholder
        except Exception as e:
            logger.debug(f"Error finding placeholder type {placeholder_type}: {e}")
        
        return None
    
    def _fallback_add_picture(self, slide, image_path: str, position: str = "bottom_right", 
                             width: float = 4.0, height: float = 3.0, alt_text: Optional[str] = None) -> Optional[object]:
        """Consolidated fallback logic for adding pictures when no placeholder is available.
        
        Args:
            slide: PowerPoint slide object
            image_path: Path to the image file
            position: Position on slide ("bottom_right", "top_right", "center")
            width: Image width in inches
            height: Image height in inches
            alt_text: Alternative text for accessibility
            
        Returns:
            Picture shape object if successful, None otherwise
        """
        try:
            # Get slide dimensions
            try:
                prs = slide.part.package.presentation_part.presentation
                slide_width = prs.slide_width
                slide_height = prs.slide_height
            except AttributeError:
                # Fallback to default dimensions if access fails
                slide_width = Inches(10)  # Standard slide width
                slide_height = Inches(7.5)  # Standard slide height
            
            img_width = Inches(width)
            img_height = Inches(height)
            margin = Inches(0.5)
            
            # Calculate position based on position parameter
            if position == "bottom_right":
                left = slide_width - img_width - margin
                top = slide_height - img_height - margin
            elif position == "top_right":
                left = slide_width - img_width - margin
                top = Inches(1)
            elif position == "center":
                left = (slide_width - img_width) / 2
                top = (slide_height - img_height) / 2
            else:
                # Default to bottom_right
                left = slide_width - img_width - margin
                top = slide_height - img_height - margin
            
            picture_shape = slide.shapes.add_picture(str(image_path), left, top, img_width, img_height)
            if alt_text:
                picture_shape.name = alt_text
            
            logger.debug(f"Added picture as floating image at {position}: {image_path}")
            return picture_shape
            
        except Exception as e:
            logger.error(f"Failed to add fallback picture '{image_path}': {e}")
            return None
    
    def _add_chart_image(self, slide, chart_path: str, alt_text: Optional[str] = None) -> None:
        """Add chart image to slide using placeholder-first approach."""
        try:
            chart_path = Path(chart_path)
            if not chart_path.exists():
                logger.error(f"Chart file not found: {chart_path}")
                return
            
            # Use consolidated placeholder finding
            ph = self.find_placeholder(slide, PP_PLACEHOLDER.PICTURE)
            if ph:
                picture_shape = ph.insert_picture(str(chart_path))
                if alt_text:
                    picture_shape.name = alt_text
                logger.debug(f"Inserted chart into picture placeholder: {chart_path}")
            else:
                logger.warning(f"No picture placeholder found on slide using layout '{slide.slide_layout.name}'. Falling back to adding chart image as floating.")
                # Use consolidated fallback logic
                self._fallback_add_picture(slide, str(chart_path), "bottom_right", 4.0, 3.0, alt_text)
                
        except Exception as e:
            logger.error(f"Failed to add chart image '{chart_path}': {e}")

    def _add_image(self, slide, image_path: str, alt_text: Optional[str] = None) -> None:
        """Add image to slide using placeholder-first approach."""
        try:
            image_path = Path(image_path)
            if not image_path.exists():
                logger.error(f"Image file not found: {image_path}")
                return
            
            # Use consolidated placeholder finding
            ph = self.find_placeholder(slide, PP_PLACEHOLDER.PICTURE)
            if ph:
                picture_shape = ph.insert_picture(str(image_path))
                if alt_text:
                    picture_shape.name = alt_text
                logger.debug(f"Inserted image into picture placeholder: {image_path}")
            else:
                logger.warning(f"No picture placeholder found on slide using layout '{slide.slide_layout.name}'. Falling back to adding image as floating.")
                # Use consolidated fallback logic
                self._fallback_add_picture(slide, str(image_path), "top_right", 3.0, 2.0, alt_text)
                
        except Exception as e:
            logger.error(f"Failed to add image '{image_path}': {e}")

    def _add_native_chart(self, slide, chart_data) -> None:
        """Add native PowerPoint chart to slide."""
        try:
            # Import here to avoid circular imports
            from .native_chart_builder import NativeChartBuilder
            
            # Convert chart data to NativeChartData if needed
            if isinstance(chart_data, dict):
                chart_builder = NativeChartBuilder(self.template_parser)
                native_chart = chart_builder.convert_legacy_chart_data(chart_data)
                
                if native_chart:
                    # Try to use chart placeholder first
                    chart_placeholder = self._find_chart_placeholder(slide)
                    if chart_placeholder:
                        chart_builder.create_chart_in_placeholder(slide, native_chart, chart_placeholder)
                    else:
                        chart_builder.create_native_chart(slide, native_chart)
                    
                    logger.debug(f"Added native chart: {native_chart.title}")
                else:
                    logger.warning("Failed to convert chart data to native format")
            
            elif isinstance(chart_data, NativeChartData):
                chart_builder = NativeChartBuilder(self.template_parser)
                chart_builder.create_native_chart(slide, chart_data)
                logger.debug(f"Added native chart: {chart_data.title}")
            
            else:
                logger.warning(f"Unsupported chart data type: {type(chart_data)}")
                
        except Exception as e:
            logger.error(f"Failed to add native chart: {e}")

    def _add_process_flow(self, slide, flow_path: str, alt_text: Optional[str] = None) -> None:
        """Add process flow SVG to slide."""
        try:
            flow_path = Path(flow_path)
            if not flow_path.exists():
                logger.error(f"Process flow file not found: {flow_path}")
                return
            
            # Convert SVG to image for insertion
            # For now, just add as image - could be enhanced to embed SVG directly
            self._add_image(slide, flow_path, alt_text=alt_text) # Pass alt_text
            logger.debug(f"Added process flow: {flow_path}")
            
        except Exception as e:
            logger.error(f"Failed to add process flow '{flow_path}': {e}")

    def _find_chart_placeholder(self, slide) -> Optional[object]:
        """Find a chart placeholder on the slide."""
        # Try chart placeholder first (type 12)
        chart_ph = self.find_placeholder(slide, 12)  # CHART placeholder (no PP_PLACEHOLDER.CHART constant available)
        if chart_ph:
            return chart_ph
        
        # Fallback to object placeholder that can hold charts
        object_ph = self.find_placeholder(slide, PP_PLACEHOLDER.OBJECT)
        if object_ph:
            return object_ph
        
        return None

    def _add_speaker_notes(self, slide, notes: str) -> None:
        """Add speaker notes to slide."""
        try:
            if hasattr(slide, 'notes_slide'):
                notes_slide = slide.notes_slide
                if hasattr(notes_slide, 'notes_text_frame'):
                    text_frame = notes_slide.notes_text_frame
                    text_frame.text = notes
                    logger.debug(f"Added speaker notes: {notes[:50]}...")
                    
        except Exception as e:
            logger.error(f"Failed to add speaker notes: {e}")

    def _add_fallback_slide(self, prs: Presentation, slide_plan: SlidePlan) -> None:
        """Add a basic fallback slide when normal slide creation fails."""
        try:
            # Use the first available layout
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            
            # Just add the title
            if hasattr(slide.shapes, 'title'):
                slide.shapes.title.text = slide_plan.title or f"Slide {slide_plan.index + 1}"
            
            logger.warning(f"Added fallback slide for slide {slide_plan.index}")
            
        except Exception as e:
            logger.error(f"Failed to add fallback slide: {e}")

    def _apply_metadata(self, prs: Presentation, outline: Outline) -> None:
        """Apply presentation metadata."""
        try:
            # Set presentation properties
            if hasattr(prs.core_properties, 'title'):
                prs.core_properties.title = outline.title
            
            if hasattr(prs.core_properties, 'subject') and outline.subtitle:
                prs.core_properties.subject = outline.subtitle
            
            if hasattr(prs.core_properties, 'author'):
                prs.core_properties.author = "Open Lilli AI"
            
            logger.debug("Applied presentation metadata")
            
        except Exception as e:
            logger.error(f"Failed to apply metadata: {e}")

    def create_slide_from_layout(
        self,
        prs: Presentation,
        layout_name: str,
        title: str,
            content: Optional[List[str]] = None,
            language: str = "en" # Add language, default to "en" for existing calls
    ) -> None:
        """
        Create a slide using a specific layout by name.
        
        Args:
            prs: Presentation object
            layout_name: Name of the layout to use
            title: Slide title
            content: Optional list of content items
            language: Language code for text alignment
        """
        try:
            layout_index = self.template_parser.get_layout_index(layout_name)
            layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(layout)
            
            # Add title
            # Assuming this internal method might not need language if it's simple title setting
            # or if it calls the main _add_title, it should pass language.
            # For now, let's assume it's a simple set, or we need to trace its callers.
            # Let's assume _add_title needs language.
            if hasattr(slide.shapes, 'title'):
                 self._add_title(slide, title, language) # Pass language
            
            # Add content if provided
            if content:
                self._add_bullet_content(slide, content, language) # Pass language
            
            logger.debug(f"Created slide with layout '{layout_name}': {title}")
            
        except Exception as e:
            logger.error(f"Failed to create slide with layout '{layout_name}': {e}")

    def analyze_slide_placeholders(self, slide) -> Dict[str, any]:
        """
        Analyze placeholders in a slide for debugging.
        
        Args:
            slide: Slide object to analyze
            
        Returns:
            Dictionary with placeholder information
        """
        placeholders_info = {
            "total_placeholders": len(slide.placeholders),
            "placeholders": []
        }
        
        for i, placeholder in enumerate(slide.placeholders):
            ph_info = {
                "index": i,
                "type": placeholder.placeholder_format.type,
                "has_text_frame": hasattr(placeholder, 'text_frame'),
                "shape_type": placeholder.shape_type if hasattr(placeholder, 'shape_type') else None
            }
            
            # Try to get placeholder name
            try:
                ph_info["name"] = placeholder.name
            except:
                ph_info["name"] = None
            
            placeholders_info["placeholders"].append(ph_info)
        
        return placeholders_info

    def get_assembly_statistics(
        self,
        slides: List[SlidePlan],
        visuals: Dict[int, Dict[str, str]]
    ) -> Dict[str, any]:
        """
        Get statistics about the assembly process.
        
        Args:
            slides: List of slides to assemble
            visuals: Dictionary of visuals
            
        Returns:
            Assembly statistics
        """
        stats = {
            "total_slides": len(slides),
            "slides_with_visuals": len([s for s in slides if s.index in visuals]),
            "total_bullets": sum(len(s.bullets) for s in slides),
            "slides_with_notes": len([s for s in slides if s.speaker_notes]),
            "slide_types": {},
            "layout_usage": {},
            "visual_types": {"charts": 0, "images": 0}
        }
        
        # Count slide types and layout usage
        for slide in slides:
            stats["slide_types"][slide.slide_type] = stats["slide_types"].get(slide.slide_type, 0) + 1
            
            layout_id = slide.layout_id or 0
            stats["layout_usage"][layout_id] = stats["layout_usage"].get(layout_id, 0) + 1
        
        # Count visual types
        for slide_visuals in visuals.values():
            if "chart" in slide_visuals:
                stats["visual_types"]["charts"] += 1
            if "image" in slide_visuals:
                stats["visual_types"]["images"] += 1
        
        return stats

    def _estimate_lines_for_bullets(
        self,
        bullets: List[BulletItem],
        placeholder_width_inches: float, # Currently unused, but good for future refinement
        base_font_size_pt: float,        # Currently unused, but good for future refinement
        avg_chars_per_line_full_width: int = 70
    ) -> int:
        """
        Estimates the total number of lines the provided bullets will occupy.

        Args:
            bullets: A list of BulletItem objects.
            placeholder_width_inches: The width of the placeholder in inches.
            base_font_size_pt: The base font size in points for level 0 bullets.
            avg_chars_per_line_full_width: Average characters that fit in a line at
                                           full placeholder width with the base font size.
                                           This is a heuristic.

        Returns:
            Total estimated number of lines.
        """
        total_estimated_lines = 0
        min_effective_chars_per_line = 20 # Minimum characters per line to avoid division by zero or overly small numbers

        if not bullets:
            return 0

        for i, bullet_item in enumerate(bullets):
            text = bullet_item.text
            level = bullet_item.level

            # Adjust avg_chars_per_line based on level.
            # Each indentation level might reduce the effective characters per line.
            # This is a simple heuristic.
            level_reduction_factor = 5 # Characters to reduce per indentation level
            effective_chars_per_line = avg_chars_per_line_full_width - (level * level_reduction_factor)
            effective_chars_per_line = max(min_effective_chars_per_line, effective_chars_per_line)

            lines_for_bullet = 0
            if not text.strip(): # Check if text is empty or only whitespace
                lines_for_bullet = 1 # Empty bullet still takes one line
            else:
                # Estimate lines based on wrapping
                # (len(text) + effective_chars_per_line - 1) // effective_chars_per_line is ceiling division
                estimated_wrapping_lines = (len(text) + effective_chars_per_line - 1) // effective_chars_per_line

                # Count explicit newline characters
                explicit_newlines = text.count('\n')

                # The number of lines is the number of explicit newlines + 1,
                # OR the estimated wrapping lines if that's more (e.g. long line with no newlines).
                # However, if wrapping lines are, say, 3, and explicit newlines are 0, it means 3 lines.
                # If wrapping lines are 1 (short text), and explicit newlines are 2 (text like "a\nb\nc"), then it's 3 lines.
                # So, lines_for_bullet should be max of (explicit_newlines + 1) and estimated_wrapping_lines.
                lines_for_bullet = max(estimated_wrapping_lines, explicit_newlines + 1)


            total_estimated_lines += lines_for_bullet

            # Basic logging for diagnostics, especially for long bullets or many lines
            if lines_for_bullet > 5 or len(text) > 150:
                logger.debug(
                    f"Bullet {i+1} (level {level}, {len(text)} chars, {explicit_newlines} newlines) "
                    f"estimated to take {lines_for_bullet} lines (effective_chars_per_line: {effective_chars_per_line})."
                )

        logger.debug(f"Total estimated lines for {len(bullets)} bullets: {total_estimated_lines}")
        return total_estimated_lines

    def validate_slides_before_assembly(self, slides: List[SlidePlan]) -> List[str]:
        """
        Validate slides before assembly and return list of issues.
        
        Args:
            slides: List of slides to validate
            
        Returns:
            List of validation issues
        """
        issues = []
        
        if not slides:
            issues.append("No slides provided")
            return issues
        
        # Check for missing titles
        for slide in slides:
            if not slide.title or not slide.title.strip():
                issues.append(f"Slide {slide.index} has no title")
        
        # Check for layout IDs
        max_layout_index = len(self.template_parser.prs.slide_layouts) - 1
        for slide in slides:
            if slide.layout_id is not None and slide.layout_id > max_layout_index:
                issues.append(f"Slide {slide.index} has invalid layout ID: {slide.layout_id}")
        
        # Check for excessive content
        for slide in slides:
            if len(slide.bullets) > 10:
                issues.append(f"Slide {slide.index} has too many bullets: {len(slide.bullets)}")
            
            if slide.title and len(slide.title) > self.max_title_length:
                issues.append(f"Slide {slide.index} title is too long: {len(slide.title)} chars")
        
        return issues

    def validate_presentation_style(self, prs: Presentation, language: Optional[str] = "en") -> None: # Add language
        """
        Validate the entire presentation against template style rules.
        
        Args:
            prs: Presentation object to validate
            language: Language code for the presentation (for font validation)
            
        Raises:
            StyleError: If style violations are found and validation is enabled
        """
        # Ensure language is available, default to "en" if not passed (e.g. from patch_existing_presentation if not updated)
        # However, assemble method should always pass it from outline.language.
        # For patch_existing_presentation, it's now a required parameter.
        # So, language should always be valid.

        if not self.validation_config.enabled or self.validation_config.mode == "disabled":
            logger.debug("Style validation disabled, skipping all checks.")
            return [] # Return empty list of feedback if disabled
        
        logger.info(f"Starting presentation style validation for language: {language}. Mode: {self.validation_config.mode}.")
        all_feedback: List[ReviewFeedback] = []
        
        # --- Step 1: Rule-based checks ---
        # These checks are performed on each slide using its properties.
        logger.debug("Performing rule-based style checks...")
        for slide_index_loop, pptx_slide_obj in enumerate(prs.slides):
            try:
                # _validate_slide_style now returns List[ReviewFeedback]
                slide_specific_feedback = self._validate_slide_style(pptx_slide_obj, slide_index_loop, language)
                all_feedback.extend(slide_specific_feedback)
            except Exception as e:
                logger.warning(f"Error during rule-based validation for slide {slide_index_loop}: {e}")
                all_feedback.append(ReviewFeedback(
                    slide_index=slide_index_loop,
                    severity="high",
                    category="rule_based_validation_error",
                    message=f"An unexpected error occurred during rule-based validation of slide {slide_index_loop + 1}: {e}",
                    suggestion="Review slide content for potential issues or report this as a bug."
                ))
        logger.debug(f"Completed rule-based checks. {len(all_feedback)} feedback items generated so far.")

        # --- Step 2: VisualProofreader checks (if enabled) ---
        if self.visual_proofreader:
            logger.debug("VisualProofreader is enabled. Preparing for LLM-based checks.")
            slide_plans_for_vp: List[SlidePlan] = []
            # Convert pptx.slide.Slide objects to SlidePlan objects for VisualProofreader.
            # This is a simplified conversion; a more robust implementation would extract more details.
            for i, ppt_slide in enumerate(prs.slides):
                title_text = ppt_slide.shapes.title.text if ppt_slide.shapes.title and ppt_slide.shapes.title.has_text_frame else f"Slide {i+1}"
                bullets_text = []
                # Basic extraction of text from body/content placeholders.
                for shape in ppt_slide.shapes:
                    if shape.is_placeholder and hasattr(shape, 'placeholder_format') and \
                       shape.placeholder_format.type in [PP_PLACEHOLDER.BODY.value, PP_PLACEHOLDER.OBJECT.value, PP_PLACEHOLDER.CONTENT.value, PP_PLACEHOLDER.TEXT_BOX.value]: # Common content placeholders
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                if para.text.strip():
                                    bullets_text.append(para.text.strip())
                    elif shape.has_text_frame and not shape.is_placeholder: # Also consider non-placeholder text boxes
                         for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                bullets_text.append(para.text.strip())

                slide_plan_for_vp = SlidePlan(
                    index=i,
                    slide_type="content", # Generic type; actual type detection from pptx is complex
                    title=title_text,
                    bullets=bullets_text, # Uses legacy flat list of bullets
                    # Note: Fields like bullet_hierarchy, image_query, chart_data are not populated here.
                    # VisualProofreader will work primarily with title and bullet text.
                )
                slide_plans_for_vp.append(slide_plan_for_vp)

            if slide_plans_for_vp:
                try:
                    logger.info(f"Running VisualProofreader on {len(slide_plans_for_vp)} slides with focus on: {self.validation_config.visual_proofreader_focus_areas or 'defaults'}.")
                    vp_result = self.visual_proofreader.proofread_slides(
                        slides=slide_plans_for_vp,
                        focus_areas=self.validation_config.visual_proofreader_focus_areas,
                        enable_corrections=self.validation_config.visual_proofreader_enable_corrections
                    )
                    if vp_result and vp_result.issues_found:
                        vp_feedback = self.visual_proofreader.convert_to_review_feedback(vp_result.issues_found)
                        all_feedback.extend(vp_feedback)
                        logger.info(f"VisualProofreader added {len(vp_feedback)} feedback items.")
                except Exception as e:
                    logger.error(f"VisualProofreader processing failed: {e}")
                    all_feedback.append(ReviewFeedback(
                        slide_index=-1, # Indicates a general VisualProofreader error
                        severity="high",
                        category="visual_proofreader_error",
                        message=f"VisualProofreader processing encountered an error: {e}",
                        suggestion="Check VisualProofreader configuration, API key, and network connectivity."
                    ))
            else:
                logger.info("No slides suitable for VisualProofreader were prepared (or no slides in presentation).")
        else:
            logger.debug("VisualProofreader is not enabled or not initialized, skipping LLM-based checks.")

        # --- Step 3: Process aggregated feedback ---
        if all_feedback:
            logger.info(f"Total feedback items generated: {len(all_feedback)}.")
            if self.validation_config.mode == "strict":
                # Convert ReviewFeedback to a list of dicts for StyleError, as it expects dicts.
                violations_for_error = [fb.model_dump() for fb in all_feedback]
                error_msg = f"Style validation failed with {len(all_feedback)} feedback items in 'strict' mode."
                style_error = StyleError(error_msg, violations=violations_for_error)
                logger.error(f"Strict mode: Style validation failed. Raising StyleError.")
                raise style_error
            elif self.validation_config.mode == "lenient":
                logger.warning(f"Style validation found {len(all_feedback)} feedback items (lenient mode):")
                for i, feedback_item in enumerate(all_feedback):
                    if i < 10: # Log details for the first 10 items
                        logger.warning(f"  - Slide {feedback_item.slide_index + 1 if feedback_item.slide_index >=0 else 'N/A'}: [{feedback_item.category}/{feedback_item.severity}] {feedback_item.message} (Suggestion: {feedback_item.suggestion or 'N/A'})")
                    elif i == 10:
                        logger.warning(f"  ... and {len(all_feedback) - 10} more feedback items (not shown in detail).")
                        break
        else:
            logger.info("Style validation passed successfully (no feedback items generated).")

        return all_feedback

    def _validate_slide_style(self, slide, slide_index: int, language: str) -> List[ReviewFeedback]:
        """
        Validate style for a single slide and return a list of ReviewFeedback objects.
        
        Args:
            slide: Slide object to validate
            slide_index: Index of the slide
            language: Language code for the presentation
            
        Returns:
            List of ReviewFeedback objects for the slide
        """
        slide_feedback: List[ReviewFeedback] = []
        
        for shape in slide.shapes:
            try:
                # _validate_shape_style will need to be refactored to return List[ReviewFeedback]
                shape_feedback = self._validate_shape_style(shape, slide, slide_index, language)
                slide_feedback.extend(shape_feedback)
            except Exception as e:
                logger.debug(f"Failed to validate shape in slide {slide_index}: {e}")
                slide_feedback.append(ReviewFeedback(
                    slide_index=slide_index,
                    severity="medium",
                    category="shape_validation_error",
                    message=f"Error validating shape {shape.shape_id if hasattr(shape, 'shape_id') else 'unknown'} on slide {slide_index + 1}: {e}",
                    suggestion="Check shape properties or content."
                ))
        
        # _check_slide_placeholder_population will also be refactored
        population_feedback = self._check_slide_placeholder_population(slide, slide_index)
        slide_feedback.extend(population_feedback)

        return slide_feedback

    def _autofix_shape_alignment(self, shape, slide, slide_index: int) -> bool:
        """
        Attempts to automatically fix shape alignment issues against slide margins.
        Args:
            shape: The shape with alignment violations.
            slide: The slide object.
            slide_index: Index of the slide.
        Returns:
            True if any fix was attempted, False otherwise.
            Note: This simple version doesn't guarantee all violations are fixed,
                  it just attempts common corrections. The caller should re-check.
        """
        # autofix_alignment is checked by the caller (_check_shape_alignment_against_margins)
        # before this method is called.

        qg_config = self.validation_config.quality_gates_config
        # qg_config and qg_config.enable_alignment_check are also checked by caller.

        slide_width_emu = slide.part.package.presentation_part.presentation.slide_width
        slide_height_emu = slide.part.package.presentation_part.presentation.slide_height
        margin_top_emu = Inches(qg_config.slide_margin_top_inches)
        margin_bottom_emu = Inches(qg_config.slide_margin_bottom_inches)
        margin_left_emu = Inches(qg_config.slide_margin_left_inches)
        margin_right_emu = Inches(qg_config.slide_margin_right_inches)

        safe_left = margin_left_emu
        safe_top = margin_top_emu
        safe_right = slide_width_emu - margin_right_emu
        safe_bottom = slide_height_emu - margin_bottom_emu

        shape_name_desc = shape.name if shape.name else f"Shape ID {shape.shape_id}"
        if hasattr(shape, 'placeholder_format') and hasattr(shape.placeholder_format, 'type'):
            try:
                ph_type = shape.placeholder_format.type
                ph_type_name = self._get_placeholder_type_name(ph_type.value if hasattr(ph_type, 'value') else int(ph_type))
                shape_name_desc = f"{ph_type_name} (Idx: {shape.placeholder_format.idx}, Name: {shape.name or 'Unnamed'})"
            except Exception: pass

        action_taken = False

        # Fix left margin violations by moving
        if shape.left < safe_left:
            original_left = shape.left
            shape.left = safe_left
            logger.info(f"Autofix: Moved shape '{shape_name_desc}' on slide {slide_index} from L:{original_left:.0f} to L:{shape.left:.0f} (due to left margin).")
            action_taken = True

        # Fix top margin violations by moving
        if shape.top < safe_top:
            original_top = shape.top
            shape.top = safe_top
            logger.info(f"Autofix: Moved shape '{shape_name_desc}' on slide {slide_index} from T:{original_top:.0f} to T:{shape.top:.0f} (due to top margin).")
            action_taken = True

        # Fix right margin violations (try moving first, then resizing width)
        if (shape.left + shape.width) > safe_right:
            overflow_amount = (shape.left + shape.width) - safe_right

            new_left_candidate = shape.left - overflow_amount
            # Ensure shape has a positive width before attempting to move/resize
            if shape.width > 0 : # Check current width is positive
                if new_left_candidate >= safe_left :
                    shape.left = new_left_candidate
                    logger.info(f"Autofix: Moved shape '{shape_name_desc}' on slide {slide_index} left to fit right margin. New L:{shape.left:.0f}, R_Edge:{(shape.left + shape.width):.0f}.")
                    action_taken = True
                else:
                    new_width = shape.width - overflow_amount
                    if new_width > Inches(0.1):
                        shape.width = new_width
                        logger.info(f"Autofix: Resized width of shape '{shape_name_desc}' on slide {slide_index} to fit right margin. New W:{shape.width:.0f}, R_Edge:{(shape.left + shape.width):.0f}.")
                        action_taken = True
                    else:
                        logger.warning(f"Autofix: Failed to fit shape '{shape_name_desc}' on slide {slide_index} within right margin without making width too small.")
            else:
                 logger.warning(f"Autofix: Shape '{shape_name_desc}' on slide {slide_index} has zero or negative width, cannot fix right margin.")


        # Fix bottom margin violations (try moving first, then resizing height)
        if (shape.top + shape.height) > safe_bottom:
            overflow_amount = (shape.top + shape.height) - safe_bottom

            new_top_candidate = shape.top - overflow_amount
            if shape.height > 0 : # Check current height
                if new_top_candidate >= safe_top:
                    shape.top = new_top_candidate
                    logger.info(f"Autofix: Moved shape '{shape_name_desc}' on slide {slide_index} up to fit bottom margin. New T:{shape.top:.0f}, B_Edge:{(shape.top + shape.height):.0f}.")
                    action_taken = True
                else:
                    new_height = shape.height - overflow_amount
                    if new_height > Inches(0.1):
                        shape.height = new_height
                        logger.info(f"Autofix: Resized height of shape '{shape_name_desc}' on slide {slide_index} to fit bottom margin. New H:{shape.height:.0f}, B_Edge:{(shape.top + shape.height):.0f}.")
                        action_taken = True
                    else:
                        logger.warning(f"Autofix: Failed to fit shape '{shape_name_desc}' on slide {slide_index} within bottom margin without making height too small.")
            else:
                logger.warning(f"Autofix: Shape '{shape_name_desc}' on slide {slide_index} has zero or negative height, cannot fix bottom margin.")

        return action_taken

    def _check_shape_alignment_against_margins(self, shape, slide, slide_index: int, language: str) -> List[ReviewFeedback]:
        """
        Checks if a shape is within the defined slide margins.
        Args:
            shape: The shape to check.
            slide: The slide object the shape belongs to.
            slide_index: Index of the slide.
            language: Language (passed for consistency, not directly used in this check).
        Returns:
            A list of ReviewFeedback objects.
        """
        feedback_items: List[ReviewFeedback] = []
        if not self.validation_config.check_alignment:
            return feedback_items

        qg_config = self.validation_config.quality_gates_config
        if not qg_config or not qg_config.enable_alignment_check:
            logger.debug("Alignment check skipped as it's not enabled in QualityGatesConfig or QualityGatesConfig is missing.")
            return feedback_items

        if not slide.part or not slide.part.package or not slide.part.package.presentation_part or not slide.part.package.presentation_part.presentation or not hasattr(slide.part.package.presentation_part.presentation, 'slide_width') or not hasattr(slide.part.package.presentation_part.presentation, 'slide_height'):
            logger.warning("Cannot access presentation-level slide dimensions for alignment check. Skipping.")
            return feedback_items

        slide_width_emu = slide.part.package.presentation_part.presentation.slide_width
        slide_height_emu = slide.part.package.presentation_part.presentation.slide_height
        margin_top_emu = Inches(qg_config.slide_margin_top_inches)
        margin_bottom_emu = Inches(qg_config.slide_margin_bottom_inches)
        margin_left_emu = Inches(qg_config.slide_margin_left_inches)
        margin_right_emu = Inches(qg_config.slide_margin_right_inches)

        safe_left = margin_left_emu
        safe_top = margin_top_emu
        safe_right = slide_width_emu - margin_right_emu
        safe_bottom = slide_height_emu - margin_bottom_emu

        shape_name_desc_orig = shape.name if shape.name else f"Shape ID {shape.shape_id}" # Keep original for logs if autofix changes name (unlikely)
        current_shape_name_desc = shape_name_desc_orig
        if hasattr(shape, 'placeholder_format') and hasattr(shape.placeholder_format, 'type'):
            try:
                ph_type = shape.placeholder_format.type
                ph_type_name = self._get_placeholder_type_name(ph_type.value if hasattr(ph_type, 'value') else int(ph_type))
                current_shape_name_desc = f"{ph_type_name} (Idx: {shape.placeholder_format.idx}, Name: {shape.name or 'Unnamed'})"
            except Exception:
                pass

        if shape.width == 0 or shape.height == 0:
            logger.debug(f"Shape '{current_shape_name_desc}' has zero width or height, skipping alignment check.")
            return feedback_items

        # Store initial violations before attempting autofix
        detected_issues_before_fix: List[Tuple[str, str]] = [] # (subtype, description)

        if shape.left < safe_left:
            detected_issues_before_fix.append(("left_margin", f"Shape '{current_shape_name_desc}' (L:{shape.left:.0f}) extends beyond the left slide margin (safe L:{safe_left:.0f})."))
        if shape.top < safe_top:
            detected_issues_before_fix.append(("top_margin", f"Shape '{current_shape_name_desc}' (T:{shape.top:.0f}) extends beyond the top slide margin (safe T:{safe_top:.0f})."))
        if (shape.left + shape.width) > safe_right:
            detected_issues_before_fix.append(("right_margin", f"Shape '{current_shape_name_desc}' (R edge:{(shape.left + shape.width):.0f}) extends beyond the right slide margin (safe R:{safe_right:.0f})."))
        if (shape.top + shape.height) > safe_bottom:
            detected_issues_before_fix.append(("bottom_margin", f"Shape '{current_shape_name_desc}' (B edge:{(shape.top + shape.height):.0f}) extends beyond the bottom slide margin (safe B:{safe_bottom:.0f})."))

        if not detected_issues_before_fix:
            return feedback_items # No violations found initially

        # Attempt autofix if enabled and violations were found
        if self.validation_config.autofix_alignment:
            logger.debug(f"Alignment violations found for '{current_shape_name_desc}' on slide {slide_index}. Attempting autofix.")
            action_taken = self._autofix_shape_alignment(shape, slide, slide_index)

            if action_taken:
                logger.info(f"Autofix attempted for '{current_shape_name_desc}' on slide {slide_index}. Re-checking alignment.")
                # Re-check violations on the potentially modified shape
                current_shape_name_desc_after_fix = shape.name if shape.name else f"Shape ID {shape.shape_id}" # Update name in case it changed (unlikely here)
                if hasattr(shape, 'placeholder_format') and hasattr(shape.placeholder_format, 'type'):
                    try:
                        ph_type = shape.placeholder_format.type
                        ph_type_name_after = self._get_placeholder_type_name(ph_type.value if hasattr(ph_type, 'value') else int(ph_type))
                        current_shape_name_desc_after_fix = f"{ph_type_name_after} (Idx: {shape.placeholder_format.idx}, Name: {shape.name or 'Unnamed'})"
                    except Exception: pass

                if shape.left < safe_left:
                    feedback_items.append(ReviewFeedback(slide_index=slide_index, severity="low", category="alignment", message=f"Shape '{current_shape_name_desc_after_fix}' still extends beyond the left slide margin after autofix (L:{shape.left:.0f}, Safe L:{safe_left:.0f}).", suggestion="Manual adjustment needed."))
                if shape.top < safe_top:
                    feedback_items.append(ReviewFeedback(slide_index=slide_index, severity="low", category="alignment", message=f"Shape '{current_shape_name_desc_after_fix}' still extends beyond the top slide margin after autofix (T:{shape.top:.0f}, Safe T:{safe_top:.0f}).", suggestion="Manual adjustment needed."))
                if (shape.left + shape.width) > safe_right:
                    feedback_items.append(ReviewFeedback(slide_index=slide_index, severity="low", category="alignment", message=f"Shape '{current_shape_name_desc_after_fix}' still extends beyond the right slide margin after autofix (R edge:{(shape.left + shape.width):.0f}, Safe R:{safe_right:.0f}).", suggestion="Manual adjustment needed."))
                if (shape.top + shape.height) > safe_bottom:
                    feedback_items.append(ReviewFeedback(slide_index=slide_index, severity="low", category="alignment", message=f"Shape '{current_shape_name_desc_after_fix}' still extends beyond the bottom slide margin after autofix (B edge:{(shape.top + shape.height):.0f}, Safe B:{safe_bottom:.0f}).", suggestion="Manual adjustment needed."))

                if not feedback_items: # No violations after autofix
                    logger.info(f"All alignment issues appear fixed for '{current_shape_name_desc_after_fix}' on slide {slide_index}.")
                else: # Some violations persist
                    logger.warning(f"Autofix did not resolve all alignment issues for '{current_shape_name_desc_after_fix}' on slide {slide_index}.")
            else: # Autofix enabled but no action_taken (e.g., resize would make shape too small)
                logger.debug(f"Autofix enabled but no specific action taken for '{current_shape_name_desc}' on slide {slide_index}. Reporting original violations.")
                for subtype, desc in detected_issues_before_fix:
                    feedback_items.append(ReviewFeedback(slide_index=slide_index, severity="medium", category="alignment", message=desc, suggestion="Autofix could not resolve this. Manual adjustment needed."))
        else: # Autofix is disabled, report all initially detected issues
            for subtype, desc in detected_issues_before_fix:
                feedback_items.append(ReviewFeedback(slide_index=slide_index, severity="medium", category="alignment", message=desc, suggestion="Enable autofix_alignment or adjust manually."))

        return feedback_items

    def _estimate_text_lines(self, text_frame, shape_width_emu) -> int:
        """
        Estimates the number of lines text will occupy in a text_frame.
        This is a simplified heuristic.
        Args:
            text_frame: The TextFrame object.
            shape_width_emu: The width of the shape in EMUs. (Currently not fully utilized)
        Returns:
            Estimated number of lines.
        """
        if not text_frame.text.strip():
            return 0

        # This would ideally come from ContentFitConfig or similar
        avg_chars_per_line_estimate = 50

        text = text_frame.text

        # Estimate lines based on text length divided by average chars per line for each segment split by newline
        estimated_lines_from_length = 0
        if avg_chars_per_line_estimate > 0:
            estimated_lines_from_length = sum(
                # Calculate lines for each segment: (length + N-1) // N for ceiling division
                (len(line_segment) + avg_chars_per_line_estimate - 1) // avg_chars_per_line_estimate
                for line_segment in text.split('\n')
            )
        else: # Avoid division by zero if estimate is bad
            estimated_lines_from_length = text.count('\n') + 1 # Fallback: count explicit newlines + 1

        # If text is "Line1\nLine2", split produces ["Line1", "Line2"].
        # Loop runs twice. len("Line1")//50 (say 1), len("Line2")//50 (say 1). Sum = 2.
        # If text is "Very long line without newlines...", loop runs once.
        # This calculation correctly sums lines from wrapped segments and adds them up.

        return estimated_lines_from_length

    def _autofix_shape_text_overflow(self, shape, text_frame, slide_index: int, language: str) -> bool:
        """
        Attempts to automatically fix text overflow in a shape by reducing font size.
        Args:
            shape: The shape with overflow.
            text_frame: The TextFrame of the shape.
            slide_index: Index of the slide.
            language: Language (currently unused here but good for consistency).
        Returns:
            True if overflow was fixed, False otherwise.
        """
        # This method is called when self.validation_config.autofix_text_overflow is already true.

        logger.debug(f"Attempting to autofix text overflow in shape '{shape.name or shape.shape_id}' on slide {slide_index}")

        min_font_size_pt = 10  # Fallback, ideally from ContentFitConfig.min_font_size
        font_shrink_step_pt = 1 # How much to shrink font by each iteration
        max_fix_attempts = 5    # Max iterations for font shrinking

        original_auto_size = text_frame.auto_size
        if text_frame.auto_size != MSO_AUTO_SIZE.NONE: # Only change if not already NONE
            text_frame.auto_size = MSO_AUTO_SIZE.NONE

        fixed = False
        for attempt in range(max_fix_attempts):
            # Re-check overflow by directly evaluating conditions, not full _check_text_frame_overflow
            # This avoids recursive calls and simplifies logic here.
            shape_height_emu = shape.height
            margin_top_emu = text_frame.margin_top if text_frame.margin_top is not None else 0
            margin_bottom_emu = text_frame.margin_bottom if text_frame.margin_bottom is not None else 0
            available_text_height_emu = shape_height_emu - margin_top_emu - margin_bottom_emu

            estimated_lines = self._estimate_text_lines(text_frame, shape.width)
            if estimated_lines == 0 : # No text or bad estimate
                 fixed = True # Assume fixed if no lines/text
                 break

            avg_font_size_pt = 18.0
            try:
                if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                    first_run_font_size = text_frame.paragraphs[0].runs[0].font.size
                    if first_run_font_size:
                        avg_font_size_pt = float(first_run_font_size.pt)
            except (AttributeError, IndexError): pass

            estimated_line_height_emu = avg_font_size_pt * 12700 * 1.2
            if estimated_line_height_emu == 0: estimated_line_height_emu = 12700 * 1.2
            required_height_emu = estimated_lines * estimated_line_height_emu

            if required_height_emu <= available_text_height_emu:
                fixed = True
                logger.info(f"Overflow fixed for shape '{shape.name or shape.shape_id}' on slide {slide_index} by font size reduction (attempt {attempt+1}).")
                break

            # If still overflowing, try to reduce font size further
            fonts_shrunk_in_step = False
            for para in text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        current_size_pt = run.font.size.pt
                        if current_size_pt > min_font_size_pt:
                            new_size_pt = max(min_font_size_pt, current_size_pt - font_shrink_step_pt)
                            if new_size_pt < current_size_pt:
                                run.font.size = Pt(new_size_pt)
                                fonts_shrunk_in_step = True

            if not fonts_shrunk_in_step:
                logger.debug(f"Could not shrink fonts further for shape '{shape.name or shape.shape_id}' on slide {slide_index}.")
                break

        if original_auto_size is not None and text_frame.auto_size != original_auto_size :
            text_frame.auto_size = original_auto_size

        if fixed:
            return True
        else:
            is_title_placeholder = False
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type:
                 try:
                    is_title_placeholder = shape.placeholder_format.type == PP_PLACEHOLDER.TITLE
                 except AttributeError: # .type might not be an enum directly if it's an unknown int
                    pass


            if is_title_placeholder and len(text_frame.text) > self.max_title_length :
                 original_text = text_frame.text
                 # Basic truncation
                 # Find the last space within max_title_length - 3 for cleaner cut
                 safe_truncate_pos = (original_text[:self.max_title_length - 3]).rfind(' ')
                 if safe_truncate_pos == -1 or safe_truncate_pos < self.max_title_length // 2 : # if no space or too short
                     truncated_text = original_text[:self.max_title_length - 3] + "..."
                 else:
                     truncated_text = original_text[:safe_truncate_pos] + "..."

                 # Preserve first paragraph's formatting for the new truncated text
                 first_para_font = None
                 if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                     first_run = text_frame.paragraphs[0].runs[0]
                     first_para_font = {
                         "name": first_run.font.name,
                         "size": first_run.font.size,
                         "bold": first_run.font.bold,
                         "italic": first_run.font.italic,
                         "color": first_run.font.color.rgb if first_run.font.color and first_run.font.color.rgb else None
                     }

                 text_frame.clear()
                 p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                 run = p.add_run()
                 run.text = truncated_text

                 if first_para_font:
                     if first_para_font["name"]: run.font.name = first_para_font["name"]
                     if first_para_font["size"]: run.font.size = first_para_font["size"]
                     if first_para_font["bold"] is not None: run.font.bold = first_para_font["bold"]
                     if first_para_font["italic"] is not None: run.font.italic = first_para_font["italic"]
                     if first_para_font["color"]: run.font.color.rgb = first_para_font["color"]

                 logger.info(f"Overflow in title shape '{shape.name or shape.shape_id}' fixed by truncation.")
                 fixed = True

            if not fixed:
                 logger.warning(f"Failed to autofix text overflow for shape '{shape.name or shape.shape_id}' on slide {slide_index}.")
            return fixed

    def _check_text_frame_overflow(self, shape, slide_index: int, language: str) -> List[ReviewFeedback]:
        """
        Checks a single shape for text overflow.
        Args:
            shape: The shape to check.
            slide_index: Index of the slide.
            language: Language of the presentation.
        Returns:
            A list of ReviewFeedback objects if overflow is detected.
        """
        feedback_items: List[ReviewFeedback] = []
        if not self.validation_config.check_text_overflow or not hasattr(shape, 'text_frame') or not shape.text_frame:
            return feedback_items

        text_frame = shape.text_frame
        # text_frame.auto_size can be:
        # MSO_AUTO_SIZE.NONE (0): No auto-sizing. Text can overflow. (This is what we want to detect)
        # MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT (1): Shape resizes to fit text. (Less likely to overflow visually)
        # MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE (2): Text font size shrinks to fit shape. (Overflow is 'hidden' by shrinking)

        # For this check, we are most interested in cases where auto_size is NONE or TEXT_TO_FIT_SHAPE
        # (where text might become too small, but that's a different check - here we estimate raw overflow).

        shape_height_emu = shape.height
        margin_top_emu = text_frame.margin_top if text_frame.margin_top is not None else 0 # Ensure EMUs
        margin_bottom_emu = text_frame.margin_bottom if text_frame.margin_bottom is not None else 0 # Ensure EMUs
        available_text_height_emu = shape_height_emu - margin_top_emu - margin_bottom_emu

        ph_type_name = "Shape" # Default
        placeholder_idx_str = f"Shape ID {shape.shape_id}"
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type: # Check type before accessing value
            try:
                # Ensure placeholder_format.type is an enum before accessing .value
                current_ph_type = shape.placeholder_format.type
                if hasattr(current_ph_type, 'value'):
                    ph_type_name = self._get_placeholder_type_name(current_ph_type.value)
                else: # Is an int, direct use might be risky if not a known value
                    ph_type_name = self._get_placeholder_type_name(int(current_ph_type))
                placeholder_idx_str = f"Placeholder Type: {ph_type_name}, Index: {shape.placeholder_format.idx}"
            except Exception as e:
                logger.debug(f"Error getting placeholder type name for overflow check: {e}")
                ph_type_name = f"Placeholder Idx {shape.placeholder_format.idx}" # Fallback


        if available_text_height_emu <= 0:
            if text_frame.text.strip():
                feedback_items.append(ReviewFeedback(
                    slide_index=slide_index,
                    severity="medium",
                    category="text_overflow",
                    message=f"Shape '{shape.name if shape.name else placeholder_idx_str}' is too small for any text, but contains text (height: {shape.height}, margins: T{text_frame.margin_top} B{text_frame.margin_bottom}).",
                    suggestion="Increase shape height or reduce text/margins."
                ))
            return feedback_items

        estimated_lines = self._estimate_text_lines(text_frame, shape.width)

        if estimated_lines == 0 and not text_frame.text.strip():
            return feedback_items # No text, no overflow

        if estimated_lines == 0 and text_frame.text.strip():
            logger.warning(f"Text overflow estimator returned 0 lines for non-empty text in shape {shape.name or placeholder_idx_str} on slide {slide_index}")
            # Potentially add a specific ReviewFeedback for estimator failure if desired
            return feedback_items

        avg_font_size_pt = 18.0 # Default
        try:
            if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                # Attempt to get font size from the first run of the first paragraph
                first_run_font_size = text_frame.paragraphs[0].runs[0].font.size
                if first_run_font_size: # Check if size is not None
                    avg_font_size_pt = float(first_run_font_size.pt)
        except (AttributeError, IndexError): # No size set, no runs, or no paragraphs
            pass # Use default

        # EMUs per point: 1 point = 12700 EMUs
        # Line height factor (e.g., 1.2 for 120% line spacing)
        estimated_line_height_emu = avg_font_size_pt * 12700 * 1.2
        if estimated_line_height_emu == 0:
             estimated_line_height_emu = 12700 * 1.2

        required_height_emu = estimated_lines * estimated_line_height_emu

        if required_height_emu > available_text_height_emu:
            if self.validation_config.autofix_text_overflow:
                was_fixed = self._autofix_shape_text_overflow(shape, text_frame, slide_index, language)
                if was_fixed:
                    logger.info(f"Text overflow in shape '{shape.name or placeholder_idx_str}' on slide {slide_index} was auto-fixed.")
                    # No violation is added if fixed
                else:
                    # Add violation if autofix failed
                    feedback_items.append(ReviewFeedback(
                        slide_index=slide_index,
                        severity="medium",
                        category="text_overflow",
                        message=f"Estimated text ({estimated_lines} lines) overflows shape '{shape.name if shape.name else placeholder_idx_str}' (autofix failed). Available height: {available_text_height_emu:.0f} EMU, Required: {required_height_emu:.0f} EMU.",
                        suggestion="Manually reduce text, decrease font size, or increase shape height."
                    ))
            else:
                # Autofix disabled, just add violation
                feedback_items.append(ReviewFeedback(
                    slide_index=slide_index,
                    severity="medium",
                    category="text_overflow",
                    message=f"Estimated text ({estimated_lines} lines) overflows shape '{shape.name if shape.name else placeholder_idx_str}'. Available height: {available_text_height_emu:.0f} EMU, Required: {required_height_emu:.0f} EMU.",
                    suggestion="Enable autofix_text_overflow or manually adjust text/shape."
                ))

        return feedback_items

    def _check_slide_placeholder_population(self, slide, slide_index: int) -> List[ReviewFeedback]:
        feedback_items: List[ReviewFeedback] = []
        if not self.validation_config.check_placeholder_population:
            return feedback_items

        required_placeholder_types = {
            PP_PLACEHOLDER.TITLE.value: "Title",
            PP_PLACEHOLDER.BODY.value: "Body Content"
        }
        has_populated_title = False
        body_placeholder_exists = False
        body_placeholder_is_populated = False

        for ph in slide.placeholders:
            ph_type_val = ph.placeholder_format.type.value
            is_empty = True
            if hasattr(ph, 'text_frame') and ph.text_frame and ph.text_frame.text.strip():
                is_empty = False
            elif ph.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    if ph.image: is_empty = False
                except ValueError: pass

            if ph_type_val == PP_PLACEHOLDER.TITLE.value:
                if not is_empty: has_populated_title = True
            elif ph_type_val == PP_PLACEHOLDER.BODY.value:
                body_placeholder_exists = True
                if not is_empty: body_placeholder_is_populated = True

        title_type_name = required_placeholder_types[PP_PLACEHOLDER.TITLE.value]
        if PP_PLACEHOLDER.TITLE.value in required_placeholder_types and not has_populated_title:
            feedback_items.append(ReviewFeedback(
                slide_index=slide_index,
                severity="high", # Missing title is generally high severity
                category="placeholder_population",
                message=f"Required '{title_type_name}' placeholder is empty or missing.",
                suggestion=f"Ensure the slide has a populated '{title_type_name}' placeholder."
            ))

        body_type_name = required_placeholder_types[PP_PLACEHOLDER.BODY.value]
        if PP_PLACEHOLDER.BODY.value in required_placeholder_types and body_placeholder_exists and not body_placeholder_is_populated:
            feedback_items.append(ReviewFeedback(
                slide_index=slide_index,
                severity="medium", # Empty body might be intentional, so medium
                category="placeholder_population",
                message=f"Required '{body_type_name}' placeholder is present but empty.",
                suggestion=f"Populate the '{body_type_name}' placeholder or remove it if not needed for this slide layout."
            ))

        return feedback_items

    def _validate_shape_style(self, shape, slide, slide_index: int, language: str) -> List[ReviewFeedback]:
        """
        Validate style for a single shape.
        
        Args:
            shape: Shape object to validate
            slide: Slide object containing the shape
            slide_index: Index of the slide containing this shape
            language: Language code for the presentation
            
        Returns:
            List of ReviewFeedback objects for the shape.
        """
        feedback_items: List[ReviewFeedback] = [] # Changed from violations to feedback_items
        
        # Alignment check (applies to all shapes)
        # _check_shape_alignment_against_margins already returns List[ReviewFeedback]
        alignment_feedback = self._check_shape_alignment_against_margins(shape, slide, slide_index, language)
        feedback_items.extend(alignment_feedback)

        # Only validate text-containing shapes for further text-related checks
        if not hasattr(shape, 'text_frame') or not shape.text_frame:
            return feedback_items # Return here if no text frame, but alignment feedback is kept
        
        text_frame = shape.text_frame
        
        placeholder_type = None
        # Ensure shape.placeholder_format exists and then check for type attribute
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
            if hasattr(shape.placeholder_format, 'type'):
                 placeholder_type = shape.placeholder_format.type
        
        # Validate each paragraph
        # _validate_paragraph_style already returns List[ReviewFeedback]
        for para_index, paragraph in enumerate(text_frame.paragraphs):
            para_feedback = self._validate_paragraph_style(
                paragraph, placeholder_type, slide_index, para_index, language
            )
            feedback_items.extend(para_feedback)

        # Check for empty placeholders (This is different from _check_slide_placeholder_population, which checks specific TITLE/BODY)
        # This one seems to be about any placeholder having no text if allow_empty_placeholders is false.
        if not self.validation_config.allow_empty_placeholders and placeholder_type is not None:
            if not text_frame.text.strip():
                ph_type_name = self._get_placeholder_type_name(placeholder_type.value if hasattr(placeholder_type, 'value') else int(placeholder_type))
                feedback_items.append(ReviewFeedback(
                    slide_index=slide_index,
                    severity="low",
                    category="placeholder_content",
                    message=f"Placeholder '{ph_type_name}' (Name: '{shape.name if shape.name else 'Unnamed'}') is empty.",
                    suggestion=f"Add content to placeholder '{ph_type_name}' or set allow_empty_placeholders to True if intentional."
                ))
        
        # Check for text overflow
        # _check_text_frame_overflow already returns List[ReviewFeedback]
        overflow_feedback = self._check_text_frame_overflow(shape, slide_index, language)
        feedback_items.extend(overflow_feedback)

        # Image Accessibility Check (runs for MSO_SHAPE_TYPE.PICTURE)
        # _check_image_accessibility already returns List[ReviewFeedback]
        if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            accessibility_feedback = self._check_image_accessibility(shape, slide_index)
            feedback_items.extend(accessibility_feedback)

        return feedback_items

    def _check_image_accessibility(self, shape, slide_index: int) -> List[ReviewFeedback]:
        """
        Checks if an image shape has meaningful alt text (stored in shape.name).
        Args:
            shape: The image shape to check (must be MSO_SHAPE_TYPE.PICTURE).
            slide_index: Index of the slide.
        Returns:
            A list of accessibility violation dictionaries.
        """
        feedback_items: List[ReviewFeedback] = []
        # The field was renamed to check_alt_text_accessibility
        if not self.validation_config.check_alt_text_accessibility:
            return feedback_items

        # Generic names that suggest missing meaningful alt text.
        # This list can be expanded. Case-insensitive check.
        generic_names = [
            "picture", "image", "chart", "graph",
            "process flow", "diagram", "flowchart", "img", "pic", "graphic"
        ] # Added a few more common ones

        shape_name = shape.name
        is_missing_alt_text = False
        violation_message = ""

        if not shape_name or not shape_name.strip():
            is_missing_alt_text = True
            violation_message = f"Image shape (ID: {shape.shape_id}) is missing alt text (name property is empty)."
        else:
            lower_shape_name = shape_name.lower()
            if lower_shape_name in generic_names:
                 is_missing_alt_text = True
                 violation_message = f"Image '{shape_name}' has alt text that is too generic ('{shape_name}') and may not be descriptive."
            elif re.match(r"^(" + "|".join(generic_names) + r")(\s+\d+)?$", lower_shape_name, re.IGNORECASE):
                is_missing_alt_text = True
                violation_message = f"Image '{shape_name}' has generic alt text ('{shape_name}') that follows a default pattern and may not be descriptive."

        if is_missing_alt_text:
            feedback_items.append(ReviewFeedback(
                slide_index=slide_index,
                severity="medium", # Missing/generic alt text is important
                category="accessibility",
                message=violation_message,
                suggestion="Provide descriptive alt text for images using the 'name' property (accessible via Selection Pane in PowerPoint)."
            ))
        return feedback_items

    def _validate_paragraph_style(
        self, 
        paragraph, 
        placeholder_type: Optional[int], 
        slide_index: int, 
        para_index: int,
        language: str # Add language
    ) -> List[ReviewFeedback]:
        """
        Validate style for a single paragraph.
        
        Args:
            paragraph: Paragraph object to validate
            placeholder_type: Type of placeholder containing this paragraph
            slide_index: Index of the slide
            para_index: Index of the paragraph
            language: Language code for the presentation
            
        Returns:
            List of ReviewFeedback objects
        """
        feedback_items: List[ReviewFeedback] = []
        
        # Skip empty paragraphs
        if not paragraph.text.strip():
            return feedback_items
        
        # Get expected style for this placeholder type
        expected_font = self._get_expected_font(placeholder_type, paragraph.level, language) # Pass language
        expected_bullet = self._get_expected_bullet(placeholder_type, paragraph.level)
        
        # Validate each run in the paragraph
        for run_index, run in enumerate(paragraph.runs):
            # _validate_run_style now returns List[ReviewFeedback]
            run_feedback = self._validate_run_style(
                run, expected_font, slide_index, para_index, run_index, language # Pass language
            )
            feedback_items.extend(run_feedback)
        
        # Validate bullet style if this is a bulleted paragraph
        if expected_bullet and self.validation_config.check_bullet_styles:
            # Assuming _validate_bullet_style will be updated to return List[ReviewFeedback]
            bullet_feedback = self._validate_bullet_style(
                paragraph, expected_bullet, slide_index, para_index
            )
            feedback_items.extend(bullet_feedback)
        
        return feedback_items

    def _validate_run_style(
        self, 
        run, 
        expected_font: Optional[FontInfo], 
        slide_index: int, 
        para_index: int, 
        run_index: int,
        language: str # Add language
    ) -> List[ReviewFeedback]:
        """
        Validate style for a single text run.
        
        Args:
            run: Run object to validate
            expected_font: Expected font information
            slide_index: Index of the slide
            para_index: Index of the paragraph
            run_index: Index of the run
            language: Language code for the presentation
            
        Returns:
            List of ReviewFeedback objects
        """
        feedback_items: List[ReviewFeedback] = []
        
        if not expected_font:
            return feedback_items
        
        font = run.font
        location = f"slide {slide_index + 1}, paragraph {para_index + 1}, run {run_index + 1}" # User-friendly 1-based indexing
        
        # Validate font name
        if (self.validation_config.enforce_font_name and 
            expected_font.name and 
            font.name and 
            font.name != expected_font.name):

            is_valid_override = False
            if language and self.template_parser and self.template_parser.template_style:
                specific_font_name = self.template_parser.template_style.language_specific_fonts.get(language.lower())
                if specific_font_name and font.name == specific_font_name:
                    is_valid_override = True

            if not is_valid_override:
                feedback_items.append(ReviewFeedback(
                    slide_index=slide_index,
                    severity="low",
                    category="style_validation_font",
                    message=f"Font name mismatch at {location}. Expected: '{expected_font.name}', Actual: '{font.name}'.",
                    suggestion="Ensure font name matches the template style or language-specific settings."
                ))
        
        # Validate font size
        if expected_font.size and font.size:
            expected_size_pt = expected_font.size
            actual_size_pt = font.size.pt
            size_diff = abs(actual_size_pt - expected_size_pt)
            
            if size_diff > self.validation_config.font_size_tolerance:
                feedback_items.append(ReviewFeedback(
                    slide_index=slide_index,
                    severity="medium",
                    category="style_validation_font",
                    message=f"Font size mismatch at {location}. Expected: {expected_size_pt}pt, Actual: {actual_size_pt:.2f}pt (Tolerance: {self.validation_config.font_size_tolerance}pt).",
                    suggestion="Adjust font size to match the template style."
                ))
        
        # Validate font weight/bold
        if (self.validation_config.enforce_font_weight and 
            expected_font.weight):
            expected_bold = expected_font.weight.lower() == "bold"
            actual_bold = font.bold is True # font.bold can be None, True, False
            
            if expected_bold != actual_bold:
                feedback_items.append(ReviewFeedback(
                    slide_index=slide_index,
                    severity="low",
                    category="style_validation_font",
                    message=f"Font weight mismatch at {location}. Expected: '{expected_font.weight}', Actual: '{'bold' if actual_bold else 'normal'}'.",
                    suggestion="Adjust font weight (bold/normal) to match the template style."
                ))
        
        # Validate font color
        if (self.validation_config.enforce_color_compliance and 
            expected_font.color and 
            font.color):
            # _validate_color already returns List[ReviewFeedback]
            # Ensure the location string passed to _validate_color is consistent
            color_feedback = self._validate_color(
                font.color, expected_font.color, f"font color in {location}", slide_index
            )
            feedback_items.extend(color_feedback)
        
        return feedback_items

    def _validate_bullet_style(
        self, 
        paragraph, 
        expected_bullet: BulletInfo, 
        slide_index: int, 
        para_index: int
    ) -> List[ReviewFeedback]:
        """
        Validate bullet point style.
        
        Args:
            paragraph: Paragraph object to validate
            expected_bullet: Expected bullet information
            slide_index: Index of the slide
            para_index: Index of the paragraph
            
        Returns:
            List of ReviewFeedback objects.
        """
        feedback_items: List[ReviewFeedback] = []
        
        if not self.validation_config.enforce_bullet_characters:
            return feedback_items
        
        location_desc = f"paragraph {para_index + 1}"

        # This is a simplified check for bullet indentation level.
        # Validating actual bullet characters (e.g., '•' vs '-') is complex with current python-pptx capabilities
        # and would typically require deeper XML inspection or more advanced template parsing.
        
        if paragraph.level != expected_bullet.indent_level:
            feedback_items.append(ReviewFeedback(
                slide_index=slide_index,
                severity="low",
                category="style_validation_bullet", # Category for bullet style issues
                message=f"Bullet indentation level mismatch for {location_desc}. Expected level: {expected_bullet.indent_level}, Actual: {paragraph.level}.",
                suggestion=f"Adjust bullet level for {location_desc} to {expected_bullet.indent_level} to match template style."
            ))
        
        return feedback_items

    def _validate_color(
        self, 
        actual_color, 
        expected_color_hex: str, 
        location: str, 
        slide_index: int
    ) -> List[ReviewFeedback]:
        """
        Validate color against expected color with tolerance.
        
        Args:
            actual_color: Color object from python-pptx
            expected_color_hex: Expected color in hex format
            location: Description of where the color is used
            slide_index: Index of the slide
            
        Returns:
            List of ReviewFeedback objects
        """
        feedback_items: List[ReviewFeedback] = []
        
        try:
            actual_hex = self._color_to_hex(actual_color)
            if not actual_hex:
                # Cannot determine actual color, cannot compare. Could log or return specific feedback.
                return feedback_items
            
            if not self._colors_match(actual_hex, expected_color_hex):
                feedback_items.append(ReviewFeedback(
                    slide_index=slide_index,
                    severity="low", # Color mismatches are often minor unless severe contrast issues
                    category="style_validation_color",
                    message=f"Color mismatch at {location}. Expected: {expected_color_hex}, Actual: {actual_hex}.",
                    suggestion=f"Ensure color at {location} matches the template's expected color {expected_color_hex} (tolerance: {self.validation_config.color_tolerance})."
                ))
        
        except Exception as e:
            logger.debug(f"Failed to validate color at {location}: {e}")
            feedback_items.append(ReviewFeedback(
                slide_index=slide_index,
                severity="low",
                category="style_validation_error",
                message=f"Error validating color at {location}: {e}",
                suggestion="Check color parsing or comparison logic."
            ))

        return feedback_items

    def _color_to_hex(self, color) -> Optional[str]:
        """
        Convert python-pptx color to hex string.
        
        Args:
            color: Color object from python-pptx
            
        Returns:
            Hex color string or None if conversion fails
        """
        try:
            if hasattr(color, 'rgb') and color.rgb:
                rgb = color.rgb
                return f"#{rgb:06X}"
        except Exception:
            pass
        
        return None

    def _colors_match(self, color1_hex: str, color2_hex: str) -> bool:
        """
        Check if two hex colors match within tolerance.
        
        Args:
            color1_hex: First color in hex format
            color2_hex: Second color in hex format
            
        Returns:
            True if colors match within tolerance
        """
        if color1_hex.lower() == color2_hex.lower():
            return True
        
        # If tolerance is 0, require exact match
        if self.validation_config.color_tolerance == 0:
            return False
        
        try:
            # Parse hex colors to RGB
            r1, g1, b1 = self._hex_to_rgb(color1_hex)
            r2, g2, b2 = self._hex_to_rgb(color2_hex)
            
            # Calculate color difference as percentage of maximum possible difference
            max_diff = 255 * 3  # Maximum possible RGB difference
            actual_diff = abs(r1 - r2) + abs(g1 - g2) + abs(b1 - b2)
            
            return (actual_diff / max_diff) <= self.validation_config.color_tolerance
            
        except Exception:
            return False

    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """
        Convert hex color to RGB tuple.
        
        Args:
            hex_color: Color in hex format (e.g., "#FF0000")
            
        Returns:
            RGB tuple (r, g, b)
        """
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def _get_expected_font(self, placeholder_type: Optional[int], level: int = 0, language: Optional[str] = None) -> Optional[FontInfo]: # Add language
        """
        Get expected font for a placeholder type and indentation level.
        
        Args:
            placeholder_type: PowerPoint placeholder type number
            level: Indentation level
            language: Language code for considering specific fonts
            
        Returns:
            Expected FontInfo or None
        """
        base_font: Optional[FontInfo] = None
        if not placeholder_type:
            base_font = self.template_parser.template_style.master_font
        else:
            # For bullet points, get font from bullet style
            bullet_info = self.template_parser.get_bullet_style_for_level(placeholder_type, level)
            if bullet_info and bullet_info.font:
                base_font = bullet_info.font
            else:
                # Fall back to placeholder default font
                base_font = self.template_parser.get_font_for_placeholder_type(placeholder_type)

        if (language and self.template_parser and 
            self.template_parser.template_style and
            hasattr(self.template_parser.template_style, 'language_specific_fonts')):
            specific_font_name = self.template_parser.template_style.language_specific_fonts.get(language.lower())
            if specific_font_name and base_font:
                # Create a new FontInfo object with the name overridden
                # Assuming FontInfo is a Pydantic model, create a new one or copy and update.
                overridden_font = base_font.model_copy() # Pydantic v2 (use .copy(deep=True) for v1)
                overridden_font.name = specific_font_name
                return overridden_font

        return base_font # Original expected font

    def _get_expected_bullet(self, placeholder_type: Optional[int], level: int = 0) -> Optional[BulletInfo]:
        """
        Get expected bullet style for a placeholder type and indentation level.
        
        Args:
            placeholder_type: PowerPoint placeholder type number
            level: Indentation level
            
        Returns:
            Expected BulletInfo or None
        """
        if not placeholder_type:
            return None
        
        return self.template_parser.get_bullet_style_for_level(placeholder_type, level)

    def _apply_font_adjustments(self, slide, slide_plan: SlidePlan) -> None:
        """
        Apply font adjustments from content fit analysis.
        
        Args:
            slide: PowerPoint slide object
            slide_plan: SlidePlan with potential font adjustments
        """
        # Check if slide has font adjustment metadata
        font_adjustment = None
        if hasattr(slide_plan, '__dict__') and 'font_adjustment' in slide_plan.__dict__:
            font_adjustment = slide_plan.__dict__['font_adjustment']
        
        if not font_adjustment:
            return  # No adjustments needed
        
        logger.debug(f"Applying font adjustment to slide {slide_plan.index}: "
                    f"{font_adjustment.original_size}pt → {font_adjustment.recommended_size}pt")
        
        try:
            # Apply font adjustment to all text content
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    self._adjust_text_frame_font(shape.text_frame, font_adjustment)
                elif hasattr(shape, 'text') and shape.text:
                    # For simple text shapes
                    if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.size and run.font.size.pt == font_adjustment.original_size:
                                    run.font.size = Pt(font_adjustment.recommended_size)
        
        except Exception as e:
            logger.error(f"Failed to apply font adjustments to slide {slide_plan.index}: {e}")

    def _adjust_text_frame_font(self, text_frame, font_adjustment: FontAdjustment) -> None:
        """
        Adjust font size in a text frame.
        
        Args:
            text_frame: PowerPoint text frame object
            font_adjustment: FontAdjustment with size information
        """
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    # Only adjust if font size matches the original size
                    if run.font.size:
                        current_size = run.font.size.pt
                        # Allow some tolerance for floating point differences
                        if abs(current_size - font_adjustment.original_size) <= 1:
                            run.font.size = Pt(font_adjustment.recommended_size)
                            logger.debug(f"Adjusted font size: {current_size}pt → {font_adjustment.recommended_size}pt")
                    else:
                        # If no explicit size set, assume it's the default and adjust if it matches
                        run.font.size = Pt(font_adjustment.recommended_size)
        
        except Exception as e:
            logger.debug(f"Error adjusting text frame font: {e}")

    def _hide_empty_placeholder(self, placeholder, placeholder_name: str) -> None:
        """
        Hide an empty placeholder to avoid style validation warnings.
        
        Args:
            placeholder: Placeholder object to hide
            placeholder_name: Name for logging purposes
        """
        try:
            # Method 1: Try to make the placeholder invisible
            if hasattr(placeholder, 'visible'):
                placeholder.visible = False
                logger.debug(f"Set {placeholder_name} placeholder invisible")
                return
            
            # Method 2: Try to set shape visibility
            if hasattr(placeholder, 'element') and hasattr(placeholder.element, 'set'):
                # Try to add hidden attribute to the shape
                try:
                    placeholder.element.set('hidden', '1')
                    logger.debug(f"Set {placeholder_name} placeholder hidden via XML attribute")
                    return
                except:
                    pass
            
            # Method 3: Move placeholder off-slide (position outside visible area)
            if hasattr(placeholder, 'left') and hasattr(placeholder, 'top'):
                placeholder.left = Inches(-10)  # Move far left, outside slide
                placeholder.top = Inches(-10)   # Move far up, outside slide
                logger.debug(f"Moved {placeholder_name} placeholder off-slide to hide it")
                return
            
            # Method 4: Set placeholder size to zero
            if hasattr(placeholder, 'width') and hasattr(placeholder, 'height'):
                placeholder.width = Inches(0.01)  # Minimal size
                placeholder.height = Inches(0.01)
                logger.debug(f"Set {placeholder_name} placeholder to minimal size")
                return
            
            # Method 5: Clear any default text that might be in the placeholder
            if hasattr(placeholder, 'text_frame') and placeholder.text_frame:
                placeholder.text_frame.clear()
                if hasattr(placeholder.text_frame, 'auto_size'):
                    try:
                        placeholder.text_frame.auto_size = MSO_AUTO_SIZE.NONE
                    except:
                        pass
                logger.debug(f"Cleared {placeholder_name} placeholder text frame")
                return
            
            logger.warning(f"Could not hide {placeholder_name} placeholder - no suitable method found")
            
        except Exception as e:
            logger.warning(f"Failed to hide {placeholder_name} placeholder: {e}")

    def _cleanup_empty_subtitle_placeholders(self, secondary_placeholders: List) -> None:
        """
        Clean up any subtitle placeholders that still contain default or empty text.
        
        Args:
            secondary_placeholders: List of secondary content placeholders
        """
        try:
            for placeholder in secondary_placeholders:
                if placeholder.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                    # Check if placeholder has default or problematic text
                    if hasattr(placeholder, 'text'):
                        current_text = placeholder.text.strip().lower()
                        
                        # List of default subtitle texts that should be cleaned up
                        default_texts = ['subtitle', 'click to add subtitle', 'add subtitle', '']
                        
                        if current_text in default_texts:
                            self._hide_empty_placeholder(placeholder, "subtitle")
                            logger.debug(f"Cleaned up subtitle placeholder with default text: '{current_text}'")
                    elif hasattr(placeholder, 'text_frame') and placeholder.text_frame:
                        # Check text frame for default content
                        text_content = ""
                        for para in placeholder.text_frame.paragraphs:
                            text_content += para.text.strip()
                        
                        if text_content.lower() in ['subtitle', 'click to add subtitle', 'add subtitle', '']:
                            self._hide_empty_placeholder(placeholder, "subtitle")
                            logger.debug(f"Cleaned up subtitle placeholder text frame with default content")
                            
        except Exception as e:
            logger.warning(f"Failed to cleanup subtitle placeholders: {e}")

    def _remove_empty_placeholders_from_slide(self, slide) -> int:
        """
        Intelligently remove or hide empty placeholders from a slide, only after
        attempting to utilize them for content distribution.
        
        Args:
            slide: Slide object to process
            
        Returns:
            Number of placeholders hidden/removed
        """
        # First, check if we have utilized available placeholders effectively
        placeholders_discovered = self._discover_all_content_placeholders(slide)
        utilization_check = self._check_placeholder_utilization(placeholders_discovered)
        
        if not utilization_check:
            logger.info(f"Low placeholder utilization detected. Attempting content redistribution before hiding placeholders.")
            # TODO: Could trigger content redistribution here if needed
        
        hidden_count = 0
        
        try:
            for placeholder in slide.placeholders:
                is_empty = True
                placeholder_type_val = placeholder.placeholder_format.type.value if hasattr(placeholder.placeholder_format.type, 'value') else int(placeholder.placeholder_format.type)

                # Check for picture placeholders (PICTURE (18) or OBJECT (10))
                if placeholder.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and \
                   (placeholder_type_val == PP_PLACEHOLDER.PICTURE.value or placeholder_type_val == PP_PLACEHOLDER.OBJECT.value):
                    try:
                        if placeholder.image:  # Accessing .image throws ValueError if no image
                            is_empty = False
                            logger.debug(f"Placeholder {placeholder.name} (type: {self._get_placeholder_type_name(placeholder_type_val)}) contains an image.")
                    except ValueError:
                        # No image, it's potentially empty. For OBJECT, also check text.
                        if placeholder_type_val == PP_PLACEHOLDER.PICTURE.value:
                            is_empty = True
                            logger.debug(f"Picture Placeholder {placeholder.name} is empty (no image).")
                        elif placeholder_type_val == PP_PLACEHOLDER.OBJECT.value:
                            # For OBJECT placeholders, if no image, check for text.
                            if hasattr(placeholder, 'text_frame') and placeholder.text_frame and placeholder.text_frame.text.strip():
                                is_empty = False
                                logger.debug(f"Object Placeholder {placeholder.name} has text, not empty.")
                            # Could add checks for .has_chart or .has_table if needed in future
                            # else: is_empty remains true if no image and no text
                            else:
                                logger.debug(f"Object Placeholder {placeholder.name} is empty (no image/text).")
                
                # Check for text in other placeholder types or if it's an OBJECT placeholder not emptied by image check
                # This 'elif' ensures text check runs if it's not a picture-type placeholder,
                # or if it's an OBJECT type that was found to not contain an image (is_empty would still be true).
                elif hasattr(placeholder, 'text_frame') and placeholder.text_frame and placeholder.text_frame.text.strip():
                    is_empty = False
                    logger.debug(f"Text Placeholder {placeholder.name} (type: {self._get_placeholder_type_name(placeholder_type_val)}) has text.")

                # Fallback for very old placeholder types that might only have .text (less common)
                elif not hasattr(placeholder, 'text_frame') and hasattr(placeholder, 'text') and placeholder.text and placeholder.text.strip():
                    is_empty = False
                    logger.debug(f"Legacy Text Placeholder {placeholder.name} (type: {self._get_placeholder_type_name(placeholder_type_val)}) has text.")
                
                else:
                    # If none of the above conditions met (e.g. not pic/obj, no text_frame, or text_frame is empty)
                    # is_empty remains true by default.
                    if not (placeholder.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and \
                            (placeholder_type_val == PP_PLACEHOLDER.PICTURE.value or placeholder_type_val == PP_PLACEHOLDER.OBJECT.value)):
                         logger.debug(f"Placeholder {placeholder.name} (type: {self._get_placeholder_type_name(placeholder_type_val)}) is considered empty (no text).")


                # Hide empty placeholders (except title which should always be visible - type 1)
                if is_empty and placeholder_type_val != PP_PLACEHOLDER.TITLE.value:
                    placeholder_type_name = self._get_placeholder_type_name(placeholder_type_val)
                    self._hide_empty_placeholder(placeholder, placeholder_type_name)
                    hidden_count += 1
        
        except Exception as e:
            logger.warning(f"Error processing empty placeholders: {e}")
        
        return hidden_count

    def _get_placeholder_type_name(self, placeholder_type: int) -> str:
        """
        Get human-readable name for placeholder type.
        
        Args:
            placeholder_type: PowerPoint placeholder type constant
            
        Returns:
            Human-readable placeholder type name
        """
        type_names = {
            1: "title",
            2: "body", 
            3: "subtitle",
            4: "center_title",
            5: "center_subtitle",
            6: "date_time",
            7: "footer",
            8: "header",
            9: "slide_number",
            10: "picture",
            11: "chart",
            12: "table",
            13: "clip_art",
            14: "organization_chart",
            15: "media_clip",
            16: "vertical_object",
            17: "vertical_body",
            18: "picture"
        }
        
        return type_names.get(placeholder_type, f"unknown_type_{placeholder_type}")