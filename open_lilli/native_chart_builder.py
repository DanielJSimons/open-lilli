"""Native PowerPoint chart builder for creating editable chart objects.

This module implements T-51: Replace PNG fallback for bar/line with editable chart objects.
The generated charts are native PowerPoint chart objects that can be edited by users.
"""

import logging
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Union

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

from .models import ChartType, NativeChartData, TemplateStyle
from .template_parser import TemplateParser

logger = logging.getLogger(__name__)


class NativeChartBuilder:
    """Builds native PowerPoint chart objects that can be edited by users."""

    def __init__(self, template_parser: Optional[TemplateParser] = None):
        """
        Initialize the native chart builder.
        
        Args:
            template_parser: Optional template parser for color palette access
        """
        self.template_parser = template_parser
        self.chart_type_mapping = {
            ChartType.BAR: XL_CHART_TYPE.BAR_CLUSTERED,
            ChartType.COLUMN: XL_CHART_TYPE.COLUMN_CLUSTERED,
            ChartType.LINE: XL_CHART_TYPE.LINE,
            ChartType.PIE: XL_CHART_TYPE.PIE,
            ChartType.SCATTER: XL_CHART_TYPE.XY_SCATTER,
            ChartType.AREA: XL_CHART_TYPE.AREA,
            ChartType.DOUGHNUT: XL_CHART_TYPE.DOUGHNUT
        }
        
        logger.info("NativeChartBuilder initialized")

    def create_native_chart(
        self,
        slide,
        chart_config: NativeChartData,
        position: Optional[Tuple[Inches, Inches, Inches, Inches]] = None
    ) -> bool:
        """
        Create a native PowerPoint chart object on a slide.
        
        Args:
            slide: PowerPoint slide object
            chart_config: Chart configuration
            position: Optional (left, top, width, height) in Inches
            
        Returns:
            True if chart was created successfully, False otherwise
        """
        try:
            # Set default position if not provided
            if position is None:
                position = (Inches(1), Inches(2), Inches(8), Inches(5))
            
            left, top, width, height = position
            
            # Create chart data
            chart_data = self._create_chart_data(chart_config)
            
            # Get PowerPoint chart type
            xl_chart_type = self.chart_type_mapping.get(
                chart_config.chart_type, 
                XL_CHART_TYPE.COLUMN_CLUSTERED
            )
            
            # Add chart to slide
            chart = slide.shapes.add_chart(
                xl_chart_type, left, top, width, height, chart_data
            ).chart
            
            # Configure chart appearance
            self._configure_chart_appearance(chart, chart_config)
            
            # Apply template colors if enabled
            if chart_config.use_template_colors and self.template_parser:
                self._apply_template_colors(chart, chart_config)
            
            logger.info(f"Created native {chart_config.chart_type} chart: {chart_config.title}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to create native chart: {e}")
            return False

    def _create_chart_data(self, chart_config: NativeChartData) -> ChartData:
        """
        Create chart data object for PowerPoint.
        
        Args:
            chart_config: Chart configuration
            
        Returns:
            ChartData object for PowerPoint
        """
        if chart_config.chart_type in [ChartType.PIE, ChartType.DOUGHNUT]:
            # Pie charts use single series data
            chart_data = CategoryChartData()
            chart_data.categories = chart_config.categories
            
            # Use first series for pie chart
            if chart_config.series:
                first_series = chart_config.series[0]
                series_name = first_series.get("name", "Values")
                series_values = first_series.get("values", [])
                chart_data.add_series(series_name, series_values)
            
        else:
            # Multi-series charts (bar, column, line, etc.)
            chart_data = CategoryChartData()
            chart_data.categories = chart_config.categories
            
            # Add each data series
            for series in chart_config.series:
                series_name = series.get("name", "Series")
                series_values = series.get("values", [])
                chart_data.add_series(series_name, series_values)
        
        return chart_data

    def _configure_chart_appearance(self, chart, chart_config: NativeChartData) -> None:
        """
        Configure chart appearance settings.
        
        Args:
            chart: PowerPoint chart object
            chart_config: Chart configuration
        """
        try:
            # Set chart title
            if chart.has_title:
                chart.chart_title.text_frame.text = chart_config.title
                
                # Style the title
                title_paragraph = chart.chart_title.text_frame.paragraphs[0]
                title_run = title_paragraph.runs[0]
                title_run.font.size = Pt(14)
                title_run.font.bold = True
            
            # Configure legend
            if chart_config.has_legend and chart.has_legend:
                chart.legend.position = XL_LEGEND_POSITION.RIGHT
                chart.legend.include_in_layout = False
            elif not chart_config.has_legend and chart.has_legend:
                chart.legend.position = XL_LEGEND_POSITION.NONE
            
            # Configure data labels
            if chart_config.has_data_labels:
                for series in chart.series:
                    if hasattr(series, 'has_data_labels'):
                        series.has_data_labels = True
            
            # Configure axes for non-pie charts
            if chart_config.chart_type not in [ChartType.PIE, ChartType.DOUGHNUT]:
                self._configure_chart_axes(chart, chart_config)
            
        except Exception as e:
            logger.error(f"Failed to configure chart appearance: {e}")

    def _configure_chart_axes(self, chart, chart_config: NativeChartData) -> None:
        """
        Configure chart axes titles and formatting.
        
        Args:
            chart: PowerPoint chart object
            chart_config: Chart configuration
        """
        try:
            # Configure category axis (X-axis)
            if hasattr(chart, 'category_axis') and chart_config.x_axis_title:
                category_axis = chart.category_axis
                if hasattr(category_axis, 'has_title'):
                    category_axis.has_title = True
                    category_axis.axis_title.text_frame.text = chart_config.x_axis_title
            
            # Configure value axis (Y-axis)
            if hasattr(chart, 'value_axis') and chart_config.y_axis_title:
                value_axis = chart.value_axis
                if hasattr(value_axis, 'has_title'):
                    value_axis.has_title = True
                    value_axis.axis_title.text_frame.text = chart_config.y_axis_title
            
        except Exception as e:
            logger.error(f"Failed to configure chart axes: {e}")

    def _apply_template_colors(self, chart, chart_config: NativeChartData) -> None:
        """
        Apply template color palette to chart.
        
        Args:
            chart: PowerPoint chart object
            chart_config: Chart configuration
        """
        try:
            if not self.template_parser or not hasattr(self.template_parser, 'palette'):
                return
            
            palette = self.template_parser.palette
            
            # Get template colors
            template_colors = []
            if 'acc1' in palette:
                template_colors.append(palette['acc1'])
            if 'acc2' in palette:
                template_colors.append(palette['acc2'])
            if 'acc3' in palette:
                template_colors.append(palette['acc3'])
            if 'acc4' in palette:
                template_colors.append(palette['acc4'])
            
            # Fallback colors if template doesn't have enough accent colors
            if not template_colors:
                template_colors = ['#1F497D', '#4F81BD', '#9BBB59', '#F79646', '#8064A2']
            
            # Apply colors to chart series
            for i, series in enumerate(chart.series):
                if i < len(template_colors):
                    color_hex = template_colors[i % len(template_colors)]
                    try:
                        # Convert hex to RGB
                        rgb = self._hex_to_rgb(color_hex)
                        
                        # Apply color to series
                        if hasattr(series, 'format') and hasattr(series.format, 'fill'):
                            fill = series.format.fill
                            fill.solid()
                            fill.fore_color.rgb = rgb
                        
                    except Exception as e:
                        logger.debug(f"Could not apply color to series {i}: {e}")
            
            logger.debug(f"Applied template colors to chart: {template_colors}")
            
        except Exception as e:
            logger.error(f"Failed to apply template colors: {e}")

    def _hex_to_rgb(self, hex_color: str):
        """
        Convert hex color to RGB for PowerPoint.
        
        Args:
            hex_color: Hex color string (e.g., "#1F497D")
            
        Returns:
            RGB color object
        """
        from pptx.dml.color import RGBColor
        
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        
        return RGBColor(r, g, b)

    def find_chart_placeholder(self, slide) -> Optional[object]:
        """
        Find a chart placeholder on the slide.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            Chart placeholder or None if not found
        """
        try:
            for placeholder in slide.placeholders:
                if (hasattr(placeholder, 'placeholder_format') and 
                    placeholder.placeholder_format.type == 12):  # CHART placeholder
                    return placeholder
                
                # Check for content placeholder that can hold charts
                if (hasattr(placeholder, 'placeholder_format') and 
                    placeholder.placeholder_format.type in [7, 2]):  # OBJECT or BODY
                    return placeholder
            
        except Exception as e:
            logger.debug(f"Error finding chart placeholder: {e}")
        
        return None

    def create_chart_in_placeholder(
        self,
        slide,
        chart_config: NativeChartData,
        placeholder=None
    ) -> bool:
        """
        Create a chart in a specific placeholder.
        
        Args:
            slide: PowerPoint slide object
            chart_config: Chart configuration
            placeholder: Specific placeholder to use, or None to find automatically
            
        Returns:
            True if chart was created successfully
        """
        try:
            # Find placeholder if not provided
            if placeholder is None:
                placeholder = self.find_chart_placeholder(slide)
            
            if placeholder is None:
                logger.warning("No suitable placeholder found for chart, using default position")
                return self.create_native_chart(slide, chart_config)
            
            # Create chart data
            chart_data = self._create_chart_data(chart_config)
            
            # Get PowerPoint chart type
            xl_chart_type = self.chart_type_mapping.get(
                chart_config.chart_type, 
                XL_CHART_TYPE.COLUMN_CLUSTERED
            )
            
            # Insert chart into placeholder
            chart = placeholder.insert_chart(xl_chart_type, chart_data).chart
            
            # Configure chart appearance
            self._configure_chart_appearance(chart, chart_config)
            
            # Apply template colors if enabled
            if chart_config.use_template_colors and self.template_parser:
                self._apply_template_colors(chart, chart_config)
            
            logger.info(f"Created native chart in placeholder: {chart_config.title}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to create chart in placeholder: {e}")
            return False

    def validate_chart_config(self, chart_config: NativeChartData) -> List[str]:
        """
        Validate chart configuration and return list of issues.
        
        Args:
            chart_config: Chart configuration to validate
            
        Returns:
            List of validation issues (empty if valid)
        """
        issues = []
        
        # Check required fields
        if not chart_config.title:
            issues.append("Chart title is required")
        
        if not chart_config.categories:
            issues.append("Chart categories are required")
        
        if not chart_config.series:
            issues.append("Chart series data is required")
        
        # Validate series data
        for i, series in enumerate(chart_config.series):
            if not isinstance(series, dict):
                issues.append(f"Series {i}: Must be a dictionary")
                continue
            
            if "name" not in series:
                issues.append(f"Series {i}: Missing 'name' field")
            
            if "values" not in series:
                issues.append(f"Series {i}: Missing 'values' field")
                continue
            
            values = series["values"]
            if not isinstance(values, list):
                issues.append(f"Series {i}: 'values' must be a list")
                continue
            
            # Check that values count matches categories count
            if len(values) != len(chart_config.categories):
                issues.append(
                    f"Series {i}: Values count ({len(values)}) "
                    f"doesn't match categories count ({len(chart_config.categories)})"
                )
        
        # Chart type specific validations
        if chart_config.chart_type in [ChartType.PIE, ChartType.DOUGHNUT]:
            if len(chart_config.series) > 1:
                issues.append(f"{chart_config.chart_type} charts support only one data series")
        
        return issues

    def get_supported_chart_types(self) -> List[str]:
        """
        Get list of supported chart types.
        
        Returns:
            List of supported chart type strings
        """
        return [chart_type.value for chart_type in ChartType]

    def convert_legacy_chart_data(self, legacy_data: Dict) -> Optional[NativeChartData]:
        """
        Convert legacy chart data format to NativeChartData.
        
        Args:
            legacy_data: Old chart data format from visual generator
            
        Returns:
            NativeChartData object or None if conversion fails
        """
        try:
            # Map legacy type to new enum
            legacy_type = legacy_data.get("type", "bar").lower()
            chart_type_map = {
                "bar": ChartType.COLUMN,  # PowerPoint "bar" is actually column
                "column": ChartType.COLUMN,
                "line": ChartType.LINE,
                "pie": ChartType.PIE,
                "scatter": ChartType.SCATTER
            }
            
            chart_type = chart_type_map.get(legacy_type, ChartType.COLUMN)
            
            # Extract categories
            categories = legacy_data.get("categories", legacy_data.get("x", ["Category"]))
            
            # Extract series data
            series = []
            if "values" in legacy_data or "y" in legacy_data:
                # Single series format
                values = legacy_data.get("values", legacy_data.get("y", [1]))
                # Flatten nested lists if present
                values = self._flatten_values(values)
                series = [{"name": "Series 1", "values": values}]
            elif "series" in legacy_data:
                # Multi-series format
                series = []
                for s in legacy_data["series"]:
                    series_data = {"name": s.get("name", "Series"), "values": self._flatten_values(s.get("values", []))}
                    series.append(series_data)
            
            # Create NativeChartData
            native_chart = NativeChartData(
                chart_type=chart_type,
                title=legacy_data.get("title", "Chart"),
                categories=categories,
                series=series,
                x_axis_title=legacy_data.get("xlabel"),
                y_axis_title=legacy_data.get("ylabel"),
                has_legend=len(series) > 1,
                has_data_labels=False,
                use_template_colors=True
            )
            
            logger.debug(f"Converted legacy chart data to native format: {chart_type}")
            return native_chart
            
        except Exception as e:
            logger.error(f"Failed to convert legacy chart data: {e}")
            return None
    
    def _flatten_values(self, values: List) -> List[float]:
        """
        Flatten nested lists and convert to floats.
        
        Args:
            values: List that may contain nested lists
            
        Returns:
            Flattened list of float values
        """
        flattened = []
        for item in values:
            if isinstance(item, list):
                # If item is a list, take the first element (or average if multiple)
                if len(item) > 0:
                    if len(item) == 1:
                        flattened.append(float(item[0]))
                    else:
                        # Take average of multiple values
                        avg = sum(item) / len(item)
                        flattened.append(float(avg))
                else:
                    flattened.append(0.0)
            else:
                # Item is already a single value
                try:
                    flattened.append(float(item))
                except (ValueError, TypeError):
                    flattened.append(0.0)
        return flattened